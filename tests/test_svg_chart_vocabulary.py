"""圖表 SVG 詞彙擴充的契約（add-deck-delivery-line tasks 2.1b）。

## 為什麼要它

圖表原本是「截成 PNG 貼進投影片」，改為**原生繪製**後，圖表 SVG 的元素要和
版面元素一樣走窄轉換器。基礎詞彙（`<rect>`／`<text>`／`<image>`）在
`test_svg_to_pptx_converter.py`；本檔只管**圖表多出來的**那些。

## 詞彙來自實掃，不是憑印象

2026-08-14 掃 `output/report_trial_20260811_094014/` 的 14 張真實圖表 SVG，
連**屬性值的形式**一起掃——上一次只掃元素標籤，漏了 `width="100%"`，
而那會讓 13/14 張在第一行就炸。

| 元素 | 次數 | 對應 |
|---|---|---|
| `<circle>` | 108 | `MSO_SHAPE.OVAL` |
| `<line>` | 19 | ✅ 全為水平／垂直（0 斜線）→ 細 rect 退化形 |
| `<polyline>` | 2 | `build_freeform()` |
| `<defs>`＋`<pattern>` | 1 | 🔴 hatch，見下 |

🔴 **hatch 不可退化成純色**：`chart_runner.py:826` 寫明「顏色分段＝申請結構，
**斜紋疊加＝已轉讓**」——第二個視覺通道承載獨立資訊。退化會讓「已轉讓」消失，
而且不會有任何東西報錯。

⚠ 本檔用合成 SVG，但屬性形式**逐一照抄實掃結果**（`fill="white"`、
`width="100%"`、`text-anchor="middle"`、`stroke-dasharray="6 4"`…）。
另有一組整合測試直接吃真實檔案，檔案不在時 skip。
"""
from __future__ import annotations

import importlib.util
import sys
import unittest
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = PROJECT_ROOT / "skills" / "html-report-to-deck" / "scripts"
REAL_CHARTS = Path(r"D:\力山\專案\專利_ppt自動\output\report_trial_20260811_094014")

DPI = 96


def _load():
    path = SCRIPTS / "svg_to_pptx.py"
    spec = importlib.util.spec_from_file_location("svg_to_pptx", path)
    module = importlib.util.module_from_spec(spec)
    sys.modules["svg_to_pptx"] = module
    spec.loader.exec_module(module)
    return module


def _svg(body: str, w: int = 949, h: int = 460, extra: str = "") -> str:
    """圖表尺寸預設用實測最常見的 949×460。"""
    return (f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'width="{w}" height="{h}" viewBox="0 0 {w} {h}"{extra}>{body}</svg>')


def _emu_px(px: float) -> int:
    return round(px / DPI * 914400)


class ChartElementTests(unittest.TestCase):
    """四種圖表元素的映射。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    def test_circle_becomes_oval(self):
        """`<circle cx cy r>` → OVAL，且 cx/cy 是**圓心**不是左上角。"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = self.mod.build([_svg('<circle cx="100" cy="200" r="12" fill="#006DF5"/>')])
        shapes = list(prs.slides[0].shapes)
        self.assertEqual(len(shapes), 1)
        shape = shapes[0]
        self.assertEqual(shape.shape_type, MSO_SHAPE_TYPE.AUTO_SHAPE)
        # 圓心 (100,200) 半徑 12 → 外框左上 (88,188)、邊長 24
        self.assertEqual(shape.left, _emu_px(88))
        self.assertEqual(shape.top, _emu_px(188))
        self.assertEqual(shape.width, _emu_px(24))
        self.assertEqual(shape.height, _emu_px(24))

    def test_horizontal_line_becomes_thin_rect(self):
        """水平 `<line>` → 細 rect（既有原則：線是矩形的退化形）。"""
        prs = self.mod.build([_svg('<line x1="50" y1="100" x2="250" y2="100" '
                                   'stroke="#00094A" stroke-width="2"/>')])
        shape = list(prs.slides[0].shapes)[0]
        self.assertEqual(shape.left, _emu_px(50))
        self.assertEqual(shape.width, _emu_px(200))
        self.assertEqual(shape.height, _emu_px(2))
        # 線以中心對齊：y=100、粗 2 → top = 99
        self.assertEqual(shape.top, _emu_px(99))

    def test_vertical_line_becomes_thin_rect(self):
        prs = self.mod.build([_svg('<line x1="80" y1="20" x2="80" y2="220" '
                                   'stroke="#00094A" stroke-width="4"/>')])
        shape = list(prs.slides[0].shapes)[0]
        self.assertEqual(shape.width, _emu_px(4))
        self.assertEqual(shape.height, _emu_px(200))
        self.assertEqual(shape.left, _emu_px(78))

    def test_diagonal_line_fails_loud(self):
        """⚠ 實掃 19 條全是水平／垂直。真出現斜線＝引擎變了，要 fail 讓人知道，
        不可用細 rect 硬畫成水平（那會靜默畫錯）。"""
        with self.assertRaises(self.mod.UnsupportedElement):
            self.mod.build([_svg('<line x1="0" y1="0" x2="100" y2="100" '
                                 'stroke="#00094A"/>')])

    def test_polyline_becomes_freeform(self):
        prs = self.mod.build([_svg('<polyline points="10,20 30,40 50,25" '
                                   'fill="none" stroke="#006DF5" stroke-width="2"/>')])
        self.assertEqual(len(list(prs.slides[0].shapes)), 1)


class AttributeFormTests(unittest.TestCase):
    """🔴 五個屬性值形式缺口——每一個都會當場炸或靜默畫錯。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    def test_percent_size_resolves_against_viewbox(self):
        """13/14 張的背景 rect 用 `width="100%"`。`_px()` 只吃絕對值會炸。"""
        prs = self.mod.build([_svg('<rect width="100%" height="100%" fill="white"/>',
                                   w=949, h=460)])
        shape = list(prs.slides[0].shapes)[0]
        self.assertEqual(shape.width, _emu_px(949))
        self.assertEqual(shape.height, _emu_px(460))

    def test_named_color_white(self):
        """`fill="white"` ×14。`_color()` 只收 #RRGGBB 會 raise。"""
        from pptx.dml.color import RGBColor

        prs = self.mod.build([_svg('<rect x="0" y="0" width="10" height="10" '
                                   'fill="white"/>')])
        shape = list(prs.slides[0].shapes)[0]
        self.assertEqual(shape.fill.fore_color.rgb, RGBColor(0xFF, 0xFF, 0xFF))

    def test_unknown_named_color_still_fails_loud(self):
        """⚠ 只放行實際用到的關鍵字，不引進整套 CSS 顏色表。"""
        with self.assertRaises(self.mod.UnsupportedElement):
            self.mod.build([_svg('<rect x="0" y="0" width="10" height="10" '
                                 'fill="rebeccapurple"/>')])

    def test_fill_opacity_applied(self):
        """`fill-opacity` ×22。忽略它會讓疊圖變成不透明，遮住底下的東西。

        ⚠ python-pptx 1.0.2 的 `FillFormat` **沒有** `transparency` 屬性
        （2026-08-14 實測），alpha 要自己寫進 DrawingML：`<a:alpha val="35000"/>`
        掛在 `<a:srgbClr>` 底下。故這裡直接驗 XML。
        """
        prs = self.mod.build([_svg('<rect x="0" y="0" width="10" height="10" '
                                   'fill="#93C5FD" fill-opacity="0.35"/>')])
        shape = list(prs.slides[0].shapes)[0]
        xml = shape.fill.fore_color._xFill.xml
        self.assertIn('<a:alpha val="35000"/>', xml)

    def test_text_anchor_middle_centers_box(self):
        """🔴 `text-anchor` 198 次（middle 152／end 46）。不做則全部偏移。"""
        prs = self.mod.build([_svg('<text x="100" y="50" font-size="16" '
                                   'text-anchor="middle" fill="#00094A">ABC</text>')])
        box = list(prs.slides[0].shapes)[0]
        # 錨點 x=100 是**中心**：框左緣必須在 100 左邊
        self.assertLess(box.left, _emu_px(100))
        self.assertAlmostEqual(box.left + box.width / 2, _emu_px(100), delta=_emu_px(1))

    def test_text_anchor_end_right_aligns(self):
        prs = self.mod.build([_svg('<text x="200" y="50" font-size="16" '
                                   'text-anchor="end" fill="#00094A">ABC</text>')])
        box = list(prs.slides[0].shapes)[0]
        self.assertAlmostEqual(box.left + box.width, _emu_px(200), delta=_emu_px(1))

    def test_text_rotate_transform(self):
        """`transform="rotate(-90 26 230)"` ×2（縱向軸標）。"""
        prs = self.mod.build([_svg('<text x="26" y="230" font-size="14" '
                                   'transform="rotate(-90 26 230)" '
                                   'fill="#00094A">件數</text>')])
        box = list(prs.slides[0].shapes)[0]
        self.assertAlmostEqual(box.rotation, 270.0, places=1)

    def test_stroke_dasharray_makes_dashed_line(self):
        """`stroke-dasharray="6 4"` ×2（參考線）。"""
        from pptx.enum.dml import MSO_LINE_DASH_STYLE

        prs = self.mod.build([_svg('<line x1="0" y1="50" x2="200" y2="50" '
                                   'stroke="#869FB2" stroke-width="1" '
                                   'stroke-dasharray="6 4"/>')])
        shape = list(prs.slides[0].shapes)[0]
        self.assertEqual(shape.line.dash_style, MSO_LINE_DASH_STYLE.DASH)


class HatchPatternTests(unittest.TestCase):
    """🔴 斜紋＝已轉讓，是獨立的視覺通道，不可退化成純色。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    HATCH_DEFS = ('<defs><pattern id="hatch" width="6" height="6" '
                  'patternUnits="userSpaceOnUse" patternTransform="rotate(45)">'
                  '<line x1="0" y1="0" x2="0" y2="6" stroke="#00094A" '
                  'stroke-width="2" stroke-opacity="0.55"/></pattern></defs>')

    def test_hatch_fill_becomes_patterned_not_solid(self):
        prs = self.mod.build([_svg(self.HATCH_DEFS
                                   + '<rect x="10" y="10" width="40" height="20" '
                                     'fill="url(#hatch)"/>')])
        shape = list(prs.slides[0].shapes)[0]
        # 必須是圖樣填滿，不是純色
        self.assertEqual(shape.fill.type, 2, "hatch 退化成純色會讓「已轉讓」消失")

    def test_unknown_url_reference_fails_loud(self):
        with self.assertRaises(self.mod.UnsupportedElement):
            self.mod.build([_svg('<rect x="0" y="0" width="10" height="10" '
                                 'fill="url(#nope)"/>')])


class IgnoredMarkupTests(unittest.TestCase):
    """互動用標記要**明確忽略**並有測試守住，不得靜默略過。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    def test_title_child_is_tooltip_not_drawn(self):
        """`<title>` 65 個掛在 circle/rect 底下，是 tooltip。"""
        prs = self.mod.build([_svg('<circle cx="50" cy="50" r="8" fill="#006DF5">'
                                   '<title>台達電 12 件</title></circle>')])
        self.assertEqual(len(list(prs.slides[0].shapes)), 1)

    def test_data_attributes_ignored(self):
        prs = self.mod.build([_svg('<circle cx="50" cy="50" r="8" fill="#006DF5" '
                                   'data-value-band="高" pointer-events="none"/>')])
        self.assertEqual(len(list(prs.slides[0].shapes)), 1)

    def test_style_element_supplies_inherited_font(self):
        """⚠ `font-family` 是繼承來的——`<style>` 或 `<svg>` 根元素。
        現行 `_add_text` 只讀元素自己的屬性，會拿不到而靜默用 pptx 預設字型。"""
        body = ("<style>text{font-family:'Noto Sans TC','Segoe UI',sans-serif}</style>"
                '<text x="10" y="20" font-size="16" fill="#00094A">中文</text>')
        prs = self.mod.build([_svg(body)])
        box = list(prs.slides[0].shapes)[0]
        run = box.text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.name, "Noto Sans TC")

    def test_root_font_family_inherited(self):
        prs = self.mod.build([_svg('<text x="10" y="20" font-size="16" '
                                   'fill="#00094A">中文</text>',
                                   extra=' font-family="Noto Sans TC, sans-serif"')])
        run = list(prs.slides[0].shapes)[0].text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.name, "Noto Sans TC")


class RealChartIntegrationTests(unittest.TestCase):
    """直接吃 14 張真實圖表 SVG——合成測試漏掉的形式由它兜底。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    def _real_svgs(self):
        if not REAL_CHARTS.is_dir():
            self.skipTest(f"真實圖表目錄不存在：{REAL_CHARTS}")
        files = [p for p in sorted(REAL_CHARTS.glob("*.svg"))
                 if not p.name.endswith(".web.svg")]
        if not files:
            self.skipTest("目錄內沒有圖表 SVG")
        return files

    def test_every_real_chart_converts(self):
        """🔴 A 層元素對帳：每張都要轉得出來，且 shape 數與可繪元素數相符。"""
        import xml.etree.ElementTree as ET

        drawable = {"rect", "text", "circle", "line", "polyline"}
        for path in self._real_svgs():
            with self.subTest(chart=path.name):
                svg = path.read_text(encoding="utf-8")
                prs = self.mod.build([svg], [path.parent])
                root = ET.fromstring(svg)
                expected = sum(1 for el in root.iter()
                               if el.tag.split("}")[-1] in drawable
                               and not self._inside_defs(root, el))
                self.assertEqual(len(list(prs.slides[0].shapes)), expected)

    @staticmethod
    def _inside_defs(root, target) -> bool:
        """`<defs>` 底下的元素是樣式定義，不直接畫。"""
        for defs in root.iter():
            if defs.tag.split("}")[-1] == "defs":
                if any(el is target for el in defs.iter()):
                    return True
        return False


if __name__ == "__main__":
    unittest.main()
