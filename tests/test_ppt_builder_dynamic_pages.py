"""PPT builder 動態頁與版型覆寫契約（v3，2026-07-30 更新）。

守兩件事：
- 基礎大綱沒列到、但本次有資料的報表，要自動插在**第一個附錄頁之前**（靠
  `is_appendix` 顯式旗標找錨點，不用頁碼魔術數字）。
- 使用者從 `layout_overrides` 挑的版型要真的套用，且頁碼重新連號後仍對得上。

`position_overrides`（拖曳座標）已於 v3 移除：座標唯一來源是 `theme.json`，
再開一條「每頁各自存一份座標」的路等於把座標分岔重新引回來。上游 runner／前端
仍可送這個 key，builder 直接忽略，不會壞。
"""

from __future__ import annotations

import importlib.util
import json
import sys
from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation

SKILL_DIR = Path(__file__).resolve().parents[1] / "skills" / "patent-report-ppt"
BUILDER_PATH = SKILL_DIR / "scripts" / "build_ppt.py"

VERSION = "report_trial_20260723_000000"
MINIMAL_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="400" height="300">'
    '<rect width="400" height="300" fill="#516CEE"/></svg>'
)


def _load_builder():
    """載入 skill 內的 build_ppt.py。"""
    spec = importlib.util.spec_from_file_location("build_ppt_dynamic", BUILDER_PATH)
    module = importlib.util.module_from_spec(spec)
    sys.modules["build_ppt_dynamic"] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture()
def report_dir(tmp_path: Path) -> Path:
    """基礎大綱有的報表 ＋ 一個不在大綱內的額外報表（owner_year_matrix）。"""
    path = tmp_path / "output" / VERSION
    path.mkdir(parents=True)
    data = {
        "parameters": {"version": VERSION},
        "reports": {
            "application_trend": {
                "label_zh": "申請趨勢",
                "report_type": "aggregate",
                "rows": [{"application_year": 2025, "patent_count": 3}],
            },
            "country_distribution": {
                "label_zh": "地域分布",
                "report_type": "aggregate",
                "rows": [{"country_code": "TW", "patent_count": 3}],
            },
            "cluster_topic_table": {
                "label_zh": "全分類技術指標總表",
                "report_type": "table",
                "rows": [{"topic": "收繩機構", "patent_count": 3}],
            },
            "owner_year_matrix": {
                "label_zh": "專利權人年份矩陣",
                "report_type": "aggregate",
                "rows": [{"current_assignee_display_name": "Rexon", "application_year": 2025, "patent_count": 3}],
            },
        },
        "family_reports": {},
    }
    (path / "report_data.json").write_text(json.dumps(data, ensure_ascii=False), encoding="utf-8")
    (path / "narratives.json").write_text(json.dumps({"reports": {}}, ensure_ascii=False), encoding="utf-8")
    (path / "artifact_manifest.json").write_text(
        json.dumps(
            {
                "metadata": {"version": VERSION},
                "artifacts": [
                    {"file": "annual_trend.svg", "artifact_type": "chart_svg",
                     "report_names": ["application_trend"], "sha256": "0" * 64},
                    {"file": "owner_year_matrix.svg", "artifact_type": "chart_svg",
                     "report_names": ["owner_year_matrix"], "sha256": "0" * 64},
                ],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    for name in ("annual_trend.svg", "owner_year_matrix.svg"):
        (path / name).write_text(MINIMAL_SVG, encoding="utf-8")
    return path


@pytest.fixture()
def approvals(tmp_path: Path) -> Path:
    path = tmp_path / "approvals.json"
    path.write_text(
        json.dumps(
            {"report_version": VERSION, "slots": {"cover.title": "手動標題"}},
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    return path


def test_build_ppt_expands_extra_report_pages_before_appendix_anchor(report_dir, approvals, tmp_path):
    """額外報表要插在第一個 is_appendix 頁之前，不用頁碼魔術數字。"""
    builder = _load_builder()
    result = builder.build_ppt(
        report_dir=report_dir, approvals_path=approvals, output_dir=tmp_path / "ppt"
    )
    pages = result["manifest"]["pages"]

    owner_index = next(i for i, page in enumerate(pages) if page["report_keys"] == ["owner_year_matrix"])
    first_appendix_index = next(i for i, page in enumerate(pages) if page["is_appendix"])

    assert owner_index < first_appendix_index
    assert pages[first_appendix_index]["report_keys"] == ["cluster_topic_table"]


def test_build_ppt_applies_layout_overrides(report_dir, approvals, tmp_path):
    """使用者挑的版型要真的套用，且套用後的頁面數與頁碼一致。"""
    builder = _load_builder()
    base = builder.build_ppt(
        report_dir=report_dir, approvals_path=approvals, output_dir=tmp_path / "ppt0"
    )
    target = next(
        page for page in base["manifest"]["pages"] if page["report_keys"] == ["owner_year_matrix"]
    )
    assert target["kind"] != "table_with_points"

    approvals.write_text(
        json.dumps(
            {
                "report_version": VERSION,
                "slots": {"cover.title": "手動標題"},
                "layout_overrides": {str(target["page"]): "table_with_points"},
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    result = builder.build_ppt(
        report_dir=report_dir, approvals_path=approvals, output_dir=tmp_path / "ppt1"
    )
    pages = result["manifest"]["pages"]
    changed = next(page for page in pages if page["report_keys"] == ["owner_year_matrix"])

    assert changed["kind"] == "table_with_points"
    assert [page["page"] for page in pages] == list(range(1, len(pages) + 1))
    tables = [shape for shape in Presentation(result["pptx_path"]).slides[changed["page"] - 1].shapes if shape.has_table]
    assert tables


def test_invalid_layout_override_is_ignored_not_fatal(report_dir, approvals, tmp_path):
    """無效版型名稱只忽略，不讓整份 PPT 產不出來。"""
    builder = _load_builder()
    approvals.write_text(
        json.dumps(
            {
                "report_version": VERSION,
                "slots": {},
                "layout_overrides": {"3": "not_a_real_layout"},
                "position_overrides": {"3.table": {"left_in": 1.25}},
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    result = builder.build_ppt(
        report_dir=report_dir, approvals_path=approvals, output_dir=tmp_path / "ppt"
    )
    assert Path(result["pptx_path"]).exists()
    assert result["manifest"]["pages"]
