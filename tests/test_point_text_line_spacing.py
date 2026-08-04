"""要點文字 16pt、行距 1.65，而且行距要**真的寫進 PPTX**（2026-08-04 使用者定案）。

## 為什麼要這支

`theme.qa.line_height_ratio`（原 1.35）**只是估算用**——程式從來沒有設定過
`paragraph.line_spacing`，PowerPoint 實際用的是預設行距。

⚠ 所以「把 line_height_ratio 調大」不會讓畫面變寬，只會讓容量估得更保守：
字更少、版面更空，看起來像倒退。要拓寬行距必須**同時**寫進段落屬性。

## 定案

| 項目 | 值 |
|---|---|
| 要點正文字級 | **16pt**（原 15） |
| 要點行距 | **1.65**（原估算 1.35、實際是 PowerPoint 預設） |

⚠ 行距只套用在**要點文字**上——標題、表格、封面沒有設段落行距，
把 `line_height_ratio` 全域改大只會讓它們的容量估算失準（估得下卻寫更少）。
"""
from __future__ import annotations

import sys
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

_SKILL = Path(__file__).resolve().parents[1] / "skills" / "patent-report-ppt" / "scripts"
sys.path.insert(0, str(_SKILL))

import build_ppt as bp  # noqa: E402


class PointTextSpacingTests(unittest.TestCase):
    def test_theme_point_text_is_16pt(self):
        self.assertEqual(bp.Theme.load().size("point_text_pt"), 16.0)

    def test_theme_declares_point_line_spacing(self):
        """行距值要有**單一定義處**，估算與渲染都讀它。"""
        theme = bp.Theme.load()
        self.assertEqual(float(theme.qa["point_line_height_ratio"]), 1.65)

    def test_line_spacing_is_written_into_the_pptx(self):
        """⚠ 只改估算值不算數——段落必須真的帶 line_spacing。"""
        theme = bp.Theme.load()
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        blocks = [("現況", "本頁共 8 件", "ink", False),
                  ("意涵", "技術集中單一大類", "ink", False)]
        bp._add_number_bold_text(slide, theme, blocks,
                                 left=1.0, top=1.0, width=3.0, height=2.0,
                                 size=theme.size("point_text_pt"))
        frame = next(sh.text_frame for sh in slide.shapes if sh.has_text_frame)
        spacings = [p.line_spacing for p in frame.paragraphs if p.text.strip()]
        self.assertTrue(spacings, "沒抓到要點段落")
        for value in spacings:
            self.assertEqual(value, 1.65, f"段落行距是 {value}，沒寫進 PPTX")

    def test_capacity_uses_the_point_line_ratio(self):
        """容量估算要用要點專屬行距，不是全域那個。"""
        theme = bp.Theme.load()
        size = theme.size("point_text_pt")
        width_in, height_in, _ = bp._points_area(theme, "chart_hero")
        _, lines = bp._text_capacity(theme, width_in=width_in, height_in=height_in,
                                     size_pt=size, line_ratio=float(theme.qa["point_line_height_ratio"]))
        expected = int(height_in / (size / 72.0 * 1.65) + 1e-6)
        self.assertEqual(lines, expected, "容量估算沒吃要點行距")


if __name__ == "__main__":
    unittest.main()
