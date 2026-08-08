"""SlidePlan 的 narrative 必須真的畫進投影片（2026-08-09 實機驗收批）。

## 為什麼要這支測試

`APPROVED_LAYOUT_PRESETS` 與 `RENDERERS` 已有雙向一致性測試，`exec_summary`／
`walls_gaps`／`reading_guide` 三個版型都「有 renderer、名稱也對得上」，測試
全綠——**但實機轉圖出來是三張一模一樣的空框**。

根因：那三支當時轉呼叫 `_render_direction`，而它的內容來自固定 slot
`direction.body`（研發方向頁專用），讀不到 SlidePlan 的 narrative。

⚠ 教訓：一致性測試驗得到「有沒有 renderer」，驗不到「renderer 畫了什麼」。
內容要有內容檢查——所以這裡直接讀回產出的 pptx 檢查文字。
"""
from __future__ import annotations

import importlib.util
import json
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

SKILL_DIR = Path(__file__).resolve().parents[1] / "skills" / "patent-report-ppt"


def _load_builder():
    spec = importlib.util.spec_from_file_location(
        "build_ppt_plan_narrative", SKILL_DIR / "scripts" / "build_ppt.py")
    module = importlib.util.module_from_spec(spec)
    sys.modules.setdefault("build_ppt_plan_narrative", module)
    spec.loader.exec_module(module)
    return module


bp = _load_builder()

# 每個無圖版型放一句可辨識的要點，產完回頭找它在不在頁面上。
_MARKERS = {
    "exec_summary": "結論標記甲乙丙",
    "walls_gaps": "牆與空白標記甲乙丙",
    "reading_guide": "判讀說明標記甲乙丙",
}


def _report_data() -> dict:
    slides = [{"slide_id": "s1", "layout_preset": "cover", "purpose": "封面",
               "chart_identities": [], "narrative": []}]
    for index, (preset, marker) in enumerate(_MARKERS.items(), start=2):
        slides.append({
            "slide_id": f"s{index}", "layout_preset": preset,
            "purpose": f"{preset} 這頁要回答什麼", "chart_identities": [],
            "narrative": [{"text": marker, "evidence_ref": "e1"}],
        })
    return {
        "version": "plan_narrative_case",
        "period": "2011–2026",
        "slide_plan": {"plan_id": "plan-test", "slides": slides},
        "chart_rows": {},
        "sections": [],
    }


class PlanNarrativeReachesSlidesTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls._tmp = tempfile.TemporaryDirectory()
        root = Path(cls._tmp.name)
        report_dir = root / "report"
        report_dir.mkdir()
        (report_dir / "report_data.json").write_text(
            json.dumps(_report_data(), ensure_ascii=False), encoding="utf-8")
        (report_dir / "artifact_manifest.json").write_text(
            json.dumps({"artifacts": []}, ensure_ascii=False), encoding="utf-8")
        result = bp.build_ppt(report_dir=str(report_dir), output_dir=str(root / "out"))
        cls.deck = Presentation(result["pptx_path"])

    @classmethod
    def tearDownClass(cls):
        cls._tmp.cleanup()

    def _page_text(self, page: int) -> str:
        parts = []
        for shape in self.deck.slides[page - 1].shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        return "\n".join(parts)

    def test_each_textual_preset_renders_its_narrative(self):
        """三個無圖版型都要印出自己的要點——空框視為未實作。"""
        for index, (preset, marker) in enumerate(_MARKERS.items(), start=2):
            with self.subTest(preset=preset):
                self.assertIn(marker, self._page_text(index),
                              f"{preset} 沒有畫出 SlidePlan 的 narrative（空框）")

    def test_panel_title_does_not_repeat_page_title(self):
        """面板標題不得等於頁標題——首版實測一頁印了兩次同一句。"""
        for index, preset in enumerate(_MARKERS, start=2):
            with self.subTest(preset=preset):
                text = self._page_text(index)
                self.assertEqual(text.count(f"{preset} 這頁要回答什麼"), 1,
                                 "頁標題重複出現在面板標題")


if __name__ == "__main__":
    unittest.main()
