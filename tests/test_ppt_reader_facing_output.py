"""簡報給讀者看的內容不得有程式殘骸（2026-07-31 實機驗收批 1）。

實測產出 `report_trial_20260731_060837.pptx`（19 頁）抓到三類缺陷：

| # | 缺陷 | 實測位置 |
|---|---|---|
| B-1 | 物件序列化外洩：`{'name': '祺驊', ...` 直接印在表格裡 | 第 9、10、18 頁（附錄 1 整張表 11 列全毀）|
| B-2 | 內部英文欄名外洩：`recent_assignee_display_nam…` 當表頭 | 第 19 頁（附錄 2）|
| B-3 | 判讀限制警語被切在句中：「…多為新案審…」 | 第 3、5 頁 |

## 為什麼網頁報表頁沒事、PPT 有事

引擎 `chart_runner._humanize_cell()` 會把 `list[dict]` 轉成「名稱 數字」，
`_data_table_html` 出的網頁表格因此乾淨；**PPT 端沒有等價處理**，`_add_table`
碰到 list 一律「、」串接，串到 dict 就印出 Python repr。

⚠ 修法**不是**在 `build_ppt` 抄第二份 `_humanize_cell`——本 repo 的欄名對照表已經
因為「同一資訊兩處落點」各自漂移過一輪。呈現規則與欄名同源：由引擎寫進
`report_data.json["table_display"]`，PPT 端讀它（skill 會被 Installer 打包到
使用者電腦，不得 import backend，故只能走檔案傳遞）。

⚠ B-2 的實際原因與原先推測不同：`recent_assignee_display_names`（複數，
`string_agg_distinct_nonblank_excl_group` 的輸出別名）**根本沒有登記中文欄名**，
`DATA_COLUMN_LABELS` 只有單數的 `recent_assignee_display_name`。不是排除規則沒掛到
第二個 report_key——它本來就不該被排除（使用者定案「後面欄都列出公司了」指的正是這欄）。
"""
from __future__ import annotations

import importlib.util
import sys
import unittest
from pathlib import Path

from pptx import Presentation

SKILL_DIR = Path(__file__).resolve().parents[1] / "skills" / "patent-report-ppt"


def _load_builder():
    spec = importlib.util.spec_from_file_location(
        "build_ppt_reader_facing", SKILL_DIR / "scripts" / "build_ppt.py")
    module = importlib.util.module_from_spec(spec)
    sys.modules.setdefault("build_ppt_reader_facing", module)
    spec.loader.exec_module(module)
    return module


bp = _load_builder()

# 實機 rows 的真實形狀（取自 report_trial_20260731_060837 的 cluster_topic_table）。
TOPIC_ROWS = [
    {
        "topic_code": "T001",
        "label": "拉繩回收",
        "source_field": "wips_independent_claims",
        "patent_count": 5,
        "top_applicants": [{"name": "祺驊", "count": 3},
                           {"name": "Brett Unsworth", "count": 1}],
        "leading_applicants_involved": ["祺驊", "Brett Unsworth"],
    },
    {
        "topic_code": "T002",
        "label": "阻力控制",
        "source_field": "effect_summary",
        "patent_count": 3,
        "top_applicants": [{"name": "廈門帝瑪", "count": 2}],
        "leading_applicants_involved": [],
    },
]


class EngineEmitsCellDisplayTextTests(unittest.TestCase):
    """B-1 引擎端：`table_display` 要一併輸出「值怎麼呈現」，不只欄名與排除欄。"""

    def setUp(self):
        from backend.app.reports import chart_runner

        self.chart_runner = chart_runner

    def test_table_display_spec_carries_display_rows(self):
        """含 list[dict] 的報表要有人類化後的字串，且不得留下 Python repr。"""
        spec = self.chart_runner.table_display_spec(
            {"cluster_topic_table": {"rows": TOPIC_ROWS}})
        rows = (spec.get("display_rows") or {}).get("cluster_topic_table")
        self.assertIsNotNone(rows, "table_display 未輸出 display_rows，PPT 端無從得知呈現規則")
        self.assertNotIn("{'", rows[0]["top_applicants"], "仍是物件序列化結果")
        self.assertEqual(rows[0]["top_applicants"], "祺驊 3；Brett Unsworth 1")
        self.assertEqual(rows[0]["leading_applicants_involved"], "祺驊、Brett Unsworth")

    def test_display_text_matches_web_table(self):
        """同一份呈現規則：display_rows 的值必須等於網頁表格用的 `_humanize_cell`。"""
        spec = self.chart_runner.table_display_spec(
            {"cluster_topic_table": {"rows": TOPIC_ROWS}})
        rows = spec["display_rows"]["cluster_topic_table"]
        for index, row in enumerate(TOPIC_ROWS):
            for column, value in row.items():
                self.assertEqual(
                    rows[index][column], self.chart_runner._humanize_cell(value),
                    f"{column} 的呈現與網頁報表頁不一致（兩處落點又漂移）")

    def test_scalar_only_reports_omitted(self):
        """純量報表不必重複一份（report_data.json 已是 124KB 等級，不做無謂膨脹）。"""
        spec = self.chart_runner.table_display_spec(
            {"application_trend": {"rows": [{"application_year": 2022, "patent_count": 15}]}})
        self.assertNotIn("application_trend", spec.get("display_rows") or {})

    def test_existing_keys_unchanged(self):
        """沿用同一區塊擴充：欄名對照與排除欄仍在原鍵，PPT 既有讀法不變。"""
        spec = self.chart_runner.table_display_spec({})
        self.assertIn("column_labels", spec)
        self.assertIn("excluded_columns", spec)


class PptUsesEngineDisplayTextTests(unittest.TestCase):
    """B-1 PPT 端：讀引擎那份，不自建第二套值轉換。"""

    @staticmethod
    def _ctx(display_rows=None):
        report_data = {
            "reports": {"cluster_topic_table": {"rows": TOPIC_ROWS, "label_zh": "主題分類統計表"}},
            "table_display": {"column_labels": {}, "excluded_columns": {}},
        }
        if display_rows is not None:
            report_data["table_display"]["display_rows"] = display_rows
        return {"report_data": report_data}

    def _spec(self, **kwargs):
        return bp.PageSpec(page=9, kind="table", title="t", topic="t",
                           report_keys=("cluster_topic_table",), **kwargs)

    def test_first_rows_prefers_engine_display_rows(self):
        from backend.app.reports.chart_runner import table_display_spec

        display = table_display_spec({"cluster_topic_table": {"rows": TOPIC_ROWS}})["display_rows"]
        rows = bp._first_rows(self._spec(), self._ctx(display))
        self.assertEqual(rows[0]["top_applicants"], "祺驊 3；Brett Unsworth 1")

    def test_row_filter_still_applies_to_display_rows(self):
        """依通道拆頁的 row_filter 不得因為改讀 display_rows 而失效。"""
        from backend.app.reports.chart_runner import table_display_spec

        display = table_display_spec({"cluster_topic_table": {"rows": TOPIC_ROWS}})["display_rows"]
        spec = self._spec(row_filter={"source_field": "effect_summary"})
        rows = bp._first_rows(spec, self._ctx(display))
        self.assertEqual([r["label"] for r in rows], ["阻力控制"])

    def test_rendered_table_has_no_object_repr(self):
        """畫進投影片的儲存格文字不得出現 `{'`。"""
        from backend.app.reports.chart_runner import table_display_spec

        display = table_display_spec({"cluster_topic_table": {"rows": TOPIC_ROWS}})["display_rows"]
        rows = bp._first_rows(self._spec(), self._ctx(display))
        texts = _render_table_cells(rows)
        self.assertFalse([t for t in texts if "{'" in t],
                         f"儲存格仍印出物件序列化結果：{texts}")
        self.assertTrue(any("祺驊 3" in t for t in texts),
                        f"前三大申請人未以可讀形式呈現：{texts}")


class TableColumnLabelCoverageTests(unittest.TestCase):
    """B-2：PPT 會排成表格的報表，其輸出欄一律要有中文欄名——規則寫成檢查。

    ⚠ 只驗「PPT 真的會印成表格」的報表：那才是讀者看得到英文欄名的地方。
    """

    def test_no_internal_column_name_reaches_a_table_page(self):
        from backend.app.reports.chart_runner import (
            DATA_COLUMN_LABELS,
            DATA_TABLE_EXCLUDED_COLUMNS,
        )
        from backend.app.reports.report_definitions import REPORT_DEFINITIONS

        table_keys = {
            key
            for spec in bp.PAGE_LAYOUT if spec.kind in ("table", "table_with_points")
            for key in spec.report_keys
        }
        self.assertIn("applicant_ranking", table_keys, "附錄 2 掛的報表變了，本測試需同步")
        missing: list[str] = []
        for key in sorted(table_keys):
            definition = REPORT_DEFINITIONS.get(key)
            if definition is None:      # 分群類報表由 cluster_analytics 產，不在定義表
                continue
            excluded = DATA_TABLE_EXCLUDED_COLUMNS.get(key, ())
            columns = list(definition.columns or ())
            columns += [alias for _, _, alias in (definition.aggregates or ())]
            missing += [f"{key}.{c}" for c in columns
                        if c not in excluded and c not in DATA_COLUMN_LABELS]
        self.assertEqual(missing, [], f"這些欄會以內部英文欄名印在簡報上：{missing}")


# 🔴 2026-08-04：原本這裡有 CaveatNotTruncatedTests——驗「判讀限制」的容量與截斷。
# 使用者定案把判讀限制整個移除（「判讀限制不要出現了，作用不大」），
# 規格沒了測試就失去存在理由，故整個類別刪除而非改寫。
# ⚠ 改寫成「驗別的事」會留下一個守不住任何意圖的空殼測試。

def _render_table_cells(rows) -> list[str]:
    """把 rows 畫進一張空白投影片，回傳所有儲存格文字（讀者真正看到的字）。"""
    theme = bp.Theme.load()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    g = theme.geometry["table"]
    bp._add_table(slide, theme, rows,
                  left=g["left_in"], top=g["top_in"], width=g["width_in"],
                  height=g["height_in"], row_height=g["row_height_in"],
                  max_columns=int(g["max_columns"]), cell_margin_in=g["cell_margin_in"],
                  cell_inset_in=g["cell_inset_in"], labels={}, excluded=set())
    table = next(shape for shape in slide.shapes if shape.has_table).table
    return [cell.text for row in table.rows for cell in row.cells]


if __name__ == "__main__":
    unittest.main()
