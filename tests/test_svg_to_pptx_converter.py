"""窄 SVG→DrawingML 轉換器的契約（add-deck-delivery-line tasks 2.1）。

## 為什麼要它

B 案把排版決定權從 PowerPoint 收回引擎：每頁先組 SVG（文字由引擎**逐行斷好**、
絕對定位），再轉成原生 PPTX。轉換器**窄**是關鍵——只支援本 skill 的元素詞彙，
不支援 SVG 標準全集。詞彙外的東西一律 fail loud，不靜默略過。

## 詞彙＝現行組版實際畫出來的東西（2026-08-13 自 `deck_layout.py` 反解）

它只用三個 python-pptx 原生元素：

| SVG | pptx | 來源 |
|---|---|---|
| `<rect>` | `add_shape(RECTANGLE / ROUNDED_RECTANGLE)` | `deck_layout.rect()` |
| `<text>` | `add_textbox()` | `deck_layout.textbox()` |
| `<image>` | `add_picture()` | `chart_stack()` |

⚠ **「線」不是獨立元素**，是矩形的退化形（`RULE_W = 0.014in` 的細 rect）。
不要為它另立詞彙——那會讓同一件事有兩個表示法。

## 座標系＝96 dpi

`deck_layout.py:73` 的註記寫死了這個前提：「0.008in 在 **96dpi** 下只有 0.77px，
次像素會讓部分細線整條消失，而且是隨機幾條」。故 SVG 以 96 px/in 為準，
13.333×7.5 in 的投影片 ＝ 1280×720 的 SVG。
"""
from __future__ import annotations

import importlib.util
import sys
import unittest
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = PROJECT_ROOT / "skills" / "html-report-to-deck" / "scripts"

DPI = 96
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5


def _load():
    path = SCRIPTS / "svg_to_pptx.py"
    spec = importlib.util.spec_from_file_location("svg_to_pptx", path)
    module = importlib.util.module_from_spec(spec)
    sys.modules["svg_to_pptx"] = module
    spec.loader.exec_module(module)
    return module


def _svg(body: str, w: int = 1280, h: int = 720) -> str:
    return (f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'xmlns:xlink="http://www.w3.org/1999/xlink" '
            f'width="{w}" height="{h}" viewBox="0 0 {w} {h}">{body}</svg>')


class VocabularyTests(unittest.TestCase):
    """詞彙內的三種元素要正確映射；詞彙外要 fail loud。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    def test_rect_becomes_autoshape(self):
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = self.mod.build([_svg('<rect x="96" y="48" width="192" height="96" '
                                   'fill="#DAE4F2"/>')])
        shapes = list(prs.slides[0].shapes)
        self.assertEqual(len(shapes), 1)
        self.assertEqual(shapes[0].shape_type, MSO_SHAPE_TYPE.AUTO_SHAPE)

    def test_text_becomes_textbox(self):
        prs = self.mod.build([_svg('<text x="96" y="60" font-size="21.33">結論句</text>')])
        shapes = list(prs.slides[0].shapes)
        self.assertEqual(len(shapes), 1)
        self.assertTrue(shapes[0].has_text_frame)
        self.assertEqual(shapes[0].text_frame.text, "結論句")

    def test_image_becomes_picture(self):
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        png = SCRIPTS.parent / "regression_baseline" / "slide01.png"
        self.assertTrue(png.is_file(), "測試素材缺失")
        prs = self.mod.build([_svg(
            f'<image x="96" y="96" width="480" height="270" '
            f'xlink:href="{png.as_uri()}"/>')])
        shapes = list(prs.slides[0].shapes)
        self.assertEqual(len(shapes), 1)
        self.assertEqual(shapes[0].shape_type, MSO_SHAPE_TYPE.PICTURE)

    def test_relative_image_path_resolves_against_svg_dir(self):
        """🔴 相對路徑要相對於 **SVG 檔所在目錄**，不是當前工作目錄。

        ⚠ 2026-08-13 映射煙霧測試發現：SVG 用 `file://` 絕對 URI 引圖時，
        Chromium 以 `set_content` 載入會因缺 base URL 而**破圖**（COM 轉圖正常）。
        正解是 SVG 存檔、圖用相對路徑、Chromium 用 `goto` 載入——那條路徑要成立，
        轉換器就必須以 SVG 檔的位置解析相對路徑。
        以 cwd 解析會在「runner 從別的目錄呼叫」時靜默找不到圖。
        """
        import shutil
        import tempfile

        src = SCRIPTS.parent / "regression_baseline" / "slide01.png"
        with tempfile.TemporaryDirectory() as tmp:
            work = Path(tmp)
            shutil.copyfile(src, work / "chart.png")
            svg_path = work / "page01.svg"
            svg_path.write_text(_svg('<image x="0" y="0" width="480" height="270" '
                                     'xlink:href="chart.png"/>'), encoding="utf-8")
            prs = self.mod.build_file([svg_path], work / "out.pptx")
        self.assertEqual(prs, 1)

    def test_unknown_element_fails_loud(self):
        """⚠ 詞彙外**不得靜默略過**——那會讓版面少一塊而沒人發現。"""
        for body in ('<circle cx="10" cy="10" r="5"/>',
                     '<path d="M0 0 L10 10"/>',
                     '<polygon points="0,0 10,0 5,10"/>'):
            with self.subTest(body=body):
                with self.assertRaises(self.mod.UnsupportedElement):
                    self.mod.build([_svg(body)])


class GeometryTests(unittest.TestCase):
    """座標換算：SVG px @96dpi → EMU。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    def test_slide_size_matches_16x9(self):
        from pptx.util import Inches

        prs = self.mod.build([_svg("")])
        self.assertEqual(prs.slide_width, Inches(SLIDE_W_IN))
        self.assertEqual(prs.slide_height, Inches(SLIDE_H_IN))

    def test_px_maps_to_inches_at_96dpi(self):
        from pptx.util import Inches

        prs = self.mod.build([_svg('<rect x="96" y="192" width="288" height="48"/>')])
        shape = list(prs.slides[0].shapes)[0]
        self.assertEqual(shape.left, Inches(1.0))
        self.assertEqual(shape.top, Inches(2.0))
        self.assertEqual(shape.width, Inches(3.0))
        self.assertEqual(shape.height, Inches(0.5))

    def test_hairline_rect_survives(self):
        """⚠ `deck_layout.py:73`：0.008in（96dpi 下 0.77px）會讓細線**隨機整條消失**。

        轉換器不得把細矩形四捨五入成 0 高——那正是「目視誤判成本來就沒分隔」的成因。
        `RULE_W = 0.014in` ＝ 1.344px。
        """
        prs = self.mod.build([_svg('<rect x="96" y="96" width="288" height="1.344"/>')])
        shape = list(prs.slides[0].shapes)[0]
        self.assertGreater(shape.height, 0, "細線被壓成 0 高＝畫不出來")


class TextFidelityTests(unittest.TestCase):
    """B 案的核心：文字逐行定位、PowerPoint 零重排自由。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    def test_word_wrap_disabled(self):
        """🔴 關 wrap 是 B 案成立的前提——開著就等於把決定權還給 PowerPoint。"""
        prs = self.mod.build([_svg('<text x="96" y="60" font-size="21.33">很長的一句話</text>')])
        shape = list(prs.slides[0].shapes)[0]
        self.assertIs(shape.text_frame.word_wrap, False)

    def test_autosize_disabled(self):
        """字級鎖死：不得讓 PowerPoint 自動縮放（`deck_layout` 早已定案 24／16）。"""
        from pptx.enum.text import MSO_AUTO_SIZE

        prs = self.mod.build([_svg('<text x="96" y="60" font-size="21.33">內文</text>')])
        tf = list(prs.slides[0].shapes)[0].text_frame
        self.assertIn(tf.auto_size, (None, MSO_AUTO_SIZE.NONE))

    def test_each_text_element_is_its_own_line(self):
        """引擎已逐行斷好：三個 `<text>` ＝ 三個獨立文字框，不是一段三行。

        ⚠ 合成一個多行段落就等於讓 PowerPoint 重新決定斷點——B 案就白做了。
        """
        body = "".join(f'<text x="96" y="{60 + i * 30}" font-size="21.33">第{i}行</text>'
                       for i in range(3))
        prs = self.mod.build([_svg(body)])
        self.assertEqual(len(list(prs.slides[0].shapes)), 3)

    def test_font_size_px_to_pt(self):
        """SVG px @96dpi → pt（1pt = 1/72 in）。16pt 內文 ＝ 21.33px。"""
        from pptx.util import Pt

        prs = self.mod.build([_svg('<text x="96" y="60" font-size="21.33">內文</text>')])
        run = list(prs.slides[0].shapes)[0].text_frame.paragraphs[0].runs[0]
        self.assertAlmostEqual(run.font.size.pt, Pt(16).pt, delta=0.1)

    def test_font_family_preserved(self):
        """字型必須帶過去——缺了會 fallback 且不報錯（design 4-0b 第 6 項）。"""
        prs = self.mod.build([_svg(
            '<text x="96" y="60" font-size="21.33" font-family="Noto Sans TC">內文</text>')])
        run = list(prs.slides[0].shapes)[0].text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.name, "Noto Sans TC")


class MultiPageTests(unittest.TestCase):
    """一頁 SVG ＝ 一張投影片。"""

    @classmethod
    def setUpClass(cls):
        cls.mod = _load()

    def test_one_svg_per_slide(self):
        prs = self.mod.build([_svg('<rect x="0" y="0" width="10" height="10"/>'),
                              _svg('<text x="10" y="20" font-size="21.33">第二頁</text>')])
        self.assertEqual(len(prs.slides), 2)


if __name__ == "__main__":
    unittest.main()
