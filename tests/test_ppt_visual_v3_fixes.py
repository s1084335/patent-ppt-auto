"""PPT skill v3 驗收後修正批次（規劃 P1；2026-07-31 使用者放行）。

實機驗收（report_trial_20260730_164919，19 頁）抓到的缺陷與使用者定案的版面升級，
唯一來源：`.agents/context/ppt-visual-rework-spec.md` 六之三節。本檔逐項釘住契約：

| # | 契約 |
|---|---|
| P1-1 | `chart_hero` 版型：大圖＋固定區註解卡＋底部核心結論條；內容頁新預設 |
| P1-2 | 機會評估頁接得到解讀（narratives 掛在 `cluster_topic_table:opportunity_*`，加 alias）|
| P1-3 | 主題分布依通道拆成兩張表格頁（rows 有資料就該用表格，不是降級大數字卡）|
| P1-4 | `table` 頁不再誤報 `narrative_missing`（表格頁本來就不配解讀）|
| P1-5 | 表格中文欄名；`topic_code` 不入畫面；`source_field` 值轉「技術／功效」|
| P1-6 | 研發方向建議移**最後、附錄之前**；動態插頁在它之前 |
| P1-7 | 研發方向合併版：三步色塊流程＋題目卡＋核心結論條；契約結構化 JSON |
| P1-8 | 封面主標用 workspace 名稱；`cover.title` AI slot 退場 |
| P1-10 | 來源註不得再寫「未經人工覆核」（預覽閘門就是人工覆核）|
"""
from __future__ import annotations

import importlib.util
import json
import sys
import tempfile
import unittest
from pathlib import Path

SKILL_DIR = Path(__file__).resolve().parents[1] / "skills" / "patent-report-ppt"
THEME_PATH = SKILL_DIR / "theme.json"


def _load_builder():
    spec = importlib.util.spec_from_file_location("build_ppt_v3fix", SKILL_DIR / "scripts" / "build_ppt.py")
    module = importlib.util.module_from_spec(spec)
    sys.modules.setdefault("build_ppt_v3fix", module)
    spec.loader.exec_module(module)
    return module


bp = _load_builder()


def _spec_by_topic(layout, topic_fragment):
    return [s for s in layout if topic_fragment in (s.topic or s.title)]


def _report_entry(rows, label_zh="X", report_type="ranking"):
    return {"label": label_zh, "label_zh": label_zh, "report_type": report_type,
            "rows": rows, "row_count": len(rows)}


def _minimal_report_data(**extra):
    data = {
        "parameters": {"version": "v-test", "workspace_name": "拉繩訓練機"},
        "reports": {
            "application_trend": _report_entry(
                [{"application_year": 2020 + i, "patent_count": 3} for i in range(4)],
                "專利申請趨勢", "trend"),
            # 附錄2 需要它才會出頁——附錄存在測試的頁序斷言才有意義。
            "applicant_ranking": _report_entry(
                [{"applicant_display_name": "TSMC", "patent_count": 9}],
                "主要申請人排名"),
        },
        "family_reports": {},
    }
    data["reports"].update(extra)
    return data


CLUSTER_ROWS = [
    {"topic_code": "T001", "label": "拉繩回收", "source_field": "wips_independent_claims",
     "patent_count": 15, "applicant_count": 13},
    {"topic_code": "T001", "label": "提升成效", "source_field": "effect_summary",
     "patent_count": 9, "applicant_count": 7},
]


class DirectionLastTests(unittest.TestCase):
    """P1-6：研發方向建議＝結論，必須在所有內容頁之後、附錄之前。"""

    def _layout(self, data):
        return bp._expand_page_layout(data, None)

    def test_direction_after_content_before_appendix(self):
        layout = self._layout(_minimal_report_data())
        kinds = [s.kind for s in layout]
        direction_at = kinds.index("direction")
        first_appendix = next(i for i, s in enumerate(layout) if s.is_appendix)
        self.assertGreater(direction_at, 0, "研發方向不得在封面前")
        content_pages = [i for i, s in enumerate(layout)
                        if s.kind not in {"cover", "direction"} and not s.is_appendix]
        self.assertTrue(all(i < direction_at for i in content_pages),
                        "還有內容頁排在研發方向之後——結論必須壓軸")
        self.assertLess(direction_at, first_appendix, "研發方向要在附錄之前")

    def test_dynamic_pages_insert_before_direction(self):
        """動態插頁也算證據，必須插在結論（direction）之前。"""
        data = _minimal_report_data(
            lifecycle=_report_entry([{"year": 2020, "patent_count": 2}], "專利生命週期"))
        layout = self._layout(data)
        direction_at = next(i for i, s in enumerate(layout) if s.kind == "direction")
        lifecycle_at = next(i for i, s in enumerate(layout)
                            if "lifecycle" in s.report_keys)
        self.assertLess(lifecycle_at, direction_at,
                        "動態插頁跑到結論之後——證據必須在結論前")


class ClusterTopicTableTests(unittest.TestCase):
    """P1-3：主題分布有 rows 就用表格，依通道拆兩頁；不再降級大數字卡。"""

    def test_split_two_table_pages_by_channel(self):
        data = _minimal_report_data(
            cluster_topic_table=_report_entry(CLUSTER_ROWS, "主題分類統計表", "cluster"))
        layout = bp._expand_page_layout(data, None)
        cluster_pages = [s for s in layout
                         if "cluster_topic_table" in s.report_keys and not s.is_appendix]
        self.assertEqual(len(cluster_pages), 2, "技術／功效應各一頁")
        kinds = {s.kind for s in cluster_pages}
        self.assertEqual(kinds, {"table_with_points"},
                         f"主題分布應為表格版型，實際 {kinds}")
        filters = {json.dumps(s.row_filter, sort_keys=True) for s in cluster_pages}
        self.assertEqual(len(filters), 2, "兩頁應帶不同的通道 row_filter")

    def test_row_filter_applied(self):
        data = _minimal_report_data(
            cluster_topic_table=_report_entry(CLUSTER_ROWS, "主題分類統計表", "cluster"))
        layout = bp._expand_page_layout(data, None)
        page = next(s for s in layout if s.row_filter
                    and "wips_independent_claims" in json.dumps(s.row_filter))
        rows = bp._first_rows(page, {"report_data": data})
        self.assertEqual(len(rows), 1)
        self.assertEqual(rows[0]["label"], "拉繩回收")


class TableColumnTests(unittest.TestCase):
    """P1-5：表格欄名中文化、topic_code 不入畫面、source_field 值轉通道。"""

    def test_column_label_map(self):
        for raw, expected in (("patent_count", "專利件數"), ("applicant_count", "申請人家數"),
                              ("label", "主題"), ("source_field", "通道")):
            with self.subTest(raw=raw):
                self.assertEqual(bp.TABLE_COLUMN_LABELS.get(raw), expected)

    def test_topic_code_excluded(self):
        self.assertIn("topic_code", bp.TABLE_EXCLUDED_COLUMNS,
                      "topic_code 供機制識別，不得入畫面（既有定案）")

    def test_source_field_value_mapped(self):
        self.assertEqual(bp.TABLE_VALUE_LABELS["source_field"]["wips_independent_claims"], "技術")
        self.assertEqual(bp.TABLE_VALUE_LABELS["source_field"]["effect_summary"], "功效")


class ChartHeroTests(unittest.TestCase):
    """P1-1：chart_hero 版型存在且為內容頁預設。"""

    def test_renderer_registered(self):
        self.assertIn("chart_hero", bp.RENDERERS)

    def test_default_kind_for_single_chart(self):
        report = {"report_type": "ranking"}
        self.assertEqual(bp._kind_for_report(report, 1), "chart_hero",
                         "單圖內容頁預設應為 chart_hero（大圖＋結論條）")

    def test_chart_dependent(self):
        """chart_hero 沒圖要降級，不得留空框。"""
        self.assertIn("chart_hero", bp.CHART_DEPENDENT_KINDS)
        self.assertIn("chart_hero", bp.SINGLE_CHART_KINDS)

    def test_conclusion_pick(self):
        """核心結論條：優先 emphasis 那條，其次「意涵」，再退 headline。"""
        points = [{"label": "現況", "text": "A", "emphasis": False},
                  {"label": "意涵", "text": "B", "emphasis": False},
                  {"label": "後續", "text": "C", "emphasis": True}]
        self.assertEqual(bp._conclusion_text("H", points), "C")
        points[2]["emphasis"] = False
        self.assertEqual(bp._conclusion_text("H", points), "B")
        self.assertEqual(bp._conclusion_text("H", []), "H")


class NarrativeAliasTests(unittest.TestCase):
    """P1-2：機會評估的解讀掛在 cluster_topic_table 的 opportunity_* 變體。"""

    NARRATIVES = {"reports": {"cluster_topic_table": {"variants": {
        "topic_table": {"headline": "主題集中", "points": [{"label": "現況", "text": "x"}], "text": "t"},
        "opportunity_tech": {"headline": "技術面機會", "points": [{"label": "意涵", "text": "y"}], "text": "t"},
        "opportunity_effect": {"headline": "功效面機會", "points": [{"label": "意涵", "text": "z"}], "text": "t"},
    }}}}

    def test_opportunity_page_finds_narrative(self):
        spec = bp.PageSpec(page=8, kind="comparison", title="機會評估", topic="機會評估",
                           report_keys=("opportunity_quadrant",),
                           charts=("opportunity_quadrant_tech.svg", "opportunity_quadrant_effect.svg"))
        matched, variant = bp._narrative_entry(self.NARRATIVES, bp._narrative_candidates(spec))
        self.assertTrue(variant, "機會評估頁找不到解讀（alias 未生效）")
        self.assertIn("機會", str(variant.get("headline")))

    def test_topic_table_page_not_hijacked(self):
        """⚠ 主題分布頁仍應拿 topic_table 變體，不得被 opportunity 搶走。"""
        spec = bp.PageSpec(page=6, kind="table_with_points", title="技術主題分布",
                           topic="技術主題分布", report_keys=("cluster_topic_table",))
        _, variant = bp._narrative_entry(self.NARRATIVES, bp._narrative_candidates(spec))
        self.assertEqual(variant.get("headline"), "主題集中")


class SplitLevelNarrativeTests(unittest.TestCase):
    """F-4：同一報表拆成 L4／L5 兩頁時，各頁必須取自己那個變體。

    🔴 2026-08-02 實機（22 頁）：p8 IPC Level 4 與 p9 IPC Level 5 的標題與四條
    判讀要點**逐字相同**，p10／p11 CPC 兩頁亦然。讀者會以為投影片貼重複。

    根因不在解讀端——`narratives.json` 確實產了 L4／L5 兩段不同內容。是組版端
    `_narrative_candidates` 只把圖檔主檔名 `ipc_main_distribution_L4` 當候選鍵，
    而 narratives 的鍵是 `ipc_main_distribution`（變體收在 `variants` 內），
    於是退回 report_key 那條線，`_narrative_entry` 取 `variants` 的**第一個**，
    兩頁都拿到 L4。缺的是把「主檔名＋後綴」翻成 `report_key:variant` 語法。
    """

    NARRATIVES = {"reports": {"ipc_main_distribution": {"variants": {
        "L4": {"headline": "A63B次分類達47件", "points": [{"label": "現況", "text": "L4"}], "text": "t"},
        "L5": {"headline": "A63B-069達19件", "points": [{"label": "現況", "text": "L5"}], "text": "t"},
    }}}}

    def _variant_of(self, chart: str) -> dict:
        spec = bp.PageSpec(page=9, kind="chart_hero", title="技術分類布局", topic="技術分類布局",
                           report_keys=("ipc_main_distribution",), charts=(chart,))
        _, variant = bp._narrative_entry(self.NARRATIVES, bp._narrative_candidates(spec))
        return variant

    def test_level4_page_takes_l4(self):
        self.assertEqual(self._variant_of("ipc_main_distribution_L4.svg").get("headline"), "A63B次分類達47件")

    def test_level5_page_takes_l5(self):
        """🔴 這條是實機錯誤的直接重現：L5 頁拿到了 L4 的解讀。"""
        self.assertEqual(self._variant_of("ipc_main_distribution_L5.svg").get("headline"), "A63B-069達19件")

    def test_two_pages_do_not_share_one_variant(self):
        l4 = self._variant_of("ipc_main_distribution_L4.svg")
        l5 = self._variant_of("ipc_main_distribution_L5.svg")
        self.assertNotEqual(l4.get("headline"), l5.get("headline"), "拆出來的兩頁共用了同一段解讀")


class CoverStatSizeTests(unittest.TestCase):
    """F-15：封面四張統計卡的主數字**字級不一致**。

    🔴 實機 p1：「60」「20」是大字，「2011–2026」因為字長被降到小字，
    四張卡並排時像三張重要、一張次要——但它們是同一層級的指標。
    ⚠ 分級規則本身沒錯（避免撐出卡片），錯在**逐張各算各的**。
    """

    def test_all_cards_share_one_size(self):
        theme = bp.Theme.load(THEME_PATH)
        stats = [("60", "件", "專利總數"), ("39 | 9", "CN | TW", "地域分布"),
                 ("2011–2026", "年", "年份區間"), ("20", "家", "申請人家數")]
        sizes = {bp._cover_stat_size(theme, stats, value) for value, _, _ in stats}
        self.assertEqual(len(sizes), 1, f"四張卡出現 {len(sizes)} 種字級：{sizes}")

    def test_shared_size_is_the_smallest_needed(self):
        """共用字級＝最長那個值需要的級數，不能大到把它撐出卡片。"""
        theme = bp.Theme.load(THEME_PATH)
        long_stats = [("60", "件", "x"), ("2011–2026", "年", "y")]
        short_stats = [("60", "件", "x"), ("20", "家", "y")]
        self.assertLess(bp._cover_stat_size(theme, long_stats, "60"),
                        bp._cover_stat_size(theme, short_stats, "60"))


class FrameworkBarFitsTests(unittest.TestCase):
    """F-15：分析框架條被截成「…等共 12 項…」——連收尾都被切掉。

    🔴 實機 p1。⚠ 根因是列數寫死 5 項：主題名長短不一，五個長名串起來就爆行。
    收尾那句「等共 N 項分析」是**資訊**（讀者才知道還有多少沒列），
    被截掉等於這行只剩半句話。
    """

    class _Spec:
        def __init__(self, topic):
            self.topic, self.title, self.kind, self.is_appendix = topic, topic, "chart_hero", False

    LONG = [_Spec.__init__ and None]  # placeholder replaced in setUp

    def setUp(self):
        names = ["申請趨勢", "專利生命週期", "保護地域分布", "公司×國家交叉表",
                 "國家佈局（現有保護）", "技術分類布局", "技術主題分布", "功效主題分布",
                 "競爭者佈局", "申請人年度專利分布矩陣", "專利權人年度布局矩陣", "機會評估"]
        self.layout = [self._Spec(n) for n in names]

    def test_tail_is_never_truncated(self):
        text = bp._framework_text(self.layout, budget_chars=60)
        self.assertLessEqual(len(text), 60)
        self.assertTrue(text.endswith("項分析"), f"收尾被截掉了：{text!r}")

    def test_fewer_items_when_names_are_long(self):
        tight = bp._framework_text(self.layout, budget_chars=40)
        loose = bp._framework_text(self.layout, budget_chars=120)
        self.assertLess(tight.count("→"), loose.count("→"))

    def test_all_items_listed_when_they_fit(self):
        short = [self._Spec(f"主題{i}") for i in range(3)]
        text = bp._framework_text(short, budget_chars=120)
        self.assertNotIn("等共", text, "全部列得下就不該出現「等共 N 項」")

    def test_total_count_is_honest(self):
        text = bp._framework_text(self.layout, budget_chars=60)
        self.assertIn(f"共 {len(self.layout)} 項", text)


class TextLanguageTagTests(unittest.TestCase):
    """F-14：中文字被 PowerPoint 的拼字檢查畫紅色波浪底線。

    🔴 使用者實機截圖 p20 可見紅色波浪線。⚠ 內容本身沒錯——是 run 沒有標語言，
    PowerPoint 拿**預設的英文校對**去檢查中文。轉圖看不到（proofing marks 不進圖），
    但客戶開檔會以為滿頁錯字。

    修法：run 上標 `lang="zh-TW"`（東亞另有 `altLang`），並關閉該 run 的拼字檢查。
    """

    def _run_xml(self, **kwargs):
        from pptx import Presentation
        from pptx.util import Inches
        theme = bp.Theme.load(THEME_PATH)
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bp._add_text(slide, theme, "同業加速進場、技術處成長期",
                     left=1.0, top=1.0, width=5.0, height=1.0, size=14.0, **kwargs)
        shape = slide.shapes[-1]
        return shape.text_frame.paragraphs[0].runs[0]._r.xml

    def test_run_declares_traditional_chinese(self):
        self.assertIn('lang="zh-TW"', self._run_xml())

    def test_run_disables_spell_check(self):
        """⚠ 只標語言不夠：PowerPoint 仍可能用簡體或其他詞庫標記專有名詞。"""
        self.assertIn('noProof="1"', self._run_xml())


class PanelHeightFitsContentTests(unittest.TestCase):
    """F-10：判讀面板高度固定，內容少時下方空掉三到四成。

    🔴 實機 p3／p4／p6／p12／p13／p16／p18／p19 八頁皆然——面板只有 2 條要點
    卻仍畫滿宣告高度。⚠ 版型算的是**框**不是**內容**，這也是 F-1（17/22 頁
    背景佔比 73–80%）的來源之一。
    """

    def _theme(self):
        return bp.Theme.load(THEME_PATH)

    def test_two_points_need_less_height_than_seven(self):
        theme = self._theme()
        few = [("現況", "A63B達47件", "", False), ("意涵", "布局集中訓練器材", "", False)]
        many = few + [("後續", f"第{i}項待查", "", False) for i in range(5)]
        g = theme.geometry["points_panel"]
        width = g["width_in"] - g["text_inset_right_in"]
        h_few = bp._points_panel_height(theme, few, width_in=width)
        h_many = bp._points_panel_height(theme, many, width_in=width)
        self.assertLess(h_few, h_many)

    def test_height_never_exceeds_declared_maximum(self):
        """⚠ 宣告高度是上限：內容再多也不能撐出版面。"""
        theme = self._theme()
        g = theme.geometry["points_panel"]
        width = g["width_in"] - g["text_inset_right_in"]
        huge = [("現況", "很長的要點" * 20, "", False) for _ in range(12)]
        self.assertLessEqual(bp._points_panel_height(theme, huge, width_in=width),
                             g["height_in"] + 1e-6)

    def test_height_has_a_floor(self):
        """空內容也要留得住標題列，不能塌成一條線。"""
        theme = self._theme()
        g = theme.geometry["points_panel"]
        width = g["width_in"] - g["text_inset_right_in"]
        self.assertGreater(bp._points_panel_height(theme, [], width_in=width), 0.5)


class SplitPageTitleTests(unittest.TestCase):
    """🔴 2026-08-03 使用者：**IPC/CPC 標題沒寫，看的人會搞混**。

    這是 F-8 的後遺症——我把 SVG 內建標題（「IPC 主分類分布 - Level 4」）移除時
    判斷「headline 已經能區分」，實測**不能**：
      p7「技術分類布局：A63B次分類達47件」← 完全沒說這是 IPC
      p8「技術分類布局：特殊訓練器械19件居首」← 也沒有，階層也不見了
    四頁併排時讀者分不出 IPC/CPC、也分不出 subclass/main group。

    ⚠ 移除重複資訊時，必須確認**剩下的那份真的講得完整**——
    我當時把「取捨待驗收時看」寫進紀錄就放行，那是把驗證延後而不是做掉。

    修法：拆頁時頁標題帶上引擎的 section 標題與 variant 顯示名（唯一來源，
    不在組版端另寫一份對照表）。
    """

    SECTIONS = [{
        "title": "IPC 主分類分布",
        "report_key": "ipc_main_distribution",
        "variants": [
            {"label": "4 階 · Subclass", "file": "ipc_main_distribution_L4.svg", "variant_key": "L4"},
            {"label": "5 階 · Main Group", "file": "ipc_main_distribution_L5.svg", "variant_key": "L5"},
        ],
    }, {
        "title": "CPC 主分類分布",
        "report_key": "cpc_main_distribution",
        "variants": [
            {"label": "4 階 · Subclass", "file": "cpc_main_distribution_L4.svg", "variant_key": "L4"},
            {"label": "5 階 · Main Group", "file": "cpc_main_distribution_L5.svg", "variant_key": "L5"},
        ],
    }]

    def _topic_of(self, chart: str) -> str:
        spec = bp.PageSpec(page=1, kind="chart_hero", title="技術分類布局",
                           topic="技術分類布局", report_keys=("ipc_main_distribution",),
                           charts=(chart,))
        return bp._chart_page_topic(spec, {"sections": self.SECTIONS})

    def test_ipc_page_says_ipc(self):
        self.assertIn("IPC", self._topic_of("ipc_main_distribution_L4.svg"))

    def test_cpc_page_says_cpc(self):
        self.assertIn("CPC", self._topic_of("cpc_main_distribution_L4.svg"))

    def test_level_is_visible(self):
        l4 = self._topic_of("ipc_main_distribution_L4.svg")
        l5 = self._topic_of("ipc_main_distribution_L5.svg")
        self.assertIn("4", l4)
        self.assertIn("5", l5)
        self.assertNotEqual(l4, l5, "兩個階層的頁標題一模一樣，讀者分不出來")

    def test_unknown_chart_keeps_original_topic(self):
        """⚠ 對不到 section 時維持原標題，不得產生空白或半截標題。"""
        spec = bp.PageSpec(page=1, kind="chart_hero", title="競爭者佈局", topic="競爭者佈局",
                           report_keys=("applicant_ranking",), charts=("applicant_ranking.svg",))
        self.assertEqual(bp._chart_page_topic(spec, {"sections": self.SECTIONS}), "競爭者佈局")

    def test_no_chart_keeps_original_topic(self):
        spec = bp.PageSpec(page=1, kind="table", title="附錄", topic="附錄",
                           report_keys=("cluster_topic_table",))
        self.assertEqual(bp._chart_page_topic(spec, {"sections": self.SECTIONS}), "附錄")


class ConclusionInPanelTests(unittest.TestCase):
    """🔴 2026-08-03 使用者：「判讀區塊那裡要能帶出核心結論，**還有不是每頁都要有核心結論**」。

    原本每頁底部固定一條「核心結論：…」橫幅，且結論那條會**從判讀面板移除**
    （`listed = [p for p in points if text != conclusion]`）——結論與依據被拆到兩處。

    改法：結論留在判讀面板（標 `emphasis`），底部橫幅取消。
    ⚠ 附帶效果正是使用者要的「圖表大一點」：橫幅讓出 0.5in＋間距，圖框得以加高。
    """

    def _theme(self):
        return bp.Theme.load(THEME_PATH)

    def test_conclusion_stays_in_the_points_list(self):
        """結論不得被排除在判讀面板之外——它是要點之一，不是附註。"""
        points = [{"label": "現況", "text": "A63B達47件", "emphasis": False},
                  {"label": "意涵", "text": "布局集中訓練器材", "emphasis": True}]
        kept = bp._points_for_panel(points)
        self.assertEqual(len(kept), 2, "結論那條被從面板拿掉了")
        self.assertIn("布局集中訓練器材", [p["text"] for p in kept])

    def test_conclusion_band_is_gone_from_chart_pages(self):
        """圖表頁的底部橫幅取消——不是每頁都要有核心結論。

        ⚠ 只針對圖表頁：研發方向建議頁**本身就是結論**，那條保留。
        """
        source = Path(bp.__file__).read_text(encoding="utf-8")
        start = source.index("def _render_chart_hero")
        body = source[start:source.index("def _render_chart_with_points")]
        self.assertNotIn("核心結論：", body,
                         "圖表頁仍畫結論橫幅；結論應由判讀面板的 emphasis 那條承擔")

    def test_chart_box_is_taller_than_before(self):
        """橫幅讓出的空間要真的給圖——否則只是少了東西、圖沒變大。"""
        g = self._theme().geometry["chart_hero"]
        self.assertGreaterEqual(g["image_height_in"], 4.9,
                                f"圖框仍是 {g['image_height_in']} in，沒吃到橫幅讓出的空間")

    def test_panel_height_follows_the_chart(self):
        """右側判讀面板要跟著加高，否則兩邊不齊。"""
        g = self._theme().geometry["chart_hero"]
        self.assertAlmostEqual(g["panel_height_in"], g["image_height_in"], places=2)


class ColumnPriorityAppliedTests(unittest.TestCase):
    """G-2：欄位放不下時要**砍尾巴**，不是照 rows 的鍵順序砍。

    🔴 實機 p11／p12／p20：`status`（技術狀態）排在 rows 的第 7 位，
    被 `max_columns=6` 擋掉——S2 整輪的重點功能一格都沒顯示。
    ⚠ 而且頁尾寫「完整欄位見附錄」，附錄也只有 6 欄——承諾了不存在的去處。

    引擎已用 `priority_columns` 宣告重要性順序；組版端必須照它取欄。
    """

    ROWS = [{
        "label": "拉繩捲輪回收機構", "patent_count": 15, "applicant_count": 13,
        "top3_share": 33, "top_applicants": "祺驊 3", "status": "成長技術",
        "representative": "CN123456 拉繩捲輪回收",
    }]
    LABELS = {"label": "主題標籤", "patent_count": "專利件數", "applicant_count": "申請人家數",
              "top3_share": "前三大占比(%)", "top_applicants": "前三大申請人",
              "status": "技術狀態", "representative": "代表專利"}
    PRIORITY = ("label", "patent_count", "applicant_count", "status",
                "top3_share", "top_applicants", "representative")

    def test_status_survives_when_columns_are_cut(self):
        """🔴 只放得下 4 欄時，技術狀態必須還在——它排在優先序第 4。"""
        cols = bp._ordered_columns(self.ROWS, excluded=set(), priority=self.PRIORITY, limit=4)
        self.assertIn("status", cols)

    def test_order_follows_priority_not_dict_order(self):
        cols = bp._ordered_columns(self.ROWS, excluded=set(), priority=self.PRIORITY, limit=7)
        self.assertEqual(cols[:4], ["label", "patent_count", "applicant_count", "status"])

    def test_unlisted_columns_come_after_priority(self):
        """⚠ 沒列進優先序的欄位不得消失——排在後面，一樣受 limit 約束。"""
        rows = [{**self.ROWS[0], "未列的欄": "x"}]
        cols = bp._ordered_columns(rows, excluded=set(), priority=self.PRIORITY, limit=8)
        self.assertEqual(cols[-1], "未列的欄")

    def test_excluded_still_wins(self):
        """排除清單優先於優先序——兩者衝突時不顯示。"""
        cols = bp._ordered_columns(self.ROWS, excluded={"status"},
                                   priority=self.PRIORITY, limit=7)
        self.assertNotIn("status", cols)

    def test_no_priority_falls_back_to_row_order(self):
        cols = bp._ordered_columns(self.ROWS, excluded=set(), priority=(), limit=3)
        self.assertEqual(cols, list(self.ROWS[0])[:3])


class NoTruncatedInformationTests(unittest.TestCase):
    """🔴 使用者定案（2026-08-03）：**資訊不能有被截斷的**。

    原本的處理是「放不下就切掉加『…』」——那是把問題丟給讀者：
    「祺驊 3：Brett Unswo…」讀者既不知道被切掉什麼，也無從查證。

    正解是**換呈現方式**：欄寬依內容分配（已做）、放不下就換行、
    列高依內容長、真的塞不下就少顯示幾列並標明完整版在附錄——
    但**已經顯示出來的每一格都必須完整**。
    """

    COLUMNS = ["主題標籤", "專利件數", "前三大申請人"]
    ROWS = [
        {"主題標籤": "拉繩捲輪回收機構", "專利件數": 15,
         "前三大申請人": "祺驊 3：Brett Unsworth 1：MOTIOFY AB 1"},
        {"主題標籤": "阻力調節拉繩機構", "專利件數": 8,
         "前三大申請人": "廈門帝瑪斯健康科技 2：MOTIOFY AB 1：OXEFIT, INC. 1"},
    ]

    def _table(self, width=12.0, height=4.0, row_height=0.4):
        from pptx import Presentation
        from pptx.util import Inches
        theme = bp.Theme.load(THEME_PATH)
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bp._add_table(slide, theme, self.ROWS, left=0.6, top=1.5, width=width, height=height,
                      row_height=row_height, max_columns=6, cell_margin_in=0.04,
                      cell_inset_in=0.06, labels={}, excluded=set())
        return next(s for s in slide.shapes if s.has_table).table

    def test_no_ellipsis_in_any_cell(self):
        table = self._table()
        for r, row in enumerate(table.rows):
            for c, _ in enumerate(table.columns):
                text = table.cell(r, c).text
                self.assertNotIn("…", text, f"儲存格 ({r},{c}) 被截斷：{text!r}")

    def test_long_value_is_kept_whole(self):
        table = self._table()
        shown = {table.cell(r, c).text for r in range(len(table.rows)) for c in range(len(table.columns))}
        self.assertIn("祺驊 3：Brett Unsworth 1：MOTIOFY AB 1", shown)

    def test_header_is_kept_whole(self):
        table = self._table()
        headers = [table.cell(0, c).text for c in range(len(table.columns))]
        self.assertEqual(headers, self.COLUMNS)

    def test_row_height_grows_for_wrapped_content(self):
        """放不下就換行——換了行就要有對應的列高，否則字會被切在框外。"""
        narrow = self._table(width=6.0)
        wide = self._table(width=12.0)
        narrow_h = sum(r.height for r in narrow.rows)
        wide_h = sum(r.height for r in wide.rows)
        self.assertGreater(narrow_h, wide_h, "窄表沒有加高列高，換行的字會被切掉")


class TableColumnWidthTests(unittest.TestCase):
    """F-3＋F-16：表格欄寬一律等分，兩個方向都出錯。

    🔴 F-3 實機 p21 十處截斷：「祺驊 3：Brett Unswo…」「廈門帝瑪斯健康科技 2…」。
    🔴 F-16 實機 p22：「最新受讓人名單」14 列只有 2 列有值、「受讓取得」全是 0/1，
    兩欄合計佔掉 **50%** 寬度。

    ⚠ 截斷是症狀，**等分配置才是根因**——`width / len(columns)` 讓「專利件數」
    這種兩位數的欄和「前三大申請人」拿一樣寬。
    """

    COLUMNS = ["主題標籤", "專利件數", "申請人家數", "前三大占比(%)", "前三大申請人"]
    ROWS = [
        {"主題標籤": "拉繩捲輪回收機構", "專利件數": 15, "申請人家數": 13,
         "前三大占比(%)": 33, "前三大申請人": "祺驊 3：Brett Unsworth 1：MOTIOFY AB 1"},
        {"主題標籤": "阻力調節拉繩機構", "專利件數": 8, "申請人家數": 7,
         "前三大占比(%)": 50, "前三大申請人": "廈門帝瑪斯健康科技 2：MOTIOFY AB 1"},
    ]

    def _widths(self, total=12.0):
        return bp._column_widths(self.COLUMNS, self.ROWS, {}, total,
                                 size_pt=11.0, inset_in=0.06)

    def test_widths_sum_to_the_table_width(self):
        self.assertAlmostEqual(sum(self._widths(12.0)), 12.0, places=4)

    def test_numeric_column_is_narrower_than_text_column(self):
        """🔴 這是 F-3／F-16 的核心：兩位數的欄不該和長名單一樣寬。"""
        widths = dict(zip(self.COLUMNS, self._widths()))
        self.assertLess(widths["專利件數"], widths["前三大申請人"])
        self.assertLess(widths["申請人家數"], widths["前三大申請人"])

    def test_header_still_fits_in_narrow_columns(self):
        """⚠ 窄欄也要放得下自己的欄頭，否則變成欄頭被截（比內容被截更糟）。"""
        widths = dict(zip(self.COLUMNS, self._widths()))
        for name in self.COLUMNS:
            need = bp._display_width(name) * (11.0 / 72.0) + 0.06 * 2
            self.assertGreaterEqual(widths[name] + 1e-6, min(need, 12.0 / len(self.COLUMNS)),
                                    f"{name} 欄放不下自己的欄頭")

    def test_mostly_empty_column_gets_less_room(self):
        """F-16：整欄幾乎沒值的欄不該佔滿寬度。"""
        columns = ["申請人", "專利件數", "最新受讓人名單", "受讓取得"]
        rows = [{"申請人": "廈門帝瑪斯健康科技", "專利件數": 13, "最新受讓人名單": "", "受讓取得": 1},
                {"申請人": "孟喬", "專利件數": 5, "最新受讓人名單": "億軒", "受讓取得": 0},
                {"申請人": "祺驊", "專利件數": 5, "最新受讓人名單": "", "受讓取得": 0}]
        widths = dict(zip(columns, bp._column_widths(columns, rows, {}, 12.0,
                                                     size_pt=11.0, inset_in=0.06)))
        self.assertLess(widths["受讓取得"], widths["申請人"])
        self.assertLess(widths["最新受讓人名單"], widths["申請人"])

    def test_single_column_takes_everything(self):
        self.assertEqual(bp._column_widths(["只有一欄"], [{"只有一欄": "x"}], {}, 9.0,
                                           size_pt=11.0, inset_in=0.06), [9.0])

    def test_labels_are_used_for_header_width(self):
        """欄頭寬要照**顯示名**算，不是照內部欄名。"""
        widths = bp._column_widths(["patent_count", "top3_applicants"],
                                   [{"patent_count": 15, "top3_applicants": "甲 3：乙 1"}],
                                   {"patent_count": "專利件數", "top3_applicants": "前三大申請人"},
                                   12.0, size_pt=11.0, inset_in=0.06)
        self.assertLess(widths[0], widths[1])


class PercentageBarRatioTests(unittest.TestCase):
    """F-7：佔比條的條長必須用「佔全體比例」，不得用「相對第一名」。

    🔴 2026-08-02 實機 p5：CN 39 件（65%）畫**滿格**，TW/US 9 件（15%）畫成短條，
    同一張圖裡兩種基準——條長分母是 `top_value`（第一名），右側標的百分比分母卻是
    `total`。讀者看到 CN 佔滿整條軌道，會讀成 100%，而字寫 65%。

    軌道（bar_track）本身就是 100% 基準，條長用真佔比才對得起來：
    CN 畫到 65%，剩下的 35% 留白正是「還有其他國家」的資訊。
    """

    def test_ratio_uses_total_not_max(self):
        """CN 39／總 60 → 0.65，不是 39/39＝1.0。"""
        self.assertAlmostEqual(bp._bar_fill_ratio(39, 60), 0.65)
        self.assertAlmostEqual(bp._bar_fill_ratio(9, 60), 0.15)
        self.assertAlmostEqual(bp._bar_fill_ratio(3, 60), 0.05)

    def test_top_bar_is_not_full_when_share_below_100(self):
        """第一名只要不是 100%，就不得畫滿。"""
        self.assertLess(bp._bar_fill_ratio(39, 60), 1.0)

    def test_ratios_sum_to_one(self):
        """全部條長加起來等於整條軌道——這是「佔比」的定義。"""
        total = 60
        self.assertAlmostEqual(sum(bp._bar_fill_ratio(v, total) for v in (39, 9, 9, 3)), 1.0)

    def test_zero_total_does_not_raise(self):
        self.assertEqual(bp._bar_fill_ratio(0, 0), 0.0)


class TablePageNoNarrativeWarningTests(unittest.TestCase):
    """P1-4：純表格頁（明細）不配解讀，不得誤報 narrative_missing。"""

    def test_no_warning_for_table_kind(self):
        data = _minimal_report_data(
            family_quality_detail=_report_entry(
                [{"a": 1}], "家族完整性明細", "detail"))
        with tempfile.TemporaryDirectory() as tmp:
            report_dir = Path(tmp)
            (report_dir / "report_data.json").write_text(
                json.dumps(data, ensure_ascii=False), encoding="utf-8")
            (report_dir / "artifact_manifest.json").write_text(
                json.dumps({"artifacts": []}), encoding="utf-8")
            result = bp.build_ppt(report_dir=report_dir, output_dir=report_dir / "out")
        bad = [w for w in result["manifest"]["warnings"]
               if w.get("type") == "narrative_missing"
               and "family_quality_detail" in str(w.get("report_key"))]
        self.assertEqual(bad, [], "明細表格頁誤報 narrative_missing")


class CoverWorkspaceTests(unittest.TestCase):
    """P1-8：封面主標用 workspace 名稱；cover.title AI slot 退場。"""

    def test_slot_keys_only_direction(self):
        self.assertEqual(bp.all_slot_keys(), ["direction.body"],
                         "cover.title 應退場——封面標題由 workspace 名確定性組成")

    def test_cover_title_from_workspace(self):
        """2026-07-31 起補上「專利分析」：單一 workspace 名不像簡報標題。"""
        self.assertEqual(bp._cover_title(_minimal_report_data(), {}), "拉繩訓練機專利分析")

    def test_cover_title_not_doubled(self):
        """workspace 名本身已含「專利分析」時不得再接一次。"""
        data = _minimal_report_data()
        data["parameters"]["workspace_name"] = "拉繩訓練機專利分析"
        self.assertEqual(bp._cover_title(data, {}), "拉繩訓練機專利分析")

    def test_cover_title_fallback(self):
        data = _minimal_report_data()
        data["parameters"].pop("workspace_name")
        self.assertEqual(bp._cover_title(data, {}), "專利情報整合分析")


class FootnoteTests(unittest.TestCase):
    """P1-10：來源註只留可追溯資訊；「未經人工覆核」退場。"""

    def test_no_unverified_disclaimer(self):
        src = (SKILL_DIR / "scripts" / "build_ppt.py").read_text(encoding="utf-8")
        self.assertNotIn("未經人工覆核", src,
                         "預覽閘門就是人工覆核——每頁再蓋此章等於否定定稿流程")
        self.assertIn("報表版本", src.split("def _render_footnote")[1][:1200],
                      "來源註應含報表版本（可追溯）")


class DirectionStructuredTests(unittest.TestCase):
    """P1-7：direction.body 契約結構化（合併版）。"""

    def test_parse_structured(self):
        payload = json.dumps({
            "situation": ["2022 起放量"], "opportunity": ["省空間區 7 件 6 家"],
            "direction": ["阻力調節差異化"],
            "topics": [{"name": "可調阻力", "basis": "必守核心區", "action": "迴避設計"}],
            "conclusion": "現在是投入窗口",
        }, ensure_ascii=False)
        parsed = bp._parse_direction_body(payload)
        self.assertEqual(parsed["conclusion"], "現在是投入窗口")
        self.assertEqual(parsed["topics"][0]["name"], "可調阻力")

    def test_parse_legacy_text(self):
        """⚠ 舊純文字要能過渡：回 None 走舊版面，不炸。"""
        self.assertIsNone(bp._parse_direction_body("整體態勢：…一段長文"))

    def test_content_rules_updated(self):
        rules = (SKILL_DIR / "report_ppt_content_rules.md").read_text(encoding="utf-8")
        for key in ("situation", "opportunity", "direction", "topics", "conclusion"):
            with self.subTest(key=key):
                self.assertIn(key, rules, f"content_rules 缺 direction 結構欄位 {key}")
        # ⚠ cover.title 允許以「退場沿革」形式出現（防止改回去），但不得再是要產的 slot。
        self.assertIn("只產一個 slot", rules, "規則應明示只剩 direction.body 一個 slot")
        self.assertIn("退場", rules, "cover.title 退場沿革要留在規則裡")


class OpportunitySplitTests(unittest.TestCase):
    """2026-07-31 使用者二輪回饋：機會評估同頁比較圖太小——改**分頁**（一象限圖一頁）。"""

    def test_opportunity_splits_into_hero_pages(self):
        spec = bp.PageSpec(page=7, kind="comparison", title="機會評估", topic="機會評估",
                           report_keys=("opportunity_quadrant",),
                           charts=("opportunity_quadrant_tech.svg", "opportunity_quadrant_effect.svg"))
        pages = bp._split_pairs_by_policy([spec], None)
        self.assertEqual(len(pages), 2, "技術面／功效面應各一頁")
        self.assertEqual({s.kind for s in pages}, {"chart_hero"},
                         "拆頁後應用大圖版型，不是要點框版")
        self.assertEqual([len(s.charts) for s in pages], [1, 1])

    def test_ipc_stays_comparison(self):
        """⚠ IPC/CPC 維持同頁比較——只有機會評估改分頁。"""
        self.assertNotIn("ipc_main_distribution", bp.SPLIT_PAIR_REPORTS)
        self.assertIn("opportunity_quadrant", bp.SPLIT_PAIR_REPORTS)


class YearMatrixMainOnlyTests(unittest.TestCase):
    """2026-07-31 使用者二輪回饋：年度矩陣只用前 10 名主表那張，「更多」圖不上 PPT。"""

    def test_more_chart_filtered(self):
        files = ("owner_year_matrix.svg", "owner_year_matrix_more.svg")
        kept = bp._filter_report_charts(("owner_year_matrix",), files)
        self.assertEqual(kept, ("owner_year_matrix.svg",))

    def test_other_reports_unaffected(self):
        files = ("ipc_main_distribution_L4.svg", "ipc_main_distribution_L5.svg")
        self.assertEqual(bp._filter_report_charts(("ipc_main_distribution",), files), files)


if __name__ == "__main__":
    unittest.main()
