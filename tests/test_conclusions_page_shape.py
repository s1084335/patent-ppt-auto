"""結論頁改版：數據降為小字、判讀與行動為主體（2026-08-18 使用者裁決）。

## 為什麼改

使用者：「這頁數據要少點，更多的是根據數據判讀後最終要轉到專利行動。」

原本四欄 `主題｜發現｜研發意涵｜專利行動`，「發現」是純統計
（`10件/9家｜最大持有 20%｜申請成長`）而且佔 3.35in——**結論頁在重述前面已經
講過的統計**。這正是 deepen 主軸要打的：「往獨立項挖深，把統計壓薄」。

## 但「發現」不能刪

它是全頁**唯一被機械驗證**的欄（逐字比對引擎字串），刪掉就沒有錨點擋 CLI 亂編。
故**降級不刪除**：從一整欄變成主題底下的灰小字，逐字比對照舊。

## 行動分組不得固定五組

使用者：「這些是選項，但不能每次都固定全部行動都有。」
⚠ 渲染時只出現**真的用到**的分組。五組固定列出會讓 CLI 覺得每組都要填
——那是形式鎖逼出硬湊（v5／v7／v9 同型）。

⚠ 反過來也**不設**「不得五組全有」的閘門：十個主題真的可能用滿五種行動，
禁止它會逼 CLI 把合理的行動改掉。這是判斷不是規則。
"""
from __future__ import annotations

import inspect
import json
import sys
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SKILL = ROOT / "skills" / "html-report-to-deck"
sys.path.insert(0, str(SKILL / "scripts"))


def _row(topic, action, reading="判讀句", finding="10件/5家｜成長", pending=0):
    return {"topic": topic, "finding": finding, "reading": reading,
            "action": action, "pending_count": pending}


class ThreeColumnsTests(unittest.TestCase):
    """四欄 → 三欄：主題（含數據小字）｜判讀｜專利行動。"""

    def test_header_says_reading_not_implication(self):
        import deck_layout

        from tests.source_assertions import executable_source

        src = executable_source(inspect.getsource(deck_layout.slide_conclusions))
        self.assertIn("判讀", src, "欄名沒改成「判讀」")
        self.assertNotIn(
            "研發意涵", src,
            "仍用「研發意涵」——欄名就是對 CLI 的指令，"
            "「意涵」容易寫成感想，「判讀」要求從資料推到結論")

    def test_three_columns_not_four(self):
        import deck_layout

        from tests.source_assertions import executable_source

        src = executable_source(inspect.getsource(deck_layout.slide_conclusions))
        self.assertNotIn(
            '"發現"', src,
            "「發現」仍是獨立欄——應降為主題底下的小字，把版面讓給判讀與行動")

    def test_finding_still_rendered_somewhere(self):
        """⚠ 降級不是刪除：逐字驗證的錨點必須還在頁面上。"""
        import deck_layout

        from tests.source_assertions import executable_source

        src = executable_source(inspect.getsource(deck_layout.slide_conclusions))
        self.assertIn(
            "finding", src,
            "發現字串完全不見了——那是唯一被機械驗證的內容，"
            "不印在頁上就等於驗證了一個沒人看到的東西")


class OnlyUsedActionGroupsTests(unittest.TestCase):
    """🔴 使用者裁決：不能每次都固定全部行動都有。"""

    def _group_labels(self, rows):
        """分組是純函式——直接驗它，不要側錄畫布。

        ⚠ 側錄 `textbox` 會把「有沒有畫出來」與「畫在哪、什麼樣式」綁在一起，
        版面一動測試就假紅。分組是可獨立驗證的邏輯，就該獨立驗。
        """
        import deck_layout

        return [verb for verb, _rows in deck_layout.conclusion_groups(rows)]

    def test_unused_groups_are_not_rendered(self):
        labels = self._group_labels([
            _row("A", "佈局"), _row("B", "佈局"), _row("C", "迴避設計"),
        ])
        group_labels = [l for l in labels]
        self.assertIn("佈局", group_labels)
        self.assertIn("迴避設計", group_labels)
        for unused in ("追蹤", "細讀比對", "暫不投入"):
            with self.subTest(verb=unused):
                self.assertNotIn(
                    unused, group_labels,
                    f"沒有任何主題用「{unused}」卻仍畫出該分組——"
                    "五組固定列出會讓 CLI 覺得每組都要填（形式鎖逼出硬湊）")

    def test_all_five_is_allowed_when_genuinely_used(self):
        """⚠ 不設「不得五組全有」的閘門：十個主題真的可能用滿五種。"""
        rows = [_row(f"T{i}", v) for i, v in enumerate(
            ("佈局", "追蹤", "迴避設計", "細讀比對", "暫不投入"))]
        labels = self._group_labels(rows)
        for verb in ("佈局", "追蹤", "迴避設計", "細讀比對", "暫不投入"):
            self.assertIn(verb, labels, f"真的用到的 {verb} 沒被畫出來")


class TemplateUsesReadingTests(unittest.TestCase):
    def test_template_row_uses_reading_key(self):
        data = json.loads(
            (SKILL / "references" / "content-template.json").read_text(encoding="utf-8"))
        row = ((data.get("conclusions") or {}).get("rows") or [{}])[0]
        self.assertIn("reading", row, "範本仍用 implication——欄名要改成判讀")
        self.assertNotIn("implication", row)

    def test_template_does_not_list_all_five_actions_as_rows(self):
        """⚠ 範本不得示範五列各一種行動——CLI 會照抄成「每種都要有」。"""
        data = json.loads(
            (SKILL / "references" / "content-template.json").read_text(encoding="utf-8"))
        rows = (data.get("conclusions") or {}).get("rows") or []
        self.assertLessEqual(
            len(rows), 2,
            "範本列了太多結論列——CLI 會照抄數量。行動是選項不是清單")



class MultiActionTests(unittest.TestCase):
    """🔴 使用者裁決（2026-08-19）：行動不是單選。

    「一種策略也許適用多個主題，也可能只適用一個主題……但也不是說行動就只能
    選擇一種而已。」前半（一策略對多主題）分組本來就成立；後半**不成立**——
    `action` 是單一字串，等於強迫一主題一行動。實務上「先迴避設計、同時追蹤
    對手審查中的案」是同一主題的兩個動作，寫不出來就會被迫二選一，
    而被丟掉的那個不會留下任何痕跡（缺席型偏差）。

    ⚠ 一併**不加**「至多／至少 N 種」的數量鎖：那是 v5／v7／v9 形式鎖同型，
    會逼 CLI 湊數。約束留在既有那條——每個動詞都要能從同列「判讀」推得出來。
    """

    def test_single_action_string_still_works(self):
        """⚠ 舊格式必須照舊：既有 content.json 與範例都是字串。"""
        import deck_layout

        groups = dict(deck_layout.conclusion_groups([_row("A", "佈局")]))
        self.assertEqual([r["topic"] for r in groups["佈局"]], ["A"])

    def test_row_appears_under_every_declared_action(self):
        import deck_layout

        groups = dict(deck_layout.conclusion_groups([
            _row("A", ["迴避設計", "追蹤"]),
            _row("B", "追蹤"),
        ]))
        self.assertEqual(sorted(groups), sorted(["追蹤", "迴避設計"]))
        self.assertEqual([r["topic"] for r in groups["迴避設計"]], ["A"])
        self.assertEqual(
            sorted(r["topic"] for r in groups["追蹤"]), ["A", "B"],
            "宣告兩個行動的主題只出現在其中一組——另一個行動被靜默丟掉")

    def test_row_actions_is_the_single_definition_point(self):
        """畫的（deck_layout）與擋的（check_content）必須讀同一份解析。"""
        import deck_layout

        self.assertEqual(deck_layout.row_actions({"action": "佈局"}), ["佈局"])
        self.assertEqual(
            deck_layout.row_actions({"action": ["佈局", "追蹤"]}), ["佈局", "追蹤"])
        self.assertEqual(deck_layout.row_actions({}), [])

    def test_gate_rejects_bad_verb_inside_a_list(self):
        """⚠ 多值最容易漏的就是「只驗第一個」。"""
        import check_content

        self.assertTrue(check_content._bad_actions({"action": ["追蹤", "立即提出申請"]}))
        self.assertFalse(check_content._bad_actions({"action": ["追蹤", "佈局"]}))

class FontSizeContractTests(unittest.TestCase):
    """🔴 SVG 與 PPTX 兩條渲染路徑必須接受同一份 `size`。

    2026-08-19 實測：`{"size": 11}`（裸 int）在 SVG 路徑完全正常——它讀
    `default_size_pt` 用的是數值；PPTX 路徑把 11 當成 **11 EMU**，換算後是
    0 centipoints，python-pptx 直接 `ValueError: must be in range 100 to
    400000`。整個 deck 側單元測試全綠，只有半真機械鏈（真的產 .pptx）才炸。

    ⚠ 這是「同一份知識兩個落點」的典型：字級單位在 SVG 端是 pt 數值、在 PPTX
    端是 Length 物件，兩邊各自解讀而沒有人對帳。修法不是「記得都寫 Pt()」
    ——那要靠人記得；而是在 `_set_font`（唯一落點）把 int 視為 pt。
    """

    def test_int_size_is_accepted_by_pptx_path(self):
        from pptx import Presentation
        from pptx.util import Inches

        import deck_layout

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        s = prs.slides.add_slide(prs.slide_layouts[6])
        deck_layout.textbox(s, 1, 1, 3, 0.4, [("小字", {"size": 11})])
        run = s.shapes[0].text_frame.paragraphs[0].runs[0]
        self.assertEqual(
            run.font.size.pt, 11,
            "裸 int 沒被當成 pt——SVG 端會過、PPTX 端會炸，且只有真產檔才看得到")

    def test_pt_object_is_not_converted_twice(self):
        """⚠ `Pt(16)` **是** int 的子類別——第一版修法用 isinstance(size, int)
        判斷，於是既有的 `T_SIZE`／`B_SIZE` 被再換算一次（同一個 ValueError
        換個數字）。判斷必須用 `Length`。
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt

        import deck_layout

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        s = prs.slides.add_slide(prs.slide_layouts[6])
        deck_layout.textbox(s, 1, 1, 3, 0.4, [("內文", {"size": Pt(16)})])
        run = s.shapes[0].text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.size.pt, 16, "Pt 物件被二次換算")

class FontScaleSingleSourceTests(unittest.TestCase):
    """🔴 字級白名單只能有一個定義處。

    `deck_layout` 定字級（`T_SIZE`／`B_SIZE`），`audit_deck` 又用字面
    `"24,16"` 定一次白名單。兩份各自演進而**不一致本身不會報錯**——
    2026-08-19 實測：§7d 加了 11pt 灰小字，deck 側單元測試全綠、PPTX 也產得
    出來，直到 audit 才紅，而且訊息是「字級分布」看不出根因。

    改為 audit 預設讀 `deck_layout.ALLOWED_SIZES`：加字級只改一處，
    第二處自動跟上。
    """

    def test_audit_default_comes_from_deck_layout(self):
        import audit_deck
        import deck_layout

        from tests.source_assertions import executable_source

        src = executable_source(inspect.getsource(audit_deck))
        self.assertNotIn(
            '"24,16"', src,
            "audit 仍自帶字級字面——它會與 deck_layout 分岔且不會報錯")
        self.assertIn(
            "ALLOWED_SIZES", src, "audit 沒有讀 deck_layout 的字級清單")
        self.assertEqual(
            sorted(s.pt for s in deck_layout.ALLOWED_SIZES), [11.0, 16.0, 24.0],
            "字級層級不是三層——標題 24／內文 16／註記小字 11")

    def test_small_size_constant_is_used_not_literal(self):
        """⚠ 小字用具名常數，不要散落 `\"size\": 11` 字面。"""
        import deck_layout

        from tests.source_assertions import executable_source

        src = executable_source(inspect.getsource(deck_layout.slide_conclusions))
        self.assertNotIn(
            '"size": 11', src,
            "小字寫成字面 11——第四處用到時沒有人知道要一起改")


if __name__ == "__main__":
    unittest.main()
