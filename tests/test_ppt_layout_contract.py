"""patent-report-ppt v3 版型與內容契約測試（2026-07-30 重建）。

重建動機（實機驗收不合格）：每頁 400–500 字單段字牆、圖檔對不上顯示佔位、
圖片溢出版面 3.5 吋、封面標題壓字、主題色 accent 完全沒用到。

本檔守的是**重建後不得再退回去**的硬契約：
1. 版型庫九種 kind 齊備，成對報表可同頁左右並排、也可分頁，但不得合成同一張圖。
2. 圖檔一律走 artifact_manifest 反查，禁止用 `{report_key}.svg` 猜檔名。
3. 找不到圖降級 stat_callout，且 PPT 內不得出現「待產出」這類佔位文字。
4. narrative 缺 headline／points（舊格式只有 text）要 fallback 成條列並寫 warnings，不得靜默。
5. PPT 階段 AI 只產 cover.title 與 direction.body 兩個 slot。
6. 字體全微軟正黑體（含 latin／ea 兩處 typeface），字級下限 12pt。
"""

from __future__ import annotations

import importlib.util
import json
import sys
from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation
from pptx.util import Inches

SKILL_DIR = Path(__file__).resolve().parents[1] / "skills" / "patent-report-ppt"
BUILDER_PATH = SKILL_DIR / "scripts" / "build_ppt.py"
THEME_PATH = SKILL_DIR / "theme.json"

VERSION = "report_trial_20260730_000000"

# 兩張尺寸不同的 SVG：寬扁圖與高瘦圖，用來驗證等比縮放後不會溢出版面。
WIDE_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="1600" height="500">'
    '<rect width="1600" height="500" fill="#516CEE"/></svg>'
)
TALL_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="500" height="1600">'
    '<rect width="500" height="1600" fill="#63DBF5"/></svg>'
)


def _load_builder():
    """以檔案路徑載入可攜的 build_ppt.py（不進主專案 import 路徑）。"""
    spec = importlib.util.spec_from_file_location("build_ppt_v3", BUILDER_PATH)
    module = importlib.util.module_from_spec(spec)
    sys.modules["build_ppt_v3"] = module
    spec.loader.exec_module(module)
    return module


def _artifact(file_name: str, report_names: list[str]) -> dict:
    return {
        "file": file_name,
        "artifact_type": "chart_svg",
        "report_name": report_names[0] if len(report_names) == 1 else None,
        "report_names": report_names,
        "sha256": "0" * 64,
    }


@pytest.fixture()
def report_dir(tmp_path: Path) -> Path:
    """實機資料的縮影：圖檔名與 report_key 不同名、IPC 有 L4/L5 兩張、
    機會矩陣有技術/功效兩張、cluster_topic_table 有資料但沒有圖。"""
    path = tmp_path / "cache" / VERSION
    path.mkdir(parents=True)

    report_data = {
        "parameters": {
            "version": VERSION,
            "generated_at": "2026-07-30T06:05:05",
            "topic_run_id": "topic-run-001",
            "topic_state_version": "state-v3",
        },
        "reports": {
            "application_trend": {
                "label_zh": "專利申請趨勢",
                "report_type": "aggregate",
                "rows": [
                    {"application_year": 2021, "patent_count": 8},
                    {"application_year": 2022, "patent_count": 15},
                    {"application_year": 2023, "patent_count": 6},
                ],
            },
            "country_distribution": {
                "label_zh": "專利受理國分布",
                "report_type": "aggregate",
                "rows": [
                    {"country_code": "CN", "patent_count": 40},
                    {"country_code": "US", "patent_count": 9},
                    {"country_code": "TW", "patent_count": 3},
                ],
            },
            "ipc_main_distribution": {
                "label_zh": "IPC 主分類分布",
                "report_type": "aggregate",
                "rows": [
                    {"Orig. IPC(Main)": "A63B", "patent_count": 47},
                    {"Orig. IPC(Main)": "F03G", "patent_count": 2},
                ],
            },
            "opportunity_quadrant": {
                "label_zh": "機會評估矩陣",
                "report_type": "aggregate",
                "rows": [
                    {"topic": "拉繩捲輪回收機構", "patent_count": 15, "applicant_count": 13},
                ],
            },
            "cluster_topic_table": {
                "label_zh": "全分類技術指標總表",
                "report_type": "table",
                "rows": [
                    {"topic": "拉繩捲輪回收機構", "patent_count": 15, "applicant_count": 13},
                    {"topic": "馬達捲繩自鎖機構", "patent_count": 8, "applicant_count": 4},
                ],
            },
            "applicant_ranking": {
                "label_zh": "主要申請人排名",
                "report_type": "aggregate",
                "rows": [{"applicant_display_name": "廈門帝瑪斯", "patent_count": 13}],
            },
        },
        "family_reports": {},
    }
    (path / "report_data.json").write_text(
        json.dumps(report_data, ensure_ascii=False), encoding="utf-8"
    )

    artifacts = {
        "metadata": {"version": VERSION},
        "artifacts": [
            _artifact("annual_trend.svg", ["application_trend", "publication_trend"]),
            _artifact("jurisdiction_distribution.svg", ["country_distribution"]),
            _artifact("ipc_main_distribution_L4.svg", ["ipc_main_distribution"]),
            _artifact("ipc_main_distribution_L5.svg", ["ipc_main_distribution"]),
            _artifact("opportunity_quadrant_tech.svg", ["opportunity_quadrant"]),
            _artifact("opportunity_quadrant_effect.svg", ["opportunity_quadrant"]),
            _artifact("applicant_ranking.svg", ["applicant_ranking"]),
        ],
    }
    (path / "artifact_manifest.json").write_text(
        json.dumps(artifacts, ensure_ascii=False), encoding="utf-8"
    )

    for name in (
        "annual_trend.svg",
        "jurisdiction_distribution.svg",
        "ipc_main_distribution_L4.svg",
        "ipc_main_distribution_L5.svg",
        "opportunity_quadrant_tech.svg",
        "applicant_ranking.svg",
    ):
        (path / name).write_text(WIDE_SVG, encoding="utf-8")
    # 高瘦圖：驗證等比縮放不會撐破版面高度。
    (path / "opportunity_quadrant_effect.svg").write_text(TALL_SVG, encoding="utf-8")

    narratives = {
        "based_on_version": VERSION,
        "reports": {
            # 新格式：有 headline 與 points。
            "annual_trend": {
                "variants": {
                    "default": {
                        "headline": "2021–2023 為布局高峰",
                        "points": [
                            {"label": "現況", "text": "2022 年 15 件為單年高峰", "emphasis": True},
                            {"label": "意涵", "text": "屬近年才熱起來的布局"},
                            {"label": "後續", "text": "追蹤 2025 年後公告進度"},
                        ],
                        "text": "長文供報表頁使用。",
                    }
                }
            },
            # 舊格式：只有 text，必須 fallback 且寫 warning。
            "ipc_main_distribution": {
                "variants": {
                    "default": {
                        "text": (
                            "觀察：技術幾乎全數集中於 A63B 健身器材，"
                            "5 階以拉繩 19 件與阻力調節 17 件兩群領先。"
                            "意涵：重心落在拉繩與阻力調節主線、由多家共同布局、集中度低。"
                            "決策提醒：就多方主群覆核代表專利確認布局密集度。"
                        )
                    }
                }
            },
        },
    }
    (path / "narratives.json").write_text(
        json.dumps(narratives, ensure_ascii=False), encoding="utf-8"
    )
    return path


@pytest.fixture()
def approvals(tmp_path: Path) -> Path:
    path = tmp_path / "approvals.json"
    path.write_text(
        json.dumps(
            {
                "report_version": VERSION,
                "slots": {
                    "cover.title": "拉繩阻力健身器材專利佈局分析",
                    "direction.body": "整體態勢：申請量自 2021 年放量。\n技術重心：集中於 A63B。",
                },
                "layout_overrides": {},
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    return path


@pytest.fixture()
def built(report_dir: Path, approvals: Path, tmp_path: Path) -> dict:
    builder = _load_builder()
    return builder.build_ppt(
        report_dir=report_dir, approvals_path=approvals, output_dir=tmp_path / "ppt"
    )



def _content_page(manifest: dict, report_key: str) -> dict:
    """取引用某 report_key 的**內容頁**；封面為了統計卡也會引用同一批 report_key，
    但它不擺報表圖也不套 narrative 標題，拿它比對會驗錯對象。"""
    return next(
        page
        for page in manifest["pages"]
        if page["kind"] not in {"cover", "direction", "section_divider"}
        and report_key in page["report_keys"]
    )

def _all_text(prs: Presentation) -> str:
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


# --------------------------------------------------------------------------
# 1. 版型庫
# --------------------------------------------------------------------------
def test_layout_library_covers_all_nine_kinds():
    """九種版型都要有 renderer，缺一種就沒辦法讓使用者挑。"""
    builder = _load_builder()
    expected = {
        "cover",
        "section_divider",
        "chart_with_points",
        "comparison",
        "stat_callout",
        "percentage_bars",
        "table",
        "table_with_points",
        "direction",
    }
    assert expected <= set(builder.RENDERERS)


def test_ppt_stage_ai_only_owns_direction_body():
    """PPT 階段 AI 只產一個 slot（P1-8：封面主標改由 workspace 名確定性組成，
    cover.title 退場）；其餘頁面文字一律來自 narratives。"""
    builder = _load_builder()
    assert builder.all_slot_keys() == ["direction.body"]


def test_paired_reports_render_side_by_side_not_merged(built):
    """IPC L4/L5 預設同頁左右並排、兩張圖都要在；機會矩陣（2026-07-31 二輪定案）
    改預設**分頁**——象限圖資訊密度高，同頁比較太小看不清。兩者都不得合成一張圖。"""
    manifest = built["manifest"]
    ipc = _content_page(manifest, "ipc_main_distribution")
    assert ipc["kind"] == "comparison"
    assert len(ipc["charts"]) == 2
    assert ipc["charts"] == ["ipc_main_distribution_L4.svg", "ipc_main_distribution_L5.svg"]

    quadrant_pages = [p for p in manifest["pages"]
                      if "opportunity_quadrant" in p["report_keys"] and not p["is_appendix"]]
    assert len(quadrant_pages) == 2, "機會矩陣應技術面／功效面各一頁"
    assert all(p["kind"] == "chart_hero" for p in quadrant_pages)
    assert all(len(p["charts"]) == 1 for p in quadrant_pages)


def test_paired_reports_can_be_split_to_separate_pages_by_override(report_dir, approvals, tmp_path):
    """使用者可經 layout_overrides 把成對報表改成分頁；分頁後兩張圖各自出頁。"""
    builder = _load_builder()
    base = builder.build_ppt(
        report_dir=report_dir, approvals_path=approvals, output_dir=tmp_path / "ppt0"
    )
    ipc_page = _content_page(base["manifest"], "ipc_main_distribution")["page"]

    approvals.write_text(
        json.dumps(
            {
                "report_version": VERSION,
                "slots": {"cover.title": "t", "direction.body": "b"},
                "layout_overrides": {str(ipc_page): "chart_with_points"},
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    result = builder.build_ppt(
        report_dir=report_dir, approvals_path=approvals, output_dir=tmp_path / "ppt1"
    )
    ipc_pages = [
        p
        for p in result["manifest"]["pages"]
        if "ipc_main_distribution" in p["report_keys"] and p["kind"] != "cover"
    ]
    assert len(ipc_pages) == 2
    assert all(p["kind"] == "chart_with_points" for p in ipc_pages)
    assert [p["charts"][0] for p in ipc_pages] == [
        "ipc_main_distribution_L4.svg",
        "ipc_main_distribution_L5.svg",
    ]


# --------------------------------------------------------------------------
# 2. 圖檔對照
# --------------------------------------------------------------------------
def test_charts_resolve_through_artifact_manifest_not_filename_guessing(built):
    """圖檔一律由 artifact_manifest 反查；猜 `{report_key}.svg` 必錯，不得出現。"""
    manifest = built["manifest"]
    trend = _content_page(manifest, "application_trend")
    assert trend["charts"] == ["annual_trend.svg"]

    used = {name for page in manifest["pages"] for name in page["charts"]}
    assert "application_trend.svg" not in used
    assert "country_distribution.svg" not in used


def test_missing_chart_degrades_to_stat_callout_without_placeholder_text(built):
    """主題分布**有 rows 就用表格**（P1-3；降級大數字卡會雙通道加總又印代碼），
    且任何頁面不得印佔位文字。"""
    manifest = built["manifest"]
    page = next(
        p
        for p in manifest["pages"]
        if p["report_keys"] == ["cluster_topic_table"] and not p["is_appendix"]
    )
    assert page["kind"] == "table_with_points"
    assert not page["degraded_from"]

    prs = Presentation(built["pptx_path"])
    text = _all_text(prs)
    for placeholder in ("圖檔待產出", "待產出", "尚未產生", "資料待補", "待確認"):
        assert placeholder not in text


def test_every_page_has_a_visual_element(built):
    """每頁都要有視覺元素（圖片、表格或色塊），不能只有一堆字。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    visual_types = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TABLE}
    prs = Presentation(built["pptx_path"])
    for index, slide in enumerate(prs.slides, start=1):
        has_visual = any(
            shape.has_table or shape.shape_type in visual_types for shape in slide.shapes
        )
        assert has_visual, f"第 {index} 頁沒有任何視覺元素（只有文字框）"


# --------------------------------------------------------------------------
# 3. narrative 契約與容錯
# --------------------------------------------------------------------------
def test_headline_becomes_interpretive_page_title(built):
    """標題＝{報表主題}：{headline}，不是裸名詞。"""
    manifest = built["manifest"]
    trend = _content_page(manifest, "application_trend")
    assert "2021–2023 為布局高峰" in trend["title"]
    assert "：" in trend["title"]


def test_points_render_as_bullets_with_emphasis(built):
    """points 逐條 bullet；emphasis 用警示紅，關鍵數字粗體。"""
    builder = _load_builder()
    theme = json.loads(THEME_PATH.read_text(encoding="utf-8"))
    alert = theme["color"]["alert"]

    prs = Presentation(built["pptx_path"])
    manifest = built["manifest"]
    trend_page = _content_page(manifest, "application_trend")
    slide = prs.slides[trend_page["page"] - 1]

    runs = [
        run
        for shape in slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    ]
    # P1-1 chart_hero：emphasis 那條升級成底部「核心結論」深藍條（白字粗體），
    # 不再以紅字混在要點裡——最重要的一句要有自己的位置。
    texts = " ".join(run.text for run in runs)
    assert "核心結論" in texts, "hero 頁缺底部核心結論條"
    assert any(run.font.bold and any(ch.isdigit() for ch in run.text) for run in runs), (
        "關鍵數字沒有粗體"
    )
    assert builder is not None


def test_legacy_narrative_without_points_falls_back_and_warns(built):
    """舊格式只有 text → fallback 條列＋截斷，並寫進 manifest warnings，不得靜默。"""
    manifest = built["manifest"]
    warnings = manifest["warnings"]
    fallback = [w for w in warnings if w["type"] == "narrative_fallback"]
    assert fallback, "narrative fallback 沒有寫進 warnings"
    assert any("ipc_main_distribution" in str(w.get("report_key", "")) for w in fallback)

    prs = Presentation(built["pptx_path"])
    ipc_page = _content_page(manifest, "ipc_main_distribution")
    slide_text = "\n".join(
        shape.text_frame.text
        for shape in prs.slides[ipc_page["page"] - 1].shapes
        if shape.has_text_frame
    )
    # fallback 後必須被切成短條，不能整段 400 字塞成一塊。
    longest_block = max((len(block) for block in slide_text.split("\n")), default=0)
    assert longest_block <= 120, f"fallback 仍是字牆（最長區塊 {longest_block} 字）"


# --------------------------------------------------------------------------
# 4. 產後 QA
# --------------------------------------------------------------------------
def test_post_build_qa_reports_no_overflow_or_overlap(built):
    """產後自檢：不得有超界、邊距不足或文字疊文字。"""
    manifest = built["manifest"]
    blocking = [
        w
        for w in manifest["warnings"]
        if w["type"] in {"out_of_bounds", "margin_violation", "text_overlap"}
    ]
    assert blocking == [], f"版面自檢不過：{blocking}"


def test_images_are_scaled_inside_their_box(built):
    """圖片一律等比縮放塞進版型框，不得溢出（舊版曾溢出 3.5 吋）。"""
    theme = json.loads(THEME_PATH.read_text(encoding="utf-8"))
    slide_w = Inches(theme["slide"]["width_in"])
    slide_h = Inches(theme["slide"]["height_in"])
    margin = Inches(theme["qa"]["min_margin_in"])

    prs = Presentation(built["pptx_path"])
    pictures = [
        shape for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13
    ]
    assert pictures, "沒有任何圖片被插入"
    for pic in pictures:
        assert pic.left >= margin - Inches(theme["qa"]["bounds_tolerance_in"])
        assert pic.top >= margin - Inches(theme["qa"]["bounds_tolerance_in"])
        assert pic.left + pic.width <= slide_w - margin + Inches(theme["qa"]["bounds_tolerance_in"])
        assert pic.top + pic.height <= slide_h - margin + Inches(theme["qa"]["bounds_tolerance_in"])


def test_qa_detects_injected_overflow(tmp_path):
    """自檢器自身的有效性：故意放一個超界形狀，必須被抓到。

    防的是「檢查器寫壞了永遠回傳空清單」造成的假性通過。
    """
    builder = _load_builder()
    from pptx import Presentation as _P
    from pptx.util import Inches as _I

    theme = builder.Theme.load(THEME_PATH)
    prs = _P()
    prs.slide_width = _I(theme.slide["width_in"])
    prs.slide_height = _I(theme.slide["height_in"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(_I(12.0), _I(1.0), _I(4.0), _I(1.0))
    box.text_frame.text = "超出右邊界"

    warnings = builder.audit_layout(prs, theme)
    assert any(w["type"] == "out_of_bounds" for w in warnings)


# --------------------------------------------------------------------------
# 5. 外觀硬規範
# --------------------------------------------------------------------------
def test_all_fonts_are_jhenghei_and_at_least_12pt(built):
    """字體全微軟正黑體（latin 與 ea 都要設），字級下限 12pt。"""
    theme = json.loads(THEME_PATH.read_text(encoding="utf-8"))
    family = theme["font"]["family"]
    min_pt = theme["font"]["min_pt"]

    prs = Presentation(built["pptx_path"])
    checked = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    checked += 1
                    assert run.font.name == family, f"字體不是 {family}：{run.font.name}"
                    assert run.font.size is not None
                    assert run.font.size.pt >= min_pt, f"字級小於下限：{run.font.size.pt}"
                    ea = run._r.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
                    ).find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
                    assert ea is not None and ea.get("typeface") == family, (
                        "沒有設 East Asian typeface，PowerPoint 會把中文退回新細明體"
                    )
    assert checked > 0


def test_accent_color_is_actually_used(built):
    """主題色 accent 必須真的用到（舊版完全沒用到）。"""
    theme = json.loads(THEME_PATH.read_text(encoding="utf-8"))
    accent = theme["color"]["accent"]

    prs = Presentation(built["pptx_path"])
    used = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.fill.type is not None and shape.fill.type == 1:
                    used.add(str(shape.fill.fore_color.rgb))
            except (AttributeError, TypeError, ValueError):
                pass
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.type is not None:
                            used.add(str(run.font.color.rgb))
    assert accent in used, f"accent {accent} 沒有被用到；實際用色 {sorted(used)}"


def test_every_content_page_has_source_footnote(built):
    """每頁底部要有資料來源註（資料來源＋統計期間）。"""
    prs = Presentation(built["pptx_path"])
    manifest = built["manifest"]
    for page in manifest["pages"]:
        if page["kind"] in {"cover", "section_divider"}:
            continue
        texts = [
            shape.text_frame.text
            for shape in prs.slides[page["page"] - 1].shapes
            if shape.has_text_frame
        ]
        joined = "\n".join(texts)
        assert "資料來源" in joined, f"第 {page['page']} 頁缺資料來源註"
        assert "統計期間" in joined, f"第 {page['page']} 頁缺統計期間"


def test_cover_has_stat_cards_and_framework_banner(built):
    """封面：主標＋期間副標＋統計卡＋分析框架條；缺就少一格，不硬湊。"""
    prs = Presentation(built["pptx_path"])
    cover = prs.slides[0]
    text = "\n".join(s.text_frame.text for s in cover.shapes if s.has_text_frame)
    assert "拉繩阻力健身器材專利佈局分析" in text
    assert "專利總數" in text
    assert "年份區間" in text
    assert "分析框架" in text

    manifest = built["manifest"]
    assert 3 <= manifest["pages"][0]["stat_cards"] <= 4


def test_manifest_carries_full_diagnostics(built):
    """manifest 要含 warnings／missing_slots／missing_reports／pages 完整資訊。"""
    manifest = built["manifest"]
    for key in ("warnings", "missing_slots", "missing_reports", "pages", "sha256"):
        assert key in manifest, f"manifest 缺 {key}"
    assert isinstance(manifest["warnings"], list)
    assert len(manifest["sha256"]) == 64
