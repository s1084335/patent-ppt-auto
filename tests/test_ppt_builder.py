"""PPT builder 基礎契約測試（v3，2026-07-30 重建）。

本檔守的是**跨版本都不能退的地基**，與版型／外觀無關（那些在
`test_ppt_layout_contract.py`）：

- 選擇驅動出頁：`report_data` 沒有資料的報表不出頁。
- 缺漏只寫 manifest，不印進 PPT（不得回到浮水印時代）。
- 封面統計卡帶年份區間。
- 同版本重跑不覆蓋既有 PPT。
- 同一張 SVG 只轉一次。
- 版型由 `PAGE_LAYOUT` 單一來源驅動，附錄靠 `is_appendix` 錨定不用頁碼魔術數字。
"""

from __future__ import annotations

import importlib.util
import json
import sys
from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation

SKILL_DIR = Path(__file__).resolve().parents[1] / "skills" / "patent-report-ppt"
BUILDER_PATH = SKILL_DIR / "scripts" / "build_ppt.py"

VERSION = "report_trial_20260723_000000"

MINIMAL_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" width="400" height="300">'
    '<rect width="400" height="300" fill="#516CEE"/></svg>'
)

# 圖檔名與 report_key 刻意不同名——這正是必須走 artifact_manifest 反查的原因。
CHART_FILES = {
    "annual_trend.svg": ["application_trend", "publication_trend"],
    "jurisdiction_distribution.svg": ["country_distribution"],
    "applicant_ranking.svg": ["applicant_ranking"],
    "owner_ranking.svg": ["owner_ranking"],
}


def _load_builder():
    """直接載入 skill 內的 deterministic PPT builder。"""
    spec = importlib.util.spec_from_file_location("build_ppt", BUILDER_PATH)
    module = importlib.util.module_from_spec(spec)
    sys.modules["build_ppt"] = module
    spec.loader.exec_module(module)
    return module


def _write_report_dir(path: Path, *, reports: dict, parameters: dict | None = None) -> Path:
    """建立 builder 測試用報表版本目錄（含 artifact_manifest）。"""
    path.mkdir(parents=True)
    payload = {
        "parameters": {
            "version": VERSION,
            "generated_at": "2026-07-23T00:00:00",
            "topic_run_id": "topic-run-001",
            "topic_state_version": "state-v3",
            **(parameters or {}),
        },
        "reports": reports,
        "family_reports": {},
    }
    (path / "report_data.json").write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    (path / "narratives.json").write_text(
        json.dumps({"based_on_version": VERSION, "reports": {}}, ensure_ascii=False), encoding="utf-8"
    )
    (path / "artifact_manifest.json").write_text(
        json.dumps(
            {
                "metadata": {"version": VERSION},
                "artifacts": [
                    {
                        "file": name,
                        "artifact_type": "chart_svg",
                        "report_names": targets,
                        "sha256": "0" * 64,
                    }
                    for name, targets in CHART_FILES.items()
                ],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    for name in CHART_FILES:
        (path / name).write_text(MINIMAL_SVG, encoding="utf-8")
    return path


@pytest.fixture()
def full_report_dir(tmp_path: Path) -> Path:
    """資料完整的報表版本。"""
    reports = {
        "application_trend": {
            "label_zh": "專利申請趨勢",
            "report_type": "aggregate",
            "rows": [
                {"application_year": 2023, "patent_count": 10},
                {"application_year": 2024, "patent_count": 25},
                {"application_year": 2025, "patent_count": 18},
            ],
        },
        "country_distribution": {
            "label_zh": "專利受理國分布",
            "report_type": "aggregate",
            "rows": [{"country_code": "CN", "patent_count": 40}, {"country_code": "US", "patent_count": 13}],
        },
        "cluster_topic_table": {
            "label_zh": "全分類技術指標總表",
            "report_type": "table",
            "rows": [{"topic": "收繩機構", "patent_count": 8}],
        },
        "applicant_ranking": {
            "label_zh": "主要申請人排名",
            "report_type": "aggregate",
            "rows": [{"applicant_display_name": "Rexon", "patent_count": 30}],
        },
        "owner_ranking": {
            "label_zh": "現專利權人排名",
            "report_type": "aggregate",
            "rows": [{"current_assignee_display_name": "Rexon", "patent_count": 22}],
        },
    }
    return _write_report_dir(tmp_path / "output" / VERSION, reports=reports)


@pytest.fixture()
def partial_approvals(tmp_path: Path) -> Path:
    """只填部分文案，用來驗證缺漏記錄但不浮水印。"""
    path = tmp_path / "approvals.json"
    path.write_text(
        json.dumps(
            {"report_version": VERSION, "slots": {"cover.title": "專利情報整合分析"}},
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    return path


def _slide_texts(prs: Presentation) -> list[str]:
    """擷取投影片文字供契約測試檢查。"""
    return [
        "\n".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)
        for slide in prs.slides
    ]


def test_manifest_metadata_traces_back_to_source_version(full_report_dir, partial_approvals, tmp_path):
    """manifest 要帶來源版本、hash 與分群追溯 metadata。"""
    builder = _load_builder()
    result = builder.build_ppt(
        report_dir=full_report_dir, approvals_path=partial_approvals, output_dir=tmp_path / "ppt"
    )

    assert Path(result["pptx_path"]).exists()
    manifest = json.loads(Path(result["manifest_path"]).read_text(encoding="utf-8"))
    assert manifest["source_report_version"] == VERSION
    assert len(manifest["sha256"]) == 64
    assert len(manifest["pages"]) == len(Presentation(result["pptx_path"]).slides._sldIdLst)
    assert manifest["metadata"]["topic_run_id"] == "topic-run-001"
    assert manifest["metadata"]["topic_state_version"] == "state-v3"


def test_missing_slots_do_not_print_watermark_but_manifest_records_them(
    full_report_dir, partial_approvals, tmp_path
):
    """缺文案只進 manifest；PPT 本體不得出現浮水印或系統狀態文字。"""
    builder = _load_builder()
    result = builder.build_ppt(
        report_dir=full_report_dir, approvals_path=partial_approvals, output_dir=tmp_path / "ppt"
    )

    prs = Presentation(result["pptx_path"])
    assert "待確認" not in "\n".join(_slide_texts(prs))

    manifest = json.loads(Path(result["manifest_path"]).read_text(encoding="utf-8"))
    assert "direction.body" in manifest["missing_slots"]
    direction = next(p for p in manifest["pages"] if p["kind"] == "direction")
    assert direction["missing_slots"] == ["direction.body"]


def test_cover_stats_include_year_range(full_report_dir, partial_approvals, tmp_path):
    """封面統計卡要顯示 application_trend 的年份區間。"""
    builder = _load_builder()
    result = builder.build_ppt(
        report_dir=full_report_dir, approvals_path=partial_approvals, output_dir=tmp_path / "ppt"
    )

    cover_text = _slide_texts(Presentation(result["pptx_path"]))[0]
    assert "2023" in cover_text
    assert "2025" in cover_text
    assert "年份區間" in cover_text


def test_unselected_or_empty_reports_do_not_create_their_pages(tmp_path, partial_approvals):
    """非 cover/direction 頁面必須由實際有資料的 report_key 驅動。"""
    builder = _load_builder()
    report_dir = _write_report_dir(
        tmp_path / "output" / VERSION,
        reports={
            "application_trend": {
                "label_zh": "專利申請趨勢",
                "report_type": "aggregate",
                "rows": [{"application_year": 2025, "patent_count": 3}],
            },
            "country_distribution": {
                "label_zh": "專利受理國分布",
                "report_type": "aggregate",
                "rows": [{"country_code": "TW", "patent_count": 3}],
            },
        },
    )

    result = builder.build_ppt(
        report_dir=report_dir, approvals_path=partial_approvals, output_dir=tmp_path / "ppt"
    )
    manifest = json.loads(Path(result["manifest_path"]).read_text(encoding="utf-8"))

    # P1-6：研發方向（結論）壓軸、附錄之前；P1-1：單圖內容頁預設 chart_hero。
    assert [p["kind"] for p in manifest["pages"]] == [
        "cover",
        "chart_hero",
        "percentage_bars",
        "direction",
    ]
    assert all("cluster_topic_table" not in p["report_keys"] for p in manifest["pages"])


def test_manifest_records_missing_report_keys_for_always_on_pages(tmp_path, partial_approvals):
    """cover/direction 恆出時，缺少的資料 key 要進 missing_reports。"""
    builder = _load_builder()
    report_dir = _write_report_dir(tmp_path / "output" / VERSION, reports={})

    result = builder.build_ppt(
        report_dir=report_dir, approvals_path=partial_approvals, output_dir=tmp_path / "ppt"
    )
    assert len(Presentation(result["pptx_path"]).slides._sldIdLst) == 2

    manifest = json.loads(Path(result["manifest_path"]).read_text(encoding="utf-8"))
    assert set(manifest["missing_reports"]) >= {"country_distribution", "application_trend"}


def test_rerun_same_version_does_not_overwrite(full_report_dir, partial_approvals, tmp_path):
    """同版本重跑要產生 _r2，不覆蓋舊 PPT。"""
    builder = _load_builder()
    out_dir = tmp_path / "ppt"

    first = Path(
        builder.build_ppt(
            report_dir=full_report_dir, approvals_path=partial_approvals, output_dir=out_dir
        )["pptx_path"]
    )
    first_bytes = first.read_bytes()
    second = Path(
        builder.build_ppt(
            report_dir=full_report_dir, approvals_path=partial_approvals, output_dir=out_dir
        )["pptx_path"]
    )

    assert second != first
    assert first.exists() and first.read_bytes() == first_bytes
    assert second.exists()


def test_page_layout_is_table_driven_with_appendix_flag():
    """版型仍由 PAGE_LAYOUT 單一來源驅動，附錄用顯式旗標而非頁碼魔術數字。"""
    builder = _load_builder()

    kinds = {spec.kind for spec in builder.PAGE_LAYOUT}
    assert {"cover", "direction"} <= kinds
    assert any(spec.is_appendix for spec in builder.PAGE_LAYOUT)
    trend = next(spec for spec in builder.PAGE_LAYOUT if "application_trend" in spec.report_keys and spec.kind != "cover")
    assert trend.topic


def test_svg_conversion_is_cached(full_report_dir, partial_approvals, tmp_path, monkeypatch):
    """同一張 SVG 不重複轉檔。"""
    builder = _load_builder()
    calls: list[Path] = []
    original = builder.rasterize_svg

    def counting(svg_path: Path, cache_dir: Path):
        calls.append(svg_path)
        return original(svg_path, cache_dir)

    monkeypatch.setattr(builder, "rasterize_svg", counting)
    builder.build_ppt(
        report_dir=full_report_dir, approvals_path=partial_approvals, output_dir=tmp_path / "ppt"
    )

    assert len(calls) == len(set(calls))


def test_approval_template_only_offers_ppt_stage_slots(tmp_path):
    """確認槽範本＝all_slot_keys()；cover.title 退場後只剩研發方向建議。"""
    builder = _load_builder()
    path = builder.write_approval_template(tmp_path / "approvals.json")
    payload = json.loads(Path(path).read_text(encoding="utf-8"))

    assert set(payload["slots"]) == {"direction.body"}
    assert "layout_overrides" in payload
