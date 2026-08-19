"""簡報版面引擎：一頁一圖、淺色系、標題 24／內文 16。

核心設計：**字級被鎖死就不能靠 autofit，只能先量文字、再配空間。**
每個文字區在畫之前先用 text_h() 估高，判讀帶依需求撐高，剩下的高度才給圖表；
圖表永遠等比縮放置中，絕不裁切。build() 會印出「可用 vs 需求」裕度表，
任何一列出現「溢出」就代表那一頁會爆版，必須改文字或改配置後重跑。

內容與版面分離：呼叫端只給 content dict（見 SKILL.md 的 schema），本檔不含任何報表內容。

用法：
    from deck_layout import build
    build(content, png_dir, out_pptx)
"""
from __future__ import annotations

import json
import math
import sys
from pathlib import Path

# 🔴 字型走 `chart_sizing` 的唯一定義處（2026-08-13 使用者裁決選項 A）。
# ⚠ 本 skill 與 backend 一起佈署（design 4-0b 環境清單），故 import 得到它；
#   判準是「改圖表字型時簡報原生文字也要跟著改」＝同一份知識，不能兩處各寫。
_REPO_ROOT = Path(__file__).resolve().parents[3]
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))
# ⚠ 別名匯入：本檔已有一個 `PALETTE`（建議卡的標籤色 cyan/blue/…，是
#   content.json 的公開契約），與色票登記表**同名不同物**。直接匯入會被遮蔽，
#   而遮蔽不會報錯——只會讓後寫的那個靜默生效。
from backend.app.reports.chart_sizing import FONT_FAMILY  # noqa: E402
from backend.app.reports.chart_sizing import PALETTE as COLOR_TOKENS  # noqa: E402

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Length, Pt
from PIL import Image

def rgb(value: str) -> RGBColor:
    """`#RRGGBB` → `RGBColor`（🔴 換算的唯一定義處）。

    ⚠ 色票以 hex 字串存（同時服務 SVG／CSS 與 PPTX），PPTX 端需要 `RGBColor`。
    這個換算每個用到的地方各寫一次的話，會在大小寫與 `#` 前綴上分岔，
    而分岔的症狀是顏色悄悄變成別的值，不是報錯。
    """
    v = value.lstrip("#")
    return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))


# ── 淺色系（非純白）：頁面淺藍灰，圖表卡白色以襯出報表原圖 ─────────
BG       = RGBColor(0xE9, 0xEE, 0xF6)
BG_PANEL = RGBColor(0xDA, 0xE4, 0xF2)
CYAN     = RGBColor(0x02, 0x6A, 0xA7)
BLUE     = RGBColor(0x43, 0x3B, 0xC4)
AMBER    = RGBColor(0xA1, 0x53, 0x07)
ROSE     = RGBColor(0xB0, 0x12, 0x3C)
GREEN    = RGBColor(0x04, 0x6B, 0x4E)
#: 🔴 頁面主文字色走色票唯一定義處（§6.3b）。
#: ⚠ 它與 `chart_runner.COLOR_TEXT`（`#00094A`）是「改一邊就得改另一邊」的
#:   同一份知識——§6.2 的媒介對照表就是在描述這個關係。兩個值都保留
#:   （使用者裁決「都留但不得同頁」），但**定義只能有一處**。
TEXT     = rgb(COLOR_TOKENS["TEXT_ON_PAGE"].hex)
MUTED    = RGBColor(0x53, 0x69, 0x8B)
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
CARD_ED  = RGBColor(0xC2, 0xD1, 0xE6)
PALETTE  = {"cyan": CYAN, "blue": BLUE, "amber": AMBER, "rose": ROSE, "green": GREEN}
TAG_COLOR = {"風險": ROSE, "機會": AMBER, "行動": GREEN, "依據": CYAN}

FONT = FONT_FAMILY
T_SIZE, B_SIZE = Pt(24), Pt(16)          # 規格鎖死：標題 24、內文 16
#: 註記小字（2026-08-18 §7d）：唯一用途是把**已在別頁講過**的數據降級掛在主體
#: 底下。⚠ 不是「塞不下就縮字」的後門——版面不夠一律拆頁，字級照舊只能量不能縮。
S_SIZE = Pt(11)
#: 🔴 字級白名單的唯一定義處。`audit_deck` 讀這一份，不自帶字面。
#: 加一級字級只改這裡，輸出前閘門自動跟上——兩處各自維護時，不一致不會報錯，
#: 只會在 audit 印出一張看不出根因的「字級分布」（2026-08-19 實際踩到）。
ALLOWED_SIZES = (T_SIZE, B_SIZE, S_SIZE)
LS = 1.22          # 設給 PowerPoint 的 line_spacing
# ⚠ **量高度不能用 LS**。LS 是「單行行距的倍數」，而中文字型的單行行距本身就大於
#   字級。2026-08-11 用 PowerPoint COM 的 TextRange.BoundHeight 量出「約 1.40 倍」，
#   當時寫成單一常數 LS_RENDER。
#
# 🔴 2026-08-13 換 Noto Sans TC 重量時發現**模型形狀是錯的**，不只是數值。
#   實測（COM，同一套方法）：
#       字級 行數  實測高   等效倍率
#        16    1   21.39   1.337
#        16    2   44.82   1.401
#        16    3   68.24   1.422
#        24    1   32.08   1.337
#        24    3  102.35   1.422
#        24    4  137.49   1.432
#   等效倍率**隨行數上升** → 固定倍率不可能同時吻合各種行數。拆開就一致了：
#   **首行 1.337、後續每行 1.464**（16pt 與 24pt 誤差 <0.1%）。
#
# ⚠ 舊值 1.40 在 2 行剛好吻合、**3 行以上一路低估**（3 行差 0.065em、4 行 0.129em）。
#   舊註記寫「段數一多就會把最後一行切掉，而且裕度表不會叫」——**根因就在這裡**：
#   不是數值不夠大，是用一個常數去逼近一條斜線。
LS_FIRST = 1.337   # 首行：含字型的 ascent＋descent，不含行距
LS_NEXT = 1.464    # 後續每行：多了段內行距


def lines_height_pt(lines: int, size_pt: float) -> float:
    """n 行文字的實際渲染高度（pt）。估高的**唯一入口**。

    ⚠ 不要在別處寫 `行數 × 字級 × 倍率`——那正是被本次修正推翻的模型。
    """
    if lines <= 0:
        return 0.0
    return size_pt * (LS_FIRST + (lines - 1) * LS_NEXT)

SW, SH = 13.333, 7.5
ML = 0.5
CW = SW - 2 * ML
CHART_TOP = 1.24
BAND_TOP, BAND_BOT = 6.02, 7.06

MIN_CHART_PT = 9.0        # 單圖頁：圖內字級低於此值視為不可讀（見 SKILL.md 優先順序）
MIN_CHART_PT_MULTI = 12.0  # 雙圖頁：門檻更高——擠成兩圖是自找的，達不到就拆成兩頁

CHART_PAD, CHART_GAP = 0.14, 0.18   # chart_stack 與預測共用，勿各寫一份

# 標籤欄式文字頁（`layout: "label"`）的幾何。
# ⚠ 存在理由是**行長**：純文字頁滿版每行約 50 個全形字，遠超中文舒適區間，
#   而字級鎖死 16pt 無法靠縮字解決，只能切掉寬度。左欄吃掉 2.2in 後，
#   內文行長降到約 40 個全形字，同時把「哪一家／哪一代」變成可掃讀的索引。
# ⚠ 不採硬分兩欄：兩欄各只剩約 23 單位，而「CN 121754861、CN 121754862、
#   CN 223248694」這類並列專利號單串就超過一行，會被迫斷行，反而更難讀。
# ⚠ ROW_PAD 不是美觀留白，是**防撞**：textbox 每段預設帶 6pt space_after，而
#   text_h 以 gap=0 估高，兩者差額會讓列間細線壓在上一列的最後一行文字上
#   （2026-08-11 實物驗收抓到）。列高一律多留這一段。
# ⚠ RULE_W 不得再細：0.008in 在 96dpi 下只有 0.77px，次像素會讓部分細線整條消失，
#   而且是隨機幾條——目視時很容易誤判成「這兩列本來就沒分隔」。
LABEL_W, LABEL_GAP, ROW_GAP, ROW_PAD, RULE_W = 2.2, 0.28, 0.14, 0.05, 0.014

_REPORT: list[tuple[str, float, float]] = []


# ── 文字量測（字級鎖死時的唯一防線） ────────────────────────────
def units(t: str) -> float:
    """字寬單位：CJK 全形算 1.0，半形算 0.55。"""
    return sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in t)


def est_lines(t: str, w_in: float, pt: int = 16) -> int:
    """換行行數。🔴 **直接數實際切出來的行**，不另外估。

    ⚠ 2026-08-13 之前這裡是獨立的數學式（`ceil(units × pt × 1.06 ÷ 寬)`）。
    加入不可分割詞組保護後兩者分岔了——保護會讓某些行提早斷，實際行數比
    純數學估的多一行，於是**估高一套、排版另一套**，版面偶爾就會溢出。
    改為委派給 `wrap_lines`：估的就是實際會切出來的東西，天然不會漂。

    安全邊界（1.06）仍在，它在 `_per_line` 裡——兩邊共用同一個容量。
    效能：一份簡報幾千字，切分成本可忽略。
    """
    return max(1, len(wrap_lines(t, w_in, pt)))


def text_h(paras, w_in: float) -> float:
    """paras = [(全文, 字級pt, space_after pt)]，回傳所需高度（英吋）。

    ⚠ 走 `lines_height_pt`，不自己乘倍率——見該函式的說明（模型是兩段式，
    不是固定倍率）。
    """
    return sum(lines_height_pt(est_lines(t, w_in, s), s) + gap
               for t, s, gap in paras) / 72


def note(region: str, avail: float, need: float):
    _REPORT.append((region, avail, need))


def _per_line(w_in: float, pt: int = 16) -> float:
    """該寬度下一行放得下幾個字寬單位（與 est_lines 同一套係數）。"""
    return w_in * 72 / (pt * 1.06)


# ── 引擎自行斷行（B 案：把排版決定權從 PowerPoint 收回） ──────────
# 行首禁則：這些字元不得出現在行首（中文排版通則）。
NO_LINE_START = "，。、；：？！）」』】》〉…—·%〞”’,.;:?!)]}"
# 行尾禁則：開括號類不得留在行尾，否則下一行開頭會孤零零一個引號。
NO_LINE_END = "（「『【《〈〝“‘([{"


def _is_token_char(char: str) -> bool:
    """算不算英數詞組的一部分（ASCII 英數與號碼常見符號）。"""
    return char.isascii() and (char.isalnum() or char in "/-.")


def _protected_ranges(text: str) -> list[tuple[int, int]]:
    """找出**不可從中間斷開**的詞組範圍 `[start, end)`。

    ⚠ 回傳「範圍」而不是「空格索引」：只記空格的話，`_inside_protected` 得再
    往兩側掃一次算邊界，那是同一份計算的第二個落點——而且我第一版就漏算了
    左側詞，導致斷點退到空格之前、詞組照樣被拆成「US」＋「 12345678」。

    🔴 2026-08-13 逐頁目視 slide05 抓到：「CN 223248696」被從空格拆開，
    行尾留下孤立的「CN」。避頭尾只管標點，管不到這個。
    ⚠ 它是「只有目視看得到」的一類——程式化檢查全綠、字寬也沒超。

    判準：空格兩側都非 CJK、且**至少一側含數字**。
    - `CN 223248696`／`A63B 069/18`／`US 12345678` → 保護
    - `11 件`／`2020 年` → 不保護（CJK 側本來就是合理斷點）
    - `the quick` → 不保護（兩側都沒數字，是一般英文）

    ⚠ 詞組邊界用「非 ASCII 英數」而不是「空白」：中英文之間沒有空格，
    「…一下，CN」往左取到空白會把整串中文吃進來，於是誤判成「含 CJK」
    而不保護（2026-08-13 實測踩到）。
    """
    ranges: list[tuple[int, int]] = []
    for index, char in enumerate(text):
        if char != " ":
            continue
        start = index
        while start > 0 and _is_token_char(text[start - 1]):
            start -= 1
        end = index + 1
        while end < len(text) and _is_token_char(text[end]):
            end += 1
        left, right = text[start:index], text[index + 1:end]
        if not left or not right:
            continue                        # 有一側不是英數詞：可斷
        if not any(c.isdigit() for c in left + right):
            continue                        # 純英文詞組：一般斷行即可
        ranges.append((start, end))
    return ranges


def wrap_lines(text: str, w_in: float, pt: int = 16) -> list[str]:
    """把一段文字切成一行一行，供 SVG 逐行絕對定位。

    🔴 **與 `est_lines` 是同一套係數**（都走 `_per_line`）。兩者若各算各的，
    就是估高一套、排版另一套——不一致不會報錯，只會讓版面偶爾溢出或留白。
    `tests/test_engine_line_breaking.py` 直接鎖住 `len(wrap_lines) == est_lines`。

    避頭尾以**回推**處理：把前一個字移到下一行，讓標點不落在行首。
    ⚠ 不用懸掛（讓標點突出右邊界）——那是排版慣例，但本 skill 是絕對定位，
    突出去會撞到右側元素（標籤欄頁的右欄就緊貼邊界）。回推的代價是該行少一個字。

    空字串回 `[""]` 而非 `[]`：下游按行數配位置，少一行會讓整段上移。
    """
    if not text:
        return [""]
    capacity = _per_line(w_in, pt)
    return _apply_kinsoku(_greedy_wrap(text, capacity), capacity)


def _greedy_wrap(text: str, capacity: float) -> list[str]:
    """貪婪斷行，遇不可分割詞組往前退。**斷行邏輯的唯一落點**。

    ⚠ `_apply_kinsoku` 推字後也需要重切，必須呼叫本函式而不是自己寫一份
    ——2026-08-13 就是因為那裡另寫了一段（沒有詞組保護），讓保護在推字後失效。
    同一份斷行知識只能有一個落點。
    """
    if not text:
        return [""]
    ranges = _protected_ranges(text)
    lines: list[str] = []
    start = 0
    used = 0.0

    for index, char in enumerate(text):
        width = 1.0 if ord(char) > 0x2E80 else 0.55
        if index > start and used + width > capacity + 1e-9:
            cut = index
            # 斷點切斷受保護詞組 → 往前退到詞組之前。
            while cut > start and _inside_protected(cut, ranges):
                cut -= 1
            # 🔴 英文**單字**內部不可斷（2026-08-14 SVG 基準目視 slide08 抓到
            #   「claim char／t」）：斷點兩側都是 ASCII 英數＝落在單字中間，
            #   回退到段首。⚠ 空格處仍可斷——保的是單字，不是片語；
            #   舊基準走 PowerPoint 斷行沒這問題，引擎自斷後才現形。
            while (start < cut < len(text)
                   and _is_token_char(text[cut]) and _is_token_char(text[cut - 1])):
                cut -= 1
            # ⚠ 退不動就照原斷點走（fail open）：詞組／單字本身超過行寬時
            #   硬保護會無限迴圈或整行溢出。
            if cut == start:
                cut = index
            lines.append(text[start:cut])
            start = cut
            used = sum(1.0 if ord(c) > 0x2E80 else 0.55 for c in text[start:index])
        used += width

    if start < len(text):
        lines.append(text[start:])
    return lines or [""]


def _token_start_at_end(line: str) -> int | None:
    """若這一行**結束在受保護詞組內**，回傳該詞組的起點索引；否則 None。

    ⚠ 只看行尾：`_apply_kinsoku` 是從行尾借字，切點永遠在那裡。
    回傳起點（而非只回傳 bool）是為了讓呼叫端能把**整個詞組**移走——
    只知道「不能借」的話，兩條規則就只能二選一。
    """
    if not line or not _is_token_char(line[-1]):
        return None
    right_start = len(line) - 1
    while right_start > 0 and _is_token_char(line[right_start - 1]):
        right_start -= 1
    if right_start == 0 or line[right_start - 1] != " ":
        return None
    left_end = right_start - 1
    left_start = left_end
    while left_start > 0 and _is_token_char(line[left_start - 1]):
        left_start -= 1
    left, right = line[left_start:left_end], line[right_start:]
    if not left or not any(c.isdigit() for c in left + right):
        return None
    return left_start


def _inside_protected(cut: int, ranges: list[tuple[int, int]]) -> bool:
    """在 `cut` 斷行會不會把某個受保護詞組切成兩半。

    ⚠ 在詞組**之前**（cut == start）或**之後**（cut == end）斷都是合法的，
    只有落在中間才算切斷。
    """
    return any(start < cut < end for start, end in ranges)


def _apply_kinsoku(lines: list[str], capacity: float) -> list[str]:
    """避頭尾修正：行首禁則字回推、行尾禁則字下推。

    ⚠ 逐行往後掃並就地修正，因為修一行會影響下一行的開頭；
    一次掃完再修會漏掉連鎖情形（實測「）。」這種連續禁則字）。
    """
    fixed = [line for line in lines]
    index = 0
    while index < len(fixed) - 1:
        current, following = fixed[index], fixed[index + 1]
        moved = False

        # 下一行以禁則字開頭 → 把本行最後一個字推過去，讓它不在行首。
        # 🔴 若行尾正在**不可分割詞組**內，改推**整個詞組**：2026-08-13 實測，
        #   借一個字會把「CN 223248696」切成「…CN 22324869」＋「6）。」
        #   ——避頭尾修好了、詞組卻被拆了。
        #   ⚠ 兩條規則本身不衝突，衝突的是「只借一個字」這個手段。
        #   推整個詞組兩者都滿足；只有整行都是詞組時才真的無解（那時停手）。
        while following and following[0] in NO_LINE_START and len(current) > 1:
            token_start = _token_start_at_end(current)
            # 🔴 英文單字（不含數字、不在詞組保護內）同樣不能只借一個字——
            #   2026-08-14 目視 slide08：「claim chart」＋行首「，」，借一字變成
            #   「claim char」＋「t，」。單字結尾就整字推（與詞組同一手段）。
            if token_start is None and _is_token_char(current[-1]):
                word_start = len(current) - 1
                while word_start > 0 and _is_token_char(current[word_start - 1]):
                    word_start -= 1
                token_start = word_start
            if token_start is not None:
                if token_start <= 0:
                    break                   # 整行都是那個詞組，推不動
                current, following = current[:token_start], current[token_start:] + following
            else:
                current, following = current[:-1], current[-1] + following
            moved = True
            # 推過去後若本行行尾又變成禁則字，下面那個迴圈會處理。

        # 本行以開括號結尾 → 把它推到下一行。
        while current and current[-1] in NO_LINE_END and len(current) > 1:
            if units(current[-1] + following) > capacity + 1e-9:
                break                       # 推不過去（下一行滿了），維持原狀
            current, following = current[:-1], current[-1] + following
            moved = True

        fixed[index], fixed[index + 1] = current, following
        if not moved:
            index += 1
            continue
        # 推字後下一行可能超寬，重切它與其後所有內容。
        # ⚠ 走 `_greedy_wrap`，**不要在這裡另寫一份**——2026-08-13 就是因為
        #   這裡有第二份斷行邏輯（沒有詞組保護），讓保護在推字後失效。
        if units(following) > capacity + 1e-9:
            fixed[index + 1:] = _greedy_wrap("".join(fixed[index + 1:]), capacity)
        index += 1
    return fixed


# 目視截圖倍率的**唯一定義處**（design 4-0c）。
# ⚠ 不得在 runner、skill 文件或規格另寫一個數字——那是「同一份知識多個落點」，
#   各自演進而不會報錯。要調就調這裡。
#
# 判準是**能力不是像素**：CLI 要能辨識行首中文標點（目視清單裡最細的一項，
# 過得了它，溢出／重疊／裁切必然過得了）。
# ⚠ 值待 tasks 2.3 實測定案：造一頁刻意植入行首標點的樣本，由低到高試，
#   取「CLI 能穩定指出該問題」的**最小值**。⚠ 不是越大越好——目視迴圈每輪都要
#   讀一整份，CLI 讀圖有 token 成本。
# 暫定 2：1280×720 的兩倍＝2560×1440，內文 16pt 在圖上約 43px。
VISUAL_SCALE = 2.0


def visual_shot_size() -> tuple[int, int]:
    """目視截圖的像素尺寸——由版面尺寸與倍率推導，呼叫端不自己算。

    ⚠ 用 `round` 不用 `int`：`SW = 13.333` 本身是 13⅓ 的截斷值
    （精確值 ×96 = 1280），再截一次會變 2559，截圖就比版面少 1px。
    """
    from svg_canvas import DPI

    return round(SW * DPI * VISUAL_SCALE), round(SH * DPI * VISUAL_SCALE)


def _text_page_lines() -> int:
    """純文字頁的內文區放得下幾個顯示行——由版面幾何算出，不寫死。"""
    avail_pt = (BAND_BOT - CHART_TOP - 0.44) * 72
    # 每則是獨立段落（各自首行），故用單行高＋8pt 段距。
    return int(avail_pt // (lines_height_pt(1, 16) + 8))


def budget() -> dict[str, float]:
    """各欄位的字寬單位上限——由版面幾何推導，**不要在別處另寫一份**。

    撰稿時先用 units() 量，超過就重寫；等到 build() 才發現溢出要重跑整條流程。
    """
    gw = (CW - 0.22) / 2
    rw = (CW - 0.44) / 3
    return {
        "deck_title":   _per_line(8.2, 24),
        "cover_line":   _per_line(8.1),          # subtitle / meta / stats_note
        "boundary":     _per_line(8.1) * 2 - 3,
        "cover_panel":  _per_line(3.18) * 5,     # 舊 schema 封面說明相容預算
        "stat_label":   _per_line(1.6),
        "page_title":   _per_line(CW - 0.4, 24),
        "takeaway":     _per_line(CW - 0.1) - 1,
        # 判讀帶會依內容自動撐高，代價是圖表變小：每多一行約吃掉 0.27in 高度、
        # 圖內字級掉約 0.6pt。2 行是建議值，4 行是還能接受的上限。
        "band_total":   _per_line(CW - 0.52) * 2,
        "band_max":     _per_line(CW - 0.52) * 4,
        "rec_title":    _per_line(gw - 0.52),
        "rec_line":     _per_line(gw - 0.52) * 2,
        # 純文字頁（charts 留空）：由實際幾何推出放得下幾行，**不要寫死倍數**。
        # ⚠ 曾經寫死 14，改用實測行高（LS_RENDER）後實際只放得下 12 行，
        #   上限變成高估——寫死的倍數不會跟著幾何調整走（2026-08-11）。
        "text_page":    _per_line(CW - 0.52) * _text_page_lines(),
        "roadmap_item": _per_line(rw - 0.48) * 3,
        "limit_line":   _per_line(CW - 0.52),
    }


# ── 繪圖低階工具 ─────────────────────────────────────────────────
def _set_font(run, *, size, bold=False, color=TEXT):
    # ⚠ 字級單位在兩條渲染路徑各有解讀：SVG 端讀 pt **數值**，PPTX 端要
    #   `Length` 物件。裸 int 在 SVG 完全正常，到 PPTX 會被當成 EMU
    #   （11 EMU → 0 centipoints → ValueError），而且只有真的產 .pptx 才炸
    #   ——deck 側單元測試全綠。2026-08-19 被半真機械鏈抓到。
    #   在這個唯一落點把 int 視為 pt，兩邊才對得起來；靠「記得寫 Pt()」是人治。
    #   ⚠ 判斷要用 `Length` 不是 `int`：`Pt(16)` **是** int 的子類別，
    #     用 `isinstance(size, int)` 會把已經換算好的值再換算一次
    #     （Pt(16) → 20320000 centipoints，同一個 ValueError 換個數字）。
    f = run.font
    f.size = size if isinstance(size, Length) else Pt(size)
    f.bold, f.name = bold, FONT
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("ea", "cs"):     # 中文必須另外指定 east-asian typeface
        rPr.append(rPr.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag,
            {"typeface": FONT}))


def _svg_canvas(w_in: float = SW, h_in: float = SH):
    """建一張 SVG 畫布，注入本模組的字型／行高／量測函式。

    ⚠ 延遲 import：`svg_canvas` 不得反向 import 本模組（會循環），
    它需要的東西一律由這裡注入——字型與行高的唯一定義處在本檔。
    """
    from svg_canvas import SvgCanvas

    return SvgCanvas(w_in, h_in, font=FONT, ls_first=LS_FIRST, ls_next=LS_NEXT,
                     unit_width=units, wrap_lines=wrap_lines)


def _is_svg(sl) -> bool:
    """畫布是 SVG 還是 pptx slide。

    用有沒有 `to_svg` 判斷，不用 isinstance——避免本模組為了型別檢查而 import
    `svg_canvas`（那會製造循環相依，而且只為了一個判斷不值得）。
    """
    return hasattr(sl, "to_svg")


def rect(sl, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    if _is_svg(sl):
        sl.rect(x, y, w, h, fill=fill, line=line,
                radius=radius if shape == MSO_SHAPE.ROUNDED_RECTANGLE else None)
        return None
    s = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    return s


def picture(sl, path, x, y, w, h):
    """貼圖原語。⚠ 原本各頁型直接呼叫 `s.shapes.add_picture`——包成原語才能分派。

    SVG 端寫**相對路徑**（`Path.name`）：圖檔與 SVG 同目錄，Chromium 用 `goto`
    載入時才抓得到（2026-08-13 實測 `file://` 絕對 URI ＋ `set_content` 會破圖）。
    """
    if _is_svg(sl):
        sl.picture(Path(path).name, x, y, w, h)
        return None
    return sl.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def textbox(sl, x, y, w, h, blocks, *, anchor=MSO_ANCHOR.TOP, space_after=6):
    if _is_svg(sl):
        sl.text_block(x, y, w, h, blocks,
                      anchor_middle=(anchor == MSO_ANCHOR.MIDDLE),
                      space_after=space_after,
                      default_size_pt=B_SIZE.pt, default_color=TEXT)
        return None
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, (txt, opt) in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opt.get("align", PP_ALIGN.LEFT)
        p.line_spacing = LS
        p.space_after = Pt(opt.get("space_after", space_after))
        p.space_before = Pt(opt.get("space_before", 0))
        # ⚠ 中文禁則處理（避頭尾）：python-pptx 產的段落沒有這兩個屬性，PowerPoint
        #   就不套東亞換行規則，句號／逗號／右括號會單獨掉到下一行的行首。
        #   2026-08-11 實測：關閉時第二行開頭是「：風磁複合」，開啟後「：」掛在行尾。
        #   ⚠ 這個瑕疵在 v1–v16 每一版都存在，逐頁目視了十幾輪都沒看出來——
        #   它不是溢出也不是重疊，眼睛會自動把標點讀回上一行。
        pPr = p._p.get_or_add_pPr()
        pPr.set("eaLnBrk", "1")        # 套用東亞換行規則
        pPr.set("hangingPunct", "1")   # 行尾標點外掛（不佔行寬）
        for seg in (txt if isinstance(txt, list) else [(txt, {})]):
            s_txt, s_opt = seg if isinstance(seg, tuple) else (seg, {})
            o = {**opt, **s_opt}
            r = p.add_run()
            r.text = s_txt
            _set_font(r, size=o.get("size", B_SIZE), bold=o.get("bold", False),
                      color=o.get("color", TEXT))
            # 沒有 lang，PowerPoint 會拿拉丁規則斷行，上面兩個屬性等於白設
            r._r.get_or_add_rPr().set("lang", "zh-TW")
    return tb


def src_line(c: dict, spec: dict | None = None) -> str | None:
    """來源行（design §7.1，機械）：「資料來源：<version>／<report_key>」。

    `_source_version` 由 make_deck 自 report.json 蓋章（CLI 不參與，
    content.json 自帶也會被蓋掉）；沒蓋章（開發側直跑舊素材）就不印。
    ⚠ 顯示時去掉恆定前綴 `report_trial_`——固定字首去掉不損回溯性，
    留著來源行 4in 起跳會擠壓 footer。
    """
    version = c.get("_source_version")
    if not version:
        return None
    shown = str(version).removeprefix("report_trial_")
    keys = "、".join(spec.get("charts") or []) if spec else ""
    return f"資料來源：{shown}／{keys}" if keys else f"資料來源：{shown}"


def base(prs, footer, page=None, source=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, fill=BG)
    rect(s, 0, 0, SW, 0.055, fill=CYAN)
    rect(s, 0, SH - 0.34, SW, 0.34, fill=BG_PANEL)
    if page is not None:
        textbox(s, SW - 1.5, SH - 0.29, 1.0, 0.26,
                [("%02d" % page, {"color": MUTED, "align": PP_ALIGN.RIGHT})])
        textbox(s, ML, SH - 0.29, 8.5, 0.26, [(footer, {"color": MUTED})])
    if source:
        # 右對齊、頁碼左側。⚠ 與 footer 左框共享 band：footer 文字超過 ~5in
        # 才會撞（現行 footer 約 3.4in）；來源行去前綴後約 4.3in 內。
        y = SH - 0.29
        textbox(s, ML + 5.2, y, SW - 1.6 - (ML + 5.2), 0.26,
                [(source, {"color": MUTED, "align": PP_ALIGN.RIGHT})])
    return s


def header(s, title, takeaway):
    """標題 24pt ＋ 下方一行結論句（不用色塊，把高度讓給圖表）。"""
    rect(s, ML, 0.26, 0.085, 0.44, fill=CYAN)
    textbox(s, ML + 0.24, 0.24, CW - 0.4, 0.52, [(title, {"size": T_SIZE, "bold": True})])
    note(f"結論句 {title[:8]}", 1.0, text_h([("▍" + takeaway, 16, 0)], CW - 0.1))
    textbox(s, ML, 0.82, CW, 0.34, [([("▍", {"color": AMBER}), (takeaway, {})], {})])


def stack_width(inv: list[float], h: float) -> float:
    """堆疊 N 張圖時的共同寬度（英吋）。chart_stack 與 predict_chart_pt 共用同一式。"""
    n = len(inv)
    return min(CW, (h - 2 * CHART_PAD * n - CHART_GAP * (n - 1)) / sum(inv))


def band_height(lines, tag) -> float:
    """判讀帶會撐到多高——決定圖表能拿到多少高度。"""
    full = ((tag + "｜") if tag else "") + " ".join(lines)
    return max(BAND_BOT - BAND_TOP, text_h([(full, 16, 0)], CW - 0.52) + 0.34)


def predict_chart_pt(sizes_px, fonts, lines, tag) -> list[float]:
    """組版前預測每張圖在投影片上的字級。

    sizes_px = [(png 寬, png 高)]（PNG 為 3 倍解析），fonts = 各圖的圖內字級。
    讓 check_content.py 能在還沒排版時就擋掉「兩張圖擠一頁會看不清楚」的配置。
    """
    inv = [hh / ww for ww, hh in sizes_px]
    h = (BAND_BOT - band_height(lines, tag)) - 0.10 - CHART_TOP
    w = stack_width(inv, h)
    return [f * (w * 72) / (ww / 3) for f, (ww, _hh) in zip(fonts, sizes_px)]


def chart_stack(s, pngs, y, h):
    """圖表等寬堆疊置中；共同寬度取「不超出區塊高度」的最大值，永不裁切。"""
    pad, gap = CHART_PAD, CHART_GAP
    inv = [Image.open(p).size[1] / Image.open(p).size[0] for p in pngs]
    w = stack_width(inv, h)
    total = sum(w * r + 2 * pad for r in inv) + gap * (len(pngs) - 1)
    # 圖很扁時（占不到區塊 6 成）改成靠下，讓圖與下方判讀帶成為一組，
    # 留白集中到標題下方；置中會變成上下各一塊空白，版面像沒收好。
    cy = y + (h - total) * (1.0 if total < h * 0.6 else 0.5)
    for p, r in zip(pngs, inv):
        ih = w * r
        rect(s, ML + (CW - w) / 2 - pad, cy, w + 2 * pad, ih + 2 * pad,
             fill=CARD, line=CARD_ED, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.04)
        picture(s, p, ML + (CW - w) / 2, cy + pad, w, ih)
        cy += ih + 2 * pad + gap
    return w


def band(s, page, lines, tag):
    """底部判讀帶：先量文字要多高，再決定帶高，剩下的全給圖表。"""
    full = ((tag + "｜") if tag else "") + " ".join(lines)
    need = text_h([(full, 16, 0)], CW - 0.52)
    h = max(BAND_BOT - BAND_TOP, need + 0.34)
    top = BAND_BOT - h
    rect(s, ML, top, CW, h, fill=BG_PANEL, line=CARD_ED,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    # ⚠ 同樣不要查「h - 0.34 vs need」——判讀帶一旦撐高，h 就是 need + 0.34，
    #   avail 恆等於 need。判讀帶會**吃掉圖表高度**，所以真正該查的是
    #   「它有沒有吃光整個圖表區」（字級是否仍可讀由 build() 另外逐張報）。
    note(f"判讀帶 P{page}", BAND_BOT - CHART_TOP - 0.10, h)
    segs = ([(tag + "｜", {"color": TAG_COLOR[tag], "bold": True})] if tag else [])
    segs.append((" ".join(lines), {}))
    textbox(s, ML + 0.26, top + 0.17, CW - 0.52, h - 0.34, [(segs, {})],
            anchor=MSO_ANCHOR.MIDDLE)
    return top


# ── 頁型 ─────────────────────────────────────────────────────────
def page_vars(c: dict) -> dict:
    """封面可用的頁碼變數——用 {chart_first} 這類佔位符，別手寫頁碼。

    ⚠ 手寫過一次就會忘了改：曾經寫「第 3–14 頁一頁一圖；第 15 頁收斂為行動」，
       但實際圖表頁是 3–15、路線圖在 16，封面直接錯一頁。
    """
    n = len(c["pages"])
    return {"chart_first": 3, "chart_last": n + 2, "last": n + 3, "total": n + 3,
            "n_charts": sum(len(p.get("charts") or []) for p in c["pages"])}


def slide_cover(prs, c):
    s = base(prs, c["footer"], source=src_line(c))
    V = page_vars(c)
    cover_panels = [(c[k], y0, col) for k, y0, col in
                    (("read_me", 1.78, CYAN), ("chart_rule", 4.22, AMBER))
                    if k in c]
    if cover_panels:
        rect(s, 9.05, 1.15, 3.78, 5.25, fill=BG_PANEL, line=CARD_ED,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        for _panel, y0, col in cover_panels:
            rect(s, 9.35, y0 - 0.26, 3.18, 0.05, fill=col)
    rect(s, ML, 1.32, 0.09, 1.55, fill=CYAN)
    textbox(s, ML + 0.28, 1.28, 8.2, 0.4, [(c["eyebrow"], {"color": CYAN})])
    textbox(s, ML + 0.28, 1.74, 8.2, 1.0,
            [(c["deck_title"], {"size": T_SIZE, "bold": True}),
             (c["subtitle"], {"color": MUTED, "space_before": 6})])
    textbox(s, ML + 0.28, 3.00, 8.1, 0.9,
            [(c["meta"][0], {}), (c["meta"][1], {"color": MUTED, "space_before": 4})])
    for i, (num, label) in enumerate(c["stats"]):
        x = ML + i * 2.07
        rect(s, x, 4.20, 1.92, 1.02, fill=CARD, line=CARD_ED,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
        rect(s, x, 4.20, 1.92, 0.05, fill=CYAN if i % 2 == 0 else BLUE)
        textbox(s, x + 0.18, 4.38, 1.6, 0.42,
                [(num, {"size": T_SIZE, "bold": True, "color": CYAN})])
        textbox(s, x + 0.18, 4.84, 1.6, 0.3, [(label, {"color": MUTED})])
    textbox(s, ML, 5.42, 8.1, 0.35, [(c["stats_note"], {"color": MUTED})])
    textbox(s, ML, 6.00, 8.1, 0.9,
            [([("邊界　", {"color": ROSE, "bold": True}),
               (c["boundary"], {"color": MUTED})], {})])
    for (head, body), y0, col in cover_panels:
        textbox(s, 9.35, y0, 3.18, 0.35, [(head, {"bold": True, "color": col})])
        textbox(s, 9.35, y0 + 0.42, 3.18, 1.6,
                [(body.format(**V), {"color": MUTED if col is AMBER else TEXT})])
    return s


def slide_rec(prs, c):
    s = base(prs, c["footer"], 2, source=src_line(c))
    header(s, c["rec_title"], c["rec_takeaway"])
    top, bot = CHART_TOP, 7.06
    gw = (CW - 0.22) / 2
    tw_ = gw - 0.52
    gh = max(text_h([(r["title"] + r["tag"], 16, 10)] + [(t, 16, 8) for t in r["lines"]], tw_)
             for r in c["recommendations"]) + 0.36
    # ⚠ 不要查「單張卡的 gh - 0.26 vs 該卡 need」——gh 是由**最高那張卡的 need**
    #   加固定 padding 算出來的，avail 恆 ≥ need，那個檢查永遠不會失敗
    #   （與 slide_roadmap 同一個坑，2026-08-11）。真正的約束是兩排卡的總高。
    avail_total, need_total = bot - top, 2 * gh + 0.18
    note("建議頁總高", avail_total, need_total)
    top += max(0.0, (avail_total - need_total) / 2)
    for i, r in enumerate(c["recommendations"]):
        x = ML + (i % 2) * (gw + 0.22)
        y = top + (i // 2) * (gh + 0.18)
        color = PALETTE[r["color"]]
        rect(s, x, y, gw, gh, fill=BG_PANEL, line=CARD_ED,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        rect(s, x, y, 0.065, gh, fill=color)
        blocks = [([(r["title"], {"bold": True, "color": color}),
                    ("　" + r["tag"], {"color": MUTED})], {"space_after": 10})]
        blocks += [([("▍", {"color": color}), (t, {})], {"space_after": 8}) for t in r["lines"]]
        textbox(s, x + 0.26, y + 0.13, tw_, gh - 0.26, blocks)
    return s


def slide_chart(prs, c, page, spec, png_dir):
    s = base(prs, c["footer"], page, source=src_line(c, spec))
    header(s, spec["title"], spec["takeaway"])
    full = ((spec.get("tag") + "｜") if spec.get("tag") else "") + " ".join(spec["lines"])
    bh = max(BAND_BOT - BAND_TOP, text_h([(full, 16, 0)], CW - 0.52) + 0.34)
    w = chart_stack(s, [png_dir / f"{n}.png" for n in spec["charts"]],
                    CHART_TOP, (BAND_BOT - bh) - 0.10 - CHART_TOP)
    band(s, page, spec["lines"], spec.get("tag"))
    return w


def _label_rows(lines):
    """把 `主體｜內容` 拆成（標籤, 內文, 需求高）；沒有分隔號的整列橫跨全寬。

    橫跨全寬的用途是「不屬於任何一個主體」的列（警語、小結）——強迫它們也擠進
    右欄只會讓左欄出現一堆空格，讀者反而要判斷那格為什麼是空的。
    """
    tw_full = CW - 0.52
    tw_body = tw_full - LABEL_W - LABEL_GAP
    rows = []
    for t in lines:
        head, sep, rest = t.partition("｜")
        if sep:
            # 左右欄各自換行，列高取兩者較高者，否則長邊會被截掉。
            h = max(text_h([("▍" + head, 16, 0)], LABEL_W),
                    text_h([(rest, 16, 0)], tw_body))
            rows.append((head, rest, h + ROW_PAD))
        else:
            rows.append((None, t, text_h([("▍" + t, 16, 0)], tw_full) + ROW_PAD))
    return rows, tw_full, tw_body


def label_page_fit(lines) -> dict:
    """標籤欄文字頁的事前量測，供 check_content.py 使用。

    ⚠ 走的是**與組版完全相同**的 `_label_rows`——不要在 check_content 另寫一份估算，
    兩份幾何一分岔就會出現「檢查過了但組版爆版」。
    """
    rows, _tw_full, _tw_body = _label_rows(lines)
    cap = _per_line(LABEL_W)
    return {
        "need": sum(h for *_, h in rows) + ROW_GAP * (len(rows) - 1),
        "avail": BAND_BOT - CHART_TOP - 0.44,
        "label_cap": cap,
        # 左欄放不下的標籤：會換行擠掉右欄對齊，必須縮短
        "over_labels": [(head, units("▍" + head)) for head, _b, _h in rows
                        if head is not None and units("▍" + head) > cap],
        "n_labelled": sum(1 for head, _b, _h in rows if head is not None),
    }


def roadmap_page_overflow(c: dict):
    """路線圖頁（卡片＋限制框）會不會撐爆整頁；不會就回 None。

    ⚠ 與 `slide_roadmap` 用**同一組式子**，供 check_content 事前擋——否則秒級閘門
    放行、組版才擋，使用者白等一輪。回傳 (可用, 需求, 卡片單行單位)。
    """
    rw = (CW - 0.44) / 3
    tw = rw - 0.48
    card_need = max(text_h([(r["label"], 16, 10)] + [(t, 16, 8) for t in r["items"]], tw)
                    for r in c["roadmap"])
    lim_need = text_h([(c.get("limits_title", "分析限制與適用邊界"), 16, 6)]
                      + [(t, 16, 4) for t in c["limits"]], CW - 0.52)
    need = (card_need + 0.50) + (lim_need + 0.34) + 0.20
    avail = BAND_BOT - CHART_TOP
    return None if need <= avail else (avail, need, _per_line(tw))


def _slide_text_label(s, page, spec, top, bot):
    """左標籤欄＋右內文欄的純文字頁（見 LABEL_W 檔頭說明）。"""
    rows, tw_full, tw_body = _label_rows(spec["lines"])
    need = sum(h for _, _, h in rows) + ROW_GAP * (len(rows) - 1)
    avail = bot - top - 0.44
    note(f"文字頁 P{page}（標籤欄）", avail, need)
    rect(s, ML, top, CW, bot - top, fill=BG_PANEL, line=CARD_ED,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    x0 = ML + 0.26
    y = top + 0.22 + max(0.0, (avail - need) / 2)
    for i, (head, body, h) in enumerate(rows):
        if i:      # 列間細線：有了左欄之後，橫向對位需要一條基準才掃得動
            rect(s, x0, y - ROW_GAP / 2, tw_full, RULE_W, fill=CARD_ED)
        if head is None:
            textbox(s, x0, y, tw_full, h,
                    [([("▍", {"color": MUTED}), (body, {})], {})], space_after=0)
        else:
            textbox(s, x0, y, LABEL_W, h,
                    [([("▍", {"color": CYAN}),
                       (head, {"bold": True, "color": CYAN})], {})], space_after=0)
            textbox(s, x0 + LABEL_W + LABEL_GAP, y, tw_body, h, [(body, {})],
                    space_after=0)
        y += h + ROW_GAP
    return s


# ── 圖形文法（design §7.4）：參數化元件，組合活、渲染死 ──────────
class FigureGrammarError(ValueError):
    """figure 宣告不在文法內。🔴 fail loud——靜默略過會讓版面少一塊而沒人發現，
    也等於偷偷開了「CLI 自由畫圖」的後門（§7.4 明令第一版不開）。"""


#: 各型節點數上限（容量閘門；check_content 引用，唯一定義處）。
FIGURE_MAX_NODES = {"flow": 6, "cycle": 6, "contrast": 8,
                    "hierarchy": 7, "parallel": 6, "timeline": 8}
#: 單節點文字上限（字寬單位）：固定卡寬放得下、不換第三行。
FIGURE_NODE_UNITS = 16.0
#: 插圖區高度（in）。固定值：figure 與文字共享頁面，量測交由裕度表。
FIG_H = 3.0


def _fig_card(s, x, y, w, h, text, color=None):
    """節點卡：圓角矩形＋置中文字（渲染死——CLI 管不到樣式）。"""
    rect(s, x, y, w, h, fill=CARD, line=color or CARD_ED,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
    textbox(s, x + 0.08, y, w - 0.16, h, [(text, {"align": PP_ALIGN.CENTER})],
            anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def _fig_glyph(s, x, y, w, h, glyph, color):
    textbox(s, x, y, w, h, [(glyph, {"bold": True, "color": color,
                                     "align": PP_ALIGN.CENTER})],
            anchor=MSO_ANCHOR.MIDDLE, space_after=0)


def draw_figure(s, fig: dict, x: float, y: float, w: float, h: float):
    """六型分派。⚠ 箭頭一律用字符（→ ↓ vs）不用多邊形——維持窄詞彙
    （rect／text／image），SVG 端才不用擴 `svg_to_pptx` 的元素表。"""
    ftype = str(fig.get("type") or "")
    nodes = [str(n) for n in fig.get("nodes") or []]
    if ftype not in FIGURE_MAX_NODES:
        raise FigureGrammarError(
            f"figure type {ftype!r} 不在文法內（可用：{sorted(FIGURE_MAX_NODES)}）"
            "——文法不足時以實例擴元件型，走 openspec 留痕，不得自由畫圖")
    if not nodes:
        raise FigureGrammarError("figure 沒有任何節點")
    {"flow": _fig_flow, "cycle": _fig_cycle, "contrast": _fig_contrast,
     "hierarchy": _fig_hierarchy, "parallel": _fig_parallel,
     "timeline": _fig_timeline}[ftype](s, nodes, x, y, w, h)


def _fig_flow(s, nodes, x, y, w, h):
    """流程：橫排卡片＋→。"""
    n = len(nodes)
    arrow_w = 0.45
    card_w = (w - arrow_w * (n - 1)) / n
    card_h = min(0.9, h - 0.4)
    cy = y + (h - card_h) / 2
    for i, t in enumerate(nodes):
        cx = x + i * (card_w + arrow_w)
        _fig_card(s, cx, cy, card_w, card_h, t, CYAN)
        if i < len(nodes) - 1:
            _fig_glyph(s, cx + card_w, cy, arrow_w, card_h, "→", CYAN)


def _fig_cycle(s, nodes, x, y, w, h):
    """循環：卡片繞橢圓排列，中央標「循環」。方向感由順時針排列承擔
    ——旋轉箭頭需要 transform，不在窄詞彙內。"""
    import math as _m

    n = len(nodes)
    card_w, card_h = min(2.4, w / 3), 0.62
    rx_ = (w - card_w) / 2
    ry_ = (h - card_h) / 2
    cx0, cy0 = x + w / 2, y + h / 2
    for i, t in enumerate(nodes):
        ang = -_m.pi / 2 + 2 * _m.pi * i / n     # 從頂端開始順時針
        cx = cx0 + rx_ * _m.cos(ang) - card_w / 2
        cy = cy0 + ry_ * _m.sin(ang) - card_h / 2
        _fig_card(s, cx, cy, card_w, card_h, t, BLUE)
    _fig_glyph(s, cx0 - 0.6, cy0 - 0.2, 1.2, 0.4, "循環", MUTED)


def _fig_contrast(s, nodes, x, y, w, h):
    """對比：左右兩欄（前半 vs 後半），中央 vs。"""
    half = (len(nodes) + 1) // 2
    panel_w = (w - 0.7) / 2
    for side, (items, color) in enumerate(
            ((nodes[:half], CYAN), (nodes[half:], AMBER))):
        px = x + side * (panel_w + 0.7)
        rect(s, px, y, panel_w, h - 0.2, fill=BG_PANEL, line=color,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        blocks = [([("▍", {"color": color}), (t, {})], {"space_after": 8})
                  for t in items]
        if blocks:
            textbox(s, px + 0.2, y + 0.15, panel_w - 0.4, h - 0.5, blocks,
                    anchor=MSO_ANCHOR.MIDDLE)
    _fig_glyph(s, x + panel_w, y, 0.7, h - 0.2, "vs", ROSE)


def _fig_hierarchy(s, nodes, x, y, w, h):
    """階層：首節點為根置頂，其餘一排在下；細 rect 畫幹與橫軌。"""
    root, children = nodes[0], nodes[1:]
    card_h = 0.62
    root_w = min(3.2, w / 2)
    _fig_card(s, x + (w - root_w) / 2, y, root_w, card_h, root, CYAN)
    if not children:
        return
    n = len(children)
    card_w = min(2.6, (w - 0.3 * (n - 1)) / n)
    total = card_w * n + 0.3 * (n - 1)
    x0 = x + (w - total) / 2
    rail_y = y + card_h + (h - 2 * card_h) / 2
    rect(s, x + w / 2 - RULE_W / 2, y + card_h, RULE_W, rail_y - y - card_h, fill=CARD_ED)
    rect(s, x0 + card_w / 2, rail_y, total - card_w, RULE_W, fill=CARD_ED)
    for i, t in enumerate(children):
        cx = x0 + i * (card_w + 0.3)
        rect(s, cx + card_w / 2 - RULE_W / 2, rail_y, RULE_W,
             y + h - card_h - rail_y, fill=CARD_ED)
        _fig_card(s, cx, y + h - card_h, card_w, card_h, t, BLUE)


def _fig_parallel(s, nodes, x, y, w, h):
    """並列：等寬卡片一排，無連接——並列本身就是語意。"""
    n = len(nodes)
    card_w = (w - 0.3 * (n - 1)) / n
    card_h = min(1.0, h - 0.4)
    cy = y + (h - card_h) / 2
    for i, t in enumerate(nodes):
        _fig_card(s, x + i * (card_w + 0.3), cy, card_w, card_h, t, GREEN)


def _fig_timeline(s, nodes, x, y, w, h):
    """時間線：橫軌＋節點方標，標籤上下交錯（避免相鄰互撞）。"""
    n = len(nodes)
    rail_y = y + h / 2
    rect(s, x, rail_y - RULE_W / 2, w, RULE_W * 2, fill=CARD_ED)
    step = w / n
    for i, t in enumerate(nodes):
        cx = x + step * i + step / 2
        rect(s, cx - 0.05, rail_y - 0.05, 0.1, 0.1, fill=CYAN)
        label_y = rail_y - 0.75 if i % 2 == 0 else rail_y + 0.18
        textbox(s, cx - step / 2 + 0.05, label_y, step - 0.1, 0.55,
                [(t, {"align": PP_ALIGN.CENTER})], anchor=MSO_ANCHOR.MIDDLE,
                space_after=0)


def slide_text(prs, c, page, spec):
    """純文字頁：`charts` 留空時使用（中場小結、方法說明、附錄等）。

    頁數不受圖表數限制——需要更多頁就多加這種頁；它一樣走量測，超出會被裕度表抓到。
    `layout: "label"` 改走左標籤欄版面（`主體｜內容` 的定義列表型內容適用）。
    """
    s = base(prs, c["footer"], page, source=src_line(c, spec))
    header(s, spec["title"], spec["takeaway"])
    top, bot = CHART_TOP, BAND_BOT
    if spec.get("layout") == "label":
        return _slide_text_label(s, page, spec, top, bot)
    # 圖形文法（design §7.4）：頁 spec 宣告 figure 時，上半渲染插圖、文字下移。
    if spec.get("figure"):
        draw_figure(s, spec["figure"], ML + 0.26, top + 0.1, CW - 0.52, FIG_H)
        top += FIG_H + 0.24
    tw_ = CW - 0.52
    need = text_h([(t, 16, 8) for t in spec["lines"]], tw_)
    note(f"文字頁 P{page}", bot - top - 0.44, need)
    rect(s, ML, top, CW, bot - top, fill=BG_PANEL, line=CARD_ED,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
    # 「主體｜內容」寫法會把主體加粗上色，做成可掃讀的定義列表（競爭者技術構型頁用）
    blocks = []
    for t in spec["lines"]:
        head, sep, rest = t.partition("｜")
        blocks.append((
            [("▍", {"color": CYAN})] +
            ([(head + sep, {"bold": True, "color": CYAN}), (rest, {})] if sep
             else [(t, {})]),
            {"space_after": 8}))
    textbox(s, ML + 0.26, top + 0.22, tw_, bot - top - 0.44, blocks,
            anchor=MSO_ANCHOR.MIDDLE)
    return s


#: 🔴 內容頁版型的唯一定義處（2026-08-18，§7a；比照下方 `ACTION_VERBS`）。
#:
#: 為什麼要這份清單：原本「能畫的」（本檔）、「會擋的」（check_content）、
#: 「CLI 照抄的」（content-template.json）各自維護，於是 `conclusions` 有畫法、
#: 有閘門，**範本裡卻沒有**——CLI 不宣告就靜默不產出那一頁，沒有人會發現。
#:
#: ⚠ 三處必須涵蓋本清單全部鍵（`tests/test_deck_layout_registry.py` 擋著）。
#: 新增版型時漏在任何一處都會紅，而不是等實機發現 CLI 從來不用它。
#:
#: ⚠ 值是**用途說明**不是裝飾：沒寫用途，CLI 不知道何時該用，結果就是永遠不用
#: ——那是缺席型偏差，目視兜不住。
LAYOUTS: dict[str, str] = {
    "chart": "有圖：`charts` 放一張 PNG，配 takeaway 與 lines。圖能表達趨勢或分布時用。",
    "text": "無圖純文字頁：`charts` 留空。適合結構化敘述與逐點論證。",
    "label": "無圖的標題卡片式：`lines` 每條是一張卡，適合並列數個對象的短敘述。",
    "table": "表格頁：`table.headers` ＋ `table.rows`（每列一個字串陣列）。"
             "多欄對照（逐家時序）、矩陣（主題×年）這類用它，不要硬塞成圖。"
             "欄寬由欄數自動算，可用 `table.weights` 指定相對寬。",
}

#: 🔴 結論頁「專利行動」的有限動詞表（design §7.7，唯一定義處；check_content 引用）。
#: CLI 只能選不能自創——避免 AI 寫出「應立即提出申請」這類法律／商業承諾，
#: 越過「AI 只輔助、不做正式決策」的界線。擴充動詞走 openspec 留痕，不在 prompt 放寬。
ACTION_VERBS = ("佈局", "追蹤", "迴避設計", "細讀比對", "暫不投入")

#: 表格版型的水平內距（左右各一），沿用結論頁既有值。
TABLE_PAD = 0.26
#: 一欄至少要有的寬度（in）。窄於此放不下中文，畫出來也讀不了。
TABLE_MIN_COL = 1.0


def table_col_widths(
    n_cols: int,
    weights: tuple[float, ...] | None = None,
) -> tuple[float, ...]:
    """依欄數（與可選的相對權重）算欄寬（2026-08-18，§7c）。

    ⚠ **欄寬計算的唯一定義處**。原本只有結論頁那份寫死的 `CONCL_COLS`；
    抽出來之後結論頁也改用它，避免兩份各自演進（改一邊另一邊不動，不會報錯）。

    ⚠ 總和一律 ≤ `CW - 2*TABLE_PAD`。超寬不會拋錯，只會**在轉圖後才發現被切掉**
    ——那時已經產完一份簡報。故在這裡算，不在呼叫端各自估。

    欄數過多時 `ValueError`：與其畫出每欄窄到讀不了的表，不如當場炸。
    """
    if n_cols < 1:
        raise ValueError(f"欄數要 ≥ 1（收到 {n_cols}）")
    usable = CW - 2 * TABLE_PAD
    if n_cols * TABLE_MIN_COL > usable:
        raise ValueError(
            f"{n_cols} 欄放不下：每欄至少 {TABLE_MIN_COL}in，可用寬只有 {usable:.2f}in"
            "——請減欄或拆表，不要縮到讀不了")
    if weights is None:
        weights = tuple(1.0 for _ in range(n_cols))
    if len(weights) != n_cols:
        raise ValueError(f"權重數 {len(weights)} 與欄數 {n_cols} 不符")
    if any(w <= 0 for w in weights):
        raise ValueError("權重必須為正")
    total = float(sum(weights))
    widths = [usable * (w / total) for w in weights]
    # 權重極端時仍要守住最小寬：不足者拉到下限，其餘按比例讓出。
    if min(widths) < TABLE_MIN_COL:
        fixed = [i for i, w in enumerate(widths) if w < TABLE_MIN_COL]
        rest = usable - TABLE_MIN_COL * len(fixed)
        rest_total = sum(weights[i] for i in range(n_cols) if i not in fixed) or 1.0
        widths = [TABLE_MIN_COL if i in fixed
                  else rest * (weights[i] / rest_total) for i in range(n_cols)]
    # ⚠ 無條件**捨去**到 0.001in，不用 round：四捨五入會讓總和超出可用寬
    #   （實測 2 欄得 11.814 > 11.813）。1‰ 的溢出不會報錯，只會在轉圖後
    #   發現最右欄被切掉——那時已經產完一份簡報。寧可少千分之一。
    import math

    return tuple(math.floor(w * 1000) / 1000 for w in widths)


#: 結論頁行動分組標題的高度（in）。
CONCL_GROUP_HEAD_H = 0.34

#: 結論頁四欄（主題｜發現｜研發意涵｜專利行動）的欄寬（in）。
#: ⚠ 2026-08-18 起結論頁改三欄（主題含發現小字｜判讀｜專利行動），本常數
#:   已不被 `slide_conclusions` 使用，保留供追溯與舊 content 對照。
#: ⚠ 總和必須 ≤ CW - 2×0.26（左右內距）；行動欄最窄——內容是動詞表單選。
CONCL_COLS = (2.35, 3.35, 4.45, 1.35)


def slide_conclusions(prs, c):
    """綜合結論頁（design §7.7）：一主題一列、三欄帶（§7.2 的應用）。

    發現＝機械（intake `topic_facts.json` 的字串，check_content 逐字閘門盯著）；
    判讀＝CLI；專利行動＝CLI 從 `ACTION_VERBS` 挑，可多選（見 `row_actions`）。
    ⚠ `conclusions` 存在時本頁**取代**建議頁（§7.10 頁數帳：取代不新增）。
    """
    cc = c["conclusions"]
    s = base(prs, c["footer"], 2, source=src_line(c))
    header(s, cc["title"], cc["takeaway"])
    top, bot = CHART_TOP, 7.06
    x0 = ML + TABLE_PAD
    # 🔴 2026-08-18 使用者裁決：「這頁數據要少點，更多的是根據數據判讀後
    #    最終要轉到專利行動。」原本四欄，「發現」是純統計卻佔 3.35in——
    #    結論頁在重述前面已經講過的統計，正是 deepen 主軸要打的「統計鋪陳」。
    #    ⚠ 但「發現」是全頁**唯一被機械驗證**的內容（逐字比對引擎字串），
    #      刪掉就沒有錨點擋 CLI 亂編。故**降級不刪除**：從一整欄變成主題底下
    #      的灰小字，逐字比對照舊；釋出的版面全給判讀與行動。
    #    ⚠ 欄名用「判讀」不用「研發意涵」：欄名就是對 CLI 的指令，
    #      「意涵」容易寫成感想，「判讀」要求從資料推到結論。
    heads = ("主題", "判讀", "專利行動")
    widths = table_col_widths(len(heads), weights=(3.0, 6.4, 2.4))
    y = top + 0.06
    x = x0
    for head, w in zip(heads, widths):
        textbox(s, x, y, w - 0.14, 0.3, [(head, {"bold": True, "color": CYAN})])
        x += w
    y += 0.42
    groups = conclusion_groups(cc["rows"])
    total_rows = sum(len(rs) for _v, rs in groups)
    # 列高只能量不能縮（字級鎖死）。第一欄要容納主題＋發現兩行。
    heights: list[float] = []
    for verb, rs in groups:
        for r in rs:
            first = (text_h([(r.get("topic", ""), 16, 0)], widths[0] - 0.14)
                     + text_h([(r.get("finding", ""), S_SIZE.pt, 0)], widths[0] - 0.14))
            act = (text_h([(verb, 16, 0)], widths[2] - 0.14)
                   + text_h([(_also_actions(verb, r), S_SIZE.pt, 0)], widths[2] - 0.14))
            h = max(first,
                    text_h([(r.get("reading", ""), 16, 0)], widths[1] - 0.14),
                    act)
            heights.append(h + ROW_PAD)
    note("結論頁列高總和", bot - y,
         sum(heights) + ROW_GAP * max(total_rows - 1, 0)
         + CONCL_GROUP_HEAD_H * len(groups))
    idx = 0
    for verb, rs in groups:
        # ⚠ 只畫**真的用到**的分組（使用者：「不能每次都固定全部行動都有」）。
        #   五組固定列出會讓 CLI 覺得每組都要填——那是形式鎖逼出硬湊。
        textbox(s, x0, y, CW - 2 * TABLE_PAD, CONCL_GROUP_HEAD_H,
                [(verb, {"bold": True, "color": GREEN})], space_after=0)
        y += CONCL_GROUP_HEAD_H
        for r in rs:
            h = heights[idx]
            idx += 1
            x = x0
            textbox(s, x, y, widths[0] - 0.14, h,
                    [(r.get("topic", ""), {"bold": True})], space_after=0)
            textbox(s, x, y + 0.26, widths[0] - 0.14, h,
                    [(r.get("finding", ""), {"color": MUTED, "size": S_SIZE})],
                    space_after=0)
            x += widths[0]
            textbox(s, x, y, widths[1] - 0.14, h,
                    [(r.get("reading", ""), {})], space_after=0)
            x += widths[1]
            # 行動格印**本分組**的動詞（不是整個 list）——同一列出現在兩組時，
            # 兩處各自顯示自己那個動詞，另一個掛灰小字，讀者不會以為分組錯了。
            textbox(s, x, y, widths[2] - 0.14, h,
                    [(verb, {"bold": True, "color": GREEN})], space_after=0)
            if (also := _also_actions(verb, r)):
                textbox(s, x, y + 0.26, widths[2] - 0.14, h,
                        [(also, {"color": MUTED, "size": S_SIZE})], space_after=0)
            y += h + ROW_GAP
            rect(s, x0, y - ROW_GAP / 2, CW - 2 * TABLE_PAD, RULE_W, fill=CARD_ED)
    return s


def row_actions(row: dict) -> list[str]:
    """一列結論宣告的專利行動，一律回 list（🔴 解析的唯一定義處）。

    2026-08-19 使用者裁決：「不是說行動就只能選擇一種而已。」一個主題可能同時
    要「迴避設計」與「追蹤」——原本 `action` 是單一字串，會逼 CLI 二選一，
    而被丟掉的那個**不會留下痕跡**（缺席型偏差，比多寫一個危險）。

    ⚠ 字串與 list 都收：既有 content.json 與範例都是字串，改成只收 list 會讓
    舊檔靜默變成「沒有行動」。
    ⚠ 畫的（`conclusion_groups`）與擋的（`check_content._bad_actions`）必須讀
    這一份——分成兩份解析後，list 在一邊被展開、在另一邊被 `str()` 成
    `"['迴避設計', '追蹤']"` 當成一個假動詞，而且不會報錯。
    """
    raw = row.get("action")
    if raw is None:
        return []
    items = raw if isinstance(raw, (list, tuple)) else [raw]
    return [s for item in items if (s := str(item).strip())]


def _also_actions(verb: str, row: dict) -> str:
    """該列除了本分組動詞以外還宣告了什麼（給行動格的灰小字）。

    沒有其他動詞時回空字串——單一行動的列（多數情況）版面完全不變。
    """
    rest = [v for v in row_actions(row) if v != verb]
    return f"同時：{'、'.join(rest)}" if rest else ""


def conclusion_groups(rows: list[dict]) -> list[tuple[str, list[dict]]]:
    """結論列依專利行動分組，**只回真的用到的分組**（2026-08-18 使用者裁決）。

    組序照 `ACTION_VERBS` 宣告序（不另定第二份順序）；組內依外部訊號
    `pending_count`（他人審查中件數）由多到少——那是對手給的時間壓力、可查證，
    取代原本 CLI 自己編的「短期 0–3 個月」。

    ⚠ **不設「不得五組全有」的檢查**：十個主題真的可能用滿五種行動，
    禁止它會逼 CLI 把合理的行動改掉。這是判斷不是規則。
    """
    order = {v: i for i, v in enumerate(ACTION_VERBS)}
    buckets: dict[str, list[dict]] = {}
    for r in rows:
        # 一列可宣告多個行動 → 在每個分組底下各出現一次（2026-08-19）。
        for verb in row_actions(r):
            buckets.setdefault(verb, []).append(r)
    out: list[tuple[str, list[dict]]] = []
    for verb in sorted(buckets, key=lambda v: (order.get(v, len(ACTION_VERBS)), v)):
        out.append((verb, sorted(
            buckets[verb],
            key=lambda r: (-int(r.get("pending_count") or 0), str(r.get("topic") or "")))))
    return out


def slide_table(prs, c, page, spec):
    """通用表格頁（2026-08-18，§7c）：參數化欄數與欄寬。

    ⚠ 畫法與結論頁**同一份**——欄寬走 `table_col_widths`、列高走實測而非估算。
    表格繪製本來就會（結論頁四欄表），只是綁死在那一頁；抽出來之後 CLI 才用得到。

    內容契約：
      `table.headers`  欄頭字串陣列
      `table.rows`     每列一個字串陣列，長度須等於 headers
      `table.weights`  （選填）相對欄寬

    ⚠ 列高「只能量不能縮」：字級鎖死，量出來放不下就是要拆頁，不是縮字
    （沿用結論頁的紀律；縮字會讓轉圖後才發現讀不了）。
    """
    tbl = spec.get("table") or {}
    heads = [str(h) for h in (tbl.get("headers") or [])]
    rows = [[str(cell) for cell in row] for row in (tbl.get("rows") or [])]
    s = base(prs, c["footer"], page, source=src_line(c))
    header(s, spec["title"], spec["takeaway"])
    if not heads:
        return s
    weights = tbl.get("weights")
    widths = table_col_widths(len(heads), tuple(weights) if weights else None)
    top, bot = CHART_TOP, 7.06
    x0 = ML + TABLE_PAD
    y = top + 0.06
    x = x0
    for head, w in zip(heads, widths):
        textbox(s, x, y, w - 0.14, 0.3, [(head, {"bold": True, "color": CYAN})])
        x += w
    y += 0.42
    heights = []
    for row in rows:
        h = max(text_h([(t, 16, 0)], w - 0.14) for t, w in zip(row, widths))
        heights.append(h + ROW_PAD)
    note(f"表格頁 P{page} 列高總和", bot - y,
         sum(heights) + ROW_GAP * max(len(rows) - 1, 0))
    for row, h in zip(rows, heights):
        x = x0
        for t, w in zip(row, widths):
            textbox(s, x, y, w - 0.14, h, [(t, {})], space_after=0)
            x += w
        y += h + ROW_GAP
        rect(s, x0, y - ROW_GAP / 2, CW - 2 * TABLE_PAD, RULE_W, fill=CARD_ED)
    return s


def slide_roadmap(prs, c, page):
    s = base(prs, c["footer"], page, source=src_line(c))
    header(s, c["roadmap_title"], c["roadmap_takeaway"])
    top = CHART_TOP
    cw_ = (CW - 0.44) / 3
    tw_ = cw_ - 0.48
    lim_title = c.get("limits_title", "分析限制與適用邊界")
    card_need = max(text_h([(r["label"], 16, 10)] + [(t, 16, 8) for t in r["items"]], tw_)
                    for r in c["roadmap"])
    lim_need = text_h([(lim_title, 16, 6)] + [(t, 16, 4) for t in c["limits"]], CW - 0.52)
    lim_h = lim_need + 0.34
    card_h = card_need + 0.50
    # ⚠ 這一頁的兩個區塊高度都是「內容需求 ＋ 固定 padding」算出來的，所以
    #   分開檢查各自的 avail vs need **恆等於通過**——那是套套邏輯，抓不到任何東西。
    #   2026-08-11 實測：限制頁加到 5 則時最後一則被頁尾切掉，裕度表仍報「溢出 0」。
    #   真正的約束是**整頁總高**，只能在這裡查。
    avail_total = BAND_BOT - CHART_TOP
    need_total = card_h + lim_h + 0.20
    note("路線圖頁總高", avail_total, need_total)
    # 放不下時靠上排（不要用負偏移把內容推到頁面外，那會連頁尾一起蓋掉）
    top += max(0.0, (avail_total - need_total) / 2)
    for i, r in enumerate(c["roadmap"]):
        x = ML + i * (cw_ + 0.22)
        color = PALETTE[r["color"]]
        rect(s, x, top, cw_, card_h, fill=BG_PANEL, line=CARD_ED,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
        rect(s, x, top, cw_, 0.06, fill=color)
        blocks = [(r["label"], {"bold": True, "color": color, "space_after": 10})]
        blocks += [([("▍", {"color": color}), (t, {})], {"space_after": 8}) for t in r["items"]]
        textbox(s, x + 0.24, top + 0.26, tw_, card_h - 0.44, blocks)
    y = top + card_h + 0.20
    rect(s, ML, y, CW, lim_h, fill=BG_PANEL, line=ROSE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    blocks = [(lim_title, {"bold": True, "color": ROSE, "space_after": 6})]
    blocks += [([("・", {"color": ROSE}), (t, {"color": MUTED})], {"space_after": 4})
               for t in c["limits"]]
    textbox(s, ML + 0.26, y + 0.17, CW - 0.52, lim_h - 0.34, blocks)
    return s


class _SvgDeck:
    """`Presentation` 的替身，讓 `base()` 與各頁型**一行都不用改**。

    它們呼叫的是 `prs.slides.add_slide(prs.slide_layouts[6])`——這裡照樣提供，
    只是回傳 SVG 畫布而非 pptx slide。
    """

    def __init__(self) -> None:
        self.pages: list = []
        self.slide_layouts = [None] * 7
        self.slides = self
        self.slide_width = self.slide_height = None

    def add_slide(self, _layout):
        canvas = _svg_canvas()
        self.pages.append(canvas)
        return canvas


def _compose(deck, content: dict, png_dir: Path) -> dict:
    """跑完整份簡報的頁型。⚠ **頁型呼叫的唯一落點**——pptx 與 SVG 兩個輸出端
    共用這一份，否則兩邊會各自演進（少一頁、順序不同都不會報錯）。
    """
    slide_cover(deck, content)
    # §7.7：conclusions 宣告時綜合結論頁**取代**建議頁（取代不並存——
    # 兩頁並存＝同一問題答兩次；§7.10 頁數帳也靠這條成立）。
    if content.get("conclusions"):
        slide_conclusions(deck, content)
    else:
        slide_rec(deck, content)
    widths = {}
    for i, spec in enumerate(content["pages"], start=3):
        if spec.get("layout") == "table":
            slide_table(deck, content, i, spec)     # 表格頁：不佔圖表寬度
        elif spec.get("charts"):
            widths[i] = slide_chart(deck, content, i, spec, png_dir)
        else:
            slide_text(deck, content, i, spec)      # 純文字頁：不佔圖表寬度
    # 🔴 2026-08-18（§7d）：路線圖頁**併入結論頁**、期程整個拿掉。
    #    結論頁四欄每欄都有引擎來源，唯獨路線圖的時間桶沒有任何資料支撐——
    #    系統不知道人力、預算與產品排程，`短期 0–3 個月` 是 CLI 憑空填的。
    #    拿掉期程後兩頁功能重疊（結論頁已有 主題｜發現｜意涵｜行動），故合併。
    #    `slide_roadmap` 暫留供舊 content 追溯，不再由 _compose 呼叫。
    return widths


def build_svg(content: dict, png_dir, out_dir) -> list[Path]:
    """B 案輸出端：每頁一個 SVG 檔，回傳檔案清單。

    ⚠ PNG 一併複製到 `out_dir`：`picture()` 在 SVG 端寫的是**相對路徑**，
    Chromium 用 `goto` 載入 SVG 時才抓得到同目錄的圖
    （2026-08-13 實測 `file://` 絕對 URI ＋ `set_content` 會破圖）。
    """
    import shutil

    png_dir, out_dir = Path(png_dir), Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    _REPORT.clear()

    deck = _SvgDeck()
    _compose(deck, content, png_dir)

    for png in png_dir.glob("*.png"):
        shutil.copyfile(png, out_dir / png.name)

    written = []
    for index, canvas in enumerate(deck.pages, start=1):
        written.append(canvas.save(out_dir / f"page{index:02d}.svg"))
    return written


def build(content: dict, png_dir, out_path) -> int:
    """組裝簡報並印出裕度表與圖表在投影片上的實際字級。回傳溢出區域數（必須為 0）。"""
    png_dir, out_path = Path(png_dir), Path(out_path)
    fonts = json.loads((png_dir / "font_choice.json").read_text(encoding="utf-8"))
    _REPORT.clear()

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    # 頁型呼叫走 _compose（唯一落點，與 build_svg 共用）
    widths = _compose(prs, content, png_dir)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)

    # 0.005in ≈ 0.36pt：吃掉浮點誤差，避免「可用 0.81 / 需求 0.81」被誤判成溢出
    bad = [r for r in _REPORT if r[2] > r[1] + 0.005]
    print("── 版面裕度（需求 > 可用即溢出）──")
    for region, avail, need in bad:
        print(f"  [溢出] {region:<24} 可用 {avail:.2f} / 需求 {need:.2f}")
    print(f"  溢出區域：{len(bad)} 個")
    print("── 圖表在投影片上的實際字級 ──")
    # 🔴 字級不足分兩類（2026-08-16 #400 首跑修正）：
    #   fixable   ＝ CLI 改內容能解（雙圖頁可拆、判讀帶佔 >2 行擠掉圖）→ 計入失敗
    #   structural＝ 原圖結構所限（判讀帶已精簡、非 chip）→ **只揭露不擋**
    # ⚠ 為什麼不能一律擋：structural 那類 CLI 怎麼改內容都不會變好，逼它過閘門
    #   只能刪判讀帶（缺席型偏差），圖還是不會變大——正是 deepen design §1.2
    #   三問第 3 題要防的 v5 同型錯誤。SKILL.md 對這類的處置本來就是「揭露」。
    weak = []          # 可修的：計入回傳值，走修稿輪
    structural = []    # 結構所限：只揭露，不擋交付
    for i, spec in enumerate(content["pages"], start=3):
        for n in spec.get("charts") or []:
            src_w = Image.open(png_dir / f"{n}.png").size[0] / 3   # PNG 為 3 倍解析
            pt = fonts[n] * (widths[i] * 72) / src_w
            floor = MIN_CHART_PT_MULTI if len(spec["charts"]) > 1 else MIN_CHART_PT
            print(f"  P{i:>2} {n:<32} 圖內 {fonts[n]:>4} → 投影片 {pt:.1f}pt"
                  f"{'  ⚠ 低於 %.0fpt' % floor if pt < floor else ''}")
            if pt < floor:
                multi = len(spec.get("charts") or []) > 1
                band_lines = est_lines(" ".join(spec["lines"]), CW - 0.52)
                (weak if (multi or band_lines > 2) else structural).append(
                    (i, n, pt, spec))

    # 依 SKILL.md 的優先順序，直接算出該頁該怎麼處理——可讀性排在頁數精簡之前
    if structural:
        print("── 圖內字級不足但屬原圖結構所限：**交付時必須揭露** ──")
        for i, n, pt, spec in structural:
            print(f"  P{i:>2} {n} {pt:.1f}pt：判讀帶已精簡、非 chip 型"
                  "——改內容無法改善，不擋交付；完成回報要逐頁列出此字級")
    if weak:
        print("── 圖內字級不足的頁面：處理方式 ──")
        for i, n, pt, spec in weak:
            band_lines = est_lines(" ".join(spec["lines"]), CW - 0.52)
            if len(spec.get("charts") or []) > 1:
                act = (f"雙圖頁未達 {MIN_CHART_PT_MULTI:.0f}pt → 拆成兩頁（頁數不是限制）")
            elif band_lines > 2:
                act = (f"判讀帶佔 {band_lines} 行、壓縮了圖表 → 先依 narrative.md 濃縮到 2 行；"
                       f"內容真的必要就移到新增的純文字頁，不要縮圖")
            else:
                act = ("判讀已精簡仍不足 → 若為 chip 型圖表可跑 rebuild_chip_chart.py；"
                       "否則屬原圖結構所限，**必須在完成回報中揭露此頁字級**")
            print(f"  P{i:>2} {n} {pt:.1f}pt：{act}")
    print(f"已輸出：{out_path}  {out_path.stat().st_size // 1024} KB  "
          f"共 {len(prs.slides._sldIdLst)} 頁")
    # 🔴 字級不足**必須計入回傳值**：只 print 不擋，等於門檻沒有牙齒。
    #    make_deck 以非零 exit 短路，runner 才知道這一步失敗。
    # 🔴 只計 weak（可修的）——structural 擋下去只會逼出缺席型偏差（見上方註記）。
    return len(bad) + len(weak)
