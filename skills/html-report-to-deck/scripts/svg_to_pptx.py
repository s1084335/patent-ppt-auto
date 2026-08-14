"""窄 SVG→DrawingML 轉換器：把引擎組好的每頁 SVG 轉成原生 PPTX。

## 為什麼窄

B 案把排版決定權從 PowerPoint 收回引擎——SVG 裡的文字**已經逐行斷好、絕對定位**，
這裡只負責忠實搬運。因此只支援本 skill 的元素詞彙（`<rect>`／`<text>`／`<image>`），
不支援 SVG 標準全集。詞彙外的元素一律 raise，**不靜默略過**：略過會讓版面少一塊
而沒有任何人發現，那正是本專案反覆踩過的靜默失敗。

## 詞彙＝現行組版實際畫得出來的東西（自 `deck_layout.py` 反解）

| SVG | pptx | 對應 |
|---|---|---|
| `<rect>` | `add_shape(RECTANGLE / ROUNDED_RECTANGLE)` | `deck_layout.rect()` |
| `<text>` | `add_textbox()`（一個 `<text>` ＝ 一行 ＝ 一個框） | `deck_layout.textbox()` |
| `<image>` | `add_picture()` | `chart_stack()` |

⚠ **「線」不是獨立詞彙**，是矩形的退化形（`RULE_W = 0.014in` 的細 rect）。
另立一種表示法就會有兩個落點。

## 圖表詞彙（tasks 2.1b，2026-08-14 加）

圖表原本截成 PNG 貼上，改**原生繪製**後，圖表 SVG 的元素也走這裡。
詞彙來自實掃 14 張真實圖表（連屬性值形式一起掃——上一次只掃元素標籤，
漏了 `width="100%"`，那會讓 13/14 張在第一行就炸）。

| SVG | 次數 | pptx |
|---|---|---|
| `<circle>` | 108 | `add_shape(OVAL)`，⚠ `cx/cy` 是**圓心**不是左上角 |
| `<line>` | 19 | 細 rect。✅ 實掃全為水平／垂直，**斜線一律 raise** |
| `<polyline>` | 2 | `build_freeform()` |
| `<defs>`＋`<pattern>` | 1 | 🔴 hatch，見下 |

🔴 **hatch 不可退化成純色**：`chart_runner.py:826` 寫明「顏色分段＝申請結構，
**斜紋疊加＝已轉讓**」——第二個視覺通道承載獨立資訊。退化會讓「已轉讓」
消失，而且不會有任何東西報錯。

🔴 **`text-anchor` 不自己量字寬**：框做成**對稱於錨點**再交給 PowerPoint 的
段落對齊（middle→CENTER、end→RIGHT）。自己量就等於在這裡開第二份字寬估算，
而 `deck_layout` 已經有一份——兩份會各自演進，且不一致不會報錯。

## 座標系＝96 dpi

`deck_layout.py:73` 的註記寫死了這個前提（「0.008in 在 96dpi 下只有 0.77px」）。
13.333×7.5in 的投影片 ＝ 1280×720 的 SVG。

用法：`from svg_to_pptx import build; build([svg1, svg2, …]).save(out)`
"""
from __future__ import annotations

import re
import xml.etree.ElementTree as ET
from pathlib import Path
from urllib.parse import unquote, urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_PATTERN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"

# 座標系：SVG 以 px 表示，96 px = 1 in（見模組 docstring）。
DPI = 96.0
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5

# 詞彙表。⚠ 新增元素要同步 `tests/test_svg_to_pptx_converter.py` 的 fail-loud 測試，
# 否則等於偷偷放寬詞彙。
SUPPORTED = {"rect", "text", "image", "circle", "line", "polyline"}
# 純結構性、不畫東西的標籤：略過不算違規。
# ⚠ `title` 是 tooltip（實掃 65 個掛在 circle／rect 底下），`defs`／`pattern`
# 是樣式定義——都不畫，但也不是違規。
STRUCTURAL = {"svg", "g", "defs", "title", "desc", "style", "metadata", "pattern"}

# 顏色關鍵字：**只收實際用到的**，不引進整套 CSS 顏色表。
# 實掃 14 張圖表只有 `white`（背景 rect ×14）。放寬要連同測試一起改。
NAMED_COLORS = {"white": "FFFFFF", "black": "000000"}

# 互動用屬性：明確忽略。⚠ 列出來是為了讓「忽略」是決定而非疏漏。
IGNORED_ATTRS = {"class", "pointer-events", "data-value-band", "data-cell",
                 "data-topic", "data-on-fill", "data-role", "id"}


class UnsupportedElement(ValueError):
    """SVG 用了詞彙外的元素。窄轉換器的核心約束——不猜、不略過、直接失敗。"""


class _Ctx:
    """走訪期間的繼承狀態。

    ⚠ 為什麼需要它：`font-family` 在真實圖表裡是**繼承來的**（`<style>` 的
    `text{...}` 或 `<svg>` 根元素），不是每個 `<text>` 都有屬性。只讀元素自己的
    屬性會拿不到，然後**靜默**用 pptx 預設字型——量測與產出一起偏。
    """

    def __init__(self, vw: float, vh: float, base_dir: Path | None):
        self.vw = vw                    # viewport 寬（供 `100%` 解析）
        self.vh = vh
        self.base_dir = base_dir
        self.font: str | None = None    # 繼承字型
        self.patterns: dict[str, dict] = {}   # id → pattern 定義（供 url(#id)）


def _local(tag: str) -> str:
    """去掉 namespace 前綴。"""
    return tag.split("}", 1)[-1] if "}" in tag else tag


def _px(value: str | None, default: float = 0.0, pct_base: float | None = None) -> float:
    """SVG 長度字串 → px。只認純數字、px 與百分比。

    🔴 百分比：實掃 14 張圖表有 **13 張**的背景 rect 用 `width="100%"`，
    不支援的話第一張就炸。`pct_base` 是該軸的 viewport 長度。
    """
    if value is None:
        return default
    text = str(value).strip()
    if text.endswith("%"):
        if pct_base is None:
            raise UnsupportedElement(f"百分比 {value!r} 沒有可對照的 viewport 長度")
        try:
            return float(text[:-1]) / 100.0 * pct_base
        except ValueError as exc:
            raise UnsupportedElement(f"看不懂的百分比：{value!r}") from exc
    if text.endswith("px"):
        text = text[:-2]
    try:
        return float(text)
    except ValueError as exc:
        raise UnsupportedElement(f"看不懂的長度值：{value!r}（詞彙只收 px 與 %）") from exc


def _emu(px: float) -> int:
    """px @96dpi → EMU。

    ⚠ 走 `Emu(round(...))` 而不是 `Inches(px/DPI)`：後者對細線（1.344px）會在
    浮點轉換時丟精度，而 `deck_layout.py:73` 明記細線消失是隨機發生、目視抓不到。
    """
    return Emu(round(px / DPI * 914400))


def _color(value: str | None) -> RGBColor | None:
    """`#RRGGBB`／有限關鍵字 → RGBColor；`none`／缺值 → None。"""
    if not value or value.strip().lower() in ("none", "transparent"):
        return None
    text = value.strip()
    named = NAMED_COLORS.get(text.lower())
    if named:
        return RGBColor.from_string(named)
    text = text.lstrip("#")
    if not re.fullmatch(r"[0-9A-Fa-f]{6}", text):
        raise UnsupportedElement(
            f"看不懂的顏色：{value!r}（詞彙收 #RRGGBB 與 {sorted(NAMED_COLORS)}）")
    return RGBColor.from_string(text.upper())


def _alpha(color_format, opacity: str | None) -> None:
    """把 `fill-opacity`／`stroke-opacity` 寫成 DrawingML 的 `<a:alpha>`。

    ⚠ python-pptx 1.0.2 的 `FillFormat` **沒有** `transparency` 屬性
    （2026-08-14 實測），只能自己掛 XML。忽略透明度會讓疊圖變不透明、
    遮住底下的東西——而且看起來仍然「像一張圖」，目視不一定抓得到。
    """
    if opacity is None:
        return
    try:
        val = float(opacity)
    except ValueError as exc:
        raise UnsupportedElement(f"看不懂的透明度：{opacity!r}") from exc
    if val >= 1.0:
        return
    srgb = color_format._xFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    node = srgb.makeelement(qn("a:alpha"), {"val": str(int(round(val * 100000)))})
    srgb.append(node)


def _resolve_fill(shape, el: ET.Element, ctx: _Ctx) -> None:
    """填色：純色／無／🔴 pattern 引用。

    🔴 `url(#hatch)` **不可退化成純色**：`chart_runner.py:826` 寫明
    「顏色分段＝申請結構，**斜紋疊加＝已轉讓**」——第二個視覺通道承載獨立資訊。
    退化會讓「已轉讓」消失，且不會有任何東西報錯。
    """
    raw = (el.get("fill") or "").strip()
    if raw.startswith("url("):
        ref = raw[4:-1].lstrip("#") if raw.endswith(")") else ""
        pattern = ctx.patterns.get(ref)
        if pattern is None:
            raise UnsupportedElement(
                f"fill 引用了不存在的樣式：{raw!r}（已定義：{sorted(ctx.patterns)}）")
        shape.fill.patterned()
        shape.fill.pattern = pattern["mso"]
        shape.fill.fore_color.rgb = pattern["fore"]
        shape.fill.back_color.rgb = RGBColor.from_string("FFFFFF")
        return
    fill = _color(raw or None)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        _alpha(shape.fill.fore_color, el.get("fill-opacity"))


def _resolve_stroke(shape, el: ET.Element) -> None:
    """描邊：顏色、寬度、虛線。"""
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    stroke = _color(el.get("stroke"))
    if stroke is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = stroke
    shape.line.width = Pt(_px(el.get("stroke-width"), 1.0) / DPI * 72)
    _alpha(shape.line.color, el.get("stroke-opacity"))
    if el.get("stroke-dasharray"):
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def _add_rect(slide, el: ET.Element, ctx: _Ctx) -> None:
    """`<rect>` → autoshape。有 `rx` 就是圓角矩形（對應 `ROUNDED_RECTANGLE`）。"""
    x, y = _px(el.get("x")), _px(el.get("y"))
    w = _px(el.get("width"), pct_base=ctx.vw)
    h = _px(el.get("height"), pct_base=ctx.vh)
    rx = _px(el.get("rx")) if el.get("rx") else 0.0
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_kind, _emu(x), _emu(y), _emu(w), _emu(h))

    _resolve_fill(shape, el, ctx)
    _resolve_stroke(shape, el)

    # 陰影一律關掉：`deck_layout.rect()` 同款（`s.shadow.inherit = False`）。
    shape.shadow.inherit = False
    if rx > 0 and w > 0:
        # SVG 的 rx 是絕對長度，pptx 的 adjustment 是相對短邊的比例。
        shape.adjustments[0] = min(0.5, rx / min(w, h) if min(w, h) else 0)


def _add_circle(slide, el: ET.Element, ctx: _Ctx) -> None:
    """`<circle>` → OVAL。⚠ `cx`／`cy` 是**圓心**，pptx 要的是外框左上角。

    實掃 108 個，用在散點圖與象限圖的資料點。搞錯就整張圖的點都位移半徑，
    而圖看起來仍然「像一張散點圖」——目視抓不到。
    """
    cx, cy, r = _px(el.get("cx")), _px(el.get("cy")), _px(el.get("r"))
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, _emu(cx - r), _emu(cy - r),
                                   _emu(r * 2), _emu(r * 2))
    _resolve_fill(shape, el, ctx)
    _resolve_stroke(shape, el)
    shape.shadow.inherit = False


def _add_line(slide, el: ET.Element, ctx: _Ctx) -> None:
    """`<line>` → 細 rect（既有原則：線是矩形的退化形，不另立詞彙）。

    ✅ 實掃 19 條**全為水平／垂直，0 條斜線**，所以退化形足夠，
    不需要 `add_connector`。
    ⚠ 真出現斜線就 raise：用細 rect 硬畫會**靜默**把它擺成水平，
    而圖看起來仍然合理。
    """
    x1, y1 = _px(el.get("x1")), _px(el.get("y1"))
    x2, y2 = _px(el.get("x2")), _px(el.get("y2"))
    thickness = _px(el.get("stroke-width"), 1.0)
    horizontal, vertical = abs(y1 - y2) < 0.5, abs(x1 - x2) < 0.5
    if not horizontal and not vertical:
        raise UnsupportedElement(
            f"斜線不在詞彙內：({x1},{y1})→({x2},{y2})。實掃 19 條全為水平／垂直——"
            "出現斜線代表引擎變了，請擴充詞彙而不要在這裡硬畫成水平。")
    if horizontal:
        left, top = min(x1, x2), y1 - thickness / 2
        w, h = abs(x2 - x1), thickness
    else:
        left, top = x1 - thickness / 2, min(y1, y2)
        w, h = thickness, abs(y2 - y1)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(left), _emu(top),
                                   _emu(w), _emu(h))
    # 線的「顏色」在 SVG 是 stroke，退化成 rect 後變成填色。
    color = _color(el.get("stroke"))
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        _alpha(shape.fill.fore_color, el.get("stroke-opacity"))
    shape.line.fill.background()
    shape.shadow.inherit = False
    if el.get("stroke-dasharray"):
        # 退化成 rect 後沒有線可以虛線化——改用外框畫虛線，保留視覺區別。
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        if color is not None:
            shape.line.color.rgb = color
            shape.line.width = Pt(thickness / DPI * 72)
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            shape.fill.background()


def _add_polyline(slide, el: ET.Element, ctx: _Ctx) -> None:
    """`<polyline points="x,y x,y …">` → freeform。實掃 2 條（趨勢折線）。"""
    raw = (el.get("points") or "").replace(",", " ").split()
    if len(raw) < 4 or len(raw) % 2:
        raise UnsupportedElement(f"看不懂的 points：{el.get('points')!r}")
    pts = [(float(raw[i]), float(raw[i + 1])) for i in range(0, len(raw), 2)]
    builder = slide.shapes.build_freeform(_emu(pts[0][0]), _emu(pts[0][1]))
    builder.add_line_segments([(_emu(x), _emu(y)) for x, y in pts[1:]], close=False)
    shape = builder.convert_to_shape()
    _resolve_fill(shape, el, ctx)
    _resolve_stroke(shape, el)
    shape.shadow.inherit = False


def _add_text(slide, el: ET.Element, ctx: _Ctx) -> None:
    """`<text>` → 一個文字框 ＝ **一行**。

    🔴 SVG 的 `y` 是**基線**，pptx 的 top 是框頂。以字級推回框頂，並關掉
    內距與自動調整，讓框的位置只由座標決定。
    ⚠ 一個 `<text>` 一個框，不合併：合併就等於把斷點交還給 PowerPoint，
    B 案整個白做。

    🔴 **`text-anchor` 不自己量字寬**（實掃 198 次：middle 152／end 46）。
    做法是把框做成**對稱於錨點**（middle）或**右緣貼齊錨點**（end），
    再交給 PowerPoint 的段落對齊。自己量就等於在這裡開第二份字寬估算，
    而 `deck_layout` 已經有一份——兩份會各自演進，且不一致不會報錯。
    """
    size_px = _px(el.get("font-size"), 21.33)
    baseline_y = _px(el.get("y"))
    # 基線→框頂：約 0.8 個字高是 ascent，其餘留白由關掉的內距吸收。
    top_px = baseline_y - size_px * 0.8
    x_px = _px(el.get("x"))
    text = "".join(el.itertext())

    # 框寬給足，反正 wrap 關掉、不會換行。CJK 以一字一 em 估上界，寧可寬不可窄
    # ——框是透明的，過寬不影響視覺；過窄才會讓對齊算錯。
    box_w = max(size_px * (len(text) + 2), size_px * 4)
    anchor = (el.get("text-anchor") or "start").strip()
    if anchor == "middle":
        left_px, align = x_px - box_w / 2, PP_ALIGN.CENTER
    elif anchor == "end":
        left_px, align = x_px - box_w, PP_ALIGN.RIGHT
    elif anchor == "start":
        left_px, align = x_px, PP_ALIGN.LEFT
    else:
        raise UnsupportedElement(f"看不懂的 text-anchor：{anchor!r}")

    box = slide.shapes.add_textbox(_emu(left_px), _emu(top_px),
                                   _emu(box_w), _emu(size_px * 1.6))
    tf = box.text_frame
    tf.word_wrap = False                    # 🔴 B 案的前提：PowerPoint 零重排自由
    tf.auto_size = MSO_AUTO_SIZE.NONE       # 字級鎖死，不讓它縮放
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size_px / DPI * 72)
    # 字型：元素自己的屬性優先，其次是繼承（`<style>` 或 `<svg>` 根元素）。
    family = el.get("font-family") or ctx.font
    if family:
        run.font.name = family.split(",")[0].strip().strip("'\"")
    weight = (el.get("font-weight") or "").strip()
    run.font.bold = weight in ("bold", "600", "700", "800", "900")
    run.font.italic = (el.get("font-style") or "").strip() == "italic"
    fill = _color(el.get("fill"))
    if fill is not None:
        run.font.color.rgb = fill
        _alpha(run.font.color, el.get("fill-opacity"))

    rotation = _rotation(el.get("transform"))
    if rotation is not None:
        box.rotation = rotation


def _rotation(transform: str | None) -> float | None:
    """`transform="rotate(-90 26 230)"` → pptx 的 rotation 角度（順時針度數）。

    實掃 2 次，用在縱向的軸標題。⚠ 只收 `rotate`；其他 transform（translate／
    scale／matrix）會改變座標系，靜默忽略會讓元素整個跑掉。
    """
    if not transform:
        return None
    text = transform.strip()
    match = re.fullmatch(r"rotate\(\s*(-?[\d.]+)(?:[\s,]+[-\d.]+[\s,]+[-\d.]+)?\s*\)", text)
    if not match:
        raise UnsupportedElement(f"transform 只收 rotate(...)：{transform!r}")
    return float(match.group(1)) % 360


def _collect_defs(root: ET.Element, ctx: _Ctx) -> None:
    """把 `<defs><pattern>` 收進 ctx，供 `fill="url(#id)"` 引用。

    🔴 實掃只有一種：45° 斜線 hatch（`applicant_ranking`，代表**已轉讓**）。
    映到 MSO 的 `DARK_UPWARD_DIAGONAL`——⚠ 間距與 SVG 的 6px 不會完全一致，
    但**視覺通道保住了**（有斜紋 vs 無斜紋），那才是它承載的資訊。
    """
    for defs in root.iter():
        if _local(defs.tag) != "defs":
            continue
        for pattern in defs:
            if _local(pattern.tag) != "pattern":
                raise UnsupportedElement(f"<defs> 只收 <pattern>，遇到 <{_local(pattern.tag)}>")
            pid = pattern.get("id")
            if not pid:
                raise UnsupportedElement("<pattern> 缺 id")
            stroke = None
            for child in pattern:
                if _local(child.tag) == "line":
                    stroke = _color(child.get("stroke"))
            ctx.patterns[pid] = {
                "mso": MSO_PATTERN.DARK_UPWARD_DIAGONAL,
                "fore": stroke or RGBColor.from_string("000000"),
            }


def _apply_style(el: ET.Element, ctx: _Ctx) -> None:
    """`<style>` → 繼承字型。

    ✅ 實掃 8 張都只有 `text{font-family:...}`，**沒有 fill／stroke 規則**，
    所以窄轉換器仍然窄。
    ⚠ 但只認這一種形式：出現別的 CSS 屬性就 raise——忽略一條 `fill` 規則
    會讓整張圖的顏色靜默錯掉。
    """
    css = (el.text or "").strip()
    if not css:
        return
    match = re.fullmatch(r"text\s*\{\s*font-family\s*:\s*([^;}]+);?\s*\}", css)
    if not match:
        raise UnsupportedElement(
            f"<style> 只收 `text{{font-family:...}}`，看到：{css[:80]!r}")
    ctx.font = match.group(1).strip()


def _add_image(slide, el: ET.Element, base_dir: Path | None) -> None:
    """`<image>` → `add_picture`。只收本機檔案，不抓網路。

    🔴 相對路徑相對於 **SVG 檔所在目錄**（`base_dir`），不是當前工作目錄。
    ⚠ 2026-08-13 映射煙霧測試的教訓：SVG 用 `file://` 絕對 URI 引圖時，
    Chromium 以 `set_content` 載入會因缺 base URL 而破圖（COM 轉圖卻正常）。
    可行的目視路徑是「SVG 存檔 ＋ 圖用相對路徑 ＋ Chromium `goto`」，
    那條路徑要成立，這裡就必須以 SVG 檔的位置解析。
    以 cwd 解析則會在 runner 從別的目錄呼叫時**靜默找不到圖**。
    """
    href = el.get(f"{{{XLINK_NS}}}href") or el.get("href")
    if not href:
        raise UnsupportedElement("<image> 缺 href")
    if href.startswith("data:"):
        raise UnsupportedElement("<image> 不收 data URI（圖檔走磁碟路徑，避免 SVG 肥大）")
    parsed = urlparse(href)
    if parsed.scheme == "file":
        path = Path(unquote(parsed.path).lstrip("/"))
    else:
        path = Path(unquote(href))
        if not path.is_absolute() and base_dir is not None:
            path = base_dir / path
    if not path.is_file():
        raise UnsupportedElement(f"<image> 指向的檔案不存在：{path}")
    slide.shapes.add_picture(str(path), _emu(_px(el.get("x"))), _emu(_px(el.get("y"))),
                             _emu(_px(el.get("width"))), _emu(_px(el.get("height"))))


def _walk(slide, el: ET.Element, ctx: _Ctx) -> None:
    """深度優先走訪；詞彙外即 raise。

    ⚠ 畫得出東西的元素**不遞迴進子節點**——`<circle>`／`<rect>` 底下的
    `<title>` 是 tooltip（實掃 65 個），本來就不該畫。
    """
    tag = _local(el.tag)
    if tag == "rect":
        _add_rect(slide, el, ctx)
        return
    if tag == "text":
        _add_text(slide, el, ctx)
        return
    if tag == "image":
        _add_image(slide, el, ctx.base_dir)
        return
    if tag == "circle":
        _add_circle(slide, el, ctx)
        return
    if tag == "line":
        _add_line(slide, el, ctx)
        return
    if tag == "polyline":
        _add_polyline(slide, el, ctx)
        return
    if tag == "style":
        _apply_style(el, ctx)
        return
    if tag == "defs":
        return                              # 樣式定義，已由 `_collect_defs` 讀走
    if tag not in STRUCTURAL:
        raise UnsupportedElement(
            f"詞彙外的元素：<{tag}>。窄轉換器只收 {sorted(SUPPORTED)}——"
            f"要新增請同步更新詞彙表與 fail-loud 測試，不要在這裡放行。")
    for child in el:
        _walk(slide, child, ctx)


def _viewport(root: ET.Element) -> tuple[float, float]:
    """viewport 尺寸，供 `width="100%"` 解析。viewBox 優先，其次 width/height。"""
    box = (root.get("viewBox") or "").replace(",", " ").split()
    if len(box) == 4:
        return float(box[2]), float(box[3])
    return (_px(root.get("width"), SLIDE_W_IN * DPI),
            _px(root.get("height"), SLIDE_H_IN * DPI))


def build(svg_pages: list[str], base_dirs: list[Path] | None = None) -> Presentation:
    """一頁 SVG ＝ 一張投影片。回傳 Presentation（呼叫端自己 `.save()`）。

    `base_dirs`：各頁 SVG 所在目錄，供解析 `<image>` 的相對路徑；
    不給則相對路徑以 cwd 解（僅適用測試與單檔手動呼叫）。
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]            # 空白版面：所有元素都由我們畫
    for index, svg in enumerate(svg_pages):
        slide = prs.slides.add_slide(blank)
        base = base_dirs[index] if base_dirs and index < len(base_dirs) else None
        root = ET.fromstring(svg)
        vw, vh = _viewport(root)
        ctx = _Ctx(vw, vh, base)
        # 🔴 順序有意義：先收樣式定義（`fill="url(#id)"` 可能出現在 defs 之前），
        # 再取根元素的繼承字型，最後才走訪畫圖。
        _collect_defs(root, ctx)
        root_font = root.get("font-family")
        if root_font:
            ctx.font = root_font
        _walk(slide, root, ctx)
    return prs


def build_file(svg_paths: list[Path | str], out_path: Path | str) -> int:
    """讀檔版：給 runner 用。回傳頁數。相對圖檔以各 SVG 自己的目錄解析。"""
    paths = [Path(p) for p in svg_paths]
    pages = [p.read_text(encoding="utf-8") for p in paths]
    prs = build(pages, [p.resolve().parent for p in paths])
    prs.save(str(out_path))
    return len(pages)
