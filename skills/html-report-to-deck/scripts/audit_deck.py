"""輸出前品質閘門：檢查字級、圖表重複與頁數。

⚠ 圖片去重必須用**內容雜湊**：python-pptx 讀到的 `image.filename` 在 pptx 內會被統一
   改名成 image.png，用檔名判斷會得到「14 張全部重複」的假結果。

用法：python audit_deck.py <pptx> [允許的字級；預設取 deck_layout.ALLOWED_SIZES]
回傳碼：0 = 全部通過；1 = 有違規（字級不符或圖表重複）
"""
from __future__ import annotations

import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation


def _default_allowed() -> set[float]:
    """字級白名單取自 `deck_layout.ALLOWED_SIZES`（🔴 唯一定義處）。

    ⚠ 原本這裡自帶字面 `"24,16"`，與 deck_layout 各自維護。2026-08-19 加了
    註記小字 11pt：deck 側單元測試全綠、PPTX 也產得出來，一路到 audit 才紅，
    而訊息只有一張「字級分布」看不出根因。同一份知識放兩處，不一致不會報錯。
    """
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from deck_layout import ALLOWED_SIZES

    return {s.pt for s in ALLOWED_SIZES}


def main() -> int:
    prs = Presentation(sys.argv[1])
    allowed = ({float(x) for x in sys.argv[2].split(",")} if len(sys.argv) > 2
               else _default_allowed())
    sizes, imgs, bad = Counter(), Counter(), []

    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.__class__.__name__ == "Picture":
                imgs[sh.image.sha1] += 1
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if not r.text.strip():
                        continue
                    pt = r.font.size.pt if r.font.size else None
                    sizes[pt] += 1
                    if pt not in allowed:
                        bad.append((i, pt, r.text[:20]))

    dups = {k: v for k, v in imgs.items() if v > 1}
    print("字級分布：", dict(sizes))
    print(f"不符字級的文字 run：{len(bad)}")
    for b in bad[:10]:
        print("   P%d  %s  %s" % b)
    print(f"圖片總數 {sum(imgs.values())}／不重複 {len(imgs)}")
    print("重複使用的圖表：", f"{len(dups)} 組" if dups else "無")
    print("頁數：", len(prs.slides._sldIdLst))
    return 1 if (bad or dups) else 0


if __name__ == "__main__":
    raise SystemExit(main())
