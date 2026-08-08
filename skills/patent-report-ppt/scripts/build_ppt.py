"""專利分析報告 PPTX 產生器 v3（deterministic，不呼叫 AI）。

設計原則
--------
- **AI 只產文案，不碰數字、不碰排版**：所有數字取自 `report_data.json`，本程式不推算、
  不四捨五入、不捏造；版型由使用者經 `layout_overrides` 挑，AI 不生成版型。
- **`theme.json` 是版面座標、字級、配色的唯一來源**：renderer 內不得出現座標數字字面值
  （`tests/test_ppt_geometry_single_source.py` 以 AST 契約測試把關）。前端投影片縮圖
  預覽也讀同一份 geometry，避免座標分岔。
- **圖檔一律走 `artifact_manifest.json` 反查**：實際檔名（`ipc_main_distribution_L4.svg`、
  `jurisdiction_distribution.svg`、`annual_trend.svg`）與 report_key 不同名，用
  `{report_key}.svg` 猜必錯，故禁止猜檔名。
- **每頁必有視覺元素**：找不到圖就降級 `stat_callout`（用該報表 rows 取關鍵數字），
  不印「（圖檔待產出）」這類佔位文字。
- **成對報表不合成同一張圖**：IPC/CPC 的 L4 與 L5、機會矩陣的技術面與功效面，
  預設同頁左右並排（`comparison`），使用者可改成分頁。
- **產後自檢**：組完逐 shape 檢查超界、邊距不足、文字疊文字與文字裝不下，
  結果寫入 manifest `warnings[]`，不靜默。

可攜獨立執行（不 import 主專案任何模組）：
    uv run --no-project --with python-pptx --with pymupdf --python 3.12 \
        python build_ppt.py --report-dir <報表版本目錄> \
        --approvals <approvals.json> --output-dir <輸出目錄>
"""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
import sys
import dataclasses
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

# ⚠ 同層模組匯入：本檔有兩種載入方式——直接執行（uv run 本檔）與以檔案路徑載入
# （backend 的 ai_report_ppt_runner._load_builder 用 spec_from_file_location）。
# 後者不會把本檔所在目錄放進 sys.path，`from starfield import ...` 會 ImportError。
_SCRIPT_DIR = str(Path(__file__).resolve().parent)
if _SCRIPT_DIR not in sys.path:
    sys.path.insert(0, _SCRIPT_DIR)

from starfield import starfield_png  # noqa: E402  （須在 sys.path 補完之後）

SKILL_ROOT = Path(__file__).resolve().parent.parent
THEME_PATH = SKILL_ROOT / "theme.json"

# 圖檔副檔名白名單：artifact_manifest 內混有 report_data／report_html，只取得出圖的。
IMAGE_SUFFIXES = (".svg", ".png", ".jpg", ".jpeg")


# --------------------------------------------------------------------------
# 內容設定表（非座標；座標一律在 theme.json）
# 報表增減、成對呈現、圖例編碼、警語只改本區的表，不動 renderer。
# --------------------------------------------------------------------------
@dataclass(frozen=True)
class PageSpec:
    """單頁版型定義。

    page：頁碼（1-based，展開後重新連號）。
    kind：組版樣式，決定用哪個 renderer。
    title：實際印出的標題（判讀式，＝「{報表主題}：{headline}」）。
    topic：報表主題（標題前半，供缺 headline 時單獨使用）。
    report_keys：本頁引用的 report_data report_key，依序取用。
    charts：本頁配圖檔名，**只能來自 artifact_manifest 反查**，不得猜檔名。
    slots：本頁需 AI 產的文案槽；v3 只有 cover.title 與 direction.body。
    is_appendix：附錄段顯式旗標，動態插頁靠它找錨點，不用頁碼魔術數字。
    degraded_from：若因缺圖降級，記錄原本的 kind，寫進 manifest 供追溯。
    """

    page: int
    kind: str
    title: str
    topic: str = ""
    report_keys: tuple[str, ...] = ()
    charts: tuple[str, ...] = ()
    slots: tuple[str, ...] = ()
    subtitle: str = ""
    is_appendix: bool = False
    degraded_from: str = ""
    # 列過濾（P1-3）：成對報表以「rows 的欄位值」拆頁時（主題分布依 source_field
    # 分技術／功效兩張表），該頁只取匹配的列。空 dict＝不過濾。
    # ⚠ 用 tuple of pairs 保持 frozen dataclass 可雜湊；讀取端用 dict(spec.row_filter)。
    row_filter: tuple[tuple[str, str], ...] = ()
    #: 附錄分頁用的列切片（起、訖）。⚠ 只在附錄放不下一頁時才設；
    #: 一般頁維持 None，行為與加這個欄位之前完全相同。
    row_slice: tuple[int, int] | None = None


def _spec_with(spec: PageSpec, **changes: Any) -> PageSpec:
    """PageSpec 是 frozen，改欄位一律走這裡建新物件。

    ⚠ 欄位**自動列舉**，不手寫清單（2026-08-03）：原本是逐一列出的 dict，
    新增 `row_slice` 時忘了同步就會被靜默丟掉——後面任何一次
    `_spec_with(spec, topic=...)` 都會把它洗回預設值，而且不報錯。
    """
    fields = {f.name: getattr(spec, f.name) for f in dataclasses.fields(spec)}
    fields.update(changes)
    return PageSpec(**fields)


# 基礎大綱：報告敘事主線。其餘有資料的報表由 `_expand_page_layout` 自動插在附錄前。
# ⚠ 頁序即論證鏈（P1-6，2026-07-31 使用者定案）：範圍 → 證據（時間→空間→技術→競爭→
# 機會）→ **結論（研發方向建議）壓軸** → 附錄。研發方向是吸收全部證據後的收尾，
# 不是開場白；動態插頁也算證據，一律插在它之前。
PAGE_LAYOUT: tuple[PageSpec, ...] = (
    PageSpec(page=1, kind="cover", title="專利情報整合分析", topic="專利情報整合分析",
             report_keys=("application_trend", "country_distribution")),
    PageSpec(page=2, kind="chart_hero", title="申請趨勢", topic="申請趨勢",
             report_keys=("application_trend", "publication_trend")),
    # 🔴 2026-08-07 使用者定案：本頁改用引擎的成對長條 SVG（每國兩條相鄰
    # ——申請件數 vs 現存有效，同尺＋圖例）。原 percentage_bars 版型由組版端
    # 拿報表原始列重畫，country_distribution 改 (國×狀態) 群組後一列一條，
    # 同一國被拆成多條、兩條 bar 分家——實機驗收抓到的錯，不得回退。
    PageSpec(page=3, kind="chart_hero", title="保護地域分布", topic="保護地域分布",
             report_keys=("country_distribution",)),
    # IPC＋CPC 同頁對照（2026-07-31）：原本只掛 ipc，cpc 落到動態插頁、與 ipc 隔開
    # 好幾頁——註解寫著「IPC/CPC 維持同頁比較」但實作沒做到，這裡補齊。
    PageSpec(page=4, kind="comparison", title="技術分類布局", topic="技術分類布局",
             report_keys=("ipc_main_distribution", "cpc_main_distribution")),
    # 主題分布：rows 帶 source_field 兩通道，展開時依通道拆成兩張表格頁（P1-3）。
    PageSpec(page=5, kind="table_with_points", title="技術主題分布", topic="技術主題分布",
             report_keys=("cluster_topic_table",)),
    # 🔴 RPT-011（2026-08-06）：owner_ranking 已刪（母體 36/55；「已轉讓」由
    # 申請人排名斜紋承接）——本頁只掛申請人排名。
    PageSpec(page=6, kind="chart_hero", title="競爭者佈局", topic="競爭者佈局",
             report_keys=("applicant_ranking",)),
    PageSpec(page=7, kind="comparison", title="機會評估", topic="機會評估",
             report_keys=("opportunity_quadrant",)),
    PageSpec(page=8, kind="direction", title="研發方向建議", topic="研發方向建議",
             slots=("direction.body",)),
    PageSpec(page=9, kind="table", title="附錄1：全分類技術指標總表", topic="全分類技術指標總表",
             report_keys=("cluster_topic_table",), is_appendix=True),
    # 🔴 2026-08-04 使用者定案：附錄只保留主題分類（附錄1）——
    # 「附錄那邊做到主題分類就好，其他的要看去網頁看就好」。
    # 原附錄2（主要專利權人與申請人）移除；排名圖的截斷註記已改指向網頁報表。
)

# 論證順序（2026-07-31 使用者定案的 17 頁大綱）：範圍 → 時間 → 空間 → 技術 →
# 競爭 → 機會 → 結論 → 附錄。
#
# ⚠ 問題背景：`PAGE_LAYOUT` 原本只固定 8 個 report_key，介面上其餘 7 種報表
# 全部落入動態插頁，順序由 `_iter_report_entries` 決定＝**引擎輸出什麼順序就什麼
# 順序**，整批塞在結論前。使用者：「我沒有看到我們的內容大綱怎安排」——論證鏈
# 只有頭尾編排過，中段是未編排的填充（CPC 與 IPC 隔開、年度矩陣與競爭者佈局分家）。
#
# ⚠ 為什麼不是把這 7 張寫成 PAGE_LAYOUT 固定條目：那樣得連 `kind` 一起寫死，
# 就失去 `_kind_for_report` 依「報表型別＋實際圖檔數」判定版型的能力（同一張報表
# 有圖沒圖該用的版型不同）。故只編排**位置**，版型仍自動判定；未列在此的報表
# （日後新增的）照舊自動出頁，排在已知證據之後、結論之前，不會漏。
EVIDENCE_ORDER: tuple[str, ...] = (
    "application_trend", "publication_trend", "lifecycle",              # 時間
    "country_distribution", "applicant_country_distribution",           # 空間
    # ⚠ family_country_layout 頁已刪（2026-08-07 併入受理局合併頁）；舊報表版本
    # 帶此鍵的資料也不撿——刪除是刻意的，不是漏排。
    "ipc_main_distribution", "cpc_main_distribution", "cluster_topic_table",  # 技術
    # ⚠ owner_ranking／owner_year_matrix 已刪（RPT-011）；舊報表版本仍可能帶
    # 這兩鍵的資料，但不在此序＝動態插頁也不會撿——刪除是刻意的，不是漏排。
    "applicant_ranking",                                                # 競爭
    "applicant_year_matrix",
    "opportunity_quadrant",                                             # 機會
)

# 不進 PPT 的報表。⚠ RPT-011（2026-08-06）後 family_quality_detail 連報表都刪了
# （家族完整性併入國家佈局頁註記），本集合仍保留該鍵：**舊報表版本**的 report_data
# 還帶著它，缺了這行會被當成動態插頁撿回簡報。
# family_country_layout（2026-08-07）：獨立頁已刪、併入受理局合併頁——引擎仍會把
# 該報表資料寫進 report_data（合併頁註記要用），動態插頁不得再為它出獨立頁。
EXCLUDED_FROM_PPT = frozenset({"family_quality_detail", "family_country_layout"})

# 截斷時優先切在這些標點之後（見 `_truncate_to_width`）：斷在標點像「話沒說完」，
# 斷在字中間像「字被砍掉」，後者會讓讀者以為產檔壞了。
TRUNCATE_BREAK_MARKS = ("，", "、", "；", "：", "。", "）", ",", ";")

# 背景層（漸層底＋星空紋理）的 shape 名稱：版面自檢與 QA 用它排除全出血元素。
BACKGROUND_SHAPE_NAME = "space-background"

# 扁圖門檻（寬/高）：超過此值改用「滿寬圖＋底部要點」版型。
# ⚠ 3.5 不是拍腦袋——實測 18 張圖在兩種版型下的縮放倍率：
#   ratio ≥3.89 → 滿寬 1.19 vs 現行 0.87（+37%，明顯有感）
#   ratio 2.1–3.0 → 滿寬 0.43–0.82，**比現行更差**（高度被底部要點壓縮）
#   ratio <2.06 → 高度先滿，換滿寬毫無幫助
# 也就是說「比框扁就換」是錯的判準，要看**換完是否真的變大**。
WIDE_CHART_ASPECT_MIN = 3.5

# 向量圖 part 的流水號：`PackURI` 必須唯一，同名會讓後插入的圖覆蓋前一張。
# 逐頁遞增、順序固定，故同一份報表重跑產出的檔案仍完全一致。
_VECTOR_INDEX = 0

# 依通道拆頁的報表（P1-3）：rows 的哪個欄位分通道、各通道的顯示名。
# ⚠ 主題分布不走「多圖拆頁」（它根本沒圖）——是**依列值**拆成兩張表格頁。
# 通道 → 解讀變體鍵（2026-07-31）：主題統計表拆成技術／功效兩頁後，各取各的解讀。
# ⚠ 上游 `chart_runner` 必須宣告同名 variant，否則這裡對不到、兩頁又會共用同一段。
# 技術通道 source_field（封面漏斗的群數只算它；功效不上封面）——
# 與 CHANNEL_SPLIT_REPORTS 同一份字面來源，改欄名時一起改。
TECHNICAL_SOURCE_FIELD = "wips_independent_claims"

CHANNEL_NARRATIVE_VARIANTS: dict[str, str] = {
    "wips_independent_claims": "topic_table_tech",
    "effect_summary": "topic_table_effect",
}

CHANNEL_SPLIT_REPORTS: dict[str, tuple[str, tuple[tuple[str, str], ...]]] = {
    "cluster_topic_table": ("source_field", (
        ("wips_independent_claims", "技術主題分布"),
        ("effect_summary", "功效主題分布"),
    )),
}

# 表格欄位的顯示對照——⚠ **僅作舊報表版本的 fallback**（2026-07-31 起）。
#
# 正式來源是引擎寫進 `report_data.json["table_display"]` 的那份（欄名對照＋逐報表
# 排除欄）。本檔這份留著只為相容：`table_display` 是新加的鍵，先前已產出的報表版本
# 沒有它，缺了會讓簡報印出原始英文欄名。
#
# ⚠ 為什麼曾經對不上：這裡自建一份、引擎另有一份，兩邊各自演進。實測差異包括
#   `top3_share` 引擎叫「前三大占比(%)」這裡叫「前三大集中度(%)」、
#   `max_share` 引擎叫「最大一家(%)」這裡叫「最大占比(%)」、
#   `current_assignee_display_name` 這裡根本沒有（寫成 owner_display_name）→ 印英文；
#   排除欄更誇張：引擎排 4 欄、這裡只排 1 欄，使用者要求拿掉的「龍頭涉入」還在簡報上。
# 新增欄位請改引擎那份，不要往這裡加。
TABLE_COLUMN_LABELS: dict[str, str] = {
    "label": "主題",
    "source_field": "通道",
    "patent_count": "專利件數",
    "applicant_count": "申請人家數",
    "top3_share": "前三大集中度(%)",
    "max_share": "最大占比(%)",
    "top_applicants": "前三大申請人",
    "year_span": "年份跨度",
    "applicant_display_name": "申請人",
    "owner_display_name": "專利權人",
    "application_year": "申請年",
    "授權公告年": "授權公告年",
    "leading_applicant_count": "主要申請人涉入數",
    "leading_applicants_involved": "主要申請人名單",
    "quadrant": "象限",
}
TABLE_EXCLUDED_COLUMNS = frozenset({"topic_code"})
TABLE_VALUE_LABELS: dict[str, dict[str, str]] = {
    "source_field": {"wips_independent_claims": "技術", "effect_summary": "功效"},
}

# narrative 別名（P1-2）：解讀掛點與 report_key 不同名時的對照。
# 機會矩陣的解讀由 ai:narrative 掛在 cluster section 的 opportunity_* 變體
# （sections 結構），而 PPT 端用 report bucket 的 opportunity_quadrant——兩個
# key 空間的歷史不一致，這裡以「report:variant」語法橋接。
NARRATIVE_ALIASES: dict[str, tuple[str, ...]] = {
    "opportunity_quadrant": ("cluster_topic_table:opportunity_tech",
                             "cluster_topic_table:opportunity_effect"),
    "opportunity_quadrant_tech": ("cluster_topic_table:opportunity_tech",),
    "opportunity_quadrant_effect": ("cluster_topic_table:opportunity_effect",),
}

# 成對報表預設呈現：預設同頁左右並排；列在這裡的改成分頁。
# - cluster_topic_table：表格內容多、並排讀不動（依通道列值拆）。
# - opportunity_quadrant：2026-07-31 使用者二輪回饋「同頁比較圖太小」——象限圖
#   資訊密度高，一圖一頁（chart_hero 大圖）才看得清。IPC/CPC 維持同頁比較。
SPLIT_PAIR_REPORTS = frozenset({"cluster_topic_table", "opportunity_quadrant"})

# 只上主圖的報表（2026-07-31 使用者定案）：年度矩陣只用前 10 名主表那張，
# 「更多」長尾圖不上 PPT（報表頁仍看得到，PPT 給決策者看主要玩家就夠）。
MAIN_CHART_ONLY_REPORTS = frozenset({"applicant_year_matrix"})


def _filter_report_charts(report_keys: tuple[str, ...], files: tuple[str, ...]) -> tuple[str, ...]:
    """套用「只上主圖」規則：MAIN_CHART_ONLY_REPORTS 的 `_more` 變體圖不上 PPT。"""
    if not any(key in MAIN_CHART_ONLY_REPORTS for key in report_keys):
        return files
    return tuple(name for name in files if "_more" not in name)

# 成對圖的左右順序偏好：L4 在 L5 前、技術面在功效面前；其餘照檔名排序保持 deterministic。
CHART_ORDER_HINTS = ("_L4", "_L5", "_tech", "_effect", "_more")

# 封面 eyebrow（小字）與主標的最後 fallback。
# ⚠ 兩者原本各自寫死「專利情報整合分析」，workspace 名稱缺失時同一句印兩次。
COVER_EYEBROW = "專利情報整合分析"
COVER_TITLE_FALLBACK = "專利布局與競爭分析"

# ── 空白頁偵測（產檔後掃描用；與版型無關，換任何一批資料都適用）──
# 正文帶＝標題與註腳之間的縱向區間（比例，不寫死英吋，換版面尺寸仍成立）。
BODY_BAND_TOP_RATIO = 0.15
BODY_BAND_BOTTOM_RATIO = 0.88
# 正文帶內少於這個字數就視為空頁。⚠ 門檻取在「只有面板標題」與「有實質敘述」
# 之間：實測空框頁只有 12 字（兩個面板標題），最短的有效內容頁遠超過 60 字。
EMPTY_BODY_MIN_CHARS = 40

# 成對圖的子標，讓使用者一眼知道左右差在哪（禁止合成一張圖的配套說明）。
CHART_VARIANT_LABELS = {
    "_L4": "4 階細分類（L4）",
    "_L5": "5 階細分類（L5）",
    "_tech": "技術面",
    "_effect": "功效面",
    "_more": "（續）",
}

# 圖例編碼說明：放圖表右上一行小字，說明視覺編碼怎麼讀。
ENCODING_NOTES = {
    "application_trend": "條長＝當年申請件數｜橫軸＝申請年",
    "publication_trend": "條長＝當年公告件數｜橫軸＝公告年",
    "country_distribution": "條長＝件數（兩條同尺）｜上＝申請件數、下＝現存有效（已授權）",
    "ipc_main_distribution": "條長＝件數｜左右為不同階層，非同圖合成",
    "cpc_main_distribution": "條長＝件數｜左右為不同階層，非同圖合成",
    "opportunity_quadrant": "橫軸＝申請人家數｜縱軸＝專利件數｜點＝技術主題",
    "cluster_topic_table": "條長＝主題件數｜家數＝投入該主題的申請人數",
    "applicant_ranking": "條長＝件數｜排序＝件數由高至低",
    "applicant_country_distribution": "格值＝件數｜列＝申請人、欄＝受理國",
    "applicant_year_matrix": "格值＝件數｜列＝申請人、欄＝申請年",
    "lifecycle": "面積／條長＝當年件數｜橫軸＝申請年",
}
DEFAULT_ENCODING_NOTE = "條長＝件數｜數值取自報表引擎"

# 警語在要點清單裡的固定 label。`_trim_blocks` 認這個字保護它不被裁切，
# 故必須是常數而非各處寫死的字面（改字面時裁切保護才不會默默失效）。
CAVEAT_LABEL = "判讀限制"

# 🔴 2026-08-04 使用者定案：**判讀限制整個移除**（「判讀限制不要出現了，作用不大」）。
# 不是移到別處——連獨立灰框也不畫。要點區改為三層說明（現況／意涵／後續）。
#
# ⚠ 清空而非刪除這個 dict：`_caveat_of` 仍在（回空字串），渲染端的分支因此自然關閉，
# 不必把警語框的程式碼一併拆掉；日後若要恢復某一頁的警語，填回這裡即可。
# ⚠ `CAVEAT_LABEL` 保留給 narrative fallback 的舊資料辨識用，不再產生新的。
CAVEATS: dict[str, str] = {}

# narrative fallback：舊格式只有長文 text，依這些段落標記切成可讀條列。
NARRATIVE_MARKERS = (
    "觀察", "現況", "意涵", "後續檢視點", "決策提醒", "後續", "結論邊界", "限制", "建議",
)

# 關鍵數字粗體：把數字（含常見量詞）從敘述中切出來單獨加粗。
NUMBER_PATTERN = re.compile(r"(\d[\d,]*(?:\.\d+)?\s*(?:%|件|家|年|項|個)?)")

# P1-10（2026-07-31 使用者定案）：來源註**只留可追溯資訊**（來源／期間／版本）。
# 舊版逐頁蓋「不保證」式免責章——數字是引擎確定性統計、預覽閘門本身就是人工定稿，
# 再蓋章等於否定定稿流程；不確定性提醒收斂到判讀限制框（該頁需要才出現）。


# 一條要點在版面上預估佔幾行（含 label 那截）。用來把「可用行數」換算成「可寫幾條」。
#: 要點區固定三層：現況（數據事實）→ 意涵（代表什麼）→ 後續（下一步看什麼）。
#: 🔴 2026-08-04 使用者定案：16pt＋行距 1.65 後側欄頁只放得下 5 條 × 19 字，
#: **碎成 5 個短句反而難讀**。改為三段有層次的敘述，同樣資訊用更少字講完。
#: ⚠ 是「濃縮」不是「少講」——同類要合併，不是把後面幾條刪掉。
NARRATIVE_LAYERS = 3

#: （已停用）每條要點預設佔幾行。三層制改為由「總行數 ÷ 3」決定每段行數。
NARRATIVE_POINT_LINES = 2

#: 一條要點的標籤（含「｜」分隔）最多佔幾個字。
#: ⚠ 容量宣告要扣掉它——`_trim_blocks` 算行數時算的是「標籤｜正文」，
#: 只宣告正文字數會讓 CLI 照著寫卻放不下（2026-08-04 實機丟 5 條的根因）。
#: 取最長的標籤才對每一種都成立。三層（現況／意涵／後續）都是 2 字。
#: 🔴 2026-08-04：原本是「判讀限制」的 5 字；判讀限制移除後降為 3。
POINT_LABEL_COST = max(len(label) for label in ("現況", "意涵", "後續")) + 1


def points_budget(per_line: int, max_lines: int, columns: int) -> dict[str, int]:
    """一個要點框放得下幾條、每條正文幾個字。**容量宣告的唯一公式**。

    🔴 2026-08-04：`max_chars` 必須**扣掉標籤佔的字**。
    症狀：第五輪實機丟了 5 條要點，但同一輪 narrative 的契約警告是 0
    ——CLI 完全照容量寫，組版還是丟。
    根因：容量算的是「正文幾個字」，`_trim_blocks` 算的是
    `len(label) + 1 + len(text)`。正文寫滿 `per_line × 2`，再加「意涵｜」3 字
    就變 3 行而不是 2 行——條數一多總行數必然溢出，尾端整條被丟，
    而丟的都是排在後面的「意涵」「後續」，正是價值最高的幾條。

    ⚠ 取**最長**標籤來扣：容量宣告必須對每一種標籤都成立。
    ⚠ 這個公式只能有這一份——測試若自己重算一次，改了程式測試還是綠的。
    """
    lines_per_point = max(1, (max_lines * columns) // NARRATIVE_LAYERS)
    return {
        "max_points": NARRATIVE_LAYERS,
        "max_chars": max(1, per_line * lines_per_point - POINT_LABEL_COST),
    }


def point_line_ratio(theme: Theme) -> float:
    """要點文字的行距。**唯一定義處**——估算（`_text_capacity`）與渲染
    （`_add_number_bold_text` 的 `paragraph.line_spacing`）都讀它。

    🔴 2026-08-04：程式原本從未設定過 `line_spacing`，`qa.line_height_ratio`
    只是估算值，PowerPoint 實際用預設行距。⚠ 兩者分開就會出現「調大了估算值、
    畫面卻沒變寬」——只是把版面寫得更空。
    """
    return float(theme.qa.get("point_line_height_ratio", theme.qa["line_height_ratio"]))


def _points_area(theme: Theme, kind: str) -> tuple[float, float, int] | None:
    """該版型放要點的區域：(寬, 高, 欄數)；沒有要點區的版型回 None。

    🔴 `caveat`（H-2，2026-08-03）：有判讀限制框的頁，要點框會縮成
    `height_with_caveat_in`（3.3）而不是 `height_in`（5.0）——渲染端一直是這樣做的
    （`_render_points_panel`），但容量計算沒跟上，於是解讀 CLI 照 5.0 寫、
    塞進 3.3 的框、被截成「…」。⚠ 光扣掉「警語本身佔幾行」補不上 1.7 in 的框差。
    """
    if kind == "chart_hero":
        g = theme.geometry["chart_hero"]
        inset = g["panel_inset_in"]
        return (g["panel_width_in"] - inset - inset,
                g["panel_height_in"] - g["panel_text_top_offset_in"] - inset, 1)
    if kind == "table_with_points":
        g = theme.geometry["table_with_points"]
        inset = g["points_band_inset_in"]
        columns = int(g["points_band_columns"])
        gap = g["points_band_column_gap_in"]
        width = (g["points_band_width_in"] - inset - inset - gap * (columns - 1)) / columns
        return (width,
                g["points_band_height_in"] - g["points_band_text_top_offset_in"] - inset, columns)
    if kind == "chart_wide":
        g = theme.geometry["chart_wide"]
        inset = g["band_inset_in"]
        columns = int(g["band_columns"])
        gap = g["band_column_gap_in"]
        width = (g["band_width_in"] - inset - inset - gap * (columns - 1)) / columns
        # 高度取「圖最矮時橫幅能拿到的空間」——扁圖高度不一，取保守值才不會高估。
        height = g["band_bottom_in"] - g["band_min_top_in"] - g["band_text_top_offset_in"] - inset
        return (width, height, columns)
    if kind == "comparison":
        g = theme.geometry["comparison"]
        return (g["column_width_in"] - g["points_inset_right_in"],
                g["points_height_in"] - g["points_top_offset_in"] - g["points_bottom_pad_in"], 1)
    # ⚠ `table`（附錄）不在此列：附錄頁**不放要點**（`_render_table` 沒有要點區），
    # 誤把它算進來會用附錄的幾何覆蓋掉同一張報表在內頁的真實容量。
    if kind in {"chart_with_points", "percentage_bars", "stat_callout"}:
        g = theme.geometry["points_panel"]
        # 判讀限制移除後不再有「要為警語框讓高度」的情形，一律用全高。
        return (g["width_in"] - g["text_inset_right_in"],
                g["height_in"] - g["text_top_offset_in"] - g["text_bottom_pad_in"], 1)
    return None


def narrative_capacity(theme: Theme | None = None,
                       charts: ChartIndex | None = None,
                       report_data: dict[str, Any] | None = None) -> dict[str, dict[str, int]]:
    """每張報表的要點區實際容量：{report_key: {max_points, max_chars}}。

    ⚠ 為什麼要這個（2026-07-31）：解讀 CLI 原本只拿到一組全域上限（4–7 條 ×55 字），
    一刀切。但 `chart_hero` 的右側直欄與無圖表格頁的底部橫幅**可用空間差很多**，
    CLI 盲寫、`_trim_blocks` 事後裁切——截斷就是這樣來的。

    這裡由 `PAGE_LAYOUT`（哪張報表上哪種版型）＋ `theme.json` 幾何（那種版型的
    要點區多大）算出真實容量，**不新增任何常數**。撰寫（prompt）、驗證（validator）
    與裁切（_trim_blocks）三處因此吃同一份數字，不會再出現「說 55 字、實際只放得下
    40」的錯位。
    """
    theme = theme or Theme.load()
    size = theme.size("point_text_pt")
    capacity: dict[str, dict[str, int]] = {}
    # 🔴 K-1（2026-08-04）：能拿到 report_data 就用**最終版面**逐頁算——
    # lifecycle／公司×國家／家族佈局／年度矩陣是動態插頁，不在 PAGE_LAYOUT 裡，
    # 舊迴圈完全沒給它們容量 → 契約退回全域 55 字、版面實際 33 字 →
    # CLI 照 55 寫、組版把「後續」整條丟掉（第六輪實機 p3/p5/p6/p15）。
    # resolve_layout 已含動態插頁、政策拆頁與扁圖長寬比判定，與渲染同一份決策。
    specs: list[PageSpec]
    if report_data is not None and charts is not None:
        specs = [spec for spec in resolve_layout(report_data, charts, theme, {})
                 if not spec.is_appendix and spec.kind not in {"cover", "direction"}]
    else:
        specs = [spec for spec in PAGE_LAYOUT if not spec.is_appendix]
    for spec in specs:
        # ⚠ 用**實際渲染時**的版型算，不是宣告的版型：列在 SPLIT_PAIR_REPORTS 的
        # comparison 頁會被 `_split_pairs_by_policy` 拆成一圖一頁的 chart_hero，
        # 拿 comparison 的窄長條去算會嚴重低估（實測只算得出 1 條）。
        kind = spec.kind
        # ⚠ 這裡要**重現執行時的版型決策順序**：先拆頁（政策拆或圖數溢出），
        # 再依長寬比選滿寬版型。只做後者會被「comparison 不在 SINGLE_CHART_KINDS」
        # 的守門條件擋掉，扁圖頁的容量就會沿用窄側欄、CLI 因此寫得比實際能放的少。
        chart_names = tuple(charts.files_for(spec.report_keys)) if charts is not None else ()
        if kind == "comparison" and (
            any(key in SPLIT_PAIR_REPORTS for key in spec.report_keys)
            or len(chart_names) > len(theme.geometry["comparison"]["column_left_in"])
        ):
            kind = "chart_hero"
        # ⚠ `chart_wide` 是**執行時依圖的長寬比**決定的，不是宣告在 PAGE_LAYOUT 裡。
        # 拿不到圖檔時只能用宣告版型估——那會低估扁圖頁（滿寬雙欄橫幅比窄側欄大得多），
        # CLI 因此寫得比實際能放的少。給了 charts 就能算準。
        def _limits_for(page_kind: str) -> dict[str, int] | None:
            """某個版型下，這一頁的要點容量。

            ⚠ 判讀限制移除後（2026-08-04），不再有「警語框先佔掉高度與行數」
            這兩層扣減——框一律是全高，要點三段平分。
            """
            area = _points_area(theme, page_kind)
            if area is None:
                return None
            width_in, height_in, columns = area
            per_line, max_lines = _text_capacity(
                theme, width_in=width_in, height_in=height_in, size_pt=size,
                line_ratio=point_line_ratio(theme))
            return points_budget(per_line, max_lines, columns)

        # 🔴 I-1（2026-08-03 實機 #166）：容量必須**逐 variant** 算。
        #
        # 同一個 report_key 的不同 variant 會落在不同版型：IPC 的 L4 是扁圖
        # （chart_wide，底部雙欄橫幅，8 條 × 54 字），L5 是一般圖
        # （chart_hero，右側窄欄，7 條 × 26 字）。原本的迴圈是
        # 「只要**任一張**圖是扁的，整個 report_key 就用 chart_wide 算」，
        # 於是 CLI 拿到 8×54 照著寫，L4 放得下、**L5 每條要 3 行**、總行數爆掉
        # ——實機丟了 10 條要點，其中 p8（IPC L5）4 條、p10（CPC L5）3 條。
        #
        # ⚠ 不走「同一 key 取最小容量」：那會讓 L4 那種寬頁寫得比能放的少而變空，
        # 等於推翻 C-9（使用者：「要的是濃縮不是丟棄」）。
        variant_kinds: dict[str, str] = {}
        widest_kind = kind
        if charts is not None:
            for name in chart_names:
                name_kind = _kind_for_aspect(kind, (name,), charts)
                stem = Path(name).stem
                for suffix in CHART_ORDER_HINTS:
                    if stem.endswith(suffix):
                        variant_kinds[suffix.lstrip("_")] = name_kind
                        break
                if name_kind == "chart_wide":
                    widest_kind = "chart_wide"

        base_limits = _limits_for(widest_kind)
        if base_limits is None:
            continue
        for key in spec.report_keys:
            # report_key 層保留：沒有 variant 的報表、以及舊資料的 fallback。
            capacity[key] = base_limits
            for variant, variant_kind in variant_kinds.items():
                limits = _limits_for(variant_kind)
                if limits is not None:
                    capacity[f"{key}:{variant}"] = limits
    return capacity


def all_slot_keys() -> list[str]:
    """回傳 PPT 階段 AI 需要產的全部文案槽（唯一來源＝PAGE_LAYOUT）。

    v3 只剩 `cover.title` 與 `direction.body`：其餘頁面文字一律來自 narratives，
    不再讓 AI 為每頁另產一份，避免同一段判讀在兩處各寫一次而互相矛盾。
    """
    keys: list[str] = []
    for spec in PAGE_LAYOUT:
        keys.extend(spec.slots)
    return keys


# --------------------------------------------------------------------------
# 樣式
# --------------------------------------------------------------------------
@dataclass
class Theme:
    """外觀樣式；改配色、字級、座標只改 theme.json，不動本程式。"""

    font: dict[str, Any]
    color: dict[str, str]
    geometry: dict[str, Any]
    slide: dict[str, float]
    qa: dict[str, Any]
    # v3 深空主題：背景漸層兩色與角度、星空紋理生成參數、圖表轉色與裁切規則。
    gradient: dict[str, Any]
    starfield: dict[str, Any]
    chart_recolor: dict[str, Any]

    @classmethod
    def load(cls, path: Path | str = THEME_PATH) -> Theme:
        data = json.loads(Path(path).read_text(encoding="utf-8"))
        return cls(
            font=data["font"],
            color=data["color"],
            geometry=data["geometry"],
            slide=data["slide"],
            qa=data["qa"],
            gradient=data["gradient"],
            starfield=data["starfield"],
            chart_recolor=data["chart_recolor"],
        )

    def rgb(self, name: str) -> RGBColor:
        return RGBColor.from_string(self.color[name])

    def size(self, name: str) -> float:
        """取字級並套用 12pt 下限；下限在 theme 內宣告，程式只負責遵守。"""
        return max(float(self.font[name]), float(self.font["min_pt"]))


# --------------------------------------------------------------------------
# 文字量估算：文字框裝不裝得下要能事前判斷（fallback 截斷）與事後自檢（QA）
# --------------------------------------------------------------------------
def _text_capacity(theme: Theme, *, width_in: float, height_in: float, size_pt: float,
                   line_ratio: float | None = None) -> tuple[int, int]:
    """回傳（每行字數, 可用行數）。中文字寬約等於字級，故以 pt/72 估字寬。

    ⚠ `line_ratio` 給**有設段落行距**的文字用（目前只有要點，見
    `POINT_LINE_RATIO`）。其餘文字沒設 `line_spacing`，PowerPoint 用預設行距，
    所以只能沿用 `qa.line_height_ratio` 這個估算值——把它全域調大，
    畫面不會變寬，只會讓容量估得更保守（字更少、版面更空）。
    """
    char_in = size_pt / 72.0 * float(theme.qa["cjk_char_width_ratio"])
    ratio = line_ratio if line_ratio is not None else float(theme.qa["line_height_ratio"])
    line_in = size_pt / 72.0 * ratio
    # ⚠ 加 epsilon：1.5 / (40/72*1.35) 在浮點下是 1.9999999998，直接 int() 會少算一行，
    # 讓剛好兩行的標題被誤判成裝不下而截字（封面標題被切成「…」就是這樣來的）。
    epsilon = 1e-6
    per_line = max(1, int(width_in / char_in + epsilon))
    lines = max(1, int(height_in / line_in + epsilon))
    return per_line, lines


def _fit_text(theme: Theme, text: str, *, width_in: float, height_in: float, size_pt: float) -> tuple[str, bool]:
    """把文字截到框內裝得下，超出加「…」。回傳（文字, 是否截斷）。"""
    per_line, lines = _text_capacity(theme, width_in=width_in, height_in=height_in, size_pt=size_pt)
    budget = per_line * lines
    if len(text) <= budget:
        return text, False
    return text[: max(1, budget - 1)].rstrip("，、。；：") + "…", True


#: 半形英數的字寬（em）。
#: 🔴 2026-08-03：原本 0.55，與 `chart_runner._display_width` **各寫一份**，
#: 而後者已於 I-3 依實測改為 0.62（轉圖掃像素量到真實字寬比估算多約 13%）
#: ——同一個估算兩處落點、只改了一邊，於是表格欄寬照舊被低估，
#: 實機 p22 把專利號 `121754861` 折成兩行（本專案第 8 次兩處落點）。
#: ⚠ 兩處必須同值；有測試 `test_display_width_matches_chart_runner` 釘住。
ALNUM_EM_WIDTH = 0.62


def _display_width(text: str) -> float:
    """字串的顯示寬度（em）。中文、全形符號約 1 em，半形英數約 `ALNUM_EM_WIDTH`。

    表格欄位混排 `applicant_display_name` 與中文公司名，一律當全形算會把英文
    表頭砍成一半；一律當半形算又會讓中文撐爆欄寬。
    """
    return sum(ALNUM_EM_WIDTH if ord(ch) < 0x2E80 else 1.0 for ch in text)


#: 儲存格內視為「不可拆」的 token 分隔符。
#: ⚠ 專利號、代碼這類 token 一旦被折行，語意就毀了（`121754861` → `12175486`／`61`
#: 會被讀成兩個號碼），與一般文字換行不同——欄寬必須保障它們完整。
_TOKEN_SEPARATORS = ("、", "；", ";", " ")


def _longest_token_em(text: str) -> float:
    """字串中最長的不可拆 token 寬度（em）。"""
    parts = [text]
    for sep in _TOKEN_SEPARATORS:
        parts = [piece for part in parts for piece in part.split(sep)]
    return max((_display_width(p) for p in parts if p), default=0.0)


def _column_widths(
    columns: list[str],
    rows: list[dict[str, Any]],
    labels: dict[str, str],
    total_width_in: float,
    *,
    size_pt: float,
    inset_in: float,
) -> list[float]:
    """依內容需求分配欄寬，總和等於表寬。

    🔴 2026-08-02：原本一律 `width / len(columns)` 等分，兩個方向都出錯——
    p21「前三大申請人」被截 10 處，p22「最新受讓人名單」14 列只有 2 列有值
    卻和其他欄一樣寬。⚠ 截斷是症狀，等分才是根因：兩位數的「專利件數」
    和一長串名單拿一樣的寬度。

    分配方式：先算每欄的需求寬（欄頭與內容取較寬者），再按需求比例分配；
    ⚠ 但每欄至少保住**自己的欄頭**——欄頭被截比內容被截更難懂
    （讀者連這欄在講什麼都不知道）。欄頭需求超過等分寬時才讓步到等分寬，
    否則欄多時會把整張表擠爆。
    """
    if not columns:
        return []
    if len(columns) == 1:
        return [total_width_in]

    per_char = size_pt / 72.0
    padding = inset_in * 2
    equal = total_width_in / len(columns)

    demands: list[float] = []
    minimums: list[float] = []
    for name in columns:
        header = _display_width(labels.get(str(name), str(name))) * per_char + padding
        content = 0.0
        for row in rows:
            value = row.get(name)
            if isinstance(value, list):
                value = "、".join(str(v) for v in value)
            content = max(content, _display_width("" if value is None else str(value)) * per_char)
        demands.append(max(header, content + padding))
        # 🔴 最長不可拆 token 也是下限（2026-08-03 實機 p22）：
        # 專利號 `121754861` 被折成 `12175486`／`61`，讀者會當成兩個號碼。
        # 一般文字換行沒關係，**token 斷開語意就毀了**。
        token = 0.0
        for row in rows:
            value = row.get(name)
            if isinstance(value, list):
                value = "、".join(str(v) for v in value)
            token = max(token, _longest_token_em("" if value is None else str(value)) * per_char)
        # ⚠ 欄頭寬需求可能本身就超過等分寬（欄多時常見），此時只能讓步到等分寬，
        # 否則所有欄的下限加起來會超出表寬。token 下限同理設上限。
        minimums.append(min(max(header, token + padding), equal))

    floor_total = sum(minimums)
    spare = total_width_in - floor_total
    if spare <= 0:
        scale = total_width_in / floor_total
        return [m * scale for m in minimums]

    # 剩餘寬度按「超出下限的需求」比例分——需求大的欄拿得多，全滿足後仍有剩就平均補。
    extra = [max(0.0, d - m) for d, m in zip(demands, minimums)]
    extra_total = sum(extra)
    if extra_total <= 0:
        return [m + spare / len(columns) for m in minimums]
    ratio = min(1.0, spare / extra_total)
    widths = [m + e * ratio for m, e in zip(minimums, extra)]
    leftover = total_width_in - sum(widths)
    if leftover > 0:
        widths = [w + leftover / len(columns) for w in widths]
    return widths


def _truncate_to_width(text: str, width_in: float, size_pt: float) -> str:
    """把字串截到指定英吋寬度內（超出加「…」），用於表格儲存格避免自動換行撐高列。

    ⚠ 截點優先落在標點（2026-07-31）：原本切在剛好超寬的那個字，會產生
    「全數歸在 A…」這種斷在詞中間的句子。改成先找容量內最後一個標點，
    切在它之後——讀起來是「一句話沒講完」而不是「字被砍斷」。
    找不到標點（例如整串是公司名）才退回原本的硬切。
    """
    budget = width_in / (size_pt / 72.0)
    if _display_width(text) <= budget:
        return text
    used, kept = 0.0, []
    for ch in text:
        step = 0.55 if ord(ch) < 0x2E80 else 1.0
        if used + step > budget - 1:
            break
        used += step
        kept.append(ch)
    clipped = "".join(kept)
    cut = max((clipped.rfind(mark) for mark in TRUNCATE_BREAK_MARKS), default=-1)
    # 太靠前的標點不採用（只留半截等於沒講），門檻取容量的六成。
    if cut >= len(clipped) * 0.6:
        clipped = clipped[: cut + 1]
    return clipped + "…"


def _lines_needed(text: str, per_line: int) -> int:
    """多段文字所需行數：每段至少佔一行，不足一行不合併。"""
    return sum(max(1, math.ceil(len(line) / per_line)) for line in text.split("\n"))


# --------------------------------------------------------------------------
# 圖檔對照：唯一來源＝artifact_manifest.json（禁止用 report_key 猜檔名）
# --------------------------------------------------------------------------
class ChartIndex:
    """report_name → 圖檔清單的反查表，並負責 SVG 轉點陣與快取。

    實際檔名與 report_key 不同名（`country_distribution` 的圖叫
    `jurisdiction_distribution.svg`），且同一 report_name 可能對應多張圖
    （IPC 的 L4/L5、機會矩陣的技術面/功效面）——這正是成對呈現的來源。
    """

    def __init__(self, report_dir: Path, cache_dir: Path, manifest: dict[str, Any] | None = None,
                 theme: Theme | None = None) -> None:
        self.report_dir = report_dir
        self.cache_dir = cache_dir
        # theme 供圖表深色轉色與白邊裁切用；None＝維持原樣只轉檔（前端預覽走這條）。
        self.theme = theme
        self._by_report: dict[str, list[str]] = {}
        self._raster_cache: dict[str, Path | None] = {}
        self._aspect_cache: dict[str, float | None] = {}
        self.manifest_found = bool(manifest)
        for artifact in (manifest or {}).get("artifacts", []) or []:
            if not isinstance(artifact, dict):
                continue
            # 上游欄位曾用 file／filename、sha256／hash 兩種寫法，兩種都吃。
            name = str(artifact.get("file") or artifact.get("filename") or "")
            if not name or not name.lower().endswith(IMAGE_SUFFIXES):
                continue
            if not (self.report_dir / name).exists():
                continue
            targets = list(artifact.get("report_names") or [])
            if artifact.get("report_name"):
                targets.append(str(artifact["report_name"]))
            for target in {str(t) for t in targets if t}:
                bucket = self._by_report.setdefault(target, [])
                if name not in bucket:
                    bucket.append(name)

    @staticmethod
    def _order_key(name: str) -> tuple[int, str]:
        for index, hint in enumerate(CHART_ORDER_HINTS):
            if hint in name:
                return index, name
        return len(CHART_ORDER_HINTS), name

    def files_for(self, report_keys: tuple[str, ...]) -> tuple[str, ...]:
        """依 report_key 順序取圖檔；同一 report_key 多張圖時保持穩定排序。"""
        found: list[str] = []
        for key in report_keys:
            for name in sorted(self._by_report.get(key, []), key=self._order_key):
                if name not in found:
                    found.append(name)
        return tuple(found)

    def owners_of(self, chart_name: str, keys: tuple[str, ...]) -> tuple[str, ...]:
        """某張圖對應到候選 report_key 中的哪幾個；拆頁時用來把 report_key 一起收窄。"""
        return tuple(key for key in keys if chart_name in self._by_report.get(key, []))

    def aspect_of(self, name: str) -> float | None:
        """圖檔的原始長寬比（寬/高）；讀不到回 None。

        版型選擇要看**圖本身的形狀**：同一個框塞得下 0.78 到 7.42 的比例，
        差 9.5 倍，沒有一種框能同時服務兩端（2026-07-31 實測）。
        """
        if name in self._aspect_cache:
            return self._aspect_cache[name]
        source = self.report_dir / name
        aspect = None
        if source.exists() and source.suffix.lower() == ".svg":
            head = source.read_text(encoding="utf-8", errors="ignore")[:400]
            match = SVG_SIZE_PATTERN.search(head)
            if match:
                width, height = float(match.group(1)), float(match.group(2))
                aspect = width / height if height else None
        self._aspect_cache[name] = aspect
        return aspect

    def resolve(self, name: str) -> Path | None:
        """把圖檔轉成 python-pptx 吃得下的點陣檔；SVG 轉 PNG 只轉一次。"""
        if name in self._raster_cache:
            return self._raster_cache[name]
        source = self.report_dir / name
        if not source.exists():
            self._raster_cache[name] = None
            return None
        if source.suffix.lower() == ".svg":
            resolved = globals()["rasterize_svg"](source, self.cache_dir, self.theme)
        else:
            resolved = source
        self._raster_cache[name] = resolved
        return resolved


BACKGROUND_RECT_PATTERN = re.compile(
    r'<rect[^>]*?fill="(?:white|#fff|#ffffff|#FFF|#FFFFFF)"[^>]*?/>', re.I)
SVG_SIZE_PATTERN = re.compile(r'<svg[^>]*?width="(\d+(?:\.\d+)?)"[^>]*?height="(\d+(?:\.\d+)?)"', re.I)
RASTER_DPI = 150


def strip_chart_title(svg_text: str) -> str:
    """移除 SVG 內建的圖表標題（引擎標了 `data-role="chart-title"` 的那一行）。

    🔴 F-8：投影片上面是 narrative 的 headline、下面是 SVG 自己畫的
    「IPC 主分類分布 - Level 4」，兩行講同一件事（實機九頁皆然）。

    ⚠ 不在引擎端砍——網頁報表頁讀的是**同一份 SVG**，那裡沒有頁標題、需要它。
    故引擎只負責標記，移除是組版端的事：同一份資料、兩種呈現。
    """
    return re.sub(r'<text[^>]*data-role="chart-title"[^>]*>.*?</text>\s*', "", svg_text, flags=re.S)


def recolor_svg(svg_text: str, recolor: dict[str, Any]) -> str:
    """把引擎產的淺底圖表 SVG 換成深空配色（PPT 端轉色，不動引擎）。

    ⚠ 為什麼在這裡轉而不是讓引擎產深色：同一份 SVG 也內嵌在**網頁報表頁**
    （淺底），引擎改深色會讓網頁那邊變成深底深字。使用者 2026-07-31 選定此方案。
    ⚠ 只換顏色，不動 viewBox、不裁切內容、不拉伸——圖表形式一律照引擎產的樣子
    （使用者定案：「不能像 NotebookLM 把圖表改形式」）。
    """
    if recolor.get("strip_background"):
        # 整版白底矩形必須先拿掉，否則深色頁上會出現一塊白板。
        svg_text = BACKGROUND_RECT_PATTERN.sub("", svg_text, count=1)
    mapping = {k.upper(): v for k, v in (recolor.get("map") or {}).items()}
    for old, new in mapping.items():
        svg_text = re.sub(re.escape(f"#{old}"), f"#{new}", svg_text, flags=re.I)
    return _recolor_paired_text(svg_text, mapping)


def _recolor_paired_text(svg_text: str, mapping: dict[str, str]) -> str:
    """依**轉色後的底色**重算「畫在圖元上的文字」顏色。

    🔴 引擎本來就會自動算對比色（`_chip_text_color`／`readable_text_on`），但那是
    對**原始**淺色主題的底色算的。本模組把底色換成深空配色之後，字色沒跟著變——
    實測象限 chip 白字掉到 1.44、泡泡數字 1.24，畫面上實質看不見。

    ⚠ 單靠字串替換無從得知「這段白字疊在哪個底上」，所以由引擎在 SVG 標
    `data-on-fill="<原始底色>"`；這裡讀它、查出新底色、重算字色。
    沒有標記的文字（座標軸、標題）不動——它們畫在頁面底上，不是圖元上。
    """
    def _swap(match: re.Match[str]) -> str:
        element, source = match.group(0), match.group(1).upper().lstrip("#")
        # ⚠ 走到這裡時 data-on-fill 的值**已被前面的全域替換換成新色**，
        # 查表查不到是正常的——此時它本身就是新底色。不要因為「查不到」
        # 就當成錯誤，也不要調換兩者順序：先算字色再換底色會讓標記失效。
        new_fill = mapping.get(source, source)
        return re.sub(r'fill="#?[0-9A-Fa-f]{6}"', f'fill="{_readable_on(new_fill)}"',
                      element, count=1)

    return re.sub(r'<text[^>]*?data-on-fill="([^"]+)"[^>]*?>', _swap, svg_text)


def _readable_on(fill: str) -> str:
    """深底用亮字、淺底用深字（WCAG 相對亮度 0.4 為界，兩側皆 ≥4.5）。"""
    value = fill.lstrip("#")
    if len(value) != 6:
        return "#FFFFFF"
    channels = []
    for offset in (0, 2, 4):
        component = int(value[offset:offset + 2], 16) / 255
        channels.append(component / 12.92 if component <= 0.03928
                        else ((component + 0.055) / 1.055) ** 2.4)
    luminance = 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]
    return "#132C44" if luminance > 0.4 else "#FFFFFF"


def _crop_box(png_path: Path, padding: int) -> tuple[int, int, int, int] | None:
    """圖檔實際內容的邊界（含安全邊）；整張都是空的回 None。

    去掉白邊是使用者 2026-07-31 授權的（「裁切就是邊邊白白的可以切」）——
    這不是改圖表形式，是不再把死空間當成圖的一部分塞進版面。實測一張趨勢圖
    左右上下共有 320 px 是空的，切掉後同樣的框內圖表內容線性大 25%。
    """
    try:
        from PIL import Image  # python-pptx 本來就依賴 Pillow，不是新增相依
    except ImportError:
        return None
    with Image.open(png_path) as image:
        box = image.getbbox()  # RGBA：回傳非透明區域
        if box is None:
            return None
        left, top, right, bottom = box
        width, height = image.size
    return (max(0, left - padding), max(0, top - padding),
            min(width, right + padding), min(height, bottom + padding))


def prepare_chart(svg_path: Path, cache_dir: Path, theme: Theme) -> tuple[Path | None, Path | None]:
    """圖表 SVG → (PNG, 轉色裁切後的 SVG)；轉不動時 PNG 回 None 由呼叫端降級。

    🔴 兩個回傳值必須**同源**：PNG 是後援與預覽用、SVG 供 PowerPoint 顯示向量。
    若 PNG 用轉色版、SVG 用原始版，會出現「縮圖淺色、放大變深色」的錯位，
    而且極難察覺（多數人不會放大到觸發 SVG 渲染）。故兩者都由同一份轉色＋
    同一個裁切框產出。
    """
    cache_dir.mkdir(parents=True, exist_ok=True)
    recolor = getattr(theme, "chart_recolor", {}) or {}
    source = svg_path.read_text(encoding="utf-8")
    # F-8：上投影片前拿掉圖表自帶標題（頁標題已經在說同一件事）。
    dark = recolor_svg(strip_chart_title(source), recolor)
    digest = hashlib.sha256(dark.encode("utf-8")).hexdigest()[:16]
    png_path = cache_dir / f"{svg_path.stem}_{digest}.png"
    out_svg = cache_dir / f"{svg_path.stem}_{digest}.svg"
    if png_path.exists() and out_svg.exists():
        return png_path, out_svg

    try:
        import pymupdf  # type: ignore
    except ImportError:
        try:
            import fitz as pymupdf  # type: ignore
        except ImportError:
            return None, None

    staged = cache_dir / f"{svg_path.stem}_{digest}_src.svg"
    staged.write_text(dark, encoding="utf-8")
    try:
        doc = pymupdf.open(str(staged))
        # alpha=True：去掉白底後背景要**透明**才貼得上深色頁，也才裁得出內容邊界。
        doc[0].get_pixmap(dpi=RASTER_DPI, alpha=True).save(str(png_path))
        doc.close()
    except Exception:
        return None, None

    box = _crop_box(png_path, int(recolor.get("crop_padding_px") or 0))
    out_svg.write_text(_cropped_svg(dark, png_path, box), encoding="utf-8")
    if box is not None:
        try:
            from PIL import Image

            with Image.open(png_path) as image:
                image.crop(box).save(png_path)
        except Exception:
            pass  # 裁不動就用未裁的，畫面只是留白多一點，不該讓整份簡報失敗
    return png_path, out_svg


def rasterize_svg(svg_path: Path, cache_dir: Path, theme: Theme | None = None) -> Path | None:
    """SVG → PNG（python-pptx 不吃 SVG）。轉不動時回 None，由呼叫端降級處理。

    ⚠ 保留這個名字與單一回傳值是為了呼叫端（`ChartIndex.resolve`）與既有測試不必改。
    有 theme 時走 `prepare_chart`（深色轉色＋白邊裁切＋同時產向量版），
    沒有 theme（例如前端預覽只拿得到 report_data）就只轉檔、不改外觀。
    """
    if theme is not None:
        png_path, _ = prepare_chart(svg_path, cache_dir, theme)
        return png_path
    cache_dir.mkdir(parents=True, exist_ok=True)
    digest = hashlib.sha256(svg_path.read_bytes()).hexdigest()[:16]
    png_path = cache_dir / f"{svg_path.stem}_{digest}_plain.png"
    if png_path.exists():
        return png_path
    try:
        import pymupdf  # type: ignore
    except ImportError:
        try:
            import fitz as pymupdf  # type: ignore
        except ImportError:
            return None
    try:
        doc = pymupdf.open(str(svg_path))
        doc[0].get_pixmap(dpi=RASTER_DPI).save(str(png_path))
        doc.close()
    except Exception:
        return None
    return png_path


def _cropped_svg(svg_text: str, png_path: Path, box: tuple[int, int, int, int] | None) -> str:
    """把裁切框換算回 SVG 使用者座標並改寫 viewBox，讓向量版與 PNG 切在同一處。"""
    size = SVG_SIZE_PATTERN.search(svg_text)
    if box is None or size is None:
        return svg_text
    try:
        from PIL import Image

        with Image.open(png_path) as image:
            png_width = image.size[0]
    except Exception:
        return svg_text
    svg_width = float(size.group(1))
    svg_height = float(size.group(2))
    if not png_width or not svg_width:
        return svg_text
    scale = png_width / svg_width
    left, top, right, bottom = (value / scale for value in box)
    width = max(1.0, right - left)
    height = max(1.0, bottom - top)
    header = size.group(0)
    replaced = re.sub(r'width="[\d.]+"', f'width="{width:.1f}"', header, count=1)
    replaced = re.sub(r'height="[\d.]+"', f'height="{height:.1f}"', replaced, count=1)
    replaced = f'{replaced} viewBox="{left:.1f} {top:.1f} {width:.1f} {height:.1f}"'
    # 原本就有 viewBox 的話先移除，避免同一個標籤出現兩個。
    svg_text = re.sub(r'\s+viewBox="[^"]*"', "", svg_text, count=1)
    _ = svg_height
    return svg_text.replace(header, replaced, 1)


# --------------------------------------------------------------------------
# report_data 取值
# --------------------------------------------------------------------------
def _iter_report_entries(report_data: dict[str, Any]) -> list[tuple[str, dict[str, Any]]]:
    entries: list[tuple[str, dict[str, Any]]] = []
    for bucket in ("reports", "family_reports"):
        section = report_data.get(bucket) or {}
        if isinstance(section, dict):
            entries.extend((k, v) for k, v in section.items() if isinstance(v, dict))
    return entries


def _entry_of(report_data: dict[str, Any], report_key: str) -> dict[str, Any]:
    for bucket in ("reports", "family_reports"):
        entry = (report_data.get(bucket) or {}).get(report_key)
        if isinstance(entry, dict):
            return entry
    return {}


def _rows_of(report_data: dict[str, Any], report_key: str) -> list[dict[str, Any]]:
    rows = _entry_of(report_data, report_key).get("rows")
    return rows if isinstance(rows, list) else []


def _label_of(report_data: dict[str, Any], report_key: str, fallback: str = "") -> str:
    entry = _entry_of(report_data, report_key)
    return str(entry.get("label_zh") or entry.get("label") or fallback or report_key)


def threshold_skips(report_data: dict[str, Any]) -> list[dict[str, str]]:
    """低於出頁門檻的分類報表清單（供 manifest 記缺頁原因）。

    🔴 缺頁不得靜默（design #5「門檻與缺頁原因進 metadata」）：頁面消失而
    manifest 無痕，讀者只會以為漏產。舊版本沒有 `classification_thresholds`
    鍵 → 回空清單，行為不變。
    """
    thresholds = report_data.get("classification_thresholds") or {}
    return [
        {"type": "below_threshold_skipped", "report_key": key,
         "reason": str(info.get("reason") or "低於出頁門檻")}
        for key, info in thresholds.items()
        if isinstance(info, dict) and info.get("below_threshold")
    ]


def _below_threshold(report_data: dict[str, Any], report_key: str) -> bool:
    info = (report_data.get("classification_thresholds") or {}).get(report_key)
    return bool(isinstance(info, dict) and info.get("below_threshold"))


def _report_key_has_data(report_data: dict[str, Any], report_key: str) -> bool:
    # 🔴 IPC/CPC 出頁門檻（2026-08-05 定案「4 階沒有 3 種以上就不出現在簡報」）：
    # below_threshold 的報表視同無資料——這裡是**唯一接縫**（固定頁、動態插頁、
    # 拆頁全走本函式），別處不再各自判斷。網頁報表不受影響（引擎照產）。
    if _below_threshold(report_data, report_key):
        return False
    entry = _entry_of(report_data, report_key)
    rows = entry.get("rows")
    if isinstance(rows, list) and rows:
        return True
    try:
        return int(entry.get("row_count") or 0) > 0
    except (TypeError, ValueError):
        return False


def _actual_report_keys(report_data: dict[str, Any], keys: tuple[str, ...]) -> tuple[str, ...]:
    return tuple(key for key in keys if _report_key_has_data(report_data, key))


def _numeric_column(rows: list[dict[str, Any]]) -> str:
    """挑出 rows 內代表件數的數值欄；優先 patent_count，其次第一個純數字欄。"""
    if not rows:
        return ""
    if "patent_count" in rows[0]:
        return "patent_count"
    for name, value in rows[0].items():
        if isinstance(value, (int, float)) and not isinstance(value, bool):
            return str(name)
    return ""


def _label_column(rows: list[dict[str, Any]], numeric: str) -> str:
    """挑出 rows 內代表項目名稱的欄。

    ⚠ 優先 `label`／顯示名欄，並跳過排除欄——舊版取「第一個非數值欄」，
    分群 rows 第一欄是 topic_code，降級要點就印出「T001 15 件」（代碼入文，
    2026-07-31 實機驗收抓到）。
    """
    if not rows:
        return ""
    for preferred in ("label", "applicant_display_name", "owner_display_name"):
        if preferred in rows[0]:
            return preferred
    for name in rows[0]:
        if str(name) != numeric and str(name) not in TABLE_EXCLUDED_COLUMNS:
            return str(name)
    return str(next(iter(rows[0])))


def _as_int(value: Any) -> int:
    try:
        return int(float(value))
    except (TypeError, ValueError):
        return 0


def _statistics_period(report_data: dict[str, Any]) -> str:
    """統計期間：先讀 parameters 明示值，沒有才由趨勢報表的年份區間推得。"""
    params = report_data.get("parameters") or {}
    for key in ("period", "date_range", "statistics_period"):
        if params.get(key):
            return str(params[key])
    years: list[int] = []
    for report_key in ("application_trend", "publication_trend", "lifecycle"):
        for row in _rows_of(report_data, report_key):
            for field in ("year", "application_year", "publication_year"):
                if field in row:
                    value = _as_int(row[field])
                    if value:
                        years.append(value)
    if not years:
        return ""
    low, high = min(years), max(years)
    return str(low) if low == high else f"{low}–{high}"


# --------------------------------------------------------------------------
# narrative：headline／points 三件套，缺就 fallback 並記 warning
# --------------------------------------------------------------------------
def _narrative_entry(narratives: dict[str, Any], candidates: tuple[str, ...]) -> tuple[str, dict[str, Any]]:
    """依候選鍵找 narrative；回傳（命中的鍵, variant 內容）。

    narratives 的鍵混用 report_key 與圖檔名（申請趨勢的解讀掛在 `annual_trend`），
    故候選鍵同時含 report_key、圖檔主檔名與去掉 _L4/_tech 等變體後綴的主檔名。
    """
    reports = narratives.get("reports") or {}
    for key in candidates:
        # 「report:variant」語法（P1-2 alias）：直接指定變體，不掃整個 entry。
        if ":" in key:
            report_key, _, variant_key = key.partition(":")
            variant = ((reports.get(report_key) or {}).get("variants") or {}).get(variant_key)
            if isinstance(variant, dict) and (variant.get("points") or variant.get("text") or variant.get("headline")):
                return key, variant
            continue
        entry = reports.get(key)
        if not isinstance(entry, dict):
            continue
        variants = entry.get("variants") or {}
        for variant_key in ("default", *variants):
            variant = variants.get(variant_key)
            if isinstance(variant, dict) and (variant.get("points") or variant.get("text") or variant.get("headline")):
                return key, variant
    return "", {}


def _split_legacy_text(text: str) -> list[dict[str, Any]]:
    """把舊格式長文依「觀察／意涵／決策提醒」等段落標記切成條列。

    這是 fallback 的可讀性補救：上游還沒升級 headline/points 契約前，
    直接把 400–500 字整段塞進版面就是字牆（實機驗收不合格的主因之一）。
    切分只依 AI 自己寫的段落標記，不新增、不改寫任何內容。
    """
    pattern = "|".join(NARRATIVE_MARKERS)
    chunks = [c.strip() for c in re.split(rf"(?=(?:{pattern})\s*[：:])", text) if c.strip()]
    points: list[dict[str, Any]] = []
    for chunk in chunks:
        match = re.match(rf"^({pattern})\s*[：:]\s*(.+)$", chunk, flags=re.DOTALL)
        if match:
            points.append({"label": match.group(1), "text": match.group(2).strip()})
        else:
            points.append({"label": "", "text": chunk})
    if len(points) > 1:
        return points
    # 找不到段落標記：依句號切成最多三條，仍好過單段字牆。
    sentences = [s.strip() for s in re.split(r"(?<=[。！？])", text) if s.strip()]
    if len(sentences) <= 1:
        return [{"label": "", "text": text.strip()}]
    size = math.ceil(len(sentences) / 3)
    return [
        {"label": "", "text": "".join(sentences[i : i + size])}
        for i in range(0, len(sentences), size)
    ]


HEADLINE_MAX_CHARS = 26


def _derive_headline(points: list[dict[str, Any]]) -> str:
    """缺 headline 時，從第一條要點取一個短句當判讀式標題。

    這是**選取**而不是生成：句子完全是 narrative 自己寫的，程式只切第一個句讀
    並檢查長度，切不出夠短的就不硬湊（標題退回報表主題）。裸名詞標題讀者看不出
    這頁在講什麼，但捏造判讀更糟——所以只在有現成句子可用時才升級標題，
    且一律寫 warning 讓人追得到這個標題是推導來的。
    """
    if not points:
        return ""
    text = str(points[0].get("text") or "").strip()
    for separator in ("。", "；", "，", "、"):
        head = text.split(separator)[0].strip()
        if head and len(head) <= HEADLINE_MAX_CHARS:
            return head
    return ""


def _normalize_narrative(variant: dict[str, Any]) -> tuple[str, list[dict[str, Any]], bool]:
    """回傳（headline, points, 是否走 fallback）。缺 headline 不自行編造。"""
    headline = str(variant.get("headline") or "").strip()
    raw_points = variant.get("points")
    if isinstance(raw_points, list) and raw_points:
        points = [
            {
                "label": str(p.get("label") or "").strip(),
                "text": str(p.get("text") or "").strip(),
                "emphasis": bool(p.get("emphasis")),
            }
            for p in raw_points
            if isinstance(p, dict) and str(p.get("text") or "").strip()
        ]
        if points:
            return headline, points, False
    text = str(variant.get("text") or "").strip()
    if not text:
        return headline, [], True
    return headline, [{**p, "emphasis": False} for p in _split_legacy_text(text)], True


# --------------------------------------------------------------------------
# 組版輔助
# --------------------------------------------------------------------------
def _set_font(run, theme: Theme, *, size: float, color: str, bold: bool) -> None:
    """套字體與字級。

    ⚠ 只設 `run.font.name` 只會寫 `a:latin`，中文會被 PowerPoint 退回新細明體；
    要全字元都是微軟正黑體必須同時寫 `a:ea`（東亞）與 `a:cs`（複雜文字）。
    """
    run.font.size = Pt(max(size, float(theme.font["min_pt"])))
    run.font.bold = bold
    run.font.color.rgb = theme.rgb(color)
    run.font.name = theme.font["family"]
    rpr = run._r.get_or_add_rPr()
    # 🔴 F-14：標語言＋關拼字檢查。不標時 PowerPoint 拿**預設的英文校對**檢查中文，
    # 整頁被畫滿紅色波浪底線（2026-08-02 使用者實機截圖 p20）。
    # ⚠ 內容沒有錯，是校對語言錯——轉圖看不到（proofing marks 不進圖），
    # 但客戶開檔第一眼就是滿頁紅線，會以為報表產壞了。
    rpr.set("lang", "zh-TW")
    rpr.set("altLang", "en-US")
    rpr.set("noProof", "1")
    for tag in ("a:ea", "a:cs"):
        element = rpr.find(qn(tag))
        if element is None:
            element = rpr.makeelement(qn(tag), {})
            rpr.append(element)
        element.set("typeface", theme.font["family"])


def _new_textbox(slide, *, left: float, top: float, width: float, height: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    # 🔴 K-2（2026-08-04）：內距一律歸零。python-pptx 文字框預設左右內距各 0.1in，
    # 而容量估算（_text_capacity）假設可用寬＝傳入寬——16pt 中文一字 0.222in，
    # 等於每行少約 1 字，多段累積就溢出 1～2 行（p4/p13/p16 溢框、p3/p5/p6/p15
    # 丟「後續」的共同殘差）。歸零後估算與實排在同一個座標系。
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    return box, frame


def _add_text(
    slide,
    theme: Theme,
    text: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: float,
    color: str = "ink",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
) -> None:
    """加入純文字框；換行字元切成獨立段落，維持行距一致。"""
    _, frame = _new_textbox(slide, left=left, top=top, width=width, height=height)
    frame.vertical_anchor = anchor
    for index, line in enumerate(str(text).split("\n")):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        _set_font(run, theme, size=size, color=color, bold=bold)


def _add_number_bold_text(
    slide,
    theme: Theme,
    blocks: list[tuple[str, str, str, bool]],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    size: float,
) -> None:
    """逐條 bullet；每條由（標籤, 內文, 內文色, 標籤是否強調）構成。

    ⚠ v3 起**全部內文一律粗體**（2026-07-31 使用者定案「文字內容記得加粗體」）：
    深色底上細字會發灰，投影時尤其糊。

    但這樣一來「用粗體切出關鍵數字」的原設計就失效了——全部都粗，數字不再突出。
    故改以**顏色**承擔區分：數字改用 accent 青，非數字用原內文色。粗體管可讀性、
    顏色管重點，兩件事分開，不再互相搶同一個視覺通道。
    ⚠ 強調條（emphasis）的數字維持 alert 色不轉青，否則整條的警示語氣會被打斷。
    """
    _, frame = _new_textbox(slide, left=left, top=top, width=width, height=height)
    spacing = point_line_ratio(theme)
    for index, (label, text, color, emphasized) in enumerate(blocks):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        # 🔴 2026-08-04：行距要**真的寫進段落**。原本只有 `qa.line_height_ratio`
        # 這個估算值，PowerPoint 用的是預設行距——調大估算值畫面不會變寬，
        # 只會讓容量估得更保守（字更少、版面更空）。
        para.line_spacing = spacing
        if label:
            chip = para.add_run()
            chip.text = f"{label}｜"
            _set_font(chip, theme, size=size, color="alert" if emphasized else "accent", bold=True)
        for piece in NUMBER_PATTERN.split(text):
            if not piece:
                continue
            is_number = bool(NUMBER_PATTERN.fullmatch(piece))
            run = para.add_run()
            run.text = piece
            _set_font(run, theme, size=size,
                      color="accent" if (is_number and color != "alert") else color,
                      bold=True)


def _add_band(slide, theme: Theme, left, top, width, height, color: str, *, rounded: bool = False, line: str = "") -> None:
    """加入色塊／底板。rounded 用於圓角面板，方角用於裝飾條。

    🔴 2026-07-31 批2：**圓角一律補可見邊框**。獨立驗收實測 46 個填色面板中
    43 個是 `<a:ln><a:noFill/>`，面板底對背景只有 1.12–1.72，等於沒有邊界、
    整頁區塊糊成一片。

    ⚠ 治在**實作**而不是逐一改 43 個呼叫點：加參數會讓介面變寬、每個呼叫端與
    測試都得跟著改（同日 `rasterize_svg` 加第三個參數就是這樣害測試變紅）。
    判斷依據用既有語意——圓角＝面板（要邊界），方角＝標題底線／進度條等裝飾條
    （加框反而礙眼），不必新增旗標。
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.rgb(color)
    if not line and rounded:
        line = "hairline"
    if line:
        shape.line.color.rgb = theme.rgb(line)
        shape.line.width = Pt(theme.font["panel_border_pt"])
    else:
        shape.line.fill.background()
    if rounded:
        shape.adjustments[0] = 0.08
    shape.shadow.inherit = False
    shape.text_frame.text = ""


def _add_background(slide, theme: Theme, cache_dir: Path, *, density: float) -> None:
    """深空背景層：全出血漸層矩形 ＋ 星空紋理疊圖（v3，2026-07-31）。

    分兩層是有原因的，不是為了好看才拆：
    - **漸層走 python-pptx 原生 gradFill**：向量，放大／投影／列印都不糊。
      ⚠ 不可改用 SVG 漸層——pymupdf 不渲染 SVG 漸層，實測整片變純黑。
    - **星點走圖片疊加**：星點數以百計，若逐顆畫成 shape 會讓檔案膨脹又拖慢
      PowerPoint 開檔；疊一張透明底 PNG 便宜得多。seed 固定＝每次都一樣。

    density＝紋理密度倍率（封面 1.0、內頁較淡）。紋理產不出來時**略過該層**
    而不是中斷組版：少一層星點只是樸素些，整份簡報產不出來才是嚴重的。
    """
    g = theme.geometry["background"]
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(g["left_in"]), Inches(g["top_in"]),
        Inches(theme.slide["width_in"]), Inches(theme.slide["height_in"]),
    )
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.text_frame.text = ""

    # 先讓 python-pptx 建出 gradFill 骨架，再把 stop 換成 theme 指定的兩色。
    # 直接改 XML 是因為 python-pptx 預設會塞多個 stop，逐一設色反而更繞。
    cfg = theme.gradient
    shape.fill.gradient()
    grad = shape._element.spPr.find(qn("a:gradFill"))
    stops = grad.find(qn("a:gsLst"))
    for gs in list(stops):
        stops.remove(gs)
    stops.append(parse_xml(
        f'<a:gs {nsdecls("a")} pos="0"><a:srgbClr val="{cfg["start"]}"/></a:gs>'))
    stops.append(parse_xml(
        f'<a:gs {nsdecls("a")} pos="100000"><a:srgbClr val="{cfg["end"]}"/></a:gs>'))
    # 線性漸層；若骨架給的是放射狀（a:path）就換掉，否則角度設定會無效。
    path = grad.find(qn("a:path"))
    if path is not None:
        grad.remove(path)
    lin = grad.find(qn("a:lin"))
    if lin is None:
        grad.append(parse_xml(
            f'<a:lin {nsdecls("a")} scaled="0" ang="{int(cfg["angle_ooxml"])}"/>'))
    else:
        lin.set("ang", str(int(cfg["angle_ooxml"])))
        lin.set("scaled", "0")

    shape.name = BACKGROUND_SHAPE_NAME
    texture = starfield_png(theme.starfield, cache_dir, density=density)
    if texture is None:
        return
    picture = slide.shapes.add_picture(
        str(texture), Inches(g["left_in"]), Inches(g["top_in"]),
        width=Inches(theme.slide["width_in"]), height=Inches(theme.slide["height_in"]),
    )
    # ⚠ 命名是為了讓版面自檢認得出它：背景是**刻意全出血**的，
    # 用「所有圖片都要在安全邊界內」去檢查它一定會誤報。
    picture.name = BACKGROUND_SHAPE_NAME


def _add_oval(slide, theme: Theme, left, top, width, height, color: str) -> None:
    """圓點裝飾（封面角落圓格網）。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.rgb(color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.text_frame.text = ""


def _attach_svg(picture, svg_path: Path, index: int) -> None:
    """把向量版掛到已插入的圖片上（OOXML 的 `asvg:svgBlip` 擴充）。

    PowerPoint 2016+ 顯示這份 SVG，舊版／網頁預覽／縮圖服務自動退回 `r:embed`
    指向的 PNG。⚠ 這不是變通做法——PowerPoint 自己插入 SVG 時產生的就是這個
    結構（點陣後援＋向量擴充），這裡只是用程式重現它。

    ⚠ `PackURI` 必須唯一：同名會讓後插入的圖覆蓋前一張，整份簡報的圖全變成同一張。
    """
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.opc.package import Part
        from pptx.opc.packuri import PackURI

        slide_part = picture.part
        partname = PackURI(f"/ppt/media/chart-vector-{index}.svg")
        svg_part = Part(partname, "image/svg+xml", slide_part.package, svg_path.read_bytes())
        rel_id = slide_part.relate_to(svg_part, RT.IMAGE)
        blip = picture._element.blipFill.find(qn("a:blip"))
        if blip is None:
            return
        blip.append(parse_xml(
            f'<a:extLst {nsdecls("a", "r")}>'
            '<a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">'
            '<asvg:svgBlip xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main" '
            f'r:embed="{rel_id}"/>'
            "</a:ext></a:extLst>"
        ))
    except Exception:
        # 掛不上就維持純 PNG：圖會是點陣的（放大略糊），但簡報照樣產得出來。
        return


def _add_picture_fitted(slide, image_path: Path, *, left: float, top: float, width: float, height: float) -> None:
    """等比縮放塞進框內並置中。

    ⚠ 舊版只給 `width=` 讓高度自由伸展，遇到高瘦圖就往下溢出 3.5 吋衝出版面。
    這裡先用原生尺寸插入再依框的長寬比縮放，圖再怪也不會出框。

    ⚠ 同名 `.svg` 若存在（`prepare_chart` 產的轉色裁切版）就一併掛成向量：
    以副檔名推導而不是多傳一個參數，是為了讓每個 renderer 的呼叫點都不必改——
    兩份檔案本來就由同一次 `prepare_chart` 同時產出、同名同 digest，不會對錯。
    """
    picture = slide.shapes.add_picture(str(image_path), Inches(left), Inches(top))
    vector = image_path.with_suffix(".svg")
    if vector.exists():
        global _VECTOR_INDEX
        _VECTOR_INDEX += 1
        _attach_svg(picture, vector, _VECTOR_INDEX)
    box_w, box_h = Inches(width), Inches(height)
    if picture.width <= 0 or picture.height <= 0:
        picture.width, picture.height = box_w, box_h
        return
    scale = min(box_w / picture.width, box_h / picture.height)
    picture.width = Emu(int(picture.width * scale))
    picture.height = Emu(int(picture.height * scale))
    picture.left = Emu(int(box_w - picture.width) // 2 + int(Inches(left)))
    picture.top = Emu(int(box_h - picture.height) // 2 + int(Inches(top)))


def _variant_label(chart_name: str) -> str:
    for suffix, label in CHART_VARIANT_LABELS.items():
        if suffix in chart_name:
            return label
    return Path(chart_name).stem


def _encoding_note(spec: PageSpec, ctx: dict[str, Any] | None = None) -> str:
    """圖表編碼說明：引擎那份優先，缺鍵才用本檔 fallback。

    ⚠ 本檔的 `ENCODING_NOTES` 僅供**舊報表版本**相容——引擎自 2026-07-31 起會把
    說明寫進 `report_data.json.table_display.encoding_notes`。兩份各自演進的後果
    已實測到：`annual_trend` 是折線卻寫「條長」、`application_growth` 縱軸是
    年增率 % 卻寫「件數」、`lifecycle` 橫軸是申請人家數卻寫「申請年」。
    新增或修改說明請改引擎那份，不要往這裡加。
    """
    engine = {}
    if ctx is not None:
        engine = (ctx["report_data"].get("table_display") or {}).get("encoding_notes") or {}
    for key in spec.report_keys:
        if key in engine:
            return engine[key]
    for key in spec.report_keys:
        if key in ENCODING_NOTES:
            return ENCODING_NOTES[key]
    return DEFAULT_ENCODING_NOTE


def _caveat_of(spec: PageSpec) -> str:
    for key in spec.report_keys:
        if key in CAVEATS:
            return CAVEATS[key]
    return ""


def _points_for(spec: PageSpec, ctx: dict[str, Any]) -> list[tuple[str, str, str, bool]]:
    """取本頁要點；缺 narrative 時退回引擎 rows 的關鍵數字，仍不留空框。"""
    _, points, _ = ctx["narratives_by_page"].get(spec.page, ("", [], False))
    if points:
        return [
            (
                str(p.get("label") or ""),
                str(p.get("text") or ""),
                "alert" if p.get("emphasis") else "ink",
                bool(p.get("emphasis")),
            )
            for p in points
        ]
    return [(label, text, "ink", False) for label, text in _row_highlights(spec, ctx)]


def _row_highlights(spec: PageSpec, ctx: dict[str, Any]) -> list[tuple[str, str]]:
    """從引擎 rows 取前幾名做要點，供缺 narrative 的頁面使用（只列數字，不下判讀）。"""
    for key in spec.report_keys:
        rows = _rows_of(ctx["report_data"], key)
        if not rows:
            continue
        numeric = _numeric_column(rows)
        if not numeric:
            continue
        label_col = _label_column(rows, numeric)
        ranked = sorted(rows, key=lambda r: _as_int(r.get(numeric)), reverse=True)[:4]
        return [(str(r.get(label_col, "")), f"{_as_int(r.get(numeric))} 件") for r in ranked]
    return []


# --------------------------------------------------------------------------
# 共用頁面區塊
# --------------------------------------------------------------------------
def _render_header(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """內頁頁首：判讀式標題、accent 底線、右上頁碼。"""
    g = theme.geometry["header"]
    _add_text(slide, theme, spec.title,
              left=g["title_left_in"], top=g["title_top_in"],
              width=g["title_width_in"], height=g["title_height_in"],
              # v3：深底主題下 navy 已是面板底色，標題再用它等於隱形（實機轉圖驗到）。
              size=theme.size("title_pt"), color="ink", bold=True)
    _add_band(slide, theme, g["rule_left_in"], g["rule_top_in"],
              g["rule_width_in"], g["rule_height_in"], "accent")
    _add_text(slide, theme, f"{spec.page:02d}",
              left=g["page_number_left_in"], top=g["page_number_top_in"],
              width=g["page_number_width_in"], height=g["page_number_height_in"],
              size=theme.size("page_number_pt"), color="accent", bold=True, align=PP_ALIGN.RIGHT)


# 通道短名（`tech`／`effect`）：分群產物的檔名後綴與母體註記鍵共用同一份，
# 上游唯一定義處在 `backend/app/clustering/sources.py::SOURCE_SEGMENT_SLUGS`。
# ⚠ 本檔是可攜 skill 不能 import backend，故這裡**由既有資料推導、不另抄一份**：
# `CHANNEL_NARRATIVE_VARIANTS` 的值就是 `topic_table_<slug>`，取尾段即得。
CHANNEL_SLUGS: tuple[str, ...] = tuple(
    variant.rsplit("_", 1)[-1] for variant in CHANNEL_NARRATIVE_VARIANTS.values())


def _page_channel_slug(spec: PageSpec) -> str:
    """本頁屬於哪個分群通道；判不出回空字串。

    兩種拆頁方式各有各的線索，都要認：
    - **主題分布頁**依列值拆（`_expand_page_layout`），通道在 `row_filter`
    - **機會矩陣頁**依圖檔拆（`_split_multi_chart_page`），通道在圖名後綴

    ⚠ 不要只認其中一種：0846 修正前兩類頁面都印合併母體，正是因為
    「一個 report_key 對兩頁」這種形狀沒有被任何一層考慮到。
    """
    source_field = dict(spec.row_filter).get("source_field")
    if source_field:
        variant = CHANNEL_NARRATIVE_VARIANTS.get(source_field, "")
        return variant.rsplit("_", 1)[-1] if variant else ""
    for name in spec.charts:
        tail = name.rsplit(".", 1)[0].rsplit("_", 1)[-1]
        if tail in CHANNEL_SLUGS:
            return tail
    return ""


def _population_note(report_data: dict[str, Any], spec: PageSpec) -> str:
    """本頁的母體註記；**引擎算好寫在 `report_data["population"]`，這裡只取用**。

    ⚠ 不在此計算：本檔是會佈署到使用者機器的可攜 skill，不能 import backend
    （全域規則「跨部署單元改走一方產生、一方消費」）。

    一頁掛多張報表時只取第一個有註記的——多張的母體通常相同（同頁對照用），
    全部印出來會把頁尾撐爆。

    🔴 分群報表的鍵帶通道後綴（`cluster_topic_table:tech`）。2026-08-06 實機驗出：
    技術頁與功效頁原本都印「母體 79/55 件」＝兩通道加總，而每頁只呈現單一通道；
    79 > 55 又無過計數說明，讀者只會判定報表算錯。
    ⚠ 判不出通道時**寧可不印**——引擎已刻意不再產出合併鍵，印不出來是預期行為，
    印一個錯的母體才是災難。
    """
    population = report_data.get("population") or {}
    slug = _page_channel_slug(spec)
    for key in spec.report_keys:
        if slug and (note := population.get(f"{key}:{slug}")):
            return note
        if note := population.get(key):
            return note
    return ""


def _render_footnote(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any], extra: str = "") -> None:
    """頁底註記：母體 → 資料來源 → 統計期間。

    🔴 **母體排最前**（A3，2026-08-06）：`_fit_text` 截斷是砍尾巴，
    母體若排在來源／期間之後，版面一擠就被砍掉，而**頁面看起來完全正常**
    ——讀者不會知道少了什麼。這種靜默消失比擠版更難查。

    ⚠ 濃縮（實測 55 字 → 41 字）：`資料來源：`→`來源：`、`統計期間：`→`期間`。
    頁尾實測容量 `12.13in × 0.22in @ 12pt` ＝ 單行約 72 中文字，而 `sources`
    是變數（4 頁掛兩個 report_key、最長報表名 11 字），濃縮只降低觸發機率，
    **排序才是護欄**。
    """
    g = theme.geometry["footnote"]
    sources = "、".join(_label_of(ctx["report_data"], key) for key in spec.report_keys) or "本次報表版本"
    period = ctx["period"] or "未標示"
    # ⚠ 2026-07-31 使用者定案：頁尾**不印報表版本**（「這種報表版本這種字不要有」）。
    # 原本印 report_trial_20260731_… 這種內部識別碼，對讀者毫無意義又佔掉頁尾寬度；
    # 可追溯性由 manifest 保留，不必寫在簡報上。
    parts = [p for p in (_population_note(ctx["report_data"], spec),) if p]
    parts.append(f"來源：{sources}")
    parts.append(f"期間 {period}")
    if extra:
        parts.append(extra)
    text = "｜".join(parts)
    text, _ = _fit_text(theme, text, width_in=g["width_in"], height_in=g["height_in"],
                        size_pt=theme.size("footnote_pt"))
    _add_text(slide, theme, text,
              left=g["left_in"], top=g["top_in"], width=g["width_in"], height=g["height_in"],
              size=theme.size("footnote_pt"), color="muted")


def _points_panel_height(
    theme: Theme,
    blocks: list[tuple[str, str, str, bool]],
    *,
    width_in: float,
    max_height_in: float | None = None,
) -> float:
    """要點面板的**實際**高度：依內容行數算，不超過宣告上限。

    🔴 2026-08-02：面板高度原本是常數，內容只有 2 條時下方空掉三到四成
    （實機 p3／p4／p6／p12／p13／p16／p18／p19 八頁）。⚠ 版型算的是框不是內容，
    這也是 F-1「17/22 頁背景佔比 73–80%」的來源之一。

    ⚠ 上限仍在：內容再多也不能撐出版面。下限保住標題列，空內容不塌成一條線。
    """
    g = theme.geometry["points_panel"]
    ceiling = max_height_in if max_height_in is not None else g["height_in"]
    chrome = g["text_top_offset_in"] + g["text_bottom_pad_in"]
    size_pt = theme.size("point_text_pt")
    per_line, _ = _text_capacity(theme, width_in=width_in, height_in=ceiling, size_pt=size_pt,
                                 line_ratio=point_line_ratio(theme))
    line_in = size_pt / 72.0 * point_line_ratio(theme)
    lines = sum(
        max(1, math.ceil(((len(label) + 1 if label else 0) + len(text)) / per_line))
        for label, text, _, _ in blocks
    )
    return max(chrome + line_in, min(ceiling, chrome + lines * line_in))


def _points_band_height(
    theme: Theme,
    blocks: list[tuple[str, str, str, bool]],
    *,
    width_in: float,
    columns: int,
) -> float:
    """底部要點橫幅的**實際**高度（H-1，2026-08-03）。

    原本固定 `points_band_height_in`＝1.75，但實機只放 2 條要點時下半是空的，
    而同一頁的表格卻因為框高不足被卡掉三筆——空間分配的兩邊都錯。

    ⚠ 複用 `_points_panel_height` 算單欄，不另寫一套估算：橫幅只是「多欄的面板」，
    兩套估法遲早會分岔（本專案已因兩處落點靜默失敗六次）。
    多欄取**最高**的那一欄——欄高不齊時以最高者為準，否則短欄會壓到別的東西。
    """
    g = theme.geometry["table_with_points"]
    inset = g["points_band_inset_in"]
    gap = g["points_band_column_gap_in"]
    chrome = g["points_band_text_top_offset_in"] + inset
    if not blocks:
        return chrome
    col_width = (width_in - inset * 2 - gap * (columns - 1)) / max(columns, 1)
    per_column = max(1, math.ceil(len(blocks) / max(columns, 1)))
    chunks = [blocks[i * per_column:(i + 1) * per_column] for i in range(columns)]
    # _points_panel_height 已含它自己的 chrome（面板的 text_top_offset＋bottom_pad），
    # 這裡要的是純內容高度，故扣掉再套橫幅自己的 chrome。
    panel_g = theme.geometry["points_panel"]
    panel_chrome = panel_g["text_top_offset_in"] + panel_g["text_bottom_pad_in"]
    body = max(
        (_points_panel_height(theme, chunk, width_in=col_width,
                              max_height_in=theme.geometry["footnote"]["top_in"]) - panel_chrome
         for chunk in chunks if chunk),
        default=0.0,
    )
    return chrome + max(body, 0.0)


def _table_available_height(theme: Theme, geometry_key: str, *, band_height_in: float = 0.0) -> float:
    """表格能用到的**實際**垂直空間：頁尾上緣 − 表格上緣 − 間距 − 要點區。

    🔴 H-1（使用者：「這個所有主題都放都還能放解讀，為甚麼要卡掉」）：
    `table_with_points.height_in` 寫死 2.88（只夠 4 列宣告高），
    但 1.62 → 6.78 實際有 5.16 in。表格從一開始就沒拿到該有的空間，
    判讀面板壓上來只是後果。

    ⚠ 宣告的 `height_in` 不再當固定值用，但仍是**下限**——版型意圖是表格至少那麼高。
    """
    g = theme.geometry[geometry_key]
    # ⚠ 表格與橫幅之間用**橫幅自己的內距**（0.18），不是欄間距 `column_gap_in`（0.30）。
    # 後者是給並排欄位的水平留白，垂直方向套上去偏寬——實測技術主題頁因此差
    # 0.07 in 放不下第 5 列（表格 3.75＋橫幅 1.18＋間距 0.30 = 5.23 > 可用 5.16）。
    # 用 0.18 後總和 5.11，第 5 列進得來，而視覺上與橫幅內距一致、不會顯得擠。
    gap = float(theme.geometry["table_with_points"]["points_band_inset_in"]) if band_height_in else 0.0
    room = theme.geometry["footnote"]["top_in"] - g["top_in"] - gap - band_height_in
    # 🔴 I-2（2026-08-03 實機 p23）：附錄最後一列壓在頁尾文字上。
    # `row_height_in` 是**宣告值**，PowerPoint 列高只增不減——實測每列 0.33–0.34，
    # 15 列累積差約 0.3 in，剛好越過 footnote 上緣。
    #
    # ⚠ **不猜實際列高**：那要靠轉圖量測，字型一換就失準
    # （I-3 的字寬係數猜了三次還沒中）。改為**預留一整列的緩衝**：
    # 即使每列都比宣告高 5%，15 列累積 0.24 in 仍在一列（0.32）之內。
    # ⚠ 代價是少放一列——但少一列會誠實顯示在「顯示前 N/M 筆」，
    # 壓到頁尾則是兩段文字疊在一起、兩邊都讀不了。
    #
    # ⚠ **只在沒有要點橫幅時扣**：有 band 的頁面（`table_with_points`）表格下方
    # 還接著 band ＋ 間距，那本身就是緩衝，再扣一列會讓技術主題頁少放第 5 筆
    # ——而「5 筆全放」是使用者 2026-08-03 明確要求的（「所有主題都放都還能放解讀」）。
    # 實際壓到頁尾的是**附錄頁**（`table` 版型，表格下方直接就是 footnote）。
    if not band_height_in:
        room -= g["row_height_in"]
    # ⚠ 下限只保「表頭＋一列」，**不拿宣告高度當下限**：
    # 附錄的宣告高度 4.86 比扣掉緩衝後的可用空間 4.84 還大，
    # 用 `max(height_in, room)` 會把緩衝整個吃掉——實機 p23 壓到頁尾就是這樣來的。
    # 宣告高度的角色是「版型預期多高」，不能凌駕「實際還剩多少」。
    return max(g["row_height_in"] * 2, room)


def _render_points_panel(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """右側要點框（＋必要時的方法論警語框）。"""
    g = theme.geometry["points_panel"]
    caveat = _caveat_of(spec)
    declared = g["height_with_caveat_in"] if caveat else g["height_in"]
    # 先用宣告高度當上限裁切，再依裁切後的實際內容把面板收回來（F-10）。
    trim_height = declared - g["text_top_offset_in"] - g["text_bottom_pad_in"]
    fitted = _trim_blocks(theme, _points_for(spec, ctx),
                          width_in=g["width_in"] - g["text_inset_right_in"],
                          height_in=trim_height, size_pt=theme.size("point_text_pt"))
    panel_height = _points_panel_height(
        theme, fitted, width_in=g["width_in"] - g["text_inset_right_in"], max_height_in=declared)
    # 🔴 K-11（2026-08-04 使用者定案）：判讀要點**去框**（面板底色＋邊框移除）——
    # 字超出框線是最醜的破版形態；p19 研發方向的卡片維持有框（使用者指名例外）。
    # panel_height 仍供下方文字區高度計算使用，只是不再畫底框。
    _add_text(slide, theme, "判讀要點",
              left=g["left_in"] + g["header_inset_left_in"], top=g["top_in"] + g["header_top_offset_in"],
              width=g["width_in"] - g["text_inset_right_in"], height=g["header_height_in"],
              size=theme.size("panel_header_pt"), color="accent", bold=True)

    text_width = g["width_in"] - g["text_inset_right_in"]
    text_height = panel_height - g["text_top_offset_in"] - g["text_bottom_pad_in"]
    size = theme.size("point_text_pt")
    blocks = fitted
    _add_number_bold_text(slide, theme, blocks,
                          left=g["left_in"] + g["text_inset_left_in"], top=g["top_in"] + g["text_top_offset_in"],
                          width=text_width, height=text_height, size=size)

    if caveat:
        c = theme.geometry["caveat_panel"]
        _add_band(slide, theme, c["left_in"], c["top_in"], c["width_in"], c["height_in"], "panel", rounded=True)
        _add_text(slide, theme, "判讀限制",
                  left=c["left_in"] + c["title_inset_left_in"], top=c["top_in"] + c["title_top_offset_in"],
                  width=c["width_in"] - c["text_inset_right_in"], height=c["title_height_in"],
                  size=theme.size("caveat_title_pt"), color="accent", bold=True)
        body_width = c["width_in"] - c["text_inset_right_in"]
        body_height = c["height_in"] - c["text_top_offset_in"] - c["text_bottom_pad_in"]
        body, _ = _fit_text(theme, caveat, width_in=body_width, height_in=body_height,
                            size_pt=theme.size("caveat_text_pt"))
        _add_text(slide, theme, body,
                  left=c["left_in"] + c["text_inset_left_in"], top=c["top_in"] + c["text_top_offset_in"],
                  width=body_width, height=body_height,
                  size=theme.size("caveat_text_pt"), color="on_dark_soft")


def _visible_column_count(rows: list[dict[str, Any]], excluded: set[str]) -> int:
    """排除欄之後實際可顯示的欄數（截欄註記的分母）。"""
    if not rows:
        return 0
    return len([name for name in rows[0] if str(name) not in excluded])


def _render_points_band(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any],
                        *, top: float | None = None, height: float | None = None) -> None:
    """底部要點橫幅：無圖的表格頁專用，讓表格拿到滿版寬度（v3，2026-07-31）。

    ⚠ 橫幅容量比右側直欄小，所以**不是**把右欄內容原樣搬下來就好——
    該頁能寫幾條、每條幾字由 `narrative_capacity()` 依本區幾何算出後餵給解讀 CLI，
    上游照容量寫，這裡的 `_trim_blocks` 只當最後保底。
    """
    g = theme.geometry["table_with_points"]
    left = g["points_band_left_in"]
    top = g["points_band_top_in"] if top is None else top
    width = g["points_band_width_in"]
    # height 省略時回到宣告值——舊呼叫端（若有）行為不變。
    height = g["points_band_height_in"] if height is None else height
    inset = g["points_band_inset_in"]
    # 🔴 K-11：去框（同側欄面板）。
    _add_text(slide, theme, "判讀要點",
              left=left + inset, top=top + g["points_band_header_top_offset_in"],
              width=width - inset - inset, height=g["points_band_header_height_in"],
              size=theme.size("panel_header_pt"), color="accent", bold=True)

    columns = int(g["points_band_columns"])
    gap = g["points_band_column_gap_in"]
    col_width = (width - inset - inset - gap * (columns - 1)) / columns
    text_top = top + g["points_band_text_top_offset_in"]
    text_height = height - g["points_band_text_top_offset_in"] - inset
    size = theme.size("point_text_pt")
    # 容量以「單欄高度 × 欄數」估：橫幅是多欄排版，只用單欄高度會低估一半。
    blocks = _trim_blocks(theme, _points_for(spec, ctx),
                          width_in=col_width, height_in=text_height * columns, size_pt=size)
    # 前半進左欄、後半進右欄——依序讀比蛇行（左右交錯）自然。
    per_column = max(1, math.ceil(len(blocks) / columns)) if blocks else 1
    for index in range(columns):
        chunk = blocks[index * per_column:(index + 1) * per_column]
        if not chunk:
            continue
        _add_number_bold_text(slide, theme, chunk,
                              left=left + inset + index * (col_width + gap), top=text_top,
                              width=col_width, height=text_height, size=size)


# 被 `_trim_blocks` 整條丟掉的要點（H-2）。⚠ 用模組級收集器而非回傳值：
# `_trim_blocks` 有 7 個呼叫端，改回傳形狀等於逼每一處都跟著改，而它們全都
# 只關心「要畫哪些條」。build_ppt() 收尾時把這裡的內容併進 manifest warnings，
# 所以丟棄仍然說得出來、不是靜默。
_DROPPED_POINTS: list[tuple[int | None, tuple[str, str, str, bool]]] = []

#: 目前正在渲染的頁碼。
#: ⚠ 為什麼用模組級狀態而不是給 `_trim_blocks` 加參數：它有 7 個呼叫端，
#: 全都只關心「要畫哪些條」，為了一句警告文字讓 7 處都改簽名並不划算。
#: 頁碼在**分派處**（RENDERERS[spec.kind] 那一行前）設定一次，呼叫端零修改。
_CURRENT_PAGE: int | None = None


def dropped_points() -> list[tuple[int | None, tuple[str, str, str, bool]]]:
    """本次組版被整條丟掉的要點 [(頁碼, 要點)]（供 QA 與測試查驗）。"""
    return list(_DROPPED_POINTS)


def reset_dropped_points() -> None:
    """每次組版開始前清空——否則同一個 process 連續產兩份會互相污染。"""
    _DROPPED_POINTS.clear()


def set_current_page(page: int | None) -> None:
    """記下正在渲染哪一頁，讓丟棄警告說得出頁碼。"""
    global _CURRENT_PAGE
    _CURRENT_PAGE = page


def _trim_blocks(
    theme: Theme,
    blocks: list[tuple[str, str, str, bool]],
    *,
    width_in: float,
    height_in: float,
    size_pt: float,
) -> list[tuple[str, str, str, bool]]:
    """把要點條列裁到框內裝得下。

    ⚠ 依序填到滿為止會讓第一條吃光版面、後面的「意涵」「決策提醒」整條消失——
    等於只給讀者半個判讀。故改成**按條數分配行數**：每條至少一行，各自截斷，
    寧可每條短一點，也要讓完整的判讀結構（現況→意涵→後續）都露出來。

    ⚠ `CAVEAT_LABEL` 那條**先扣足行數、不參與均分**（2026-07-31 實機第 3、5 頁）：
    判讀限制被切成「…多為新案審…」，而那段沒有標點可切，`_truncate_to_width`
    的「切在標點」邏輯救不了。警語講一半比不講更糟，故讓要點讓路而不是讓警語斷句。
    """
    if not blocks:
        return blocks
    per_line, lines = _text_capacity(theme, width_in=width_in, height_in=height_in,
                                     size_pt=size_pt, line_ratio=point_line_ratio(theme))
    needs = [
        max(1, math.ceil(((len(label) + 1 if label else 0) + len(text)) / per_line))
        for label, text, _, _ in blocks
    ]
    if sum(needs) <= lines:
        return blocks

    protected = {index for index, block in enumerate(blocks) if block[0] == CAVEAT_LABEL}
    reserved = min(lines, sum(needs[index] for index in protected))
    # 🔴 H-2（2026-08-03 使用者：「這種卡掉的敘述不要再有」）：**一律不截字**。
    # 依序放到裝不下為止，放不下的整條不放並記進 dropped——
    # 句子斷在半路讀者看不懂，比少一條更糟；少一條至少「看到的都完整」，
    # 而且 warnings 會講出來，不是靜默。
    # ⚠ 真正的治本在上游：`narrative_capacity()` 把每頁**實際**能寫多少交給 CLI，
    # 讓它照著寫。這裡只是保底，正常情況不該觸發。
    others = [index for index in range(len(blocks)) if index not in protected]
    room = lines - reserved
    kept: set[int] = set()
    used = 0
    for index in others:
        if used + needs[index] > room:
            continue
        kept.add(index)
        used += needs[index]
    dropped = [blocks[index] for index in others if index not in kept]
    if dropped:
        _DROPPED_POINTS.extend((_CURRENT_PAGE, block) for block in dropped)
    return [block for index, block in enumerate(blocks)
            if index in protected or index in kept]


# --------------------------------------------------------------------------
# 版型 renderer
# --------------------------------------------------------------------------
def _render_cover(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """封面：主標＋統計期間副標＋統計卡＋分析框架條＋幾何裝飾（深底亮字）。

    ⚠ v3 起**不再自己畫整頁底色**：底色由 `_add_background` 的漸層＋星空負責，
    這裡若再鋪一張不透明矩形會把背景整個蓋掉（v2 是淺底主題才需要）。
    """
    g = theme.geometry["cover"]
    _add_band(slide, theme, g["accent_block_left_in"], g["accent_block_top_in"],
              g["accent_block_width_in"], g["accent_block_height_in"], "royal")

    # 角落圓格網與斜線紋：Slidesgo 的幾何裝飾語彙，純視覺、不承載資訊。
    for row in range(int(g["dots_rows"])):
        for col in range(int(g["dots_cols"])):
            _add_oval(slide, theme,
                      g["dots_left_in"] + col * g["dots_step_in"],
                      g["dots_top_in"] + row * g["dots_step_in"],
                      g["dots_size_in"], g["dots_size_in"], "accent")
    for index in range(int(g["stripe_count"])):
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(g["stripe_left_in"] + index * g["stripe_step_in"]), Inches(g["stripe_top_in"]),
            Inches(g["stripe_width_in"]), Inches(g["stripe_height_in"]),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = theme.rgb("royal")
        stripe.line.fill.background()
        stripe.shadow.inherit = False
        stripe.rotation = g["stripe_rotation_deg"]
        stripe.text_frame.text = ""

    _add_text(slide, theme, COVER_EYEBROW,
              left=g["eyebrow_left_in"], top=g["eyebrow_top_in"],
              width=g["eyebrow_width_in"], height=g["eyebrow_height_in"],
              size=theme.size("cover_subtitle_pt"), color="accent", bold=True)

    title = _cover_title(ctx["report_data"], ctx["slots"])
    title, _ = _fit_text(theme, title, width_in=g["title_width_in"], height_in=g["title_height_in"],
                         size_pt=theme.size("cover_title_pt"))
    _add_text(slide, theme, title,
              left=g["title_left_in"], top=g["title_top_in"],
              width=g["title_width_in"], height=g["title_height_in"],
              size=theme.size("cover_title_pt"), color="ink", bold=True)
    _add_band(slide, theme, g["rule_left_in"], g["rule_top_in"], g["rule_width_in"], g["rule_height_in"], "accent")

    period = ctx["period"]
    # ⚠ 2026-07-31 使用者定案：封面也不印報表版本（內部識別碼對讀者無意義）。
    #    可追溯性由 manifest 保留；頁尾同步移除，見 _render_footnote。
    # 統計期間＋設計案備註併成一行 muted 小字（A4，2026-08-06）。
    #
    # ⚠ 為什麼併進這一行而不另開一列：另開要在 theme 加一組座標，且**實測封面
    # 統計卡下方沒有餘裕**；併行則零版面改動。⚠ 也不得併進統計卡的 value
    # ——`_cover_stat_size()` 四張卡同級、級數由最長值決定，併進去會把四張卡一起縮小。
    #
    # ⚠ 備註文字由**引擎**產（`report_data["patent_kind"]["design_note"]`），
    # 本檔不自行判定設計案——判定的唯一定義處在 backend 的 `transforms/patent_kind.py`，
    # 而本檔是會佈署到使用者機器的可攜 skill，不能 import backend。
    design_note = ((ctx["report_data"].get("patent_kind") or {}).get("design_note") or "")
    subtitle = "｜".join(p for p in (f"統計期間 {period}" if period else "", design_note) if p)
    _add_text(slide, theme, subtitle,
              left=g["period_left_in"], top=g["period_top_in"],
              width=g["period_width_in"], height=g["period_height_in"],
              size=theme.size("cover_subtitle_pt"), color="muted")

    for index, (value, unit, label) in enumerate(ctx["cover_stats"]):
        left = g["stat_left_in"] + index * g["stat_gap_in"]
        _add_band(slide, theme, left, g["stat_top_in"], g["stat_width_in"], g["stat_height_in"],
                  "panel", rounded=True)
        _add_band(slide, theme, left, g["stat_top_in"], g["stat_width_in"], g["stat_accent_height_in"], "accent")
        # 🔴 F-15：四張卡**共用**同一個字級。分級規則本身沒錯（避免撐出卡片），
        # 錯在逐張各算各的——實機 p1 的「2011–2026」被降級後，四張並排看起來
        # 像三張重要、一張次要，但它們是同一層級的指標。
        value_size = _cover_stat_size(theme, ctx["cover_stats"], value)
        _add_text(slide, theme, value,
                  left=left, top=g["stat_value_top_in"],
                  width=g["stat_width_in"], height=g["stat_value_height_in"],
                  # v3：卡片底色就是 navy，數值再用 navy 等於隱形（實機轉圖驗到）。
                  size=value_size, color="on_dark", bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, theme, unit,
                  left=left, top=g["stat_unit_top_in"],
                  width=g["stat_width_in"], height=g["stat_unit_height_in"],
                  size=theme.size("stat_unit_pt"), color="blue", bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, theme, label,
                  left=left + g["stat_label_inset_in"], top=g["stat_label_top_in"],
                  width=g["stat_width_in"] - g["stat_label_inset_in"] * 2, height=g["stat_label_height_in"],
                  size=theme.size("stat_label_pt"), color="ink", align=PP_ALIGN.CENTER)

    _add_band(slide, theme, g["banner_left_in"], g["banner_top_in"],
              g["banner_width_in"], g["banner_height_in"], "panel", rounded=True)
    banner, _ = _fit_text(theme, ctx["framework_text"],
                          width_in=g["banner_text_width_in"], height_in=g["banner_text_height_in"],
                          size_pt=theme.size("cover_banner_pt"))
    _add_text(slide, theme, banner,
              left=g["banner_left_in"] + g["banner_text_inset_left_in"],
              top=g["banner_top_in"] + g["banner_text_top_offset_in"],
              width=g["banner_text_width_in"], height=g["banner_text_height_in"],
              size=theme.size("cover_banner_pt"), color="on_dark_soft")


def _render_section_divider(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """章節隔頁：深色塊大字，只做段落分隔。

    ⚠ v3 起不再自畫整頁底色（同 `_render_cover` 的理由：會蓋掉漸層與星空）。
    ⚠ 目前大綱不使用本版型——使用者定案「不加章節號、不做章節隔頁」，
    論證鏈靠頁序本身表達。保留 renderer 是因為 layout_overrides 仍可指定它。
    """
    g = theme.geometry["section_divider"]
    _add_band(slide, theme, g["accent_left_in"], g["accent_top_in"],
              g["accent_width_in"], g["accent_height_in"], "accent")
    _add_text(slide, theme, spec.title,
              left=g["title_left_in"], top=g["title_top_in"],
              width=g["title_width_in"], height=g["title_height_in"],
              size=theme.size("section_title_pt"), color="on_dark", bold=True)
    if spec.subtitle:
        _add_text(slide, theme, spec.subtitle,
                  left=g["subtitle_left_in"], top=g["subtitle_top_in"],
                  width=g["subtitle_width_in"], height=g["subtitle_height_in"],
                  size=theme.size("section_subtitle_pt"), color="on_dark_soft")


def _points_for_panel(points: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """判讀面板要顯示的要點——**結論也在裡面**。

    🔴 2026-08-03：原本結論那條會被濾掉（它被抽去底部「核心結論」橫幅），
    造成結論與支撐它的依據分處兩地。使用者定案：結論回到判讀區塊，
    且**不是每頁都要有**——沒有結論性的那條就不標 `emphasis`，不硬湊。
    """
    return list(points)


def _conclusion_text(headline: str, points: list[dict[str, Any]]) -> str:
    """核心結論條的文字：emphasis 那條 → 「意涵」那條 → headline。

    NotebookLM 範例的手法：每頁底部一句話收束。優先取 narrative 自己標記
    最重要的（emphasis），其次判讀性最強的「意涵」，都沒有才退回標題句。
    """
    for point in points:
        if point.get("emphasis"):
            return str(point.get("text") or "")
    for point in points:
        if "意涵" in str(point.get("label") or ""):
            return str(point.get("text") or "")
    return headline


def _render_chart_hero(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """大圖版型（P1-1，內容頁新預設）：圖佔主幅＋右側精簡註解卡＋底部核心結論條。

    「圖表為主、文字為輔」的落實：文字不是消失，而是（a）判讀進標題、
    （b）要點成右側小卡、（c）最重要一句進底部結論條。
    ⚠ 註解卡放固定區，不精準指向圖內資料點——不動 SVG 內部就拿不到錨點座標。
    """
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["chart_hero"]
    image = ctx["charts"].resolve(spec.charts[0]) if spec.charts else None
    # R-3：框高先讓過頁尾帶，再談縮放（長圖會撐滿框高，底注否則壓住資料來源）。
    frame_h = image_frame_height(theme, g["image_top_in"], g["image_height_in"])
    # 說明靠右對齊**圖的實際右緣**，不是框的右緣。圖填不滿框時（瘦圖只佔框寬的
    # 三分之一），靠框對齊會讓說明飄在圖右邊幾吋外的空白處（獨立驗收 p6 抓到）。
    shown_w, _ = (_fitted_size(image, g["image_width_in"], frame_h)
                  if image is not None else (g["image_width_in"], 0.0))
    edge_left = g["image_left_in"] + (g["image_width_in"] - shown_w) / 2
    _add_text(slide, theme, _encoding_note(spec, ctx),
              left=edge_left, top=g["encoding_top_in"],
              width=shown_w, height=g["encoding_height_in"],
              size=theme.size("encoding_note_pt"), color="muted", align=PP_ALIGN.RIGHT)
    if image is not None:
        _add_picture_fitted(slide, image,
                            left=g["image_left_in"], top=g["image_top_in"],
                            width=g["image_width_in"], height=frame_h)

    headline, points, _ = ctx["narratives_by_page"].get(spec.page, ("", [], False))
    # 🔴 2026-08-03 使用者：「判讀區塊那裡要能帶出核心結論，**還有不是每頁都要有
    # 核心結論**」。原本結論被抽去底部橫幅、並從面板濾掉——結論與依據拆在兩處，
    # 而且每頁都硬要有一條。改為結論留在面板（由 `emphasis` 標示），橫幅取消。
    # ⚠ 附帶效果正是「圖表要大一點」：橫幅讓出的空間全部給圖框。
    listed = _points_for_panel(points)
    if not listed and not points:
        listed = [{"label": label, "text": text, "emphasis": False}
                  for label, text in _row_highlights(spec, ctx)]
    blocks = [(str(p.get("label") or ""), str(p.get("text") or ""),
               "alert" if p.get("emphasis") else "ink", bool(p.get("emphasis")))
              for p in listed]
    caveat = _caveat_of(spec)
    if caveat:
        blocks = blocks + [(CAVEAT_LABEL, caveat, "muted", False)]
    # 🔴 K-11（2026-08-04 使用者定案）：判讀要點**去框**——字超出框線是最醜的
    # 破版形態；p19 研發方向的卡片維持有框（使用者指名例外）。
    _add_text(slide, theme, "判讀要點",
              left=g["panel_left_in"] + g["panel_inset_in"],
              top=g["panel_top_in"] + g["panel_header_top_offset_in"],
              width=g["panel_width_in"] - g["panel_inset_in"] * 2,
              height=g["panel_header_height_in"],
              size=theme.size("panel_header_pt"), color="accent", bold=True)
    text_width = g["panel_width_in"] - g["panel_inset_in"] * 2
    text_height = g["panel_height_in"] - g["panel_text_top_offset_in"] - g["panel_inset_in"]
    size = theme.size("point_text_pt")
    _add_number_bold_text(slide, theme,
                          _trim_blocks(theme, blocks, width_in=text_width,
                                       height_in=text_height, size_pt=size),
                          left=g["panel_left_in"] + g["panel_inset_in"],
                          top=g["panel_top_in"] + g["panel_text_top_offset_in"],
                          width=text_width, height=text_height, size=size)

    _render_footnote(slide, theme, spec, ctx)


def _render_chart_with_points(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """內容頁預設版型：左圖約 60% 寬，右側要點框（＋必要時警語框）。"""
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["chart_with_points"]
    image = ctx["charts"].resolve(spec.charts[0]) if spec.charts else None
    # 說明靠右對齊**圖的實際右緣**，不是框的右緣。圖填不滿框時（瘦圖只佔框寬的
    # 三分之一），靠框對齊會讓說明飄在圖右邊幾吋外的空白處（獨立驗收 p6 抓到）。
    shown_w, _ = (_fitted_size(image, g["image_width_in"], g["image_height_in"])
                  if image is not None else (g["image_width_in"], 0.0))
    edge_left = g["image_left_in"] + (g["image_width_in"] - shown_w) / 2
    _add_text(slide, theme, _encoding_note(spec, ctx),
              left=edge_left, top=g["encoding_top_in"],
              width=shown_w, height=g["encoding_height_in"],
              size=theme.size("encoding_note_pt"), color="muted", align=PP_ALIGN.RIGHT)
    if image is not None:
        _add_picture_fitted(slide, image,
                            left=g["image_left_in"], top=g["image_top_in"],
                            width=g["image_width_in"], height=g["image_height_in"])
    _render_points_panel(slide, theme, spec, ctx)
    _render_footnote(slide, theme, spec, ctx)


def _render_comparison(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """成對報表同頁左右並排：兩張圖各配子標、圖例編碼與要點。

    ⚠ 成對報表（IPC/CPC 的 L4 與 L5、機會矩陣的技術面與功效面）**禁止合成同一張圖**；
    合成會讓兩個不同母體的數字被讀成同一組，這裡只做並排比較。
    """
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["comparison"]
    note = _encoding_note(spec, ctx)
    blocks = _points_for(spec, ctx)
    for index, chart_name in enumerate(spec.charts[: len(g["column_left_in"])]):
        left = g["column_left_in"][index]
        width = g["column_width_in"]
        _add_text(slide, theme, _variant_label(chart_name),
                  left=left, top=g["caption_top_in"], width=width, height=g["caption_height_in"],
                  size=theme.size("chart_caption_pt"), color="accent", bold=True)
        _add_text(slide, theme, note,
                  left=left, top=g["encoding_top_in"], width=width, height=g["encoding_height_in"],
                  size=theme.size("encoding_note_pt"), color="muted")
        image = ctx["charts"].resolve(chart_name)
        if image is not None:
            _add_picture_fitted(slide, image,
                                left=left, top=g["image_top_in"], width=width, height=g["image_height_in"])
        # K-11：去框（同 chart_hero）。
        text_width = width - g["points_inset_right_in"]
        text_height = g["points_height_in"] - g["points_top_offset_in"] - g["points_bottom_pad_in"]
        size = theme.size("point_text_pt")
        half = blocks[index::2] or blocks[:1]
        _add_number_bold_text(slide, theme,
                              _trim_blocks(theme, half, width_in=text_width, height_in=text_height, size_pt=size),
                              left=left + g["points_inset_left_in"], top=g["points_top_in"] + g["points_top_offset_in"],
                              width=text_width, height=text_height, size=size)
    _render_footnote(slide, theme, spec, ctx)


def _render_stat_callout(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """大數字焦點頁；同時是「圖檔缺失時的降級版型」——確保每頁都有視覺元素。"""
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["stat_callout"]
    _add_band(slide, theme, g["block_left_in"], g["block_top_in"],
              g["block_width_in"], g["block_height_in"], "panel", rounded=True)

    rows: list[dict[str, Any]] = []
    report_key = ""
    for key in spec.report_keys:
        rows = _rows_of(ctx["report_data"], key)
        if rows:
            report_key = key
            break
    numeric = _numeric_column(rows)
    label_col = _label_column(rows, numeric) if numeric else ""
    total = sum(_as_int(row.get(numeric)) for row in rows) if numeric else 0

    _add_text(slide, theme, f"{total:,}" if total else str(len(rows)),
              left=g["value_left_in"], top=g["value_top_in"],
              width=g["value_width_in"], height=g["value_height_in"],
              size=theme.size("callout_value_pt"), color="on_dark", bold=True)
    _add_text(slide, theme, "件" if total else "筆",
              left=g["unit_left_in"], top=g["unit_top_in"],
              width=g["unit_width_in"], height=g["unit_height_in"],
              size=theme.size("callout_unit_pt"), color="accent", bold=True)
    _add_text(slide, theme, _label_of(ctx["report_data"], report_key, spec.topic) if report_key else spec.topic,
              left=g["label_left_in"], top=g["label_top_in"],
              width=g["label_width_in"], height=g["label_height_in"],
              size=theme.size("callout_label_pt"), color="on_dark_soft")
    _add_band(slide, theme, g["rule_left_in"], g["rule_top_in"], g["rule_width_in"], g["rule_height_in"], "accent")

    ranked = sorted(rows, key=lambda r: _as_int(r.get(numeric)), reverse=True) if numeric else []
    for index, row in enumerate(ranked[: int(g["row_max"])]):
        _add_text(slide, theme,
                  f"{row.get(label_col, '')}　{_as_int(row.get(numeric)):,} 件",
                  left=g["row_left_in"], top=g["row_top_in"] + index * g["row_step_in"],
                  width=g["row_width_in"], height=g["row_height_in"],
                  size=theme.size("callout_row_pt"), color="on_dark_soft")

    panel = theme.geometry["points_panel"]
    caveat = _caveat_of(spec)
    # K-11：去框（同 chart_hero）。
    _add_text(slide, theme, "判讀要點",
              left=g["points_left_in"] + panel["header_inset_left_in"],
              top=g["points_top_in"] + panel["header_top_offset_in"],
              width=g["points_width_in"] - panel["text_inset_right_in"], height=panel["header_height_in"],
              size=theme.size("panel_header_pt"), color="accent", bold=True)
    text_width = g["points_width_in"] - panel["text_inset_right_in"]
    text_height = g["points_height_in"] - panel["text_top_offset_in"] - panel["text_bottom_pad_in"]
    size = theme.size("point_text_pt")
    blocks = _points_for(spec, ctx)
    if caveat:
        blocks = blocks + [(CAVEAT_LABEL, caveat, "muted", False)]
    _add_number_bold_text(slide, theme,
                          _trim_blocks(theme, blocks, width_in=text_width, height_in=text_height, size_pt=size),
                          left=g["points_left_in"] + panel["text_inset_left_in"],
                          top=g["points_top_in"] + panel["text_top_offset_in"],
                          width=text_width, height=text_height, size=size)
    _render_footnote(slide, theme, spec, ctx)


def _bar_fill_ratio(value: int, total: int) -> float:
    """佔比條的條長比例——分母是**全體總數**，不是第一名。

    🔴 2026-08-02 實機 p5：條長原本用 `value / top_value`，右側標的百分比卻用
    `value / total`。CN 39 件被畫滿整條軌道、字寫 65%——同一張圖兩種基準，
    讀者把滿格讀成 100%。軌道本身就是 100% 基準，條長用真佔比才對得起來：
    CN 停在 65%，留白的 35% 正是「還有其他國家」這個資訊。
    """
    return (value / total) if total else 0.0


def _render_percentage_bars(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """佔比條列（如受理國分布）：條長＝佔全體比例，右側數值為件數與佔比。"""
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["percentage_bars"]
    rows: list[dict[str, Any]] = []
    for key in spec.report_keys:
        rows = _rows_of(ctx["report_data"], key)
        if rows:
            break
    numeric = _numeric_column(rows)
    label_col = _label_column(rows, numeric) if numeric else ""
    ranked = sorted(rows, key=lambda r: _as_int(r.get(numeric)), reverse=True) if numeric else []
    ranked = ranked[: int(g["row_max"])]
    total = sum(_as_int(r.get(numeric)) for r in ranked) or 1

    # 條數少時整組垂直置中，避免條列擠在上方、下半頁一片空白。
    body = theme.geometry
    span = len(ranked) * g["row_step_in"]
    room = body["body_top_in"] + body["body_height_in"] - g["row_top_in"]
    first_top = g["row_top_in"] + max(0.0, (room - span) / 2)

    for index, row in enumerate(ranked):
        top = first_top + index * g["row_step_in"]
        value = _as_int(row.get(numeric))
        _add_text(slide, theme, str(row.get(label_col, "")),
                  left=g["row_left_in"], top=top, width=g["label_width_in"], height=g["label_height_in"],
                  size=theme.size("percent_label_pt"), color="ink", bold=True)
        _add_band(slide, theme, g["track_left_in"], top + g["track_top_offset_in"],
                  g["track_width_in"], g["track_height_in"], "bar_track", rounded=True)
        ratio = _bar_fill_ratio(value, total)
        fill_width = max(g["track_height_in"], g["track_width_in"] * ratio)
        _add_band(slide, theme, g["track_left_in"], top + g["track_top_offset_in"],
                  # 🔴 批2：原本第一名用 royal、其餘用 blue——兩者都是**裝飾色**
                  # （色相距 accent 僅 0.6°／0.7°），使用者反映「資料看起來像裝飾」；
                  # 且 royal 比 blue 暗，等於**最大值最不顯眼**、語意反了。
                  # 改為資料暖色，第一名用主序列、其餘用淺階，明暗與大小一致。
                  fill_width, g["track_height_in"],
                  "series_primary" if index == 0 else "series_light", rounded=True)
        _add_text(slide, theme, f"{value:,} 件　{value / total:.0%}",
                  left=g["value_left_in"], top=top, width=g["value_width_in"], height=g["value_height_in"],
                  size=theme.size("percent_value_pt"), color="ink", bold=True)
    _render_points_panel(slide, theme, spec, ctx)
    _render_footnote(slide, theme, spec, ctx)


def _render_table(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """全寬表格（附錄）：直接列引擎 rows，不加解讀。"""
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["table"]
    rows = _first_rows(spec, ctx)
    labels, excluded, priority = _table_display(ctx, spec)
    # 附錄不放要點，整個可用區都給表格（H-4：原本用宣告的 height_in，13 筆只出 6 筆）。
    shown, _used = _add_table(slide, theme, rows,
                              left=g["left_in"], top=g["top_in"], width=g["width_in"],
                              height=_table_available_height(theme, "table"),
                              row_height=g["row_height_in"], max_columns=int(g["max_columns"]),
                              cell_margin_in=g["cell_margin_in"], cell_inset_in=g["cell_inset_in"],
                              labels=labels, excluded=excluded, priority=priority)
    _render_footnote(slide, theme, spec, ctx,
                     _rows_note(shown, rows, int(g["max_columns"]), _visible_column_count(rows, excluded)))


def _render_table_with_points(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """表格＋底部要點橫幅：明細與判讀同頁（v3 改滿版，2026-07-31）。

    ⚠ 這類頁面（技術／功效主題分布）**沒有圖**，右欄本來就不必讓給要點卡。
    量測：舊版 8.05 in ÷ 4 欄＝每欄 2.01 in；改滿版 12.13 in ÷ 6 欄＝每欄 2.02 in
    ——欄寬幾乎不變，欄位卻從 4 欄變 6 欄全到齊。要點改放底部橫幅雙欄。
    """
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["table_with_points"]
    rows = _first_rows(spec, ctx)
    # 🔴 K-9（2026-08-04 使用者定案）：**內頁**代表專利只放 1 件（引擎排序後的首件）。
    # 3 個號碼 × U+2011 不可拆＝3 行，把每列撐到 0.96in，可用高 3.03in 只塞得下
    # 2 列（實機功效表 2/8）。附錄全表（_render_table）維持 3 件——「內頁精選、附錄放齊」。
    rows = [
        {**row, "representative": str(row.get("representative", "")).split("、")[0]}
        if row.get("representative") else row
        for row in rows
    ]
    labels, excluded, priority = _table_display(ctx, spec)
    # 🔴 H-1（2026-08-03）：先算要點區**實際**要多高，剩下的全給表格。
    # ⚠ 順序不能反：表格能放幾列取決於剩多少空間，而剩多少取決於要點內容——
    #   要點內容是已知的（narrative 早就產好了），表格列數才是被算出來的那一邊。
    band_blocks = _points_for(spec, ctx)
    band_height = _points_band_height(theme, band_blocks,
                                      width_in=g["points_band_width_in"],
                                      columns=int(g["points_band_columns"]))
    shown, used_height = _add_table(
        slide, theme, rows,
        left=g["left_in"], top=g["top_in"], width=g["width_in"],
        height=_table_available_height(theme, "table_with_points", band_height_in=band_height),
        row_height=g["row_height_in"], max_columns=int(g["max_columns"]),
        cell_margin_in=g["cell_margin_in"], cell_inset_in=g["cell_inset_in"],
        labels=labels, excluded=excluded, priority=priority)
    # 橫幅接在表格**實際**底緣之後。
    # ⚠ 用 `used_height`（_add_table 逐列累加的真值），不是 `(shown+1) * row_height_in`
    #   ——後者是宣告列高，內容一換行就低估，橫幅會往上壓住表格（實機 p11／p12）。
    _render_points_band(slide, theme, spec, ctx,
                        top=g["top_in"] + used_height + g["points_band_inset_in"],
                        height=band_height)
    _render_footnote(slide, theme, spec, ctx,
                     _rows_note(shown, rows, int(g["max_columns"]), _visible_column_count(rows, excluded)))


#: 題目卡明細的標籤成本（「依據｜」「行動｜」各 3 字）。
DIRECTION_DETAIL_LABEL_COST = 3


def direction_capacity(theme: Theme, *, topic_text_pt: float | None = None) -> dict[str, int]:
    """研發方向頁（p19）各欄位的**實際版面容量**，由 theme 幾何推導。

    🔴 R-1（2026-08-05 實機 p19）：規範檔原本寫死「basis／action 各 ≤20 字」，
    那個 20 是 **12.5pt 時代**算出來的（卡片文字寬 3.5in ÷ (12.5/72) ＝ 20.1 字，
    剛好對上）。K-10 把 `topic_text_pt` 改 16pt 後每行只剩 15.7 字、再扣掉
    「依據｜」3 字＝12.7 字——20 字必被 `_fit_text` 截成「…」。
    ⚠ 改字級沒同步字數上限，就是「同一份知識兩處落點」的第 15 例；
    故上限一律**從幾何算**，不再寫死：日後改字級、改卡片大小都自動跟著動。

    回傳的鍵直接進 PPT 文案提示（見 ai_report_ppt_runner.build_report_ppt_payload），
    讓 CLI 拿到的數字與組版端量的是同一個。
    """
    g = theme.geometry["direction_flow"]
    size = float(topic_text_pt if topic_text_pt is not None else theme.size("topic_text_pt"))
    flow_size = theme.size("flow_text_pt")
    detail_width = g["topic_width_in"] - g["topic_inset_in"] * 2
    detail_height = g["topic_height_in"] - g["topic_text_top_offset_in"] - g["topic_inset_in"]
    detail_per_line, detail_lines = _text_capacity(
        theme, width_in=detail_width, height_in=detail_height, size_pt=size)
    step_width = g["step_width_in"] - g["step_inset_in"] * 2
    step_height = g["step_height_in"] - g["step_text_top_offset_in"] - g["step_inset_in"]
    step_per_line, step_lines = _text_capacity(
        theme, width_in=step_width, height_in=step_height, size_pt=flow_size)
    concl_width = g["conclusion_width_in"] - g["conclusion_inset_in"] * 2
    concl_per_line, _ = _text_capacity(
        theme, width_in=concl_width,
        height_in=g["conclusion_height_in"] - g["conclusion_text_top_offset_in"] * 2,
        size_pt=theme.size("conclusion_pt"))
    return {
        "topic_detail_per_line": int(detail_per_line),
        "topic_detail_lines": int(detail_lines),
        # 每段（依據／行動）各佔一段，扣掉標籤成本後的正文上限
        "topic_detail_max_chars": max(1, int(detail_per_line) * max(1, int(detail_lines) // 2)
                                      - DIRECTION_DETAIL_LABEL_COST),
        "topic_name_max_chars": int((g["topic_width_in"] - g["topic_inset_in"] * 2)
                                    / (theme.size("topic_name_pt") / 72.0)),
        "step_line_max_chars": int(step_per_line),
        "step_max_lines": int(step_lines),
        "conclusion_max_chars": max(1, int(concl_per_line) - len("核心結論：")),
        "topic_max": int(g["topic_max"]),
    }


def validate_direction_body(theme: Theme, body: dict[str, Any]) -> list[str]:
    """驗 direction.body 各欄位是否在版面容量內，回傳警告清單（合規＝空）。

    ⚠ 只寫在提示裡、沒有程式驗證的規則等於沒有規則（known-issues C-1）——
    截斷是 `_fit_text` 的最後防線，但截了讀者就看不到後半句，必須在這裡先叫出來。
    """
    cap = direction_capacity(theme)
    warnings: list[str] = []
    for index, topic in enumerate(body.get("topics") or []):
        for field in ("basis", "action"):
            text = str((topic or {}).get(field) or "")
            if len(text) > cap["topic_detail_max_chars"]:
                warnings.append(
                    f"direction.topics[{index}].{field} 超出卡片容量"
                    f"（{len(text)} 字 > {cap['topic_detail_max_chars']}）——會被截斷")
        name = str((topic or {}).get("name") or "")
        if len(name) > cap["topic_name_max_chars"]:
            warnings.append(
                f"direction.topics[{index}].name 超限（{len(name)} 字 > "
                f"{cap['topic_name_max_chars']}）")
    for key in ("situation", "opportunity", "direction"):
        for index, line in enumerate(body.get(key) or []):
            if len(str(line)) > cap["step_line_max_chars"] * cap["step_max_lines"]:
                warnings.append(f"direction.{key}[{index}] 超出色塊容量")
    conclusion = str(body.get("conclusion") or "")
    if len(conclusion) > cap["conclusion_max_chars"]:
        warnings.append(
            f"direction.conclusion 超限（{len(conclusion)} 字 > "
            f"{cap['conclusion_max_chars']}）")
    return warnings


def _parse_direction_body(body: str) -> dict[str, Any] | None:
    """解析結構化 direction.body（P1-7 合併版契約）；不是合法 JSON 就回 None。

    契約形狀（content_rules 同步定義；AI 產 JSON 字串塞進 slot）：
        {"situation": [..], "opportunity": [..], "direction": [..],
         "topics": [{"name","basis","action"}, ..], "conclusion": ".."}
    ⚠ 舊純文字要能過渡：回 None 讓 renderer 走舊的條列版面，不炸、不硬轉。
    """
    text = str(body or "").strip()
    if not text.startswith("{"):
        return None
    try:
        parsed = json.loads(text)
    except json.JSONDecodeError:
        return None
    if not isinstance(parsed, dict):
        return None
    return {
        "situation": [str(x) for x in parsed.get("situation") or []],
        "opportunity": [str(x) for x in parsed.get("opportunity") or []],
        "direction": [str(x) for x in parsed.get("direction") or []],
        "topics": [t for t in parsed.get("topics") or [] if isinstance(t, dict)],
        "conclusion": str(parsed.get("conclusion") or ""),
    }


def image_frame_height(theme: Theme, top_in: float, declared_height_in: float) -> float:
    """圖框可用高度——**不得延伸進頁尾帶**（R-3，2026-08-05 實機 p17／p18）。

    🔴 象限板長寬比 ~1.47 比圖框（8.9×5.0＝1.78）更高，於是**高度受限**、
    撐滿整個框高：框底 1.86+5.0＝6.86in 落進頁尾帶（footnote.top 6.78in），
    圖自己最後一行底注「本分析非侵權迴避(FTO)結論…」就疊在組版頁尾
    「資料來源：…」上面。⚠ 寬度受限的圖（長條 949×453＝2.1）實際高度只有 4.2in，
    本來就碰不到頁尾——所以夾限只會咬到真正會撞的那幾張，其餘零影響。

    ⚠ 用夾限而不是把宣告高度改小：頁尾位置是唯一事實來源（theme.footnote），
    日後頁尾搬家時圖框自動跟著讓位，不必再記得同步第二個數字。
    """
    footnote_top = theme.geometry["footnote"]["top_in"]
    return min(declared_height_in, max(0.0, footnote_top - top_in))


def _fitted_size_from_px(width_px: float, height_px: float,
                         box_w: float, box_h: float) -> tuple[float, float]:
    """已知像素尺寸時的等比縮放結果（供 `_fitted_size` 與測試共用同一套算法）。"""
    if not width_px or not height_px:
        return box_w, box_h
    scale = min(box_w / (width_px / 96), box_h / (height_px / 96))
    return (width_px / 96) * scale, (height_px / 96) * scale


def _fitted_size(image_path: Path, box_w: float, box_h: float) -> tuple[float, float]:
    """圖等比縮放塞進框後的**實際**尺寸（英吋）。

    元件要對齊圖的實際範圍而不是框——圖填不滿框時，靠框對齊的說明文字會飄在
    空白處（獨立驗收在扁圖頁抓到）。取不到尺寸就回框的大小，退化為原行為。
    """
    try:
        from PIL import Image

        with Image.open(image_path) as img:
            width_px, height_px = img.size
    except Exception:
        return box_w, box_h
    return _fitted_size_from_px(width_px, height_px, box_w, box_h)


def _render_chart_wide(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """扁圖版型（ratio ≥3.5）：圖佔滿寬置頂，要點在下方橫幅。

    ⚠ 為什麼另立版型：`chart_hero` 的框是 8.9×4.32（比例 2.06），扁圖塞進去會
    **寬度先滿、高度大量浪費**——IPC 四階分布（比例 6.05）實測只用到 1.5 in 高，
    下方空掉三分之一頁。改滿寬後縮放倍率 0.87→1.19（+37%，實測 6 張皆然）。
    ⚠ 門檻取 3.5 而非「比框扁就換」：比例 2.1–3.0 的圖換滿寬反而更差
    （倍率 0.43–0.82 < 現行 0.56–0.76），因為高度被底部要點壓縮。
    """
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["chart_wide"]
    image = ctx["charts"].resolve(spec.charts[0]) if spec.charts else None

    # 說明文字靠右對齊**圖的實際右緣**，不是框的右緣。
    shown_w, shown_h = (_fitted_size(image, g["image_width_in"], g["image_height_in"])
                        if image is not None else (g["image_width_in"], g["image_height_in"]))
    edge_left = g["image_left_in"] + (g["image_width_in"] - shown_w) / 2
    _add_text(slide, theme, _encoding_note(spec, ctx),
              left=edge_left, top=g["encoding_top_in"],
              width=shown_w, height=g["encoding_height_in"],
              size=theme.size("encoding_note_pt"), color="muted", align=PP_ALIGN.RIGHT)
    if image is not None:
        # ⚠ 框高直接給**圖的實際高度**：`_add_picture_fitted` 會把圖置中在框裡，
        # 框比圖高就會往下推（IPC 四階實測被推 0.87 in），底緣算式就對不上、
        # 要點橫幅直接蓋住圖表。把框縮成圖的大小，置中即成無作用。
        _add_picture_fitted(slide, image,
                            left=g["image_left_in"], top=g["image_top_in"],
                            width=g["image_width_in"], height=shown_h)

    # 要點橫幅跟著圖的**實際底緣**走；圖矮時橫幅上移，不留一大塊空白。
    band_top = max(g["band_min_top_in"], g["image_top_in"] + shown_h + g["band_gap_in"])
    _render_wide_points_band(slide, theme, spec, ctx, top=band_top)
    _render_footnote(slide, theme, spec, ctx)


def _render_wide_points_band(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any],
                             *, top: float) -> None:
    """扁圖頁的底部要點橫幅（雙欄），高度吃到頁尾之前的可用空間。"""
    g = theme.geometry["chart_wide"]
    inset = g["band_inset_in"]
    columns_n = int(g["band_columns"])
    gap_w = g["band_column_gap_in"]
    col_w = (g["band_width_in"] - inset - inset - gap_w * (columns_n - 1)) / columns_n
    size_pt = theme.size("point_text_pt")
    # ⚠ 高度依**內容**決定，`band_bottom_in` 只是上限：橫幅若一律吃到底，
    # 資料少的頁面（IPC 四階只有 2 條）下半部就是一大片空白——
    # 這正是本批要治的毛病，換個位置再犯一次沒有意義。
    lines = sum(_lines_needed(f"{label}｜{text}", _text_capacity(
        theme, width_in=col_w, height_in=g["band_bottom_in"], size_pt=size_pt,
        line_ratio=point_line_ratio(theme))[0])
        for label, text, _, _ in _points_for(spec, ctx))
    per_column_lines = math.ceil(lines / columns_n) if lines else 1
    needed = (g["band_text_top_offset_in"] + inset
              + per_column_lines * size_pt / 72.0 * point_line_ratio(theme))
    height = min(g["band_bottom_in"] - top,
                 max(g["band_header_height_in"] + inset + inset, needed))
    # 🔴 K-11：去框（同側欄面板）。
    _add_text(slide, theme, "判讀要點",
              left=g["band_left_in"] + inset, top=top + g["band_header_top_offset_in"],
              width=g["band_width_in"] - inset - inset, height=g["band_header_height_in"],
              size=theme.size("panel_header_pt"), color="accent", bold=True)

    columns, gap, col_width, size = columns_n, gap_w, col_w, size_pt
    text_top = top + g["band_text_top_offset_in"]
    text_height = height - g["band_text_top_offset_in"] - inset
    blocks = _trim_blocks(theme, _points_for(spec, ctx),
                          width_in=col_width, height_in=text_height * columns, size_pt=size)
    # 🔴 K-3（2026-08-04）：分欄改**按行數**貪婪均分。原本按條數對半（3 條→左 2 右 1），
    # 三層制後段長差異大——左欄（現況＋意涵）行數遠超右欄（後續），
    # 實機 p7 左欄溢出壓到頁尾資料來源、p9 貼底。
    per_line_band, _ = _text_capacity(theme, width_in=col_width, height_in=text_height,
                                      size_pt=size, line_ratio=point_line_ratio(theme))
    needs = [_lines_needed(f"{label}｜{text}" if label else text, per_line_band)
             for label, text, _c, _e in blocks]
    total_lines = sum(needs)
    chunks: list[list[tuple[str, str, str, bool]]] = []
    start = 0
    acc = 0
    for index, need in enumerate(needs):
        # 加了這段會超過均分行數就先換欄（⚠ 加完才斷會把長段整段堆在左欄，
        # 左欄溢出、右欄空一半——驗證時 p7 實際發生）。
        if (len(chunks) < columns - 1 and acc > 0
                and acc + need > total_lines / columns):
            chunks.append(blocks[start:index])
            start = index
            acc = 0
        acc += need
    chunks.append(blocks[start:])
    for index, chunk in enumerate(chunks[:columns]):
        if not chunk:
            continue
        _add_number_bold_text(slide, theme, chunk,
                              left=g["band_left_in"] + inset + index * (col_width + gap),
                              top=text_top, width=col_width, height=text_height, size=size)


def _render_direction_flow(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any],
                           parsed: dict[str, Any]) -> None:
    """合併版（2026-07-31 使用者經表單選定）：
    上半「態勢→機會→方向」三步色塊流程、下半研發題目卡片、底部核心結論條。
    python-pptx 畫不了原生 SmartArt——色塊＋箭頭是同等效果的確定性圖形。
    """
    g = theme.geometry["direction_flow"]
    # 🔴 2026-07-31 批2：三步原本用 royal／blue／navy，明暗完全不一致——
    # `blue`(4FC3F7) 是亮青，配 on_dark_soft 深灰字等於讀不出來（實機第 17 頁）。
    # 改為三階**都夠深**的同族藍，靠由深到淺表達流程推進，文字一律亮字。
    # ⚠ 三張卡是**並列的同階資訊**（y/w/h 完全相同），不是有大小的序階，
    # 故底色一律相同——流程推進由中間的箭頭表達。批2 初版讓三者由暗到亮遞增，
    # 獨立驗收指出那是語意錯亂，且最暗的「態勢」卡對背景只有 1.343 不達標。
    steps = (("態勢", parsed["situation"], "panel"),
             ("機會", parsed["opportunity"], "panel"),
             ("方向", parsed["direction"], "panel"))
    for index, (label, lines, color) in enumerate(steps):
        left = g["step_left_in"] + index * g["step_gap_in"]
        _add_band(slide, theme, left, g["step_top_in"], g["step_width_in"], g["step_height_in"],
                  color, rounded=True, line="royal")
        _add_text(slide, theme, label,
                  left=left + g["step_inset_in"], top=g["step_top_in"] + g["step_label_top_offset_in"],
                  width=g["step_width_in"] - g["step_inset_in"] * 2, height=g["step_label_height_in"],
                  size=theme.size("flow_label_pt"), color="accent", bold=True)
        body_width = g["step_width_in"] - g["step_inset_in"] * 2
        body_height = g["step_height_in"] - g["step_text_top_offset_in"] - g["step_inset_in"]
        body, _ = _fit_text(theme, "\n".join(lines), width_in=body_width, height_in=body_height,
                            size_pt=theme.size("flow_text_pt"))
        _add_text(slide, theme, body,
                  left=left + g["step_inset_in"], top=g["step_top_in"] + g["step_text_top_offset_in"],
                  width=body_width, height=body_height,
                  # 批2：on_dark_soft 在最亮的 panel_deep 上只有 3.28，改亮字。
                  size=theme.size("flow_text_pt"), color="on_dark")
        if index < len(steps) - 1:
            _add_text(slide, theme, "→",
                      left=left + g["step_width_in"] + g["arrow_left_offset_in"],
                      top=g["arrow_top_in"], width=g["arrow_width_in"], height=g["arrow_height_in"],
                      size=theme.size("flow_arrow_pt"), color="accent", bold=True,
                      align=PP_ALIGN.CENTER)

    topics = parsed["topics"][: int(g["topic_max"])]
    for index, topic in enumerate(topics):
        left = g["topic_left_in"] + index * g["topic_gap_in"]
        _add_band(slide, theme, left, g["topic_top_in"], g["topic_width_in"], g["topic_height_in"],
                  "panel", rounded=True)
        _add_band(slide, theme, left, g["topic_top_in"], g["topic_width_in"],
                  g["topic_accent_height_in"], "accent")
        _add_text(slide, theme, str(topic.get("name") or ""),
                  left=left + g["topic_inset_in"], top=g["topic_top_in"] + g["topic_name_top_offset_in"],
                  width=g["topic_width_in"] - g["topic_inset_in"] * 2, height=g["topic_name_height_in"],
                  size=theme.size("topic_name_pt"), color="ink", bold=True)
        detail_lines = []
        if topic.get("basis"):
            detail_lines.append(f"依據｜{topic['basis']}")
        if topic.get("action"):
            detail_lines.append(f"行動｜{topic['action']}")
        body_width = g["topic_width_in"] - g["topic_inset_in"] * 2
        body_height = g["topic_height_in"] - g["topic_text_top_offset_in"] - g["topic_inset_in"]
        body, _ = _fit_text(theme, "\n".join(detail_lines), width_in=body_width,
                            height_in=body_height, size_pt=theme.size("topic_text_pt"))
        _add_text(slide, theme, body,
                  left=left + g["topic_inset_in"], top=g["topic_top_in"] + g["topic_text_top_offset_in"],
                  width=body_width, height=body_height,
                  size=theme.size("topic_text_pt"), color="ink")

    if parsed["conclusion"]:
        _add_band(slide, theme, g["conclusion_left_in"], g["conclusion_top_in"],
                  g["conclusion_width_in"], g["conclusion_height_in"], "panel_deep", rounded=True)
        text_width = g["conclusion_width_in"] - g["conclusion_inset_in"] * 2
        body, _ = _fit_text(theme, f"核心結論：{parsed['conclusion']}", width_in=text_width,
                            height_in=g["conclusion_height_in"] - g["conclusion_text_top_offset_in"] * 2,
                            size_pt=theme.size("conclusion_pt"))
        _add_text(slide, theme, body,
                  left=g["conclusion_left_in"] + g["conclusion_inset_in"],
                  top=g["conclusion_top_in"] + g["conclusion_text_top_offset_in"],
                  width=text_width,
                  height=g["conclusion_height_in"] - g["conclusion_text_top_offset_in"] * 2,
                  size=theme.size("conclusion_pt"), color="on_dark", bold=True)


def _render_direction(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """研發方向建議（結論頁，壓軸）。

    結構化 direction.body → 合併版（色塊流程＋題目卡＋結論條）；
    舊純文字 → 走原條列版面過渡，並由 build_ppt 記 warning（不靜默）。
    """
    _render_header(slide, theme, spec, ctx)
    parsed = _parse_direction_body(ctx["slots"].get("direction.body") or "")
    if parsed is not None:
        _render_direction_flow(slide, theme, spec, ctx, parsed)
        _render_footnote(slide, theme, spec, ctx)
        return
    g = theme.geometry["direction"]
    _add_band(slide, theme, g["body_left_in"], g["body_top_in"],
              g["body_width_in"], g["body_height_in"], "panel", rounded=True)
    _add_text(slide, theme, "研發方向與具體題目",
              left=g["body_left_in"] + g["body_header_inset_left_in"],
              top=g["body_top_in"] + g["body_header_top_offset_in"],
              width=g["body_width_in"] - g["body_text_inset_right_in"], height=g["body_header_height_in"],
              size=theme.size("panel_header_pt"), color="accent", bold=True)
    text_width = g["body_width_in"] - g["body_text_inset_right_in"]
    text_height = g["body_height_in"] - g["body_text_top_offset_in"] - g["body_text_bottom_pad_in"]
    size = theme.size("body_pt")
    body = ctx["slots"].get("direction.body") or ""
    blocks = [("", line.strip(), "ink", False) for line in body.split("\n") if line.strip()]
    if not blocks:
        blocks = [(label, text, "ink", False) for label, text in _row_highlights(spec, ctx)]
    _add_number_bold_text(slide, theme,
                          _trim_blocks(theme, blocks, width_in=text_width, height_in=text_height, size_pt=size),
                          left=g["body_left_in"] + g["body_text_inset_left_in"],
                          top=g["body_top_in"] + g["body_text_top_offset_in"],
                          width=text_width, height=text_height, size=size)

    _add_band(slide, theme, g["basis_left_in"], g["basis_top_in"],
              g["basis_width_in"], g["basis_height_in"], "panel", rounded=True)
    _add_text(slide, theme, "專利地圖依據",
              left=g["basis_left_in"] + g["basis_header_inset_left_in"],
              top=g["basis_top_in"] + g["basis_header_top_offset_in"],
              width=g["basis_item_width_in"], height=g["basis_header_height_in"],
              size=theme.size("panel_header_pt"), color="accent", bold=True)
    for index, label in enumerate(ctx["included_report_labels"][: int(g["basis_item_max"])]):
        text, _ = _fit_text(theme, f"・{label}", width_in=g["basis_item_width_in"],
                            height_in=g["basis_item_height_in"], size_pt=theme.size("point_label_pt"))
        _add_text(slide, theme, text,
                  left=g["basis_left_in"] + g["basis_item_inset_left_in"],
                  top=g["basis_item_top_in"] + index * g["basis_item_step_in"],
                  width=g["basis_item_width_in"], height=g["basis_item_height_in"],
                  size=theme.size("point_label_pt"), color="on_dark_soft")
    _render_footnote(slide, theme, spec, ctx)


def _table_display(ctx: dict[str, Any], spec: PageSpec) -> tuple[dict[str, str], set[str], tuple[str, ...]]:
    """本頁表格的欄名對照、排除欄與**顯示優先序**：引擎那份優先，缺鍵才用本檔 fallback。

    排除欄是**逐報表**的（同一欄在 A 報表要藏、在 B 報表要顯示），故依本頁掛的
    report_keys 逐一併集。

    🔴 優先序（2026-08-03）：欄位放不下時砍尾巴不砍中間。上一輪 `status`
    排在 rows 第 7 位被 `max_columns` 依鍵順序切掉，整輪重點功能一格沒顯示（G-2）。
    ⚠ 哪一欄重要是**資料語意**，故順序由引擎宣告，組版端只照著取。
    """
    display = ctx["report_data"].get("table_display") or {}
    labels = {**TABLE_COLUMN_LABELS, **(display.get("column_labels") or {})}
    excluded = set(TABLE_EXCLUDED_COLUMNS)
    per_report = display.get("excluded_columns") or {}
    priority: list[str] = []
    per_report_priority = display.get("priority_columns") or {}
    for key in spec.report_keys:
        excluded.update(per_report.get(key) or ())
        for name in per_report_priority.get(key) or ():
            if name not in priority:
                priority.append(name)
    return labels, excluded, tuple(priority)


def _rows_note(shown: int, rows: list[dict[str, Any]], max_columns: int, visible_columns: int) -> str:
    """表格截列／截欄時據實說明，避免讀者把前 N 筆當成全部。

    ⚠ `visible_columns` 是**排除欄之後**的欄數：拿原始欄數去比會把引擎刻意藏起來
    的欄（topic_code、龍頭涉入等）算進分母，印出「前 6/10 欄」這種嚇人又不實的註記。
    """
    notes = []
    if rows and shown < len(rows):
        notes.append(f"顯示前 {shown}/{len(rows)} 筆")
    if visible_columns > max_columns:
        notes.append(f"前 {max_columns}/{visible_columns} 欄，完整欄位見附錄")
    return "、".join(notes)


def _first_rows(spec: PageSpec, ctx: dict[str, Any]) -> list[dict[str, Any]]:
    """取本頁表格要印的 rows；帶 row_filter 的頁（依通道拆頁）只留匹配的列。

    ⚠ 優先取引擎寫在 `table_display.display_rows` 的**呈現字串**：`top_applicants`
    這類欄的原始值是物件陣列，直接印會變成 `{'name': '祺驊', ...`（2026-07-31 實機
    第 9、10、18 頁）。呈現規則的唯一來源在引擎（`chart_runner._humanize_cell`），
    本檔不自建第二份；舊報表版本沒有這個鍵時退回原始 rows。
    """
    display = (ctx["report_data"].get("table_display") or {}).get("display_rows") or {}
    filters = dict(spec.row_filter)
    for key in spec.report_keys:
        rows = display.get(key) or _rows_of(ctx["report_data"], key)
        if filters:
            rows = [r for r in rows
                    if all(str(r.get(col)) == value for col, value in filters.items())]
        if rows:
            # 附錄分頁的列切片（2026-08-03）——⚠ 必須在**過濾之後**才切：
            # 先切再過濾會讓每頁的可見列數對不上分頁時算的那一份。
            if spec.row_slice:
                start, stop = spec.row_slice
                return rows[start:stop]
            return rows
    return []


def _ordered_columns(
    rows: list[dict[str, Any]],
    *,
    excluded: set[str],
    priority: tuple[str, ...] | list[str],
    limit: int,
) -> list[str]:
    """決定表格要顯示哪幾欄、以什麼順序——**放不下時砍尾巴，不砍中間**。

    🔴 G-2（2026-08-03 實機）：`status`（技術狀態）在 rows 裡排第 7，
    被 `max_columns=6` 依鍵順序切掉——S2 整輪的重點功能一格都沒顯示出來，
    而頁尾還寫著「完整欄位見附錄」，附錄同樣只有 6 欄。

    ⚠ 順序取自引擎的 `priority_columns`（唯一來源），組版端不自己排——
    哪一欄重要是資料語意，不是版面問題。
    ⚠ 沒列進優先序的欄位排在後面但**不消失**，否則新增欄位會被靜默吞掉。
    """
    if not rows:
        return []
    available = [name for name in rows[0] if str(name) not in excluded]
    ordered = [name for name in priority if name in available]
    ordered += [name for name in available if name not in ordered]
    return ordered[:limit]


def _table_line_plan(
    theme: Theme,
    rows: list[dict[str, Any]],
    columns: list[str],
    labels: dict[str, str],
    width: float,
    *,
    cell_inset_in: float,
) -> tuple[list[float], int, list[int]]:
    """算出欄寬、表頭行數、以及**每一列各佔幾行**。

    ⚠ 抽出來的理由（2026-08-03）：附錄分頁要在 `_expand_page_layout` 階段就知道
    「一頁放得下幾列」，而那個答案只有這段邏輯算得準。留在 `_add_table` 裡面
    等於逼分頁端另寫一套估法——本專案已因「同一資訊兩處落點」靜默失敗六次。
    """
    col_widths = _column_widths(columns, rows, labels, width,
                                size_pt=theme.size("table_body_pt"), inset_in=cell_inset_in)
    text_widths = [w - cell_inset_in * 2 for w in col_widths]
    body_pt = theme.size("table_body_pt")

    def _lines_for(row: dict[str, Any]) -> int:
        needed = 1
        for index, name in enumerate(columns):
            value = row.get(name)
            if isinstance(value, list):
                value = "、".join(str(v) for v in value)
            text = "" if value is None else str(value)
            span = _display_width(text) * (body_pt / 72.0)
            needed = max(needed, math.ceil(span / max(text_widths[index], 1e-6)))
        return needed

    header_lines = max(
        (math.ceil(_display_width(labels.get(str(n), str(n))) * (theme.size("table_header_pt") / 72.0)
                   / max(text_widths[i], 1e-6)) for i, n in enumerate(columns)), default=1)
    return col_widths, max(1, header_lines), [_lines_for(row) for row in rows]


def _appendix_rows_per_page(
    theme: Theme,
    rows: list[dict[str, Any]],
    labels: dict[str, str],
    excluded: set[str],
    priority: tuple[str, ...],
) -> int:
    """附錄一頁放得下幾列（用與渲染端同一套行數估算）。"""
    if not rows:
        return 0
    g = theme.geometry["table"]
    columns = _ordered_columns(rows, excluded=excluded, priority=priority,
                               limit=int(g["max_columns"]))
    _widths, header_lines, line_counts = _table_line_plan(
        theme, rows, columns, labels, g["width_in"], cell_inset_in=g["cell_inset_in"])
    height = _table_available_height(theme, "table")
    row_height = g["row_height_in"]
    used = header_lines * row_height
    count = 0
    for lines in line_counts:
        if used + lines * row_height > height and count:
            break
        used += lines * row_height
        count += 1
    return max(1, count)


def _add_table(
    slide,
    theme: Theme,
    rows: list[dict[str, Any]],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    row_height: float,
    max_columns: int,
    cell_margin_in: float,
    cell_inset_in: float,
    labels: dict[str, str],
    excluded: set[str],
    priority: tuple[str, ...] = (),
) -> tuple[int, float]:
    """把引擎 rows 畫成表格，回傳（實際顯示列數, **實際用掉的高度**）。

    列數依框高與**逐列實際行數**累加後截斷（PowerPoint 的表格列高只增不減）。

    ⚠ 高度一定要交出來（H-1，2026-08-03）：本函式內部本來就算過 `used_height`，
    但沒回傳，於是呼叫端拿宣告列高 `row_height_in` 自己重估一次——同一個量兩處落點，
    而且重估的那處是錯的（宣告 0.32 vs 換行後實際 0.6），底部要點橫幅因此壓住表格。
    """
    if not rows:
        g = theme.geometry["table"]
        _add_band(slide, theme, left, top, width, height, "panel", rounded=True)
        _add_text(slide, theme, "本頁報表無資料列",
                  left=g["empty_text_left_in"], top=g["empty_text_top_in"],
                  width=g["empty_text_width_in"], height=g["empty_text_height_in"],
                  size=theme.size("table_body_pt"), color="muted", align=PP_ALIGN.CENTER)
        return 0, height

    # 欄位顯示規則：排除欄與中文欄名以引擎那份為準（labels／excluded 由呼叫端備妥），
    # 欄值轉譯仍在本檔（source_field 的原始欄值不得入畫面，轉「技術／功效」）。
    columns = _ordered_columns(rows, excluded=excluded, priority=priority,
                               limit=max_columns)

    # 🔴 2026-08-03 使用者定案：**資訊不能有被截斷的**。
    # 原本是「放不下就切掉加『…』」——讀者既不知道被切掉什麼，也無從查證。
    # 改為：欄寬依內容分配 → 放不下就換行 → **列高跟著長**；
    # 真的塞不進框時少顯示幾列（完整版在附錄），但**顯示出來的每一格都完整**。
    col_widths, header_lines, row_line_counts_all = _table_line_plan(
        theme, rows, columns, labels, width, cell_inset_in=cell_inset_in)

    # 逐列累加實際高度，超過框就停——不是用 height/row_height 平均估。
    display: list[dict[str, Any]] = []
    row_line_counts: list[int] = []
    used = header_lines * row_height
    for row, lines in zip(rows, row_line_counts_all):
        if used + lines * row_height > height and display:
            break
        display.append(row)
        row_line_counts.append(lines)
        used += lines * row_height
    # 表高依實際列數收縮：宣告高度是**上限**不是固定值，列少時下半截留白會很難看
    # （主題分布通常只有 8–12 列，舊版固定 4.86 in 有一半是空的）。
    used_height = min(height, (header_lines + sum(row_line_counts)) * row_height)
    table = slide.shapes.add_table(
        len(display) + 1, len(columns), Inches(left), Inches(top), Inches(width), Inches(used_height)
    ).table
    # 列高依該列實際行數——換了行卻不加高，字就會被切在框外。
    table.rows[0].height = Inches(header_lines * row_height)
    for index, lines in enumerate(row_line_counts, start=1):
        table.rows[index].height = Inches(lines * row_height)

    for index, col_width in enumerate(col_widths):
        table.columns[index].width = Inches(col_width)
    for index, name in enumerate(columns):
        cell = table.cell(0, index)
        cell.text = labels.get(str(name), str(name))
        # v3 深空：表頭深藍底＋accent 青字（原白字在深底上與內文分不出層次）。
        _style_cell(cell, theme, size=theme.size("table_header_pt"), color="accent", bold=True,
                    fill="navy", margin_in=cell_margin_in, inset_in=cell_inset_in)
    for r, row in enumerate(display, start=1):
        for c, name in enumerate(columns):
            value = row.get(name)
            mapped = TABLE_VALUE_LABELS.get(str(name), {})
            if isinstance(value, list):
                value = "、".join(str(v) for v in value)
            value = mapped.get(str(value), value) if mapped else value
            cell = table.cell(r, c)
            cell.text = "" if value is None else str(value)
            # bold=True：v3 使用者定案「文字內容加粗體」，深底細字會發灰。
            _style_cell(cell, theme, size=theme.size("table_body_pt"), color="ink", bold=True,
                        fill="paper" if r % 2 else "panel_alt",
                        margin_in=cell_margin_in, inset_in=cell_inset_in)
    return len(display), used_height


def _set_cell_borders(cell, theme: Theme) -> None:
    """給儲存格四邊套主題細線。

    ⚠ python-pptx 沒有儲存格框線 API，不寫就會沿用 PowerPoint **預設表格樣式**的
    淺灰線（實測約 A8B3BD–DFE0E3），與全簡報統一的細線不一致——規格 S1-6 寫了
    「外框細線」卻一直沒實作，獨立驗收在 p9/p10/p18/p19 抓到。
    ⚠ 四邊都要寫：只寫左與上，相鄰儲存格之間會留下缺口。
    """
    colour = theme.color["hairline"]
    width = int(Pt(theme.font["panel_border_pt"]))
    props = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = props.find(qn(tag))
        if existing is not None:
            props.remove(existing)
        props.append(parse_xml(
            f'<{tag} {nsdecls("a")} w="{width}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{colour}"/></a:solidFill>'
            f'</{tag}>'
        ))


def _style_cell(
    cell, theme: Theme, *, size: float, color: str, bold: bool = False, fill: str = "paper",
    margin_in: float = 0.0, inset_in: float = 0.0,
) -> None:
    """套儲存格底色與字體。

    ⚠ python-pptx 的儲存格預設上下內距 0.05 in，PowerPoint 的列高只增不減——
    預設內距會讓實際列高遠大於宣告值，整張表往下溢出版面。故一律壓縮內距。
    """
    cell.fill.solid()
    cell.fill.fore_color.rgb = theme.rgb(fill)
    cell.margin_top = Inches(margin_in)
    cell.margin_bottom = Inches(margin_in)
    cell.margin_left = Inches(inset_in)
    cell.margin_right = Inches(inset_in)
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            _set_font(run, theme, size=size, color=color, bold=bold)
    _set_cell_borders(cell, theme)


# ── KP 版型（P2；範例＝滑雪機 V2 p7–p10）──────────────────────────
# ⚠ 三個都是**備選版型**：沒有那個內容就不出那一頁（2026-08-07 定案），
# 由規劃端（SlidePlan）決定出不出，組版端只負責畫得對。


def _render_kp_quadrant(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """KP 競爭定位象限：整頁一張泡泡圖（軸義與圖例都畫在 SVG 內）。

    沿用 chart_hero 的幾何與圖框邏輯——象限圖與一般大圖的版面需求相同，
    另立一套座標只會讓 theme 多一份會漂移的定義。
    """
    _render_chart_hero(slide, theme, spec, ctx)


def _render_kp_deepdive(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """單一 KP 深入：演進時間軸＋三數字卡（家族件數／技術群／布局地區）。

    ⚠ 沒有軌跡（不同申請年 <3）的公司**不該出這一頁**——規劃端把關；
    組版端若真收到無軌跡資料，仍照畫數字卡，不自行降級成別的版型。
    """
    _render_chart_with_points(slide, theme, spec, ctx)


def _render_kp_cards(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """利基／新興玩家小卡矩陣（名稱＋一句定位＋件數）。"""
    _render_table_with_points(slide, theme, spec, ctx)


# 無圖要點頁的面板標題（依版型語意，不重複頁標題）。
POINTS_PAGE_PANEL_TITLES = {
    "exec_summary": "關鍵結論",
    "walls_gaps": "要迴避的牆與可切入的空白",
    "reading_guide": "判讀說明",
}


def _render_points_page(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """無圖的純要點頁：整頁一個面板，內容＝本頁 narrative。

    🔴 2026-08-09：這三種版型原本轉呼叫 `_render_direction`，但那支的內容來自
    固定 slot `direction.body`（研發方向頁專用），**讀不到 SlidePlan 的
    narrative**——實機轉圖出來是三張一模一樣的空框（p2／p5／p9）。
    ⚠ 版型名稱對、renderer 存在、雙向一致性測試也綠，成品仍是空的：
    那種測試驗得到「有沒有 renderer」，驗不到「renderer 畫了什麼」。

    座標沿用 direction 的 body 框並延伸到 basis 框右緣（不另立一套幾何）。
    """
    _render_header(slide, theme, spec, ctx)
    g = theme.geometry["direction"]
    left = g["body_left_in"]
    width = (g["basis_left_in"] + g["basis_width_in"]) - left
    _add_band(slide, theme, left, g["body_top_in"], width, g["body_height_in"],
              "panel", rounded=True)
    # ⚠ 面板標題不能用 spec.topic——那就是頁標題，會一頁印兩次（首版實測）。
    _add_text(slide, theme, POINTS_PAGE_PANEL_TITLES.get(spec.kind, "重點"),
              left=left + g["body_header_inset_left_in"],
              top=g["body_top_in"] + g["body_header_top_offset_in"],
              width=width - g["body_text_inset_right_in"],
              height=g["body_header_height_in"],
              size=theme.size("panel_header_pt"), color="accent", bold=True)
    text_width = width - g["body_text_inset_right_in"]
    text_height = (g["body_height_in"] - g["body_text_top_offset_in"]
                   - g["body_text_bottom_pad_in"])
    size = theme.size("body_pt")
    blocks = _points_for(spec, ctx)
    if not blocks:
        blocks = [(label, text, "ink", False) for label, text in _row_highlights(spec, ctx)]
    _add_number_bold_text(
        slide, theme,
        _trim_blocks(theme, blocks, width_in=text_width, height_in=text_height, size_pt=size),
        left=left + g["body_text_inset_left_in"],
        top=g["body_top_in"] + g["body_text_top_offset_in"],
        width=text_width, height=text_height, size=size)
    _render_footnote(slide, theme, spec, ctx)


def _render_exec_summary(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """結論先行頁（範例 p2）：把三個可行動判斷放最前面。

    內容來自 SlidePlan 的 narrative（CLI 已把結論寫成具名發現）。
    """
    if spec.report_keys and spec.charts:
        _render_table_with_points(slide, theme, spec, ctx)
    else:
        _render_points_page(slide, theme, spec, ctx)


def _render_walls_gaps(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """要迴避的牆 vs 可切入的空白（範例 p3／割草機 p2）：收斂成可行動清單。"""
    _render_points_page(slide, theme, spec, ctx)


def _render_reading_guide(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """判讀說明（範例 p11）：母體口徑、可觀測性偏差、資料限制。"""
    _render_points_page(slide, theme, spec, ctx)


def _render_kp_compare(slide, theme: Theme, spec: PageSpec, ctx: dict[str, Any]) -> None:
    """兩個 Key Player 左右對照（範例 p9）：核心技術架構與布局並列。"""
    _render_comparison(slide, theme, spec, ctx)


RENDERERS = {
    "cover": _render_cover,
    "section_divider": _render_section_divider,
    "chart_hero": _render_chart_hero,
    "chart_with_points": _render_chart_with_points,
    "comparison": _render_comparison,
    "stat_callout": _render_stat_callout,
    "percentage_bars": _render_percentage_bars,
    "table": _render_table,
    "table_with_points": _render_table_with_points,
    "chart_wide": _render_chart_wide,
    "direction": _render_direction,
    # KP 版型（P2）：備選版型庫的一部分，出不出由規劃端依內容決定。
    "kp_quadrant": _render_kp_quadrant,
    "kp_deepdive": _render_kp_deepdive,
    "kp_cards": _render_kp_cards,
    "kp_compare": _render_kp_compare,
    # 敘事鏈版型（範例骨架）：結論先行、可行動清單、判讀說明。
    "exec_summary": _render_exec_summary,
    "walls_gaps": _render_walls_gaps,
    "reading_guide": _render_reading_guide,
}

# 需要圖才成立的版型：解析不到圖就降級 stat_callout，不留佔位文字。
CHART_DEPENDENT_KINDS = frozenset({"chart_hero", "chart_with_points", "comparison",
                                  "kp_quadrant", "kp_deepdive"})
# 單圖版型：被指定給多圖頁面時要拆成多頁（成對報表的「分頁」呈現）。
SINGLE_CHART_KINDS = frozenset({"chart_hero", "chart_with_points", "stat_callout"})


# --------------------------------------------------------------------------
# 頁面展開
# --------------------------------------------------------------------------
def _kind_for_report(report: dict[str, Any], chart_count: int) -> str:
    """動態頁的預設版型：明細／矩陣走表格，多圖走並排，單圖走大圖，沒圖走大數字。

    單圖預設 `chart_hero`（P1-1）：大圖＋底部核心結論條——「圖表為主、文字為輔」
    的預設呈現；要點框版（chart_with_points）保留給版型下拉切換。
    """
    if str(report.get("report_type") or "").lower() in {"detail", "table"}:
        return "table"
    if chart_count >= 2:
        return "comparison"
    if chart_count == 1:
        return "chart_hero"
    return "stat_callout"


def _page_should_render(report_data: dict[str, Any], spec: PageSpec) -> bool:
    """cover／direction 恆出；其餘頁面必須至少有一個實際有資料的 report_key。"""
    if spec.kind in {"cover", "direction", "section_divider"}:
        return True
    return bool(_actual_report_keys(report_data, spec.report_keys))


def _kind_for_aspect(kind: str, files: tuple[str, ...], charts: ChartIndex | None) -> str:
    """單圖頁若是**扁圖**，改用滿寬版型；其餘維持原版型。

    ⚠ 只對單圖頁生效：多圖頁（並排比較）本來就各自佔半寬，換滿寬沒有意義。
    """
    if charts is None or kind not in SINGLE_CHART_KINDS or len(files) != 1:
        return kind
    aspect = charts.aspect_of(files[0])
    return "chart_wide" if aspect and aspect >= WIDE_CHART_ASPECT_MIN else kind


def _evidence_rank(spec: PageSpec) -> int:
    """證據頁在論證鏈上的位置；未列名者排在已知證據之後（仍在結論之前）。

    取第一個命中 EVIDENCE_ORDER 的 report_key——成對報表（如申請人／專利權人
    排名同頁）以先列者定位即可，兩者在論證上本來就相鄰。
    """
    for key in spec.report_keys:
        if key in EVIDENCE_ORDER:
            return EVIDENCE_ORDER.index(key)
    return len(EVIDENCE_ORDER)


def _expand_page_layout(report_data: dict[str, Any], charts: ChartIndex | None = None,
                        theme: Theme | None = None) -> list[PageSpec]:
    """選擇驅動出頁：把有資料的報表展開成頁面清單，頁碼重新連號。

    `charts` 省略時（前端縮圖預覽只拿得到 report_data）僅回頁面骨架，
    圖檔欄位留空——留空是誠實的，猜檔名才是錯的。
    """
    base: list[PageSpec] = []
    for spec in PAGE_LAYOUT:
        if not _page_should_render(report_data, spec):
            continue
        keys = spec.report_keys if spec.kind in {"cover", "direction"} else _actual_report_keys(report_data, spec.report_keys)
        # 封面／隔頁／方向頁不擺報表圖，charts 留空才不會讓 manifest 反查出誤導的對照。
        files = _filter_report_charts(keys, charts.files_for(keys)) if charts and spec.kind not in {"cover", "direction", "section_divider"} else ()
        kind = spec.kind
        if kind == "comparison" and len(files) < 2:
            kind = "chart_hero"
        resolved = _spec_with(spec, report_keys=keys, charts=files, kind=kind)
        # 依通道拆頁（P1-3）：主題分布 rows 帶兩通道，各自成一張表格頁；
        # 單一通道時維持一頁、不加通道字樣。
        base.extend(_split_by_channel(resolved, report_data))

    covered = {key for spec in PAGE_LAYOUT for key in spec.report_keys}
    extra: list[PageSpec] = []
    for report_key, report in _iter_report_entries(report_data):
        if report_key in covered or report_key in EXCLUDED_FROM_PPT:
            continue
        if not _report_key_has_data(report_data, report_key):
            continue
        files = _filter_report_charts((report_key,), charts.files_for((report_key,))) if charts else ()
        topic = _label_of(report_data, report_key)
        extra.append(
            PageSpec(page=0, kind=_kind_for_report(report, len(files)), title=topic, topic=topic,
                     report_keys=(report_key,), charts=files)
        )

    # ⚠ 錨點＝direction 或第一個附錄，取先出現者：動態插頁也是證據，
    #   必須排在結論（研發方向建議）之前——結論永遠壓軸（P1-6）。
    anchor = next((i for i, spec in enumerate(base)
                   if spec.is_appendix or spec.kind == "direction"), len(base))
    # 證據段依 EVIDENCE_ORDER 重排（2026-07-31）：固定條目與動態插頁**混在一起**
    # 排序，而不是「固定的在前、動態的在後」——否則 CPC 仍會離 IPC 很遠。
    # sort 是穩定的，故未列名的報表維持引擎輸出的相對順序，只是整批落在最後。
    head = [spec for spec in base[:anchor] if spec.kind == "cover"]
    evidence = [spec for spec in base[:anchor] if spec.kind != "cover"] + extra
    evidence.sort(key=_evidence_rank)
    merged = head + evidence + base[anchor:]
    merged = _split_pairs_by_policy(merged, charts)
    # ⚠ theme 省略時不分頁：既有呼叫端（前端縮圖預覽只拿得到 report_data）行為不變。
    if theme is not None:
        merged = _paginate_appendix(merged, report_data, theme)
    # 🔴 拆頁後補回分類系統與階層（F-8 移除 SVG 標題後，頁標題必須自己講清楚
    # 這是 IPC 還是 CPC、哪一階——否則四頁併排讀者分不出誰是誰）。
    merged = [_spec_with(spec, topic=_chart_page_topic(spec, report_data)) for spec in merged]
    # ⚠ 版型的長寬比調整要在**拆頁之後**：多圖頁要拆完才知道每頁只有一張圖，
    # 且動態插頁走 `_kind_for_report` 不經基礎迴圈——放在這裡才涵蓋全部頁面。
    merged = [_spec_with(spec, kind=_kind_for_aspect(spec.kind, spec.charts, charts))
              for spec in merged]
    return [_spec_with(spec, page=index) for index, spec in enumerate(merged, start=1)]


def split_rows_evenly(total: int, *, per_page: int) -> list[int]:
    """把 total 列切成每頁不超過 per_page 的**平均**份數。

    🔴 I-5（2026-08-03 實機 p21／p22）：原本「每頁塞滿 per_page 筆」，
    最後一頁拿到餘數——8 筆、每頁 7 筆切成 **7＋1**，
    第 2 頁只有一列、整頁 90% 空白，而且兩頁欄寬還不一樣（各自算）。

    改為**先算頁數、再平均攤**：8 筆 2 頁 → 4＋4。
    ⚠ 每頁仍不得超過 `per_page`（版面放得下的量）——平均只在頁數確定後攤平，
    不會因為攤平而讓某頁塞爆。
    """
    if total <= 0:
        return []
    if total <= per_page:
        return [total]
    pages = math.ceil(total / per_page)
    base, extra = divmod(total, pages)
    # 前 extra 頁各多一列——差距最多 1，不會出現「7＋1」這種懸殊分配。
    return [base + (1 if i < extra else 0) for i in range(pages)]


def _paginate_appendix(
    specs: list[PageSpec],
    report_data: dict[str, Any],
    theme: Theme,
) -> list[PageSpec]:
    """附錄放不下一頁時切成多頁——**附錄要放齊**（2026-08-03 使用者定案）。

    ⚠ 只切附錄：內頁是「精選」，少列是刻意的；附錄的職責才是「完整」。
    ⚠ 每頁能放幾列用 `_appendix_rows_per_page`，與渲染端同一套行數估算——
    分頁端另寫一套的話，切出來的頁數與實際放得下的列數會對不起來。
    ⚠ 標題加「（N/M）」讓讀者知道還有下一頁；只有一頁時不加，避免無謂的雜訊。
    """
    ctx = {"report_data": report_data}
    out: list[PageSpec] = []
    for spec in specs:
        if not spec.is_appendix or spec.kind != "table":
            out.append(spec)
            continue
        rows = _first_rows(spec, ctx)
        if not rows:
            out.append(spec)
            continue
        labels, excluded, priority = _table_display(ctx, spec)
        per_page = _appendix_rows_per_page(theme, rows, labels, excluded, priority)
        if per_page >= len(rows):
            out.append(spec)
            continue
        sizes = split_rows_evenly(len(rows), per_page=per_page)
        total_pages = len(sizes)
        start = 0
        for index, size in enumerate(sizes):
            out.append(_spec_with(
                spec,
                title=f"{spec.title}（{index + 1}/{total_pages}）",
                row_slice=(start, start + size),
            ))
            start += size
    return out


def _split_by_channel(spec: PageSpec, report_data: dict[str, Any]) -> list[PageSpec]:
    """rows 含多通道的報表依通道拆頁（每通道一張表，**附錄也拆**）。

    🔴 H-4（2026-08-03 實機）：附錄原本不拆，於是「全分類**技術**指標總表」裡
    混進「提升訓練成效」這種功效主題——標題宣稱只有技術，內容不是。
    使用者定案：**內頁精選、附錄放齊**；放齊的前提是兩種主題各自成表。

    ⚠ 附錄的標題要**保留「附錄N」前綴**再加通道名：內頁拆頁是直接把 title 換成
    通道名，附錄照抄會丟掉附錄身分，讀者在頁序上找不到它。
    """
    if not spec.report_keys:
        return [spec]
    config = CHANNEL_SPLIT_REPORTS.get(spec.report_keys[0])
    if config is None:
        return [spec]
    column, channels = config
    rows = _rows_of(report_data, spec.report_keys[0])
    present = [(value, topic) for value, topic in channels
               if any(str(r.get(column)) == value for r in rows)]
    if len(present) <= 1:
        return [spec]
    if spec.is_appendix:
        # 「附錄1：全分類技術指標總表」＋「技術主題分布」→「附錄1：技術主題分布（全表）」
        prefix = spec.title.split("：", 1)[0]
        return [
            _spec_with(spec, topic=topic, title=f"{prefix}：{topic}（全表）",
                       row_filter=((column, value),))
            for value, topic in present
        ]
    return [
        _spec_with(spec, topic=topic, title=topic, row_filter=((column, value),))
        for value, topic in present
    ]


def theme_comparison_columns() -> list[float]:
    """並排版型的欄位左緣（欄數＝len）。⚠ 取自 theme，不在程式寫死欄數。"""
    return list(Theme.load().geometry["comparison"]["column_left_in"])


def _chart_page_topic(spec: PageSpec, report_data: dict[str, Any]) -> str:
    """拆頁後的頁標題主題——帶上分類系統與階層。

    🔴 2026-08-03 使用者：「IPC/CPC 標題沒寫，看的人會搞混」。
    F-8 移除 SVG 內建標題（「IPC 主分類分布 - Level 4」）時，我判斷
    「headline 已經能區分」——實測**不能**：p7「技術分類布局：A63B次分類達47件」
    完全沒說這是 IPC，階層也不見了。四頁併排讀者分不出誰是誰。

    ⚠ 顯示名取自引擎的 section title 與 variant label（**唯一來源**），
    不在組版端另寫一份 L4→「次分類」的對照表——那就是第二處落點。
    ⚠ 對不到 section 時回原 topic，不得產生半截標題。
    """
    if not spec.charts:
        return spec.topic
    wanted = str(spec.charts[0])
    for section in (report_data.get("sections") or []):
        for variant in (section.get("variants") or []):
            if str(variant.get("file") or "") != wanted:
                continue
            title = str(section.get("title") or "").strip()
            label = str(variant.get("label") or "").strip()
            if title and label:
                return f"{title}（{label}）"
            return title or spec.topic
    return spec.topic


def _split_pairs_by_policy(layout: list[PageSpec], charts: ChartIndex | None = None) -> list[PageSpec]:
    """成對報表改成分頁；⚠ 並排版型只有兩欄，超過兩張圖一律拆頁。

    🔴 2026-07-31：技術分類布局頁掛了 4 張圖（IPC L4/L5 ＋ CPC L4/L5），
    `comparison` 只畫得下兩欄，**CPC 兩張從來沒被畫出來**——而頁尾仍寫著
    「資料來源：IPC 主分類分布、CPC 主分類分布」。靜默丟圖比畫不下更糟：
    讀者以為看到的就是全部。故除了 SPLIT_PAIR_REPORTS 的政策拆頁外，
    只要圖數超過並排欄數就一律拆。
    """
    columns = len(theme_comparison_columns())
    result: list[PageSpec] = []
    for spec in layout:
        policy_split = any(key in SPLIT_PAIR_REPORTS for key in spec.report_keys)
        overflow = len(spec.charts) > columns
        if spec.kind == "comparison" and (policy_split or overflow):
            # 拆頁後用大圖版型（chart_hero）——分頁的動機就是圖要大。
            result.extend(_split_multi_chart_page(_spec_with(spec, kind="chart_hero"), charts))
        else:
            result.append(spec)
    return result


def _split_multi_chart_page(spec: PageSpec, charts: ChartIndex | None = None) -> list[PageSpec]:
    """單圖版型碰到多圖頁面：一圖一頁。

    拆頁時把 `report_keys` 一併收窄到該圖真正對應的報表——否則兩頁都掛著全部
    report_key，會抓到同一段 narrative、印出兩張一模一樣的標題與註腳。
    """
    if len(spec.charts) <= 1:
        return [spec]
    pages: list[PageSpec] = []
    for name in spec.charts:
        owners = charts.owners_of(name, spec.report_keys) if charts else ()
        pages.append(_spec_with(spec, charts=(name,), report_keys=owners or spec.report_keys))
    return pages


# ══ SlidePlan 消費（P2 第 5 節）══════════════════════════════════════
# 🔴 分工紅線：CLI 決定「哪一頁講什麼、用哪張圖、用哪種版型」；
# builder 決定「那種版型長什麼樣」（座標／字級／顏色一律由 theme 解析）。
# CLI 若硬塞幾何欄位，這裡**一律忽略**——不是報錯後照用，是根本不讀。


class SlidePlanError(RuntimeError):
    """SlidePlan 無法轉成版面（未知版型等）。"""


def page_specs_from_plan(plan: dict[str, Any],
                        charts: "ChartIndex | None" = None) -> list[PageSpec]:
    """把通過驗證的 SlidePlan 轉成 PageSpec 序列（頁碼依 slides 順序連號）。

    ⚠ `charts` 要給：chart_identity 只是 `report_key:variant`，實際檔名不同名
    （country_distribution 的圖叫 jurisdiction_distribution.svg）——不反查就會
    每頁都因「找不到圖」降級成 stat_callout（2026-08-09 首次串接實測）。
    """
    specs: list[PageSpec] = []
    for index, slide in enumerate(plan.get("slides") or [], start=1):
        preset = str(slide.get("layout_preset") or "")
        if preset not in RENDERERS:
            raise SlidePlanError(
                f"slide {slide.get('slide_id', '?')} 的版型 {preset!r} 不在組版端支援清單")
        identities = [str(i) for i in slide.get("chart_identities") or []]
        report_keys = tuple(i.split(":", 1)[0] for i in identities)
        # ⚠ files_for 收的是 tuple；傳單一字串會被當序列逐字元迭代而全部落空。
        files = list(charts.files_for(report_keys)) if charts is not None else []
        specs.append(PageSpec(
            page=index,
            kind=preset,
            title=str(slide.get("title") or slide.get("purpose") or ""),
            topic=str(slide.get("purpose") or ""),
            # chart_identity 形如 `report_key:variant_key`——取前段當 report_key。
            report_keys=report_keys,
            charts=tuple(dict.fromkeys(files)),
        ))
    return specs


def plan_coverage_manifest(
    plan: dict[str, Any],
    selected_identities: set[str],
) -> dict[str, Any]:
    """選圖覆蓋清單：使用者選了卻沒進 PPT 的圖必須現形。"""
    used: set[str] = set()
    for slide in plan.get("slides") or []:
        used.update(str(i) for i in slide.get("chart_identities") or [])
    return {
        "plan_id": plan.get("plan_id", ""),
        "slide_count": len(plan.get("slides") or []),
        "used_charts": sorted(used),
        "missing_selected": sorted(selected_identities - used),
        "unselected_used": sorted(used - selected_identities),
    }


def resolve_layout(report_data: dict[str, Any], charts: ChartIndex,
                   theme: Theme | None, overrides: dict[str, str]) -> list[PageSpec]:
    """從報表資料算出最終版面：展開 → 套版型覆寫（含拆頁）→ 圖檔降級 → **重算標題**。

    🔴 2026-08-04（J-2）：最後那步不能省。`_expand_page_layout` 在**還是一頁兩張圖**
    的時候就把 topic 定死了，`_apply_layout_overrides` 之後才一圖一頁拆開——
    於是 p13／p14 兩頁的標題一字不差都是「主要申請人排名（Applicants）」，
    但 p14 畫的是 `owner_ranking.svg`、內容講的是權利人。**只有標題錯**。

    ⚠ `_split_multi_chart_page` 已經很細心地把 `report_keys` 收窄到該圖真正對應的
    報表（它的 docstring 就寫著「否則兩頁會印出一模一樣的標題」），但漏了 topic 這一半。
    這裡在拆完之後重算一次——`_chart_page_topic` 依 `charts[0]` 判斷，拆完才問得到正確答案。

    ⚠ 收成單一入口是為了讓測試驗得到**完整結果**：三個步驟散在呼叫端時，
    測試只驗得到中間狀態，正是這個 bug 混過去的原因。
    """
    # P2：有通過驗證的 SlidePlan 就照它出頁；沒有就走既有固定 PAGE_LAYOUT 展開
    # ⚠ 舊路徑保留（不是換掉）——既有報告仍要能重產，且出問題有回頭路。
    plan = (report_data.get("slide_plan") or {}) if isinstance(report_data, dict) else {}
    if plan.get("slides"):
        layout = page_specs_from_plan(plan, charts)
    else:
        layout = _expand_page_layout(report_data, charts, theme)
    layout = _apply_layout_overrides(layout, overrides, charts)
    layout = _apply_chart_degradation(layout, charts)
    return [_spec_with(spec, topic=_chart_page_topic(spec, report_data)) for spec in layout]


def _clean_layout_overrides(value: Any) -> dict[str, str]:
    """只接受 renderer 支援的版型名稱，無效覆寫忽略而非讓產檔失敗。"""
    if not isinstance(value, dict):
        return {}
    return {str(page): str(kind) for page, kind in value.items() if str(kind) in RENDERERS}


def _apply_layout_overrides(
    layout: list[PageSpec], overrides: dict[str, str], charts: ChartIndex | None = None
) -> list[PageSpec]:
    """套用使用者挑的版型；單圖版型碰到多圖頁面會自動拆頁，頁碼重新連號。"""
    result: list[PageSpec] = []
    for spec in layout:
        kind = overrides.get(str(spec.page), spec.kind)
        target = _spec_with(spec, kind=kind)
        if kind in SINGLE_CHART_KINDS and len(spec.charts) > 1:
            result.extend(_split_multi_chart_page(target, charts))
        elif kind == "comparison" and len(spec.charts) < 2:
            result.append(_spec_with(target, kind="chart_with_points"))
        else:
            result.append(target)
    return [_spec_with(spec, page=index) for index, spec in enumerate(result, start=1)]


def _apply_chart_degradation(layout: list[PageSpec], charts: ChartIndex) -> list[PageSpec]:
    """圖檔缺失／轉檔失敗 → 降級 stat_callout，確保每頁都有視覺元素。"""
    result: list[PageSpec] = []
    for spec in layout:
        if spec.kind not in CHART_DEPENDENT_KINDS:
            result.append(spec)
            continue
        usable = tuple(name for name in spec.charts if charts.resolve(name) is not None)
        if usable:
            result.append(_spec_with(spec, charts=usable))
        else:
            result.append(_spec_with(spec, kind="stat_callout", charts=(), degraded_from=spec.kind))
    return result


# --------------------------------------------------------------------------
# 封面統計卡與分析框架
# --------------------------------------------------------------------------
def _cover_title(report_data: dict[str, Any], slots: dict[str, str]) -> str:
    """封面主標＝workspace 顯示名稱（P1-8，確定性組成；cover.title AI slot 已退場）。

    ⚠ slots 參數保留：使用者若日後經 approvals 明確給了標題仍尊重（人工定稿
    優先於推導），但**不再請 AI 產**——AI 只剩 direction.body 一個 slot。
    parameters 缺 workspace 名稱時退回通用標題，不硬湊。
    """
    manual = str(slots.get("cover.title") or "").strip()
    if manual:
        return manual
    # 🔴 2026-08-09：goal-driven 規劃時，CLI 會在 s1 的第一條 narrative 寫出
    # 針對這批資料的主標。⚠ 順位排在 workspace 名稱**之後**（不推翻 07-31
    # 定案），但要排在通用 fallback 之前——否則 workspace 名稱缺失時主標會
    # 退回寫死的「專利情報整合分析」，與封面 eyebrow 一字不差印兩次（實測）。
    plan_title = ""
    for slide in ((report_data.get("slide_plan") or {}).get("slides") or [])[:1]:
        for point in slide.get("narrative") or []:
            plan_title = str(point.get("text") or "").strip()
            if plan_title:
                break
    params = report_data.get("parameters") or {}
    for key in ("workspace_name", "workspace_display_name", "workspace"):
        value = str(params.get(key) or "").strip()
        if value:
            # 2026-07-31 使用者定案：「封面頁主題要顯示成 workspace 名稱配上專利分析」
            # ——單獨一個「滑雪機」不像簡報標題，補上主題詞才成句。
            return value if value.endswith("專利分析") else f"{value}專利分析"
    return plan_title or COVER_TITLE_FALLBACK


def _cover_funnel(report_data: dict[str, Any]) -> tuple[str, str, str] | None:
    """三層漏斗併 1 格（Q3，2026-08-05 使用者定案）：原始→同族合併→技術主題。

    ⚠ 單位與標籤走 unit／label 欄，**不得併進 value**——四張卡取同一級、
    級數由最長值決定，把「件」「群」寫進 value 會讓整排字縮小（規格明列風險）。
    ⚠ 技術群數只算**技術通道**（功效通道不上封面，也不得加總）。
    缺同族或分群資料時回 None：少一格，不硬湊也不寫 0。
    """
    trend_rows = _rows_of(report_data, "application_trend")
    if not trend_rows:
        return None
    total = sum(_as_int(row.get("patent_count")) for row in trend_rows)
    # 第二層＝**同族合併後件數**（distinct 家族）。⚠ 不得用各國家族數加總——
    # 跨國家族在每個國家各算一次（2026-08-07 真資料：加總 46 ≠ distinct 48）。
    family_total = _as_int((report_data.get("parameters") or {}).get("family_merged_total"))
    topic_rows = (report_data.get("chart_rows") or {}).get("cluster_topic_table") or []
    tech_groups = len({
        str(r.get("topic_code") or r.get("topic_key") or "")
        for r in topic_rows
        if str(r.get("source_field") or "") == TECHNICAL_SOURCE_FIELD
    } - {""})
    if not (total and family_total and tech_groups):
        return None
    return (f"{total}→{family_total}→{tech_groups}", "件→件→群", "原始→同族合併→技術主題")


def _cover_stats(report_data: dict[str, Any]) -> list[tuple[str, str, str]]:
    """封面統計卡；資料不足就少一格，不硬湊低價值指標。"""
    stats: list[tuple[str, str, str]] = []
    trend_rows = _rows_of(report_data, "application_trend")
    if trend_rows:
        total = sum(_as_int(row.get("patent_count")) for row in trend_rows)
        stats.append((f"{total:,}", "件", "專利總數"))
    country_rows = _rows_of(report_data, "country_distribution")
    if country_rows:
        numeric = _numeric_column(country_rows)
        label_col = _label_column(country_rows, numeric)
        # 🔴 J-4（2026-08-04）：原本只列前二國——「39｜9」加起來 48，
        # 但專利總數 60（US 9、EP 3 被丟）；且 TW／US 同 9 件時挑誰是任意的。
        # 改為 ≤4 局全列、>4 局取前 3 ＋「其他」合計：**件數總和恆等於總數**，
        # 排序（件數 desc, 代碼 asc）決定同數順序，不再任意。
        # 🔴 2026-08-07：country_distribution 改 (國×狀態) 群組後一國多列——
        # 先按國彙總再排序，否則同一國重複出現、數字變成狀態分項。
        merged: dict[str, int] = {}
        for row in country_rows:
            code = str(row.get(label_col, "-"))
            merged[code] = merged.get(code, 0) + _as_int(row.get(numeric))
        ordered = [{label_col: code, numeric: total}
                   for code, total in sorted(merged.items(), key=lambda kv: (-kv[1], kv[0]))]
        if len(ordered) > 4:
            shown = ordered[:3]
            rest = sum(_as_int(r.get(numeric)) for r in ordered[3:])
            pairs = [(str(r.get(label_col, "-")), _as_int(r.get(numeric))) for r in shown]
            pairs.append(("其他", rest))
        else:
            pairs = [(str(r.get(label_col, "-")), _as_int(r.get(numeric))) for r in ordered]
        stats.append((
            "｜".join(str(n) for _, n in pairs),
            "｜".join(code for code, _ in pairs),
            "地域分布（件數）",
        ))
    period = _statistics_period(report_data)
    if period:
        stats.append((period, "年", "年份區間"))
    funnel = _cover_funnel(report_data)
    if funnel:
        stats.append(funnel)
    # 第 4 格由資料現有欄位組成；都沒有就只出 3 格。
    for report_key, unit, label in (
        ("applicant_ranking", "家", "申請人家數"),
        ("ipc_main_distribution", "類", "IPC 主分類數"),
        ("cluster_topic_table", "群", "技術主題數"),
    ):
        rows = _rows_of(report_data, report_key)
        if rows:
            stats.append((str(len(rows)), unit, label))
            break
    return stats[:4]


FRAMEWORK_BUDGET_CHARS = 58


def _cover_stat_size(theme: Theme, stats: list[tuple[str, str, str]], value: str) -> float:
    """封面統計卡的主數字字級——**四張卡取同一級**。

    ⚠ 級數由最長的那個值決定：卡片寬度固定，最長的放得下，其餘自然放得下。
    逐張各算會讓長值被降級、短值維持大字，並排時像重要性不同（實機 p1）。
    """
    longest = max((len(str(v)) for v, _unit, _label in stats), default=len(str(value)))
    if longest <= 4:
        return theme.size("stat_value_pt")
    if longest <= 8:
        return theme.size("stat_value_medium_pt")
    return theme.size("stat_value_small_pt")


def _framework_text(layout: list[PageSpec], budget_chars: int = FRAMEWORK_BUDGET_CHARS) -> str:
    """分析框架條：固定的論證鏈分組名＋本次實際項數。

    🔴 J-11（2026-08-04）：原本用「實際頁面主題名」串動線——主題名 15–20 字，
    58 字預算只塞得下 1–2 個，實機 p1 印成「專利申請趨勢與專利授權公告趨勢
    （Trend）→ 等共 16 項分析」，讀起來像句子被切掉。
    ⚠ F-15 修的是「收尾被截」，沒修「只剩半條動線」——根因是拿長度不可控的
    主題名去湊長度固定的一行。

    改用**固定分組名**（時間→地域→技術→競爭→機會＝論證鏈的五段，
    見 EVIDENCE_ORDER 的排序意圖）：短而穩定，永遠放得下、永遠是完整句。
    """
    topics = list(dict.fromkeys(
        spec.topic or spec.title
        for spec in layout
        if spec.kind not in {"cover", "direction", "section_divider"} and not spec.is_appendix
    ))
    if not topics:
        return "分析框架：本次僅含封面與研發方向建議"
    return f"分析框架：時間趨勢 → 地域布局 → 技術分類 → 競爭者 → 機會評估，共 {len(topics)} 項分析"


# --------------------------------------------------------------------------
# 產後自檢（QA）：溢出、邊距、重疊、文字裝不下
# --------------------------------------------------------------------------
def _shape_font_pt(shape) -> float:
    sizes = [
        run.font.size.pt
        for para in shape.text_frame.paragraphs
        for run in para.runs
        if run.font.size is not None
    ]
    return max(sizes) if sizes else 0.0


def audit_layout(prs: Presentation, theme: Theme) -> list[dict[str, Any]]:
    """逐頁逐 shape 檢查版面問題，回傳 warnings（不修改簡報，只回報）。

    檢四項：超出版面邊界、邊距不足、文字疊文字、文字估算裝不下。
    最後一項用字級與框大小估算——PowerPoint 的文字溢出不會改變 shape 尺寸，
    只看座標抓不到「字太多」，而字牆正是這次重建要解決的問題。
    """
    warnings: list[dict[str, Any]] = []
    slide_w = float(theme.slide["width_in"])
    slide_h = float(theme.slide["height_in"])
    margin = float(theme.qa["min_margin_in"])
    bounds_tol = float(theme.qa["bounds_tolerance_in"])
    overlap_tol = float(theme.qa["overlap_tolerance_in"])
    slack = float(theme.qa["capacity_slack"])

    for page, slide in enumerate(prs.slides, start=1):
        boxes: list[tuple[str, float, float, float, float]] = []
        # 🔴 空白頁偵測（2026-08-09）：一頁若只有標題與註腳、正文區完全沒東西，
        # 那是「版型有 renderer 但沒畫出內容」——實機出過三張一模一樣的空框。
        # ⚠ 放在這個後置掃描而不是各 renderer 內部：這裡對**所有版型**一體適用，
        # 在 renderer 裡各記各的只保護當下改到的那幾種，換一種版型又會靜默。
        body_chars = 0
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            left = shape.left / 914400
            top = shape.top / 914400
            right = left + (shape.width or 0) / 914400
            bottom = top + (shape.height or 0) / 914400
            text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            name = f"{shape.shape_type}｜{text[:18]}" if text else str(shape.shape_type)

            # 全出血底板（封面／隔頁背景）本來就貼齊版面，不算違規。
            full_bleed = left <= bounds_tol and top <= bounds_tol and right >= slide_w - bounds_tol
            # 邊距規則保護的是「內容不要貼邊」；色塊、圓點、斜線這類裝飾出血是設計語彙，
            # 只受版面邊界（out_of_bounds）約束，不受安全區約束。
            is_content = bool(text) or shape.has_table or shape.shape_type == MSO_SHAPE_TYPE.PICTURE

            if not full_bleed and (
                left < -bounds_tol or top < -bounds_tol
                or right > slide_w + bounds_tol or bottom > slide_h + bounds_tol
            ):
                warnings.append({
                    "type": "out_of_bounds", "page": page, "shape": name,
                    "overflow_in": round(max(-left, -top, right - slide_w, bottom - slide_h), 3),
                })
            elif is_content and not full_bleed and (
                left < margin - bounds_tol or top < margin - bounds_tol
                or right > slide_w - margin + bounds_tol or bottom > slide_h - margin + bounds_tol
            ):
                warnings.append({
                    "type": "margin_violation", "page": page, "shape": name,
                    "margin_in": round(min(left, top, slide_w - right, slide_h - bottom), 3),
                })

            # 正文帶＝標題與註腳之間；圖片視為足量內容（圖表頁的正文就是圖）。
            if BODY_BAND_TOP_RATIO * slide_h < top < BODY_BAND_BOTTOM_RATIO * slide_h:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.has_table:
                    body_chars += EMPTY_BODY_MIN_CHARS
                else:
                    body_chars += len(text)

            if not text:
                continue
            boxes.append((name, left, top, right, bottom))

            size_pt = _shape_font_pt(shape)
            if size_pt:
                per_line, lines = _text_capacity(
                    theme, width_in=right - left, height_in=bottom - top, size_pt=size_pt
                )
                needed = _lines_needed(text, per_line)
                if needed > lines * slack:
                    warnings.append({
                        "type": "text_overflow_estimated", "page": page, "shape": name,
                        "lines_needed": needed, "lines_available": lines,
                    })

        if body_chars < EMPTY_BODY_MIN_CHARS:
            warnings.append({
                "type": "empty_body", "page": page, "body_chars": body_chars,
                "detail": "正文區幾乎沒有內容——版型有 renderer 但沒畫出東西，"
                          "或規劃未提供本頁的敘述",
            })

        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                overlap_w = min(a[3], b[3]) - max(a[1], b[1])
                overlap_h = min(a[4], b[4]) - max(a[2], b[2])
                if overlap_w > overlap_tol and overlap_h > overlap_tol:
                    warnings.append({
                        "type": "text_overlap", "page": page,
                        "shapes": [a[0], b[0]],
                        "overlap_in": [round(overlap_w, 3), round(overlap_h, 3)],
                    })
    return warnings


# --------------------------------------------------------------------------
# 產出
# --------------------------------------------------------------------------
def _next_available_path(output_dir: Path, version: str) -> Path:
    """同版本重跑不覆蓋舊檔，改產 `_r2`、`_r3`。"""
    base = output_dir / f"{version}.pptx"
    if not base.exists():
        return base
    index = 2
    while (output_dir / f"{version}_r{index}.pptx").exists():
        index += 1
    return output_dir / f"{version}_r{index}.pptx"


def _sha256_of(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1 << 20), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _load_json(path: Path, default: dict) -> dict:
    if not path.exists():
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return default


def _narrative_candidates(spec: PageSpec) -> tuple[str, ...]:
    """narrative 查找鍵：report_key → 圖檔主檔名 → 去後綴主檔名 → alias。

    ⚠ alias（NARRATIVE_ALIASES）解「兩個 key 空間不同名」：機會矩陣的解讀
    掛在 `cluster_topic_table:opportunity_*`（sections 結構），report bucket
    卻叫 `opportunity_quadrant`——沒有 alias 這頁永遠 narrative_missing
    （2026-07-31 實機 P10）。
    """
    # 🔴 順序＝精確度，不是「先 report_key 再圖檔」（2026-07-31 修）。
    # 症狀：機會評估拆成技術／功效兩頁後，兩頁印出同一段解讀（圖是功效、文字是技術）。
    # 根因：拆頁時 report_keys 被收窄成同一個 `opportunity_quadrant`，兩頁完全相同；
    # 而 alias 依 report_key 層先展開，它同時對到 tech 與 effect 兩個變體，
    # 於是兩頁都命中先列的 tech。圖檔主檔名（`opportunity_quadrant_effect`）才是
    # 能區分兩頁的唯一線索，必須先查。
    # 通道拆頁（主題統計表）最精確：該頁只呈現一個通道的列，解讀也該只取那個通道。
    # ⚠ 這頁沒有圖檔，所以圖檔主檔名那條線索不存在，只能靠 row_filter 認。
    channel: list[str] = []
    for column, value in spec.row_filter:
        variant = CHANNEL_NARRATIVE_VARIANTS.get(str(value))
        if not variant:
            continue
        for key in spec.report_keys:
            candidate = f"{key}:{variant}"
            if candidate not in channel:
                channel.append(candidate)

    specific: list[str] = []
    for name in spec.charts:
        stem = Path(name).stem
        if stem not in specific:
            specific.append(stem)
        for suffix in CHART_ORDER_HINTS:
            if not stem.endswith(suffix):
                continue
            base = stem[: -len(suffix)]
            # 🔴 2026-08-02：檔名後綴**就是**變體鍵，先翻成 `report_key:variant` 再查。
            # 症狀：p8 IPC Level 4 與 p9 Level 5 的標題與四條要點逐字相同（CPC 亦然）。
            # 根因：解讀端其實產了 L4／L5 兩段不同內容，但候選鍵只有主檔名
            # `ipc_main_distribution_L4`（narratives 沒這個鍵）與 base
            # `ipc_main_distribution`，於是 `_narrative_entry` 走「取 variants 第一個」
            # 那條路，兩頁都拿 L4。⚠ scoped 必須排在 base 前面，否則等於沒改。
            scoped = f"{base}:{suffix.lstrip('_')}"
            if scoped not in specific:
                specific.append(scoped)
            if base not in specific:
                specific.append(base)
    generic = [key for key in spec.report_keys if key not in specific]

    keys: list[str] = channel + list(specific) + generic
    # alias 也照同一順序展開：圖檔層 alias 一律排在 report_key 層 alias 之前。
    for key in specific + generic:
        for alias in NARRATIVE_ALIASES.get(key, ()):
            if alias not in keys:
                keys.append(alias)
    return tuple(keys)


def build_ppt(
    *,
    report_dir: Path | str,
    approvals_path: Path | str | None = None,
    output_dir: Path | str | None = None,
    theme_path: Path | str = THEME_PATH,
) -> dict[str, Any]:
    """依版型表組出報告 PPTX，回傳輸出路徑、manifest 路徑與 manifest 內容。"""
    report_dir = Path(report_dir)
    report_data = _load_json(report_dir / "report_data.json", {})
    # 門檻缺頁現形（design #5）：低於 IPC/CPC 出頁門檻的報表不出頁，
    # 原因記進 manifest warnings——沒有這一筆，讀者只會以為漏產。
    threshold_warnings = threshold_skips(report_data)
    narratives = _load_json(report_dir / "narratives.json", {})
    artifact_manifest = _load_json(report_dir / "artifact_manifest.json", {})
    approvals = _load_json(Path(approvals_path), {}) if approvals_path else {}
    slots: dict[str, str] = approvals.get("slots") or {}

    version = (
        (report_data.get("parameters") or {}).get("version")
        or approvals.get("report_version")
        or report_dir.name
    )
    output_dir = Path(output_dir) if output_dir else Path("data/report_artifacts/ppt")
    output_dir.mkdir(parents=True, exist_ok=True)

    theme = Theme.load(theme_path)
    # ⚠ 同一個 process 連續產兩份時要先清空，否則上一份丟掉的要點會算到這一份頭上。
    reset_dropped_points()
    charts = ChartIndex(report_dir, output_dir / ".cache", artifact_manifest, theme)

    warnings: list[dict[str, Any]] = list(threshold_warnings)
    if not charts.manifest_found:
        warnings.append({
            "type": "artifact_manifest_missing",
            "detail": "找不到 artifact_manifest.json；本次無法對照圖檔，含圖頁面一律降級為 stat_callout。",
        })

    direction_body = str(slots.get("direction.body") or "").strip()
    if direction_body and _parse_direction_body(direction_body) is None:
        warnings.append({
            "type": "direction_unstructured",
            "detail": "direction.body 非結構化 JSON（舊契約純文字），已用條列版面過渡；"
                      "重跑 ai:report_ppt 產新契約後將呈現色塊流程＋題目卡。",
        })

    layout = resolve_layout(report_data, charts, theme,
                            _clean_layout_overrides(approvals.get("layout_overrides")))

    # 逐頁備妥 narrative（判讀式標題＋要點），fallback 一律寫 warning，不靜默。
    narratives_by_page: dict[int, tuple[str, list[dict[str, Any]], bool]] = {}
    # plan 的 narrative（頁碼→要點）；沒有 plan 時為空 dict，行為與既有一致。
    plan_narratives: dict[int, list[dict[str, Any]]] = {}
    for index, slide in enumerate((report_data.get("slide_plan") or {}).get("slides") or [],
                                  start=1):
        points = [{"text": str(n.get("text") or ""), "label": "",
                   "emphasis": bool(n.get("emphasis"))}
                  for n in (slide.get("narrative") or []) if n.get("text")]
        if points:
            plan_narratives[index] = points
    titled: list[PageSpec] = []
    for spec in layout:
        if spec.kind in {"cover", "direction", "section_divider"} or spec.is_appendix:
            # 附錄頁沒有文案 slot、只渲染表格，標題保留「附錄N：…」的定位不套判讀式標題。
            titled.append(spec)
            continue
        if spec.kind == "table":
            # 純表格頁（明細類，如家族完整性明細）本來就不配解讀——查了也是空，
            # 誤報 narrative_missing 會讓人白跑一趟「補解讀」（P1-4，實機 P17）。
            titled.append(spec)
            continue
        # 🔴 SlidePlan 自帶 narrative 時以它為準（規劃與敘述同一份產出，
        # 不必再去 narratives.json 找——那是固定頁序時代的來源）。
        plan_points = plan_narratives.get(spec.page)
        if plan_points is not None:
            narratives_by_page[spec.page] = (spec.title, plan_points, False)
            titled.append(spec)
            continue
        matched, variant = _narrative_entry(narratives, _narrative_candidates(spec))
        headline, points, fell_back = _normalize_narrative(variant) if variant else ("", [], False)
        narratives_by_page[spec.page] = (headline, points, fell_back)
        if not variant:
            warnings.append({
                "type": "narrative_missing", "page": spec.page,
                "report_key": ",".join(spec.report_keys),
                "detail": "找不到對應 narrative；本頁要點改用引擎 rows 的關鍵數字。",
            })
        elif fell_back:
            warnings.append({
                "type": "narrative_fallback", "page": spec.page,
                "report_key": matched or ",".join(spec.report_keys),
                "detail": "narrative 缺 headline／points（舊格式只有 text），已切段落並依版面截斷。",
            })
        if not headline and points:
            headline = _derive_headline(points)
            if headline:
                warnings.append({
                    "type": "headline_derived", "page": spec.page,
                    "report_key": matched or ",".join(spec.report_keys),
                    "detail": f"narrative 未提供 headline，標題取自要點首句：「{headline}」。",
                })
        titled.append(_spec_with(spec, title=f"{spec.topic}：{headline}" if headline else spec.topic))
    layout = titled

    included_labels = [
        _label_of(report_data, key)
        for spec in layout
        for key in spec.report_keys
        if _report_key_has_data(report_data, key)
    ]
    cover_stats = _cover_stats(report_data)
    ctx: dict[str, Any] = {
        "report_data": report_data,
        "narratives": narratives,
        "narratives_by_page": narratives_by_page,
        "slots": slots,
        "version": version,
        "charts": charts,
        "period": _statistics_period(report_data),
        "cover_stats": cover_stats,
        "framework_text": _framework_text(layout),
        "included_report_labels": list(dict.fromkeys(included_labels)),
    }

    prs = Presentation()
    prs.slide_width = Inches(theme.slide["width_in"])
    prs.slide_height = Inches(theme.slide["height_in"])
    blank = prs.slide_layouts[6]

    pages: list[dict[str, Any]] = []
    inner_density = float(theme.starfield["inner_page_density"])
    for spec in layout:
        slide = prs.slides.add_slide(blank)
        # 背景先畫＝壓在最底層；封面用滿密度，內頁調淡避免星點被當成圖表資料點。
        _add_background(slide, theme, charts.cache_dir,
                        density=1.0 if spec.kind in {"cover", "section_divider"} else inner_density)
        set_current_page(spec.page)
        RENDERERS[spec.kind](slide, theme, spec, ctx)
        page_info: dict[str, Any] = {
            "page": spec.page,
            "kind": spec.kind,
            "title": spec.title,
            "topic": spec.topic,
            "report_keys": list(spec.report_keys),
            "charts": list(spec.charts),
            "is_appendix": spec.is_appendix,
            "degraded_from": spec.degraded_from,
            "filled_slots": [s for s in spec.slots if slots.get(s)],
            "missing_slots": [s for s in spec.slots if not slots.get(s)],
            "missing_reports": [k for k in spec.report_keys if not _report_key_has_data(report_data, k)],
        }
        if spec.kind == "cover":
            page_info["stat_cards"] = len(cover_stats)
        if spec.degraded_from:
            warnings.append({
                "type": "chart_missing_degraded", "page": spec.page,
                "report_key": ",".join(spec.report_keys),
                "detail": f"artifact_manifest 內找不到可用圖檔，{spec.degraded_from} 已降級為 stat_callout。",
            })
        pages.append(page_info)

    warnings.extend(audit_layout(prs, theme))

    # R-1：direction.body 超出卡片容量要說出來（截斷是最後防線，不是合格狀態）。
    _direction_body = _parse_direction_body(slots.get("direction.body") or "")
    if _direction_body:
        for message in validate_direction_body(theme, _direction_body):
            warnings.append({"type": "direction_capacity_exceeded", "detail": message})

    # 🔴 H-2：整條放不下而被丟掉的要點必須說出來。
    # ⚠ 本專案的原則是「沒有靜默的截斷」——舊做法是截字加「…」，使用者當場抓到
    # 「這種卡掉的敘述不要再有」。現在改成不截字，但**不代表可以默默少一條**：
    # 正常情況這裡應該是空的（容量已由 narrative_capacity 交給 CLI），
    # 一旦有值就代表上游容量算錯了，要當成 bug 追。
    for page, (label, text, _color, _emph) in dropped_points():
        entry: dict[str, Any] = {
            "type": "points_dropped",
            "detail": f"要點整條未放入（版面不足，未截字）：{label}｜{text}",
        }
        if page is not None:
            entry["page"] = page
        warnings.append(entry)

    pptx_path = _next_available_path(output_dir, version)
    prs.save(str(pptx_path))

    selected = (report_data.get("parameters") or {}).get("reports_selected") or []
    rendered_keys = {key for page in pages for key in page["report_keys"]}
    manifest = {
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "builder_version": "v3",
        "source_report_version": version,
        "source_report_dir": str(report_dir),
        "pptx_file": pptx_path.name,
        "sha256": _sha256_of(pptx_path),
        "slot_total": len(all_slot_keys()),
        "slot_filled": sum(len(p["filled_slots"]) for p in pages),
        "missing_slots": sorted({s for p in pages for s in p["missing_slots"]}),
        # M-5（2026-08-04）：扣掉 EXCLUDED_FROM_PPT——那些是「刻意不進 PPT」
        # （07-31 定案，如 family_quality_detail），列進 missing_reports 會把
        # 「不放」誤報成「缺料」，讓監控與驗收把正常狀態當問題追。
        "missing_reports": sorted(
            ({str(key) for key in selected if not _report_key_has_data(report_data, str(key))}
             | {key for p in pages for key in p["missing_reports"]}
             | {str(key) for key in selected if str(key) not in rendered_keys and _report_key_has_data(report_data, str(key))})
            - EXCLUDED_FROM_PPT
        ),
        "metadata": {
            key: (report_data.get("parameters") or {}).get(key)
            for key in ("topic_run_id", "topic_state_version")
            if (report_data.get("parameters") or {}).get(key) is not None
        },
        "warnings": warnings,
        "pages": pages,
    }
    manifest_path = pptx_path.with_suffix(".manifest.json")
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"pptx_path": str(pptx_path), "manifest_path": str(manifest_path), "manifest": manifest}


def write_approval_template(path: Path) -> Path:
    """產出確認槽範本供填入定稿文案；v3 只有兩個槽。"""
    payload = {
        "report_version": "<報表版本>",
        "slots": {slot: "" for slot in all_slot_keys()},
        "layout_overrides": {},
    }
    Path(path).write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return Path(path)


def main() -> None:
    parser = argparse.ArgumentParser(description="專利分析報告 PPTX 產生器 v3（deterministic）")
    parser.add_argument("--report-dir", help="報表版本目錄（含 report_data.json、artifact_manifest.json）")
    parser.add_argument("--approvals", help="確認槽定稿文案 JSON")
    parser.add_argument("--output-dir", default="data/report_artifacts/ppt", help="輸出目錄")
    parser.add_argument("--init-approvals", help="產出確認槽範本到指定路徑後結束")
    args = parser.parse_args()

    if args.init_approvals:
        print(f"approval template: {write_approval_template(Path(args.init_approvals))}")
        return
    if not args.report_dir:
        parser.error("--report-dir is required unless --init-approvals is used")

    result = build_ppt(
        report_dir=args.report_dir, approvals_path=args.approvals, output_dir=args.output_dir
    )
    manifest = result["manifest"]
    print(f"pptx: {result['pptx_path']}")
    print(f"manifest: {result['manifest_path']}")
    print(f"sha256: {manifest['sha256']}")
    print(f"pages: {len(manifest['pages'])}")
    print(f"warnings: {len(manifest['warnings'])}")
    for warning in manifest["warnings"]:
        print(f"  - [{warning['type']}] {warning.get('page', '-')} {warning.get('detail', '')}".rstrip())


if __name__ == "__main__":
    main()
